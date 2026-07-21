#!/usr/bin/env python3
"""Component audit — did this deck hand-roll a form the library already implements?

WHY THIS EXISTS, measured: across three delivered decks (45 slides) the build scripts reached for
**3 of deckkit's 59 form components**. Everything else was composed from raw `box` + `text`. That is
not a style complaint — a hand-rolled form re-introduces the exact geometry bugs the component was
written to fix (a baseline hardcoded short of the last bar; a value label floating off the bar's
centerline; a segmented band whose parts do not sum; a timeline whose end caption drifts off its dot).
SKILL.md has said "when a COMPONENT exists for the form, BUILD that component" since forever — as
PROSE, with no gate. It was violated dozens of times and detected zero times.

WHY IT IS ADVISORY, NOT A BLOCKER — and this is a deliberate design decision, not timidity:
geometry alone CANNOT distinguish a lazy hand-roll from a deliberate bespoke composition, and the
deliberate one is the *signature move* — the single thing that makes a deck not-a-template. A gate
that blocked on "you hand-rolled a bar row" would delete the best page of the deck that motivated
it. So this tool reports and directs attention; it never vetoes. What it CAN state as fact is the
usage ratio and the specific component whose geometry the cluster matches.

CLI:  python component_audit.py <build_script.py> [<deck.pptx>] [--json]
Exit: 0 = nothing to report · 2 = at least one cluster matches an unused component (advisory).
"""
import argparse
import json
import os
import re
import sys

# EMITTERS — components that actually DRAW rect geometry, so their own output could be mistaken
# for a hand-roll and must suppress a report. This list has been wrong twice in two consecutive
# edits (columns() returns rects and draws nothing; table() emits a GraphicFrame) and each time it
# silently killed real detections deck-wide, so it is now DERIVED from deckkit's source instead of
# hand-kept. A component that only returns geometry, or only adds a table/chart/picture frame, can
# never be the source of a loose-rect cluster and must never suppress one.
def _emitting_components():
    import re
    try:
        src = open(os.path.join(os.path.dirname(os.path.abspath(__file__)), "deckkit.py"),
                   encoding="utf-8").read()
    except OSError:
        return set()
    out = set()
    for m in re.finditer(r"^def ([a-z][a-z0-9_]*)\(", src, re.M):
        name, start = m.group(1), m.start()
        nxt = re.search(r"^def ", src[start + 4:], re.M)
        body = src[start: start + 4 + (nxt.start() if nxt else len(src) - start)]
        if re.search(r"\b(box|_flat|add_shape|node|icon_tile|build_freeform|_iso_poly|iso_prism)\s*\(", body):
            out.add(name)
    return out


_DRAWS_RECTS = _emitting_components()

# EMITTERS must be the intersection of "draws rects" and "is a FORM component". The derivation
# alone returns 74 names because the PRIMITIVES (box, chip, arrow, bullet, icon_tile) draw rects
# too — and every deck calls dk.box(), so using the raw set suppressed 29 of 36 real decks, which
# is worse than the bug it replaced. Calling box() IS the hand-rolling; it can never excuse it.
# Populated after FORM_GUARANTEE is defined (see below).
EMITTERS = set()

# form components + the concrete guarantee each one gives that a hand-roll does not
FORM_GUARANTEE = {
    "native_chart": "a real editable chart; axis derived from the data, non-Latin labels render",
    "segmented_bar": "parts are normalised to sum, and each segment's label is fitted to its width",
    "meter_bar": "the value label is centered ON the bar's centerline by construction",
    "stat_row": "even column pitch and one baseline for every figure",
    "scorecard": "equal tiles with one hero, from the region — not a hand-picked pitch",
    "timeline": "end markers are inset so the first/last caption stays co-centred with its dot",
    "step_list": "numbered steps with a derived pitch; no stranded gap under the header",
    "tier_stack": "one taper (funnel/pyramid) sized from the values, not clamped by a min-size floor",
    "dot_strip": "even spacing with the end dots inset",
    "eval_matrix": "options x criteria with fitted glyph cells",
    "heat_matrix": "one colour scale bound to the value range",
    "leaderboard": "ranked rows keyed to a chart, with one highlight",
    "table": "header/rule/highlight styling and column widths that respect the content",
    "org_tree": "centroid parents and a horizontal bus; raises when it cannot fit legibly",
    "position_map": "labelled 2-D positions with anti-collision labels",
    "small_multiples": "every panel pinned to ONE shared value axis",
    "iso_bars": "a FAITHFUL 2.5D bar chart — height linear in the value, zero-based",
    "iso_stack": "an isometric layered stack with labels aligned to each slab",
    "iso_prism": "one extruded isometric block with fixed one-light-source face shading",
}


EMITTERS = _DRAWS_RECTS & set(FORM_GUARANTEE)

def _script_calls(path):
    r"""Every deckkit name this script calls, however it reached for it.

    A bare `\bdk\.` regex missed `from deckkit import scorecard` and any alias other than `dk`,
    so a deck that imported components directly had its OWN component output reported as a
    hand-roll — the tool punishing exactly the behaviour it exists to encourage.
    """
    src = open(path, encoding="utf-8", errors="ignore").read()
    names = set()
    aliases = {"dk", "deckkit"}
    for m in re.finditer(r"^\s*import\s+deckkit\s+as\s+([A-Za-z_]\w*)", src, re.M):
        aliases.add(m.group(1))
    for m in re.finditer(r"^\s*from\s+deckkit\s+import\s+([^\n(]+)", src, re.M):
        for part in m.group(1).split(","):
            part = part.strip()
            if not part or part == "*":
                continue
            names.add(part.split(" as ")[-1].strip())
            names.add(part.split(" as ")[0].strip())
    for a in aliases:
        names |= {m.group(1) for m in re.finditer(r"\b%s\.([a-z][a-z0-9_]*)\s*\(" % re.escape(a), src)}
    return names


def _shapes(pptx_path):
    from pptx import Presentation
    from pptx.util import Emu
    out = []
    prs = Presentation(pptx_path)
    for n, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            try:
                x, y = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except Exception:
                continue
            txt = ""
            try:
                if sh.has_text_frame:
                    txt = sh.text_frame.text.strip()
            except Exception:
                pass
            filled = False
            try:
                filled = sh.fill.type is not None and "BACKGROUND" not in str(sh.fill.type)
            except Exception:
                pass
            if getattr(sh, "has_chart", False):
                kind = "chart"
            elif getattr(sh, "has_table", False):
                kind = "table"
            elif "PICTURE" in str(getattr(sh, "shape_type", "")):
                kind = "pic"
            elif filled:
                # a FILLED shape is rect geometry whether or not it also carries a label. The
                # earlier "text if txt else rect" test made the MOST COMMON hand-roll — box(...)
                # with its own text inside — invisible to all four detectors.
                kind = "rect"
            else:
                kind = "text" if txt else "rect"
            out.append({"slide": n, "x": x, "y": y, "w": w, "h": h, "kind": kind, "txt": txt,
                        "self_labelled": bool(txt and filled)})
    return out


def _clusters(shapes):
    """Geometry signatures that a deckkit component already implements."""
    hits = []
    by_slide = {}
    for s in shapes:
        by_slide.setdefault(s["slide"], []).append(s)

    for n, shs in by_slide.items():
        rects = [s for s in shs if s["kind"] == "rect" and s["w"] > 0.12 and s["h"] > 0.06]
        texts = [s for s in shs if s["kind"] == "text"]

        def near_text(r, pad=0.55):
            own = [r] if r.get("self_labelled") else []
            return own + [t for t in texts
                          if abs(t["y"] - r["y"]) < pad and t["x"] > r["x"] - 2.6
                          and t["x"] < r["x"] + r["w"] + 2.6]

        # BAR ROW — same x0 and h, different w, stacked vertically, with digit-bearing labels
        rows = {}
        for r in rects:
            rows.setdefault((round(r["x"], 2), round(r["h"], 2)), []).append(r)
        for (x0, h), group in rows.items():
            if len(group) < 3:
                continue
            if len({round(g["w"], 2) for g in group}) < 2:
                continue                       # all identical width -> that is a card stack, below
            if any(re.search(r"\d", t["txt"]) for g in group for t in near_text(g)):
                hits.append({"slide": n, "pattern": "bar row",
                             "detail": f"{len(group)} rects sharing x={x0}in and h={h}in with varying "
                                       f"width and numeric labels",
                             "suggest": ["native_chart", "meter_bar", "segmented_bar"]})

        # 100% BAND — >=3 rects on one y, abutting (gap ~0), spanning a wide run
        bands = {}
        for r in rects:
            bands.setdefault((round(r["y"], 2), round(r["h"], 2)), []).append(r)
        for (y0, h), group in bands.items():
            if len(group) < 3:
                continue
            g = sorted(group, key=lambda r: r["x"])
            gaps = [g[i + 1]["x"] - (g[i]["x"] + g[i]["w"]) for i in range(len(g) - 1)]
            span = (g[-1]["x"] + g[-1]["w"]) - g[0]["x"]
            if span > 3.0 and all(abs(v) < 0.06 for v in gaps):
                hits.append({"slide": n, "pattern": "abutting 100% band",
                             "detail": f"{len(group)} rects abutting across {span:.1f}in at y={y0}in",
                             "suggest": ["segmented_bar"]})

        # CARD/TILE ROW — >=3 identical rects evenly spaced on one axis
        for (y0, h), group in bands.items():
            if len(group) < 3:
                continue
            g = sorted(group, key=lambda r: r["x"])
            if len({round(r["w"], 2) for r in g}) != 1:
                continue
            pitches = [g[i + 1]["x"] - g[i]["x"] for i in range(len(g) - 1)]
            if pitches and max(pitches) - min(pitches) < 0.05 and min(pitches) > g[0]["w"]:
                hits.append({"slide": n, "pattern": "tile row",
                             "detail": f"{len(group)} identical rects on an even pitch at y={y0}in",
                             "suggest": ["scorecard", "stat_row", "columns"]})

        # MARKER ROW — >=3 small same-size shapes evenly spaced with captions under them
        dots = [r for r in rects if r["w"] < 0.4 and r["h"] < 0.4]
        if len(dots) >= 3:
            g = sorted(dots, key=lambda r: r["x"])
            ys = {round(r["y"], 1) for r in g}
            pitches = [g[i + 1]["x"] - g[i]["x"] for i in range(len(g) - 1)]
            if len(ys) == 1 and pitches and max(pitches) - min(pitches) < 0.08:
                caps = [t for t in texts if any(abs(t["x"] + t["w"] / 2 - (r["x"] + r["w"] / 2)) < 1.2
                                                and t["y"] > r["y"] for r in g)]
                if len(caps) >= 3:
                    hits.append({"slide": n, "pattern": "marker row with captions",
                                 "detail": f"{len(g)} evenly spaced markers with captions below",
                                 "suggest": ["timeline", "spaced_centers", "step_list"]})
    # one visual row must not be reported twice (a body rect and its own accent rule are two
    # rows at nearly the same y)
    seen, dedup = set(), []
    for h in sorted(hits, key=lambda x: (x["slide"], x["pattern"])):
        k = (h["slide"], h["pattern"])
        if k in seen:
            continue
        seen.add(k); dedup.append(h)
    return dedup


def audit(script_path, pptx_path=None):
    if not os.path.isfile(script_path):
        return {"used_forms": [], "form_total": len(FORM_GUARANTEE), "clusters": [],
                "actionable": [], "inspected": False,
                "why": "build script not found: %s" % script_path}
    called = _script_calls(script_path)
    used_forms = sorted(f for f in FORM_GUARANTEE if f in called)
    hits, inspected, why = [], True, ""
    if not pptx_path:
        inspected, why = False, "no deck inspected — pass the .pptx explicitly"
    elif not os.path.isfile(pptx_path):
        inspected, why = False, "deck not found: %s" % pptx_path
    else:
        try:
            hits = _clusters(_shapes(pptx_path))
        except Exception as e:                                   # noqa: BLE001
            inspected, why = False, "could not read the deck: %s" % e
    # TWO LISTS, deliberately not one:
    #   suggest  — what the report tells you to reach for (may include LAYOUT helpers such as
    #              columns(), which is the right advice even though it draws nothing);
    #   emitters — what could have DRAWN this geometry, used ONLY to suppress.
    # Conflating them silently broke the tool once: columns() went into the suppression list, and
    # since 12 of 44 real decks call it, a quarter of all tile-row detections were being killed by
    # a helper that never drew a tile. Suppression is what stops a component's OWN output being
    # reported as a hand-roll (org_tree's node rects are indistinguishable from hand-placed boxes
    # in the finished file) — so `emitters` lists every component that DRAWS, and nothing that
    # merely returns geometry.
    # Suppress when the deck called ANY component that draws rects: in a finished pptx a
    # component's own output is indistinguishable from hand-placed boxes. EMITTERS is derived
    # from deckkit's source, so it cannot drift the way the hand-kept list did (twice).
    drew = EMITTERS & called
    actionable = [h for h in hits if h.get("suggest")] if not drew else []
    suppressed_by = sorted(drew)
    return {"used_forms": used_forms, "form_total": len(FORM_GUARANTEE),
            "clusters": hits, "actionable": actionable, "suppressed_by": suppressed_by,
            "inspected": inspected, "why": why}


def main():
    ap = argparse.ArgumentParser(description="did this deck hand-roll a form the library implements?")
    ap.add_argument("script")
    ap.add_argument("pptx", nargs="?")
    ap.add_argument("--json", action="store_true", dest="as_json")
    a = ap.parse_args()
    pptx = a.pptx
    if pptx is None:
        d = os.path.dirname(os.path.abspath(a.script))
        cand = [f for f in os.listdir(d) if f.endswith(".pptx") and not f.startswith("~$")]
        pptx = os.path.join(d, cand[0]) if len(cand) == 1 else None
    r = audit(a.script, pptx)
    if a.as_json:
        print(json.dumps(r, indent=1))
        sys.exit(1 if not r.get("inspected") else (2 if r["actionable"] else 0))
    # NEVER report clean for a deck that was not opened: a wrong path, an unreadable file, or an
    # ambiguous directory used to print the success line and exit 0 — a green PRE-FLIGHT tick for
    # a check that did no work, which is the worst failure a checklist tool can have.
    if not r.get("inspected"):
        print("[components] NOT CHECKED — {}".format(r.get("why") or "unknown"))
        print("[components] rerun with the deck path: component_audit.py <build_*.py> <deck.pptx>")
        sys.exit(1)
    print("[components] this deck calls {} of the {} form components this tool can name a "
          "guarantee for (deckkit's wider form catalogue is ~59): {}".format(
              len(r["used_forms"]), r["form_total"], ", ".join(r["used_forms"]) or "none"))
    if r.get("suppressed_by"):
        print("[components] cluster reporting suppressed — this deck draws with {}, whose own "
              "output is indistinguishable from hand-placed boxes in a finished file."
              .format(", ".join(r["suppressed_by"])))
        sys.exit(0)
    if not r["actionable"]:
        print("[components] no hand-rolled cluster matches an unused component.")
        sys.exit(0)
    print("[components] {} cluster(s) look like a form the library already implements. "
          "This is ADVISORY:".format(len(r["actionable"])))
    print("             a bespoke composition IS the signature move — but a hand-rolled COMMON form")
    print("             re-inherits the geometry bugs the component fixed. Confirm each is deliberate.")
    print("             NOTE: this is about HOW to draw a form you already chose, never about WHETHER")
    print("             that form is right — form-selection.md owns that, and 'use scorecard' is not")
    print("             a licence for a card grid.")
    CAP = 8
    for h in r["actionable"][:CAP]:
        print(f"  slide {h['slide']:>2}  {h['pattern']}: {h['detail']}")
        for c in h["suggest"]:
            g = FORM_GUARANTEE.get(c)
            print(f"            deckkit.{c}() — {g}" if g else f"            deckkit.{c}()")
    if len(r["actionable"]) > CAP:
        rest = len(r["actionable"]) - CAP
        pats = sorted({h["pattern"] for h in r["actionable"][CAP:]})
        print(f"  … and {rest} more ({', '.join(pats)}) — rerun with --json for the full list.")
    sys.exit(2)


if __name__ == "__main__":
    main()

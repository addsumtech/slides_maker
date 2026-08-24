#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""run_eval.py — score a produced deck against an eval's machine-decidable assertions.

WHY THIS EXISTS
    The 20 suites in tests/ all ask "does this code work". None asks "do these 2109 lines of
    SKILL.md, 35 references and 5 agents make the DECK better" — so every claim that a change
    improved the skill has been an impression. CLAUDE.md states the failure directly: a lossless
    refactor can score 2298/2298 and still degrade the skill, because the regression is in WHEN
    content reaches context, not whether it survives.

    Twice now that has bitten for real: a refactor whose next live run read 3 of 10 reference
    files (and shipped zero icons, and printed a raw format string on a slide, and passed every
    automated gate), and an opt-in gate that reported 0 findings on the very deck whose defect
    motivated it. Both are invisible to unit tests and both are visible here.

THE SPLIT, WHICH IS THE WHOLE DESIGN
    `--score` reads an ALREADY-BUILT deck directory and decides. No model, no network, ~1s. That
    half is CI-able and is what tests/ exercises.
    `--run` invokes the pipeline to produce the deck first. That half costs minutes and real
    tokens, so it is developer- or nightly-only and never gated on a PR.

    Keeping them apart is what stops this from becoming the kind of eval nobody runs.

WHAT IT DELIBERATELY DOES NOT DO
    No taste scoring. No LLM judge, no aesthetic composite. Both paths were examined against the
    published benchmarks and rejected on the merits: reference-similarity scoring rewards
    imitation (the opposite of this skill's direction gate and diversity floors), and closed-form
    beauty scores fight its own semantic-colour and variation rules. Everything asserted here is
    something a program can be wrong about in only one direction.

USAGE
    python3 scripts/run_eval.py --list
    python3 scripts/run_eval.py --score ~/Downloads/my-deck --eval happy-source-provided \\
                               [--transcript run.jsonl] [--record]
    python3 scripts/run_eval.py --score ~/Downloads/my-deck --eval happy-source-provided --json out.json

EXIT CODES
    0 = every assertion passed · 1 = at least one failed · 2 = could not run (never "clean")
"""
import argparse
import json
import os
import re
import subprocess
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
SKILL = os.path.dirname(HERE)
EVALS = os.path.join(SKILL, "evals", "evals.json")
HISTORY = os.path.join(SKILL, "evals", "history.json")

# A number on a slide: 2.1 · 0.94 · 1,261 · 45.8% · 8x. Deliberately greedy on separators so
# "1,261.1" is ONE number rather than three — a split would make the unsourced check nonsense.
_NUM = re.compile(r"\d[\d,]*(?:\.\d+)?")


def die(msg, code=2):
    print("run_eval: " + msg, file=sys.stderr)
    sys.exit(code)


def load_evals():
    try:
        with open(EVALS, encoding="utf-8") as fh:
            return json.load(fh)["evals"]
    except (OSError, ValueError, KeyError) as exc:
        die("cannot read %s (%s)" % (EVALS, exc))


def deck_files(deck_dir):
    """The .pptx a deck folder is built around, plus its gates file."""
    if not os.path.isdir(deck_dir):
        die("%s is not a directory" % deck_dir)
    pptx = [f for f in sorted(os.listdir(deck_dir))
            if f.endswith(".pptx") and not f.startswith("~$")]
    return (os.path.join(deck_dir, pptx[0]) if pptx else None,
            os.path.join(deck_dir, ".deck-gates.json"))


def slide_texts(pptx):
    from pptx import Presentation
    out = []
    for slide in Presentation(pptx).slides:
        buf = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                buf.append(sh.text_frame.text or "")
        out.append("\n".join(buf))
    return out


_C_NS = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_CJK_RE = re.compile("[぀-ヿ㐀-䶿一-鿿豈-﫿가-힯]")


def _parts(pptx, prefix):
    """Yield (name, parsed XML) for package parts under `prefix`. Unparseable parts are skipped —
    this is an eval, and a corrupt part is a different failure that other checks already own."""
    import xml.etree.ElementTree as ET
    import zipfile
    with zipfile.ZipFile(pptx) as z:
        for name in z.namelist():
            if not (name.startswith(prefix) and name.endswith(".xml")):
                continue
            try:
                yield name, ET.fromstring(z.read(name))
            except ET.ParseError:
                continue


def signed_native_bars(pptx):
    """Native BAR/COLUMN charts carrying a NEGATIVE value — a rendering bug no gate here can see.

    Measured: LibreOffice draws a negative column/bar at its ABSOLUTE height, so a -12% bar comes
    out of the rasteriser looking identical to +12%. The .pptx itself is correct, so every check
    that reads the deck passes, and every check that reads the PNG is looking at a plausible chart
    pointing the wrong way. Nothing in this repo catches it — not lint_deck, not the visual
    contract, not a critic reading a render. The only reliable fix is to not use a native chart
    for signed data (draw it with shapes, or raster it), which makes this decidable from the
    artifact: a native bar chart with a negative cached value is the defect.

    Scoped to `c:barChart` on purpose: line and scatter render signed data correctly, and flagging
    them would train authors to route working charts through the workaround.
    """
    hits = []
    for name, root in _parts(pptx, "ppt/charts/"):
        for bar in root.iter(_C_NS + "barChart"):
            for val in bar.iter(_C_NS + "val"):
                for v in val.iter(_C_NS + "v"):
                    try:
                        if float((v.text or "").strip()) < 0:
                            hits.append(name)
                            break
                    except ValueError:
                        continue
                if hits and hits[-1] == name:
                    break
            if hits and hits[-1] == name:
                break
    return sorted(set(hits))


def cjk_runs_missing_ea(pptx):
    """Runs containing CJK text whose `a:rPr` declares no `<a:ea>` East-Asian typeface.

    PowerPoint picks the face for CJK glyphs from `<a:ea>`, not `<a:latin>`. A run that sets only
    `<a:latin>` renders its Chinese through whatever the host substitutes — a different face, a
    different width, and on a machine without the fallback, tofu. This repo has already shipped a
    geometry bug from the same asymmetry (every width measurement read `<a:latin>` while the glyphs
    came from `<a:ea>`, under-reporting Chinese runs by ~46%), which is exactly why a CJK deck
    needs an assertion that reads the file rather than a reviewer who reads the render on the one
    machine that happens to have the font.
    """
    bad = []
    for name, root in _parts(pptx, "ppt/slides/slide"):
        for run in root.iter(_A_NS + "r"):
            t = run.find(_A_NS + "t")
            if t is None or not _CJK_RE.search(t.text or ""):
                continue
            rpr = run.find(_A_NS + "rPr")
            if rpr is None or rpr.find(_A_NS + "ea") is None:
                bad.append("%s: %s" % (name.rsplit("/", 1)[-1], (t.text or "")[:18]))
    return bad


def canvas_inches(pptx):
    from pptx import Presentation
    prs = Presentation(pptx)
    return (prs.slide_width / 914400.0, prs.slide_height / 914400.0)


def notes_texts(pptx):
    from pptx import Presentation
    out = []
    for slide in Presentation(pptx).slides:
        try:
            out.append(slide.notes_slide.notes_text_frame.text or "")
        except Exception:
            out.append("")
    return out


def hard_findings(pptx):
    """The render-time lint's HARD findings. `--static` keeps a static deck from failing on the
    advisory NO BUILDS; the [warn]/[stats] lines are advisory by definition and not asserted."""
    r = subprocess.run([sys.executable, os.path.join(HERE, "lint_deck.py"), pptx, "--static"],
                       capture_output=True, text=True)
    return [l.strip() for l in (r.stdout or "").splitlines()
            if re.match(r"\s*slide \d+:", l) and "[warn]" not in l]


def check(a, ctx):
    """Return (ok, detail). Every assertion must be decidable from the artifacts alone."""
    kind = a.get("kind")
    if kind == "pptx_exists":
        return bool(ctx["pptx"]), "no .pptx in the deck folder"
    if kind == "slides_between":
        n = len(ctx["texts"])
        return a["lo"] <= n <= a["hi"], "%d slides, wanted %d-%d" % (n, a["lo"], a["hi"])
    if kind == "lint_hard_clean":
        f = hard_findings(ctx["pptx"])
        return not f, "%d hard finding(s): %s" % (len(f), (f[0][:80] if f else ""))
    if kind == "gates_present":
        g = ctx["gates"]
        if g is None:
            return False, "no .deck-gates.json — the hand-off evidence file was never written"
        missing = [k for k in a["keys"] if k not in g]
        return not missing, "missing gate block(s): %s" % ", ".join(missing)
    if kind == "text_contains":
        blob = "\n".join(ctx["texts"])
        missing = [v for v in a["values"] if v not in blob]
        return not missing, "not on any slide: %s" % ", ".join(missing)
    if kind == "text_absent":
        blob = "\n".join(ctx["texts"]).lower()
        found = [v for v in a["values"] if v.lower() in blob]
        return not found, "build residue on a slide: %s" % ", ".join(found)
    if kind == "no_unsourced_numbers":
        allow = set(a.get("allow") or [])
        # Page numbers, years and the as-of date are chrome, not claims.
        extra = set(str(i) for i in range(1, 100)) | set(str(y) for y in range(1900, 2100))
        bad = []
        for i, t in enumerate(ctx["texts"], 1):
            for m in _NUM.findall(t):
                s = m.rstrip(".")
                if s in allow or s in extra or s.replace(",", "") in allow:
                    continue
                bad.append("s%d:%s" % (i, s))
        return not bad, "figures with no source in the prompt: %s" % ", ".join(bad[:8])
    if kind == "notes_coverage":
        notes = ctx["notes"]
        if not notes:
            return False, "no slides"
        have = sum(1 for n in notes if len((n or "").strip()) > 20)
        share = have / float(len(notes))
        return share >= a.get("min_share", 0.8), "%d of %d slides carry real notes (%.0f%%)" % (
            have, len(notes), share * 100)
    if kind == "canvas_is":
        w, h = canvas_inches(ctx["pptx"])
        tol = a.get("tol", 0.05)
        good = abs(w - a["w_in"]) <= tol and abs(h - a["h_in"]) <= tol
        return good, "canvas is %.2f x %.2f in, wanted %.2f x %.2f" % (w, h, a["w_in"], a["h_in"])
    if kind == "no_signed_native_bars":
        hits = signed_native_bars(ctx["pptx"])
        return not hits, ("native bar/column chart(s) carrying negative values: %s — LibreOffice "
                          "renders these at absolute height, so the PNG shows the bar pointing the "
                          "wrong way while the .pptx is correct. Draw signed data with shapes or "
                          "raster it." % ", ".join(hits))
    if kind == "cjk_runs_have_ea_face":
        bad = cjk_runs_missing_ea(ctx["pptx"])
        return not bad, ("%d CJK run(s) declare no <a:ea> typeface (e.g. %s) — PowerPoint takes the "
                         "East-Asian face from <a:ea>, so these render through a host substitution"
                         % (len(bad), "; ".join(bad[:4])))
    if kind == "reference_reached":
        tr = ctx["transcript"]
        if tr is None:
            return None, "needs --transcript; nothing else can answer whether a file was READ"
        missing = [p for p in a["paths"] if p not in tr]
        return not missing, "never opened during the run: %s" % ", ".join(missing)
    return None, "unknown assertion kind %r" % kind


def score(deck_dir, ev, transcript=None):
    pptx, gates_path = deck_files(deck_dir)
    ctx = {"pptx": pptx, "texts": [], "notes": [], "gates": None, "transcript": None}
    if pptx:
        ctx["texts"] = slide_texts(pptx)
        ctx["notes"] = notes_texts(pptx)
    if os.path.exists(gates_path):
        try:
            with open(gates_path, encoding="utf-8") as fh:
                ctx["gates"] = json.load(fh)
        except ValueError:
            ctx["gates"] = None
    if transcript:
        try:
            with open(transcript, encoding="utf-8", errors="replace") as fh:
                ctx["transcript"] = fh.read()
        except OSError as exc:
            die("cannot read transcript %s (%s)" % (transcript, exc))

    rows = []
    for a in ev.get("assertions") or []:
        ok, detail = (False, "assertion crashed")
        try:
            ok, detail = check(a, ctx)
        except Exception as exc:                     # a crashed assertion is never a pass
            detail = "%s: %s" % (type(exc).__name__, str(exc)[:70])
        rows.append({"kind": a.get("kind"), "ok": ok, "detail": detail,
                     "why": a.get("why", "")})
    return rows


def main():
    ap = argparse.ArgumentParser(add_help=True)
    ap.add_argument("--list", action="store_true")
    ap.add_argument("--score", metavar="DECK_DIR")
    ap.add_argument("--eval", metavar="ID")
    ap.add_argument("--transcript", metavar="FILE")
    ap.add_argument("--json", metavar="OUT")
    ap.add_argument("--record", action="store_true",
                    help="append the result to evals/history.json under the current VERSION")
    args = ap.parse_args()

    evals = load_evals()
    if args.list:
        for e in evals:
            n = len(e.get("assertions") or [])
            print("%-24s %2d assertion(s)  %s" % (e["id"], n, e["prompt"][:56]))
        return 0
    if not args.score:
        ap.print_help()
        return 2
    if not args.eval:
        die("--score needs --eval <id>; run --list to see them")
    ev = next((e for e in evals if e["id"] == args.eval), None)
    if ev is None:
        die("no eval %r (have: %s)" % (args.eval, ", ".join(e["id"] for e in evals)))

    rows = score(args.score, ev, args.transcript)
    print("eval %s  ·  %s" % (ev["id"], args.score))
    hard = [r for r in rows if r["ok"] is False]
    skipped = [r for r in rows if r["ok"] is None]
    for r in rows:
        mark = "  ok  " if r["ok"] else ("  skip" if r["ok"] is None else "  FAIL")
        line = "%s %s" % (mark, r["kind"])
        if r["ok"] is not True:
            line += "  — " + r["detail"]
        print(line)
        if r["ok"] is False and r["why"]:
            print("         why it matters: " + r["why"])
    n_ok = sum(1 for r in rows if r["ok"] is True)
    print("\n%d passed, %d failed, %d skipped" % (n_ok, len(hard), len(skipped)))
    if skipped:
        print("  (a skipped assertion is NOT a pass — it could not be decided from what was given)")

    if args.json:
        with open(args.json, "w", encoding="utf-8") as fh:
            json.dump({"eval": ev["id"], "deck": args.score, "rows": rows}, fh,
                      ensure_ascii=False, indent=1)
    if args.record:
        try:
            with open(os.path.join(SKILL, "VERSION"), encoding="utf-8") as fh:
                version = fh.read().strip()
        except OSError:
            version = "unknown"
        try:
            with open(HISTORY, encoding="utf-8") as fh:
                hist = json.load(fh)
        except (OSError, ValueError):
            hist = []
        hist.append({"version": version, "eval": ev["id"], "deck": args.score,
                     "passed": n_ok, "failed": len(hard), "skipped": len(skipped),
                     "rows": [{"kind": r["kind"], "ok": r["ok"]} for r in rows]})
        with open(HISTORY, "w", encoding="utf-8") as fh:
            json.dump(hist, fh, ensure_ascii=False, indent=1)
            fh.write("\n")
        print("recorded -> %s (VERSION %s)" % (HISTORY, version))
    return 1 if hard else 0


try:                                            # console safety: a legacy code page must
    from _console import safe_stdio             # degrade a tick, never kill the report
    safe_stdio()
except Exception:
    pass


if __name__ == "__main__":
    sys.exit(main())

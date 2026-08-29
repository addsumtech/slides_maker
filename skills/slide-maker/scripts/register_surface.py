#!/usr/bin/env python3
"""Build a register's SURFACE, not just its palette — the half `presets.apply()` never had.

WHY. `apply()` calls exactly four things: `set_palette`, `set_geometry`, `set_ground` and the font
setters. That is a colourway and a corner radius. Measured, and rendered: one identical page taken
through all 18 presets produced 18 pages differing only in ground colour, ink colour, accent, one
font swap and a line weight. memphis had none of its colour bands or scattered marks, bauhaus none
of its oversized primitives, glassmorphism no glass, risograph no overprint, terminal no scanlines.
Every preset's `surface` field describes those things correctly — in prose, for an author to read
and build by hand, which in practice meant they were never built.

WHAT A KIT IS. Three things, because a register lives in all three and a background alone is a
wallpaper:

    ground(slide, register, role=...)   paints the register's own furniture and RETURNS the content
                                        rect that is left — like `title_bar()` returning its
                                        content top. Marks stay in the margins; nothing is painted
                                        into the band it hands back.
    card(slide, register, x, y, w, h)   the register's CARD FORM. Same call, different object: a
                                        memphis banded card, a bauhaus hard square, a riso sticker
                                        with a crisp offset, a glass panel, a terminal output block.
                                        This is what stops the eighteen pages being one page.
    marks                               the public primitives the grounds are built from
                                        (`halftone`, `starburst`, `boomerang`, `zigzag`, `tri`,
                                        `scanlines`, `color_band`) — reusable in a bespoke register.

DETERMINISM. Placement varies per page so a deck does not repeat one arrangement 12 times, but it
varies by INDEX, never by a random number: the same deck built twice is byte-identical, which the
render-parallel test and every diff-based check depend on.

IT OBEYS ITS OWN REGISTER. `check_register_guard` reads the same `presets.FORBIDS` these builders
are written against — bauhaus gets exactly ONE oversized primitive (never a confetti), risograph's
halftone is discrete dots rather than a gradient, terminal stays monospace, and every shape here
goes through deckkit's `_flat` + `shadow.inherit = False` so nothing arrives with a theme shadow.
The selftest runs the guard over every kit's own output, so a builder that violates the register it
claims to build fails here rather than in a delivered deck.

    python3 scripts/register_surface.py --selftest
    python3 scripts/register_surface.py --sample <out.pptx>   # one page per kit, to LOOK at
    python3 scripts/register_surface.py --list                # each kit beside its surface spec
"""
from __future__ import annotations

import argparse
import sys
import warnings
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx.enum.shapes import MSO_SHAPE            # noqa: E402
from pptx.util import Inches, Pt                  # noqa: E402

import deckkit as dk                              # noqa: E402

MARGIN = 0.45                  # marks live outside the content band, never on top of it
LOUD_MAX = 3                   # the motif budget this skill already states, applied to the ground


# ---------------------------------------------------------------------------- deterministic spread

def _h(index, salt=0):
    """A stable small integer from a page index — placement must not need a random number.

    `Math.random`-style variation makes two builds of one deck differ, which breaks the byte-identity
    the parallel-render check proves and makes every visual diff untrustworthy. Varying by index
    gives the same spread every time and still stops twelve pages sharing one arrangement.
    """
    v = (int(index) + 1) * 2654435761 + salt * 40503
    return (v ^ (v >> 13)) & 0xFFFF


def _pick(seq, index, salt=0):
    return seq[_h(index, salt) % len(seq)]


def _frac(index, salt=0):
    return (_h(index, salt) % 1000) / 1000.0


_MARKS = []                    # rects painted by the CURRENT ground() call, minus textures


def _record(x, y, w, h, texture):
    """A LOUD mark is one that must clear the content band; a texture may lie under it.

    The difference is not decorative bookkeeping. Scanlines, glows and a faint screen are grounds:
    type sits on them and the contrast checks judge the result. A triangle, a disc, a burst, an
    oversized primitive are objects: type landing on one is the collision a reader sees first, and
    this module's own docstring promised it could not happen — while the first render put a memphis
    triangle through a card corner and a bauhaus disc through a third card. A promise with nothing
    measuring it is how that ships.
    """
    if not texture:
        _MARKS.append((x, y, x + w, y + h))


def _blend(a, b, t):
    """`t` of the way from colour `a` to colour `b` — how a texture is made faint on any ground."""
    A, B = dk._as_rgb(a), dk._as_rgb(b)          # RGBColor is a bytes-like triple; the API only
    return dk.RGBColor(*[int(round(A[i] + (B[i] - A[i]) * t))   # accepts RGBColor back, not a tuple
                         for i in range(3)])


def _note(x, y, w, h):
    """Register a mark this module did NOT draw itself — a `seal`, a `ghost_numeral`, a badge.

    The invariant only ever saw shapes made by `_shape` and the mark helpers, so a loud object
    drawn through a deckkit helper was invisible to it. Measured: editorial_paper's ghost folio
    landed inside the returned content rect AND ran off the bottom edge, and the check said
    nothing — a contract that only inspects the calls you remembered to route through it is not
    a contract.
    """
    _record(x, y, w, h, False)


def _shape(slide, kind, x, y, w, h, fill=None, line=None, line_w=1.0, rot=None, texture=False):
    """One primitive, flattened the way every deckkit helper flattens: no theme style, no shadow."""
    _record(x, y, w, h, texture)
    s = dk._flat(slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h)))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = dk._as_rgb(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = dk._as_rgb(line)
        s.line.width = Pt(line_w * getattr(dk, "RULE_W_SCALE", 1.0))
    s.shadow.inherit = False
    if rot:
        s.rotation = rot
    return s


# ------------------------------------------------------------------------------------------ marks

def color_band(slide, y, h, color, *, x=0.0, w=None):
    """A full-bleed band of colour. memphis's header bands, riso's ink bar, a section rule."""
    W, _H_ = dk._slide_size(slide)
    return _shape(slide, MSO_SHAPE.RECTANGLE, x, y, (W - x) if w is None else w, h, fill=color)


def tri(slide, x, y, w, h, color, *, direction="up", line=None, line_w=1.0):
    """An isosceles triangle — bauhaus's third primitive, memphis's scatter, midcentury's wedge."""
    rot = {"up": 0, "right": 90, "down": 180, "left": 270}.get(direction, 0)
    return _shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, w, h,
                  fill=color, line=line, line_w=line_w, rot=rot)


def zigzag(slide, x, y, w, h, color, *, teeth=5, line_w=2.5, texture=False):
    """A memphis squiggle, drawn as a real polyline of hairlines rather than a picture."""
    _record(x, y, w, h, texture)
    out = []
    step = w / max(1, teeth)
    for i in range(teeth):
        x0 = x + i * step
        y0 = y + (0 if i % 2 == 0 else h)
        y1 = y + (h if i % 2 == 0 else 0)
        c = dk._flat(slide.shapes.add_connector(
            1, Inches(x0), Inches(y0), Inches(x0 + step), Inches(y1)))
        c.line.color.rgb = dk._as_rgb(color)
        c.line.width = Pt(line_w * getattr(dk, "RULE_W_SCALE", 1.0))
        c.shadow.inherit = False
        out.append(c)
    return out


def halftone(slide, x, y, w, h, color, *, cols=14, rows=8, r_max=0.11, r_min=0.02, index=0,
             texture=False):
    """A riso/print DOT FIELD that ramps across the rect — discrete dots, never a gradient.

    risograph's own guard forbids gradients (its fills are flat and its shadows hard), so the
    obvious way to fake a tonal ramp is exactly the thing the register prohibits. Real riso tone is
    a screen: the dots stay the same colour and change SIZE, which is what this draws.
    """
    _record(x, y, w, h, texture)
    out = []
    dx, dy = w / cols, h / rows
    for cx in range(cols):
        t = cx / max(1, cols - 1)
        r = r_max - (r_max - r_min) * t
        for cy in range(rows):
            if r <= 0.005:
                continue
            out.append(_shape(slide, MSO_SHAPE.OVAL,
                              x + cx * dx + (dx - r) / 2, y + cy * dy + (dy - r) / 2,
                              r, r, fill=color, texture=True))   # the FIELD was recorded, not each dot
    return out


def starburst(slide, cx, cy, r, color, *, rays=12, line_w=1.2, index=0, texture=False):
    """The midcentury atomic burst: hairlines radiating from a point, alternating length."""
    import math
    _record(cx - r, cy - r, 2 * r, 2 * r, texture)
    out = []
    for i in range(rays):
        a = (2 * math.pi * i / rays) + (_frac(index, 7) * 0.4)
        rr = r * (1.0 if i % 2 == 0 else 0.62)
        c = dk._flat(slide.shapes.add_connector(
            1, Inches(cx), Inches(cy), Inches(cx + rr * math.cos(a)), Inches(cy + rr * math.sin(a))))
        c.line.color.rgb = dk._as_rgb(color)
        c.line.width = Pt(line_w * getattr(dk, "RULE_W_SCALE", 1.0))
        c.shadow.inherit = False
        out.append(c)
    return out


def boomerang(slide, x, y, w, h, color, *, line_w=None, index=0):
    """The other midcentury mark: an open arc, drawn as a thick-stroked ARC primitive.

    Stroked rather than filled, because a filled crescent at this scale reads as a blob and the
    midcentury mark is a GESTURE — the eye should follow it, not weigh it.
    """
    s = _shape(slide, MSO_SHAPE.ARC, x, y, w, h, fill=None,
               line=color, line_w=line_w if line_w is not None else 3.0,
               rot=_pick((0, 30, 200, 250), index, 11))
    return s


def scanlines(slide, color, *, spacing=0.09, line_w=0.75, top=0.0, bottom=None):
    """terminal's phosphor texture: faint horizontal hairlines over the whole surface.

    Kept to a hairline at the register's MUTE colour — a scanline that competes with body type is
    a screen effect, and the point is a screen you can still read.
    """
    W, H = dk._slide_size(slide)
    bottom = H if bottom is None else bottom
    out, yy = [], top          # a texture by definition: type sits ON it, contrast judges it
    while yy < bottom:
        c = dk._flat(slide.shapes.add_connector(1, Inches(0), Inches(yy), Inches(W), Inches(yy)))
        c.line.color.rgb = dk._as_rgb(color)
        c.line.width = Pt(line_w)
        c.shadow.inherit = False
        out.append(c)
        yy += spacing
    return out


# ------------------------------------------------------------------------------------- the grounds

def _band_below(slide, y_top, *, side_margin=0.6, footer_gap=0.55):
    W, H = dk._slide_size(slide)
    return (side_margin, y_top, W - 2 * side_margin, H - y_top - footer_gap)


def _memphis(slide, role, index):
    """cream ground · a colour band that TITLES the page · 2–3 scattered marks in the margins."""
    W, H = dk._slide_size(slide)
    accent = _pick((dk.MAGENTA, dk.BLUE, dk.TEAL), index, 1)
    band_h = 0.85 if role in ("cover", "section") else 0.5
    color_band(slide, 0, band_h, accent)
    if role in ("cover", "section"):
        color_band(slide, band_h, 0.14, dk.DEEP)
    # Marks live OUTSIDE the band this returns — memphis scatters in the margin, it does not
    # scribble over the text. Three at most: the motif budget, applied to the furniture.
    top = band_h + 0.46
    foot = 0.62                                  # the bottom margin the marks live in
    marks = [
        lambda: zigzag(slide, W - 1.7, H - 0.5, 1.1, 0.26, dk.DEEP, teeth=4),
        lambda: tri(slide, 0.2, H - 0.56, 0.42, 0.42, dk.TEAL, line=dk.DEEP, line_w=1.5),
        lambda: _shape(slide, MSO_SHAPE.OVAL, W - 0.62, band_h + 0.06, 0.32, 0.32, fill=dk.MAGENTA),
        lambda: _shape(slide, MSO_SHAPE.PIE, 0.2, band_h + 0.04, 0.36, 0.36, fill=dk.BLUE),
    ]
    seen = set()
    for k in range(LOUD_MAX):                    # LOUD_MAX distinct marks, never the same one twice
        j = _h(index, 20 + k) % len(marks)
        while j in seen:
            j = (j + 1) % len(marks)
        seen.add(j)
        marks[j]()
    return _band_below(slide, top, side_margin=0.7, footer_gap=foot)


def _bauhaus(slide, role, index):
    """ONE oversized primitive, bleeding off an edge, in a primary — never a confetti of shapes."""
    W, H = dk._slide_size(slide)
    color = _pick((dk.MAGENTA, dk.BLUE, dk.TEAL), index, 2)
    size = 4.6 if role in ("cover", "section") else 3.4
    corner = _pick(("tr", "br", "bl"), index, 3)
    x = W - size * 0.55 if corner in ("tr", "br") else -size * 0.45
    y = -size * 0.35 if corner == "tr" else H - size * 0.6
    kind = _pick((MSO_SHAPE.OVAL, MSO_SHAPE.RECTANGLE, MSO_SHAPE.ISOSCELES_TRIANGLE), index, 4)
    _shape(slide, kind, x, y, size, size, fill=color,
           rot=(_pick((0, 180), index, 5) if kind == MSO_SHAPE.ISOSCELES_TRIANGLE else None))
    # The hairline datum bauhaus composes against. One rule, full bleed, the register's ink.
    color_band(slide, H - 0.42, 0.03, dk.DEEP)
    # Derived from where the primitive ACTUALLY landed, not from a guess at where it lands: the
    # guessed version left a 0.23in overlap that put the hero disc through the third card.
    pad = 0.3
    if corner in ("tr", "br"):
        left, right = 0.7, x - pad
    else:
        left, right = x + size + pad, W - 0.7
    top = 0.85
    bottom = (H - 0.62) if corner == "tr" else min(H - 0.62, y - pad)
    return (left, top, max(3.2, right - left), max(1.6, bottom - top))


def _risograph(slide, role, index):
    """Two flat inks, deliberately out of register, plus a screen — the whole look is the misfit."""
    W, H = dk._slide_size(slide)
    a = _pick((dk.MAGENTA, dk.BLUE), index, 6)
    b = dk.TEAL if a is not dk.TEAL else dk.MAGENTA
    bar_y = 0.0 if role in ("cover", "section") else H - 0.72
    bar_h = 1.05 if role in ("cover", "section") else 0.4
    color_band(slide, bar_y, bar_h, a)
    color_band(slide, bar_y + 0.055, bar_h, b, x=0.075, w=W - 0.15)   # the OFFSET plate
    if bar_y == 0:                                # cover: bar on top, screen in the bottom margin
        top, bottom = bar_h + 0.4, H - 0.95
        halftone(slide, W - 2.6, H - 0.82, 2.4, 0.66, a, rows=4, r_max=0.09, index=index)
    else:                                         # content: bar at the foot, screen in the top one
        top, bottom = 0.95, bar_y - 0.28
        halftone(slide, W - 2.6, 0.14, 2.4, 0.62, a, rows=4, r_max=0.09, index=index)
    return (0.65, top, W - 1.3, bottom - top)


def _terminal(slide, role, index):
    """A screen: scanlines, a prompt line for chrome, a block cursor. Everything monospace."""
    W, H = dk._slide_size(slide)
    # Faint and wide: at MUTE on 0.11in this rendered as ruled notebook paper with the body type
    # sitting on the rules. A phosphor screen is a texture you stop seeing; blend it most of the
    # way back to the ground and give it room.
    scanlines(slide, _blend(dk.GROUND, dk.MUTE, 0.3), spacing=0.2, line_w=0.6)
    dk.text(slide, 0.55, 0.3, W - 1.1, 0.3,
            [[("user@deck", 11, dk.MUTE, False, False, dk.MONO),
              (":~$ ", 11, dk.MUTE, False, False, dk.MONO),
              (("cover" if role == "cover" else "slide --{}".format(role)),
               11, dk.MAGENTA, True, False, dk.MONO)]])
    _shape(slide, MSO_SHAPE.RECTANGLE, W - 0.72, H - 0.55, 0.16, 0.26, fill=dk.MAGENTA)
    return _band_below(slide, 0.82, side_margin=0.55, footer_gap=0.8)


def _midcentury(slide, role, index):
    """Atomic-age marks on warm paper: a burst, a boomerang, one hairline datum."""
    W, H = dk._slide_size(slide)
    side = _pick(("l", "r"), index, 9)
    cx = (W - 0.95) if side == "r" else 0.95
    starburst(slide, cx, H - 0.5, 0.42, _pick((dk.MAGENTA, dk.TEAL), index, 10),
              rays=_pick((10, 12, 14), index, 12), index=index)
    boomerang(slide, (0.4 if side == "r" else W - 1.9), 0.1, 1.4, 0.62, dk.BLUE, index=index)
    color_band(slide, H - 0.98, 0.02, dk.DEEP)
    return _band_below(slide, 0.86, side_margin=0.75, footer_gap=1.05)


def _glassmorphism(slide, role, index):
    """A lit dark ground — glass only reads as glass on something with light behind it."""
    W, H = dk._slide_size(slide)
    dk.glow(slide, _frac(index, 13) * (W - 4.0) + 2.0, 1.2, 6.0, 4.4,
            _pick((dk.MAGENTA, dk.BLUE), index, 14))
    dk.glow(slide, _frac(index, 15) * (W - 3.0) + 1.5, H - 1.0, 5.0, 3.6, dk.TEAL, alpha=0.4)
    return _band_below(slide, 0.9, side_margin=0.7, footer_gap=0.7)



def _blueprint(slide, role, index):
    """A drawing sheet: faint grid, a thin cyan frame, corner ticks, a title block at the foot.

    blueprint's guard reserves the ONE coral for the focal path, so nothing here is coral — the
    ground is the sheet, and the accent is spent by the content on the thing that matters.
    """
    W, H = dk._slide_size(slide)
    step = 0.5
    x = step
    while x < W:                                   # the sheet's grid: texture, type sits on it
        c = dk._flat(slide.shapes.add_connector(1, Inches(x), Inches(0), Inches(x), Inches(H)))
        c.line.color.rgb = dk._as_rgb(_blend(dk.GROUND, dk.MAGENTA, 0.16)); c.line.width = Pt(0.5)
        c.shadow.inherit = False
        x += step
    y = step
    while y < H:
        c = dk._flat(slide.shapes.add_connector(1, Inches(0), Inches(y), Inches(W), Inches(y)))
        c.line.color.rgb = dk._as_rgb(_blend(dk.GROUND, dk.MAGENTA, 0.16)); c.line.width = Pt(0.5)
        c.shadow.inherit = False
        y += step
    # accents[0] is the CYAN line-work; accents[1] is the coral this register RESERVES for the
    # focal path, so no chrome here may touch it — apply() binds them to MAGENTA and BLUE, whose
    # names say nothing about which is which.
    line = dk.MAGENTA
    dk.catalogue_frame(slide, inset=0.3, gap=0.0, color=line, line_w=0.9)
    for cx, cy in ((0.3, 0.3), (W - 0.3, 0.3), (0.3, H - 0.3), (W - 0.3, H - 0.3)):
        _shape(slide, MSO_SHAPE.RECTANGLE, cx - 0.06, cy - 0.06, 0.12, 0.12,
               fill=None, line=line, line_w=0.9)
    # The title block a real drawing carries, bottom-right, in the register's mono chrome.
    tb_w, tb_h = 2.9, 0.46
    _shape(slide, MSO_SHAPE.RECTANGLE, W - 0.3 - tb_w, H - 0.3 - tb_h, tb_w, tb_h,
           fill=None, line=line, line_w=0.9)
    dk.text(slide, W - 0.22 - tb_w, H - 0.28 - tb_h, tb_w - 0.16, tb_h - 0.04,
            [[("SHEET ", 9, dk.MUTE, False, False, dk.MONO),
              ("{:02d}".format(index + 1), 9, line, True, False, dk.MONO),
              ("  ·  SCALE 1:1  ·  REV A", 9, dk.MUTE, False, False, dk.MONO)]])
    return (0.72, 0.85, W - 1.44, H - 0.85 - 1.0)


def _brutalist(slide, role, index):
    """A newspaper front page: a slab rule over the masthead, another under it, column hairlines."""
    W, H = dk._slide_size(slide)
    color_band(slide, 0.42, 0.09, dk.DEEP)
    color_band(slide, 0.78, 0.03, dk.DEEP)
    dk.text(slide, 0.6, 0.1, W - 1.2, 0.3,
            [[("SECTION {} ".format(index + 1), 10, dk.DEEP, True, False, dk.MONO),
              ("— " + ("EDITION" if role in ("cover", "section") else "CONTINUED"),
               10, dk.MUTE, False, False, dk.MONO)]])
    cols = 4 if role == "content" else 3
    top, bottom = 1.05, H - 0.72
    for i in range(1, cols):                       # the column grid, stated not implied
        x = 0.6 + (W - 1.2) * i / cols
        c = dk._flat(slide.shapes.add_connector(1, Inches(x), Inches(top), Inches(x), Inches(bottom)))
        c.line.color.rgb = dk._as_rgb(_blend(dk.GROUND, dk.DEEP, 0.22)); c.line.width = Pt(0.6)
        c.shadow.inherit = False
    color_band(slide, H - 0.55, 0.06, dk.DEEP)
    return (0.6, top, W - 1.2, bottom - top)


def _consulting(slide, role, index):
    """The management-consulting page: a semantic gradient rule on top, a status stamp at the foot."""
    W, H = dk._slide_size(slide)
    dk.gradient_rule(slide, 0.0, 0.0, W, dk.BLUE, dk.TEAL, h=0.11)
    dk.text(slide, 0.62, 0.26, W - 1.24, 0.26,
            [[(("EXECUTIVE SUMMARY" if role in ("cover", "section")
                else "SECTION {:02d}".format(index + 1)), 9.5, dk.MUTE, True, False)]])
    dk.hrule(slide, 0.62, H - 0.62, W - 1.24, _blend(dk.GROUND, dk.DEEP, 0.25))
    dk.text(slide, 0.62, H - 0.55, 3.0, 0.28, [[("CONFIDENTIAL", 8.5, dk.MUTE, True, False)]])
    return (0.62, 0.66, W - 1.24, H - 0.66 - 0.78)


def _dark_tech(slide, role, index):
    """A product-launch dark page: a tracked mono eyebrow, a warm->cool rule, a faint node field."""
    W, H = dk._slide_size(slide)
    for i in range(24):                            # the faint node field: texture, never content
        gx = 0.4 + (i % 8) * ((W - 0.8) / 7.0)
        gy = H - 1.5 + (i // 8) * 0.42
        _shape(slide, MSO_SHAPE.OVAL, gx, gy, 0.045, 0.045,
               fill=_blend(dk.GROUND, dk.INK if hasattr(dk, "INK") else dk.MUTE, 0.35),
               texture=True)
    dk.text(slide, 0.65, 0.34, W - 1.3, 0.28,
            [[(">_ ", 10, dk.TEAL, True, False, dk.MONO),
              (("K E Y N O T E" if role in ("cover", "section") else
                "M O D U L E   {:02d}".format(index + 1)), 10, dk.MUTE, True, False, dk.MONO)]])
    dk.gradient_rule(slide, 0.65, 0.68, 3.4, dk.MAGENTA, dk.BLUE, h=0.045)
    return (0.65, 0.95, W - 1.3, H - 0.95 - 1.7)


def _eastern_traditional(slide, role, index):
    """Warm paper, a 传统色 swatch column in the margin, one small seal — the colours tell the story."""
    W, H = dk._slide_size(slide)
    _note(0.34, 0.9, 0.2, 3 * 0.46)
    for i, col in enumerate((dk.MAGENTA, dk.BLUE, dk.TEAL)):   # the named hues, shown as themselves
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.34, 0.9 + i * 0.46, 0.2, 0.36, fill=col)
    dk.hrule(slide, 0.72, 0.72, W - 1.44, _blend(dk.GROUND, dk.DEEP, 0.3))
    _note(W - 0.82, H - 0.86, 0.42, 0.42)
    dk.seal(slide, W - 0.82, H - 0.86, 0.42,
            dk.cjk_numeral(index + 1) if hasattr(dk, "cjk_numeral") else "印",
            fill=dk.MAGENTA, shape="square", rounded=False)
    return (0.78, 0.92, W - 1.56, H - 0.92 - 1.0)


def _editorial_paper(slide, role, index):
    """A magazine spread's chrome: caps eyebrow, hairline, a big ghost folio. Colour is the photo's."""
    W, H = dk._slide_size(slide)
    dk.text(slide, 0.7, 0.42, W - 1.4, 0.26,
            [[("FEATURE", 9.5, dk.MAGENTA, True, False),
              ("   ·   PAGE {:02d}".format(index + 11), 9.5, dk.MUTE, False, False)]])
    dk.hrule(slide, 0.7, 0.74, W - 1.4, _blend(dk.GROUND, dk.DEEP, 0.35))
    # The folio lives in the BOTTOM MARGIN, under the band and over the closing rule. It used to
    # sit at H-1.35 at 1.1in tall: inside the content rect and clipped by the page edge, which the
    # render showed and nothing measured until `_note` existed.
    _note(W - 1.35, H - 0.78, 1.05, 0.5)
    dk.ghost_numeral(slide, W - 1.35, H - 0.78, 1.05, 0.5, "{:02d}".format(index + 11),
                     color=dk.MAGENTA, opacity=0.16)
    dk.hrule(slide, 0.7, H - 0.2, W - 1.4, _blend(dk.GROUND, dk.DEEP, 0.35))
    return (0.7, 0.92, W - 1.4, H - 0.92 - 0.9)


def _editorial_report(slide, role, index):
    """Dark gravitas: a roman-numeral section mark, one red hairline, a source line's worth of room."""
    W, H = dk._slide_size(slide)
    roman = ("I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X", "XI", "XII")
    dk.text(slide, 0.72, 0.4, 2.0, 0.3,
            [[(roman[index % len(roman)], 11, dk.MAGENTA, True, False),
              ("   ", 11, dk.MUTE, False, False),
              (("REPORT" if role in ("cover", "section") else "ANALYSIS"),
               11, dk.MUTE, True, False)]])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 0.74, 0.85, 0.035, fill=dk.MAGENTA)
    dk.hrule(slide, 0.72, H - 0.66, W - 1.44, _blend(dk.GROUND, dk.INK if hasattr(dk, "INK") else dk.MUTE, 0.25))
    return (0.72, 0.95, W - 1.44, H - 0.95 - 0.85)


def _ink_wash(slide, role, index):
    """留白 first: one hairline, a CJK section numeral in the margin, one chop. Nothing filled.

    ink_wash's guard is the only one that forbids something by ABSENCE — no filled corners, no
    dense grids. So this ground is the quietest in the file on purpose, and the temptation it
    resists (a wash gradient) is also on its FORBIDS list.
    """
    W, H = dk._slide_size(slide)
    num = dk.cjk_numeral(index + 1) if hasattr(dk, "cjk_numeral") else str(index + 1)
    _note(0.5, 0.85, 0.5, 0.6)
    dk.text(slide, 0.5, 0.85, 0.5, 1.2,
            [[(str(num), 22, _blend(dk.GROUND, dk.DEEP, 0.45), False, False)]])
    dk.hrule(slide, 1.15, 0.78, W - 2.0, _blend(dk.GROUND, dk.DEEP, 0.28))
    _note(W - 0.78, H - 0.82, 0.36, 0.36)
    dk.seal(slide, W - 0.78, H - 0.82, 0.36, "印", fill=dk.MAGENTA, shape="square", rounded=False)
    return (1.15, 1.0, W - 2.3, H - 1.0 - 0.95)


def _luxury_dark(slide, role, index):
    """A masthead and an issue line in champagne hairlines, and then a great deal of nothing."""
    W, H = dk._slide_size(slide)
    dk.hrule(slide, 0.85, 0.62, W - 1.7, dk.MAGENTA, weight=0.008)
    dk.text(slide, 0.85, 0.28, W - 1.7, 0.28,
            [[("M A I S O N", 10, dk.MAGENTA, True, False),
              ("        ISSUE {:02d}".format(index + 1), 10, dk.MUTE, False, False)]])
    dk.hrule(slide, 0.85, H - 0.72, W - 1.7, _blend(dk.GROUND, dk.MAGENTA, 0.45), weight=0.006)
    return (0.95, 0.95, W - 1.9, H - 0.95 - 1.0)


def _museum_memorial(slide, role, index):
    """The exhibition catalogue: a double-line frame inset from the edges, a brass year badge."""
    W, H = dk._slide_size(slide)
    dk.catalogue_frame(slide, inset=0.34, gap=0.07, color=dk.MAGENTA, line_w=0.8)
    dk.text(slide, 0.75, 0.5, W - 1.5, 0.28,
            [[("PLATE {:02d}".format(index + 1), 9.5, dk.MAGENTA, True, False),
              ("   ·   CATALOGUE", 9.5, dk.MUTE, False, False)]])
    dk.hrule(slide, 0.75, 0.82, W - 1.5, _blend(dk.GROUND, dk.MAGENTA, 0.4))
    return (0.85, 1.0, W - 1.7, H - 1.0 - 0.85)


def _swiss(slide, role, index):
    """A grid and a ghost folio. Swiss is what you DON'T draw — this ground is nearly all restraint.

    Its guard forbids rounded cards, gradients and soft shadows, and its `card()` here draws no box
    at all: a swiss column is a hairline and a measure, and boxing it is the commonest way a deck
    claims swiss while looking like every other card deck.
    """
    W, H = dk._slide_size(slide)
    top, bottom = 0.95, H - 0.75
    for i in range(1, 4):
        x = 0.75 + (W - 1.5) * i / 4.0
        c = dk._flat(slide.shapes.add_connector(1, Inches(x), Inches(top), Inches(x), Inches(bottom)))
        c.line.color.rgb = dk._as_rgb(_blend(dk.GROUND, dk.DEEP, 0.12))
        c.line.width = Pt(0.5)
        c.shadow.inherit = False
    dk.hrule(slide, 0.75, 0.78, W - 1.5, dk.DEEP)
    _note(W - 1.45, H - 0.7, 1.05, 0.52)     # the folio goes under the band, not through it
    dk.ghost_numeral(slide, W - 1.45, H - 0.7, 1.05, 0.52, "{:02d}".format(index + 1),
                     color=dk.DEEP, opacity=0.1)
    return (0.75, top, W - 1.5, bottom - top)


def _synthwave(slide, role, index):
    """The receding grid horizon and a banded sun — dosed by ROLE, which is this register's guard.

    synthwave's guard says the LOUD full horizon belongs on a cover or a divider, and that a faint
    low-edge register may repeat every slide because that is system rather than stamping. So the
    role is not decoration here: it decides which of the two this page gets.
    """
    import math
    W, H = dk._slide_size(slide)
    loud = role in ("cover", "section")
    # The horizon sits LOW even when loud: on a 10x5.63in canvas a horizon at 0.42H left a
    # 1.01in content band, which is not a page. The sun rises BEHIND the title (texture), which is
    # what the register does anyway.
    hz = H * (0.66 if loud else 0.86)
    if loud:
        for i in range(7):                         # the banded sunset, flat bands (no gradient)
            t = i / 6.0
            col = _blend(dk.MAGENTA, dk.BLUE, t)
            _shape(slide, MSO_SHAPE.RECTANGLE, W / 2 - 1.9, hz - 1.62 + i * 0.22, 3.8, 0.17,
                   fill=col, texture=True)
    for i in range(-9, 10):                        # converging lines to the vanishing point
        x_bottom = W / 2 + i * (W / 6.0)
        c = dk._flat(slide.shapes.add_connector(
            1, Inches(W / 2), Inches(hz), Inches(x_bottom), Inches(H)))
        c.line.color.rgb = dk._as_rgb(_blend(dk.GROUND, dk.TEAL, 0.55 if loud else 0.3))
        c.line.width = Pt(0.9 if loud else 0.6)
        c.shadow.inherit = False
    rows = 6 if loud else 3
    for i in range(rows):                          # ...and the receding horizontals
        yy = hz + (H - hz) * ((i + 1) / rows) ** 2.1
        c = dk._flat(slide.shapes.add_connector(1, Inches(0), Inches(yy), Inches(W), Inches(yy)))
        c.line.color.rgb = dk._as_rgb(_blend(dk.GROUND, dk.TEAL, 0.5 if loud else 0.28))
        c.line.width = Pt(0.9 if loud else 0.6)
        c.shadow.inherit = False
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, hz - 0.012, W, 0.024, fill=dk.MAGENTA, texture=True)
    return (0.8, 0.9, W - 1.6, (hz - 0.32 if loud else H - 1.35) - 0.9)


GROUNDS = {
    "blueprint": _blueprint,
    "brutalist": _brutalist,
    "consulting": _consulting,
    "dark_tech": _dark_tech,
    "eastern_traditional": _eastern_traditional,
    "editorial_paper": _editorial_paper,
    "editorial_report": _editorial_report,
    "ink_wash": _ink_wash,
    "luxury_dark": _luxury_dark,
    "museum_memorial": _museum_memorial,
    "swiss": _swiss,
    "synthwave": _synthwave,
    "memphis": _memphis,
    "bauhaus": _bauhaus,
    "risograph": _risograph,
    "terminal": _terminal,
    "midcentury": _midcentury,
    "glassmorphism": _glassmorphism,
}


# --------------------------------------------------------------------------------------- the cards

def _card_memphis(slide, x, y, w, h, label=None):
    body = dk.box(slide, x, y, w, h, fill=dk.TINT, line=dk.DEEP, line_w=1.6, round=True, r=0.12)
    head = dk.box(slide, x, y, w, 0.42, fill=dk.MAGENTA, corners="top", r=0.12)
    if label:
        dk.text(slide, x + 0.18, y + 0.05, w - 0.36, 0.32,
                [[(label, 12, dk.on(dk.MAGENTA), True, False)]])   # auto-contrast on the band
    return body, head


def _card_bauhaus(slide, x, y, w, h, label=None):
    return dk.box(slide, x, y, w, h, fill=dk.TINT, line=dk.DEEP, line_w=2.2), None


def _card_risograph(slide, x, y, w, h, label=None):
    return dk.offset_shadow(slide, x, y, w, h, dk.TINT, shadow=dk.MAGENTA,
                            line=dk.DEEP, line_w=1.6, round=False), None


def _card_terminal(slide, x, y, w, h, label=None):
    return dk.box(slide, x, y, w, h, fill=None, line=dk.MUTE, line_w=1.0), None


def _card_midcentury(slide, x, y, w, h, label=None):
    return dk.box(slide, x, y, w, h, fill=dk.TINT, line=dk.DEEP, line_w=1.0, round=True, r=0.06), None


def _card_glass(slide, x, y, w, h, label=None):
    return dk.glass_card(slide, x, y, w, h, dk.BLUE), None



def _card_blueprint(slide, x, y, w, h, label=None):
    """A panel on the sheet: thin cyan line-work, square (blueprint forbids rounded)."""
    body = dk.box(slide, x, y, w, h, fill=_blend(dk.GROUND, dk.MAGENTA, 0.10),
                  line=dk.MAGENTA, line_w=0.9)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.16, 0.03, fill=dk.MAGENTA)
    return body, None


def _card_brutalist(slide, x, y, w, h, label=None):
    """A slab: no fill, a rule so heavy it IS the design. Square by the register's own guard."""
    body = dk.box(slide, x, y, w, h, fill=None, line=dk.DEEP, line_w=3.0)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.12, fill=dk.DEEP)
    return body, None


def _card_consulting(slide, x, y, w, h, label=None):
    """A scorecard tile: white, a navy top keyline, the semantic colour left to the content."""
    body = dk.box(slide, x, y, w, h, fill=dk.TINT, line=_blend(dk.GROUND, dk.DEEP, 0.2), line_w=0.8,
                  round=True, r=0.05)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.06, fill=dk.BLUE)
    return body, None


def _card_dark_tech(slide, x, y, w, h, label=None):
    """An insight panel with the accent LEFT BAR this register uses instead of a header band."""
    body = dk.box(slide, x, y, w, h, fill=_blend(dk.GROUND, dk.TINT, 0.35),
                  line=_blend(dk.GROUND, dk.TEAL, 0.35), line_w=0.9, round=True, r=0.07)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.055, h, fill=dk.TEAL)
    return body, None


def _card_eastern(slide, x, y, w, h, label=None):
    """Paper panel, hairline, one ochre rule at the head — the hue is the message here."""
    body = dk.box(slide, x, y, w, h, fill=_blend(dk.GROUND, dk.DEEP, 0.05),
                  line=_blend(dk.GROUND, dk.DEEP, 0.28), line_w=0.7)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w * 0.34, 0.04, fill=dk.BLUE)
    return body, None


def _card_editorial_paper(slide, x, y, w, h, label=None):
    """No box: a hairline over the measure. The chrome stays neutral so a photo can carry colour."""
    dk.hrule(slide, x, y, w, _blend(dk.GROUND, dk.DEEP, 0.4))
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.5, 0.022, fill=dk.MAGENTA)
    return None, None


def _card_editorial_report(slide, x, y, w, h, label=None):
    """A dark panel with one red tick — this register spends its red once per slide."""
    body = dk.box(slide, x, y, w, h, fill=_blend(dk.GROUND, dk.TINT, 0.22),
                  line=_blend(dk.GROUND, dk.MUTE, 0.3), line_w=0.7)
    return body, None


def _card_ink_wash(slide, x, y, w, h, label=None):
    """A hairline frame and nothing else — a filled card would be the 留白 its guard protects."""
    body = dk.box(slide, x, y, w, h, fill=None, line=_blend(dk.GROUND, dk.DEEP, 0.3), line_w=0.7)
    return body, None


def _card_luxury(slide, x, y, w, h, label=None):
    """Champagne hairline, no fill. The luxury is the space around it."""
    body = dk.box(slide, x, y, w, h, fill=None, line=_blend(dk.GROUND, dk.MAGENTA, 0.55),
                  line_w=0.7)
    return body, None


def _card_museum(slide, x, y, w, h, label=None):
    """A plate: navy panel inside a brass hairline, the way a catalogue sets a reproduction."""
    body = dk.box(slide, x, y, w, h, fill=_blend(dk.GROUND, dk.TINT, 0.2),
                  line=_blend(dk.GROUND, dk.MAGENTA, 0.5), line_w=0.8)
    return body, None


def _card_swiss(slide, x, y, w, h, label=None):
    """NO BOX. A swiss column is a hairline and a measure.

    This is the one card in the file that draws almost nothing, and it is the most important one:
    boxing a column is exactly how a deck claims swiss and looks like every other card deck. The
    register's own guard forbids the rounded card; the discipline it actually asks for is to not
    draw the card at all.
    """
    dk.hrule(slide, x, y, w, dk.DEEP)
    return None, None


def _card_synthwave(slide, x, y, w, h, label=None):
    """A SOLID panel with a neon rim — its guard: keep text on a solid panel where it crosses the grid."""
    body = dk.box(slide, x, y, w, h, fill=_blend(dk.GROUND, (0, 0, 0), 0.35),
                  line=dk.BLUE, line_w=1.2, round=True, r=0.08)
    return body, None


CARDS = {
    "blueprint": _card_blueprint,
    "brutalist": _card_brutalist,
    "consulting": _card_consulting,
    "dark_tech": _card_dark_tech,
    "eastern_traditional": _card_eastern,
    "editorial_paper": _card_editorial_paper,
    "editorial_report": _card_editorial_report,
    "ink_wash": _card_ink_wash,
    "luxury_dark": _card_luxury,
    "museum_memorial": _card_museum,
    "swiss": _card_swiss,
    "synthwave": _card_synthwave,
    "memphis": _card_memphis,
    "bauhaus": _card_bauhaus,
    "risograph": _card_risograph,
    "terminal": _card_terminal,
    "midcentury": _card_midcentury,
    "glassmorphism": _card_glass,
}


# ------------------------------------------------------------------------------------ the public API

def has(register):
    """Is there a surface kit for this register? All 18 have one; anything else answers False."""
    return str(register or "").strip().lower() in GROUNDS


def ground(slide, register, *, role="content", index=0):
    """Paint the register's own surface and RETURN the content rect `(x, y, w, h)` left over.

    Call it as the FIRST thing on a slide, before any content: everything it draws is furniture,
    and furniture goes behind. The rect it returns is where content may go — nothing this function
    painted overlaps it, which is the contract that lets a builder stay hands-off about placement.

    `role` is the page's job (`cover` / `section` / `content` / `closer`): a cover gets the loud
    version of the register, a content page the quiet one, because a register applied at full
    volume on every page is the "flashy but unreadable" failure this skill keeps naming.
    """
    reg = str(register or "").strip().lower()
    if reg not in GROUNDS:
        raise KeyError(
            "no surface kit for register {!r} — kits exist for {}. Do NOT silently fall back to a "
            "plain page: a caller that asked for a register's surface and got a blank one would "
            "ship the colourway-only deck this module was written to end.".format(
                register, ", ".join(sorted(GROUNDS))))
    if dk.GROUND is None:
        raise RuntimeError(
            "the palette is not set yet — call `presets.apply({!r})` (or `deckkit.set_ground(...)` "
            "for a bespoke look) BEFORE painting its surface. Every kit reads the register's own "
            "GROUND/DEEP/accents to build from; without them the first blend fails deep inside a "
            "builder with a TypeError that says nothing about the real mistake.".format(reg))
    del _MARKS[:]
    band = GROUNDS[reg](slide, str(role or "content").lower(), int(index))
    bx, by, bw, bh = band
    clash = [m for m in _MARKS
             if min(bx + bw, m[2]) - max(bx, m[0]) > 0.02
             and min(by + bh, m[3]) - max(by, m[1]) > 0.02]
    if clash:
        raise AssertionError(
            "{}'s ground painted {} loud mark(s) INTO the content rect it returned "
            "({}) — a caller that trusts the rect would set type on top of them. Marks belong in "
            "the margins; a full-bleed ground passes texture=True. Offenders: {}".format(
                reg, len(clash), tuple(round(v, 2) for v in band),
                [tuple(round(v, 2) for v in c) for c in clash]))
    return band


def card(slide, register, x, y, w, h, *, label=None):
    """The register's CARD FORM at this rect. Same call on every register, a different object.

    This is the half that makes the pages differ in FORM rather than in colour: a memphis banded
    card, a bauhaus hard square with a heavy keyline, a riso sticker with a crisp offset plate, a
    frosted glass panel, a terminal output block with no fill at all. Returns `(body, header)`;
    `header` is None for the registers whose card has no band.
    """
    reg = str(register or "").strip().lower()
    if reg not in CARDS:
        raise KeyError("no card form for register {!r} — kits exist for {}".format(
            register, ", ".join(sorted(CARDS))))
    return CARDS[reg](slide, x, y, w, h, label)


def registers():
    """The registers that have a kit, sorted — for a gallery, a doc, or a coverage report."""
    return sorted(GROUNDS)


# ------------------------------------------------------------------------------------------ sample

def sample(out_path):
    """One page per kit, same content, so the difference is the REGISTER and nothing else."""
    import presets
    prs = None
    for i, reg in enumerate(registers()):
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            presets.apply(reg)
        if prs is None:
            prs = dk.blank_deck()
        s = dk.add_slide(prs)
        role = "cover" if i == 0 else "content"
        bx, by, bw, bh = ground(s, reg, role=role, index=i)
        dk.text(s, bx, by, bw, 0.7,
                [[(reg.replace("_", " ").upper(), 30, dk.DEEP, True, False)]])
        cw = (bw - 0.6) / 3.0
        for c in range(3):
            cx = bx + c * (cw + 0.3)
            card(s, reg, cx, by + 0.95, cw, bh - 1.15, label="CARD {}".format(c + 1))
            dk.text(s, cx + 0.22, by + 1.55, cw - 0.44, bh - 1.9,
                    [[("The same three cards, the same three sentences, on every one of these "
                       "pages.", 12, dk.DEEP, False, False)]])
    prs.save(str(out_path))
    return out_path


def _selftest():
    import presets
    import check_register_guard as guard
    import tempfile
    ok, bad = [], []
    tmp = Path(tempfile.mkdtemp(prefix="regsurf-"))

    snap = {k: getattr(dk, k) for k in dir(dk) if k.isupper()}
    try:
        for reg in registers():
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                presets.apply(reg)
            prs = dk.blank_deck()
            for i, role in enumerate(("cover", "content", "section")):
                s = dk.add_slide(prs)
                bx, by, bw, bh = ground(s, reg, role=role, index=i)
                W, H = dk._slide_size(s)
                if not (bx >= 0 and by >= 0 and bw > 3.0 and bh > 1.5
                        and bx + bw <= W + 1e-6 and by + bh <= H + 1e-6):
                    bad.append("{} {}: content rect {} is off-canvas or too small".format(
                        reg, role, (bx, by, bw, bh)))
                head = 0.6
                cw = min(3.2, bw)
                ch = max(0.7, bh - head)
                card(s, reg, bx, by + head, cw, ch, label="X")
                if ch > 0.95:
                    dk.text(s, bx + 0.2, by + head + 0.3, max(0.5, cw - 0.4), ch - 0.5,
                            [[("a card carries content", 12, dk.DEEP, False, False)]])
            p = tmp / "{}.pptx".format(reg)
            prs.save(str(p))
            viol, _facts = guard.check(p, register=reg)
            (ok if not viol else bad).append(
                "`{}`'s own surface kit obeys `{}`'s prohibitions".format(reg, reg)
                if not viol else "{} violates its own register: {}".format(
                    reg, [c for c, _m in viol]))
    finally:
        for k, v in snap.items():
            setattr(dk, k, v)

    # Determinism: the same index must give the same page, or every diff-based check is noise.
    def _fingerprint():
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            presets.apply("memphis")
        prs = dk.blank_deck()
        s = dk.add_slide(prs)
        ground(s, "memphis", role="content", index=7)
        return [(sh.shape_type, sh.left, sh.top, sh.width, sh.height) for sh in s.shapes]
    snap2 = {k: getattr(dk, k) for k in dir(dk) if k.isupper()}
    try:
        a, b = _fingerprint(), _fingerprint()
    finally:
        for k, v in snap2.items():
            setattr(dk, k, v)
    (ok if a == b else bad).append(
        "the same page index builds the same page twice — placement varies by INDEX, never by a "
        "random number, so a deck built twice stays byte-identical"
        if a == b else "non-deterministic ground")

    try:
        ground(dk.add_slide(dk.blank_deck()), "no_such_register")
        bad.append("an unknown register was silently given a blank surface")
    except KeyError:
        ok.append("a name that is NOT a register raises instead of quietly returning a plain page "
                  "— the silent fallback would ship exactly the colourway-only deck this ends")

    try:
        prs = dk.blank_deck()
        sl = dk.add_slide(prs)
        _g = dk.GROUND
        dk.GROUND = None
        try:
            ground(sl, "swiss")
            bad.append("painting a surface with no palette set did not raise")
        except RuntimeError as exc:
            (ok if "presets.apply" in str(exc) else bad).append(
                "painting a surface before the palette is set says SO — it used to die on a blend "
                "deep inside a builder with a TypeError that named nothing the caller could fix"
                if "presets.apply" in str(exc) else str(exc)[:90])
        finally:
            dk.GROUND = _g
    except Exception as exc:
        bad.append("palette-not-set probe blew up: {}".format(exc))

    for line in ok:
        print("  ok   " + line)
    for line in bad:
        print("  FAIL " + line)
    print("\n{} passed, {} failed".format(len(ok), len(bad)))
    return 1 if bad else 0


def main(argv=None):
    from _console import safe_stdio
    safe_stdio()
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--selftest", action="store_true")
    ap.add_argument("--sample", metavar="OUT.pptx")
    ap.add_argument("--list", action="store_true")
    a = ap.parse_args(argv)
    if a.selftest:
        return _selftest()
    if a.list:
        import presets
        for r in registers():
            print("  {:16s} {}".format(r, str(presets.PRESETS[r].get("surface"))[:88]))
        missing = [r for r in sorted(presets.PRESETS) if r not in GROUNDS]
        print("\n  {} of {} registers have a surface kit; still prose-only: {}".format(
            len(GROUNDS), len(presets.PRESETS), ", ".join(missing)))
        return 0
    if a.sample:
        print("wrote", sample(a.sample))
        return 0
    ap.print_help()
    return 0


if __name__ == "__main__":
    sys.exit(main())

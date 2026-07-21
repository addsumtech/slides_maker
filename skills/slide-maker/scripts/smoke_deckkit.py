#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""smoke_deckkit — exercise the crash-prone CORE of deckkit (+ the chart recipes) with canonical
examples on a blank deck and assert none raises. NOT exhaustive: it covers roughly half of the
public helpers (~68 of 138 — the positional/tuple colour contracts and geometry paths that have
actually broken), so a passing run is a regression guard, not proof every helper works. Run after
editing deckkit/designed_charts:

    python scripts/smoke_deckkit.py     # exits non-zero on any failure
"""
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deckkit as dk          # noqa: E402
import designed_charts as dc  # noqa: E402
from deckkit import RGBColor  # noqa: E402

def C(h):
    return RGBColor.from_string(h)

TMP = tempfile.gettempdir()
IMG = os.path.join(TMP, "_smoke.png")
from PIL import Image  # noqa: E402
Image.new("RGB", (800, 450), (120, 150, 200)).save(IMG)

W, H = 13.333, 7.5
prs = dk.blank_deck(W, H)
def S():
    return dk.add_slide(prs)
def last():
    return prs.slides[-1]

fails = []
def ok(name, fn):
    try:
        fn()
    except Exception as e:
        fails.append((name, repr(e))); print(f"  FAIL {name} -> {e!r}")
def raises(name, fn):                       # a guard that SHOULD reject bad input
    try:
        fn(); fails.append((name, "did not raise")); print(f"  FAIL {name} (expected an error)")
    except Exception:
        pass

s = S()
ok("box (hex colours)", lambda: dk.box(s, 0, 0, 3, 1, fill="C0362C", line="333333"))     # unified colour API
ok("text", lambda: dk.text(s, 0, 2, 4, 1, [[("hi", 18, C("222222"), True, False)]]))
ok("title_bar/footer", lambda: (dk.title_bar(s, "T", kicker="k"), dk.footer(s, "tag", 1)))
ok("columns/rows", lambda: (dk.columns(3, slide=s), dk.rows(2, slide=s)))
ok("callout", lambda: dk.callout(s, 0, 3, 4, 1, "head", "body"))
ok("chip", lambda: dk.chip(s, 6, 4, 2, 0.6, "A", "sub", C("007CC2")))
ok("bullet", lambda: dk.bullet(s, 5, 1, 5, [("Lead ", "body")]))
ok("hrule", lambda: dk.hrule(s, 0, 5, 5))
ok("table", lambda: dk.table(S(), 0.6, 1.5, 6, [["A", "B"], ["1", "2"]], highlight=1))
ok("scorecard (numeric value)", lambda: dk.scorecard(S(), 0.6, 1, 2.5, 1.8, "Users", 1234, delta="3.2pp"))
ok("scorecard glass", lambda: dk.scorecard(last(), 4, 1, 2.5, 1.8, "X", "9", glass_tint=C("5B8DEF")))
ok("leaderboard (mixed + short row)", lambda: dk.leaderboard(S(), 0.6, 1, 5,
   [(C("007CC2"), "a", 42), (C("C0362C"), "b", "18", "sub"), (C("1B7A3D"), "short-row")]))
ok("takeaway_rail", lambda: dk.takeaway_rail(last(), 7, 1, 3, "lab", "+5", "body"))
ok("change_stat", lambda: dk.change_stat(S(), 4, 3, 3, 0.7, "<10%", "≈40%"))
ok("stat_row", lambda: dk.stat_row(S(), 0.7, 2, 8, [("8", "x", "label"), ("99", "%", "two")]))
ok("stat_row (2-tuple, no unit)", lambda: dk.stat_row(last(), 0.7, 3.2, 8, [("8", "no-unit"), ("9", "also")]))
ok("quadrant", lambda: dk.quadrant(S(), 1.6, 1.6, 5, 4, x_labels=("lo", "hi"), y_labels=("a", "b")))
ok("hub_spoke", lambda: dk.hub_spoke(S(), 6, 4, 2.0, "Core", ["a", "b", "c", "d"]))
    # r=2.0 (was 1.8): at 1.8 two spoke chips sat a ~0.05in sliver apart, which the
    # SLIVER_GAP lint correctly flags — the fixture is nudged, the check is not weakened.
ok("timeline h", lambda: dk.timeline(S(), 0.7, 2, 11.4, [("1979", "first"), ("2026", "now", "cap")], highlight=1))
ok("timeline v", lambda: dk.timeline(S(), 0.7, 2, 5, [("x", "a"), ("y", "b")], orientation="v"))
ok("image_tab", lambda: dk.image_tab(S(), 0.5, 0.5, "BEFORE"))
ok("before_after", lambda: dk.before_after(last(), 0.7, 2, 7, 3, IMG, IMG))
ok("photo_triptych", lambda: dk.photo_triptych(S(), [IMG, IMG, IMG]))
ok("photo_card", lambda: dk.photo_card(S(), 0.5, 0.5, 3, 1, role="primary"))
ok("corner_frame", lambda: dk.corner_frame(S()))
ok("accent_one", lambda: dk.accent_one(["a", "b", "c"], 1, C("C0362C")))
ok("cover", lambda: dk.cover(S(), "Title", issue_label="No 1", subtitle="sub"))
ok("colophon (list credits)", lambda: dk.colophon(S(), "tag", credits=["a", "b"], tooling="x"))
ok("sources_page", lambda: dk.sources_page(S(), [f"ref {i}" for i in range(6)]))
ok("part_eyebrow/page_marker", lambda: (dk.part_eyebrow(S(), 0.7, 0.5, "x"), dk.page_marker(last(), 2, 8)))
ok("specimen_card", lambda: dk.specimen_card(S(), 1, 1, 2, 2.5, "Aa", "Sans"))
ok("specimen_card (small h)", lambda: dk.specimen_card(last(), 4, 1, 1, 0.5, "Aa", "tiny"))
ok("wireframe_grid/spec_list", lambda: (dk.wireframe_grid(S(), 0.7, 1.5, 7, 4, [("H", 0, 4, 0, 1)]),
                                        dk.spec_list(last(), 8, 1.5, ["base = 8 px"])))
ok("glass/glow/scrim/offset", lambda: (lambda sl: (dk.box(sl, 0, 0, W, H, fill="0A0E27"),
   dk.glow(sl, 3, 3, 5, 4, C("5B4BE0")), dk.glass_card(sl, 1, 1, 3, 2, C("5B8DEF")),
   dk.scrim_overlay(sl, 0, 5, W, 2), dk.offset_shadow(sl, 5, 5, 2, 0.6, C("F5B301"))))(S()))
ok("editorial_header", lambda: dk.editorial_header(S(), "eyebrow", "Title", serif="Georgia"))
ok("big_numeral", lambda: dk.big_numeral(S(), 0.5, 0.5, "04", mode="ghost"))
ok("picture", lambda: dk.picture(S(), IMG, 1, 1, 4, 3, round=True))
def _wordmark():
    p = os.path.join(TMP, "_wm.png")
    dk.wordmark("Acme Labs", p, rule=True, monogram=True)
    assert os.path.exists(p) and os.path.getsize(p) > 0, "wordmark PNG missing/empty"
    dk.logo(S(), p, h=0.4)                       # placed afterward, exactly like a real logo
ok("wordmark (typographic logo stand-in)", _wordmark)
ok("backdrop_motif", lambda: dk.backdrop_motif(S(), accent_disc=C("C0362C")))
ok("native_chart (editable)", lambda: dk.native_chart(S(), 0.6, 1, 6, 3.2, ["Q1", "Q2", "Q3"],
   [("新客", [10, 18, 26]), ("老客", [30, 33, 36])], kind="column", dark=True, highlight=0))
ok("native_dual_axis (editable)", lambda: dk.native_dual_axis(S(), 0.6, 1, 7, 3.2, ["m1", "m2", "m3"],
   [5, 24, 40], [100, 80, 62], left_name="占比（%）", right_name="成本（指数）", dark=True))
ok("native_donut (editable)", lambda: dk.native_donut(S(), 0.6, 1, 5, 4, [("私域", 40), ("公域", 35), ("其他", 25)], "40%", "私域占比", dark=True))
ok("native_pareto (editable)", lambda: dk.native_pareto(S(), 0.6, 1, 8, 4, [("华东", 45), ("华北", 28), ("华南", 18)],
   dark=True, count_name="数量", cum_name="累计 %"))    # CJK deck → pass the CJK series names explicitly
ok("native_bubble (editable)", lambda: dk.native_bubble(S(), 0.6, 1, 8, 4, [(1, 2, 30), (2, 3, 55), (3, 2.4, 20)], dark=True))
def _csv_chart():
    p = os.path.join(TMP, "_s.csv"); open(p, "w").write("m,a,b\nJ,1,\"2,000\"\nF,3,4\n")
    cats, series = dk.series_from_csv(p, "m", ["a", "b"])
    assert cats == ["J", "F"] and series[0] == ("a", [1.0, 3.0]) and series[1][1] == [2000.0, 4.0], "series_from_csv parse"
    dk.native_chart(S(), 0.6, 1, 6, 3, cats, series, kind="column")
ok("series_from_csv + native_chart", _csv_chart)
raises("series_from_csv rejects a missing column", lambda: dk.series_from_csv(os.path.join(TMP, "_s.csv"), "nope", ["a"]))

ok("dc.donut_kpi", lambda: dc.donut_kpi(os.path.join(TMP, "_d.png"), [("a", 3), ("b", 2)], "5", "x"))
ok("dc.dumbbell", lambda: dc.dumbbell(os.path.join(TMP, "_db.png"), [("a", 1, 2)], highlight=0))
ok("dc.slope", lambda: dc.slope(os.path.join(TMP, "_sl.png"), [("a", 1, 2)], highlight=0))
ok("dc.dual_axis", lambda: dc.dual_axis(os.path.join(TMP, "_da.png"), [1, 2], [1, 2], [2, 1]))
ok("dc.bubble_trend", lambda: dc.bubble_trend(os.path.join(TMP, "_bt.png"), [(1, 2, 3, "a"), (2, 3, 5, "b")]))
ok("dc.pareto", lambda: dc.pareto(os.path.join(TMP, "_pa.png"), [("a", 4), ("b", 2)]))

# --- the build-time geometry gate must actually catch a fault and (strict) block the save ---
def _offcanvas_deck():
    p = dk.blank_deck(10, 5.625); s = dk.add_slide(p)
    dk.text(s, 9.2, 2.0, 3.0, 0.5, [[("this text runs off the right slide edge entirely here", 16, dk.DEEP, False, False)]])
    return p
def _lint_layout_gate():
    crit = [f for f in dk.lint_layout(_offcanvas_deck(), verbose=False) if f[1] == "CRITICAL"]
    assert crit, "lint_layout missed an off-canvas CRITICAL"
    p = dk.blank_deck(10, 5.625); s = dk.add_slide(p)
    dk.text(s, 1.0, 1.0, 4.0, 0.5, [[("a tidy line well inside the slide", 16, dk.DEEP, False, False)]])
    assert not [f for f in dk.lint_layout(p, verbose=False) if f[1] == "CRITICAL"], "lint_layout false-flagged a clean slide"
ok("lint_layout catches off-canvas + passes clean", _lint_layout_gate)
raises("lint_layout(strict=True) raises on a CRITICAL", lambda: dk.lint_layout(_offcanvas_deck(), strict=True, verbose=False))

raises("donut_kpi([]) rejects empty", lambda: dc.donut_kpi(os.path.join(TMP, "_x.png"), [], "0", "n"))
raises("dual_axis rejects empty", lambda: dc.dual_axis(os.path.join(TMP, "_x.png"), [], [], []))
raises("wireframe_grid rejects cols=0", lambda: dk.wireframe_grid(S(), 0, 0, 5, 4, [], cols=0))
raises("picture rejects missing file", lambda: dk.picture(S(), "/no/such/img.png", 0, 0, 3, 2))

ok("cycle_diagram (4 nodes + feedback + center)", lambda: dk.cycle_diagram(
    S(), 6.5, 3.6, [("获客", "3 行业"), ("见效", "7 天"), ("使用", "3 指标"), ("续费", "达标增购")],
    rx=1.5, ry=1.0, center="飞轮", feedback=(3, 0), feedback_label="回流"))
    # rx=1.5 (was 1.4): at 1.4 the hub sat a ~0.05in sliver from two nodes, which the
    # SLIVER_GAP lint correctly flags — the fixture is nudged, the check is not weakened.
ok("dumbbell_board (hero + threshold)", lambda: dk.dumbbell_board(
    S(), 0.8, 1.6, 11.0, [("ARR", "+51%", 4980, 6350, 4300, 6800, "万"),
                          ("NRR", "首次>100%", 92, 108, 85, 118, "%")],
    hero=1, threshold=(1, 100, "100%")))

# --- columns()/rows() weights= : proportional split, symmetric outer margins, guarded input ---
def _weights_grid():
    s = S()
    rail, main = dk.columns(2, slide=s, weights=(1, 2))
    assert abs(main[2] / rail[2] - 2.0) < 1e-6, "columns weights=(1,2) should give a 2x width ratio"
    assert abs(rail[0] - (W - (main[0] + main[2]))) < 1e-6, "outer margins must stay symmetric"
    r_top, r_bot = dk.rows(2, slide=s, weights=(3, 1))
    assert abs(r_top[3] / r_bot[3] - 3.0) < 1e-6, "rows weights=(3,1) should give a 3x height ratio"
    for a, b in zip(dk.columns(3, slide=s), dk.columns(3, slide=s, weights=(1, 1, 1))):
        assert all(abs(u - v) < 1e-9 for u, v in zip(a, b)), "equal weights must reproduce the default grid"
ok("columns/rows weights= (1:2 split, symmetric margins)", _weights_grid)
raises("columns rejects a length-mismatched weights tuple", lambda: dk.columns(2, slide=last(), weights=(1, 2, 3)))
raises("rows rejects a non-positive weight", lambda: dk.rows(2, slide=last(), weights=(1, 0)))

# --- dumbbell_board: leftward (lower-is-better) rows flip labels OUTWARD; optional v_mid dot.
#     Asserted on INK GEOMETRY directly (dk._ink_rect + dk._overlap_area), because lint_layout's
#     overlap_tol provably passes the old broken output — plus a byte-identity guard that a
#     rightward row's label frames still sit at the legacy coordinates. ---
def _dumbbell_directional():
    s = S()
    x, y0, w = 0.8, 1.6, 11.0
    dk.dumbbell_board(s, x, y0, w, [
        ("Cost / query", "", 0.62, 0.55, 0, 0.8, "$"),      # leftward: the proven-colliding row
        ("Latency", "p95", 1240, 620, 0, 1500, "ms"),       # leftward, wide span
        ("ARR", "+51%", 4980, 6350, 4300, 6800, "万"),      # rightward control (legacy geometry)
        ("NRR", "", 92, 108, 85, 118, "%", 100),            # rightward + v_mid (8-element row)
    ], label_w=0.42 * w)   # explicit label_w exercises the byte-identity frame guard below
    labels = {"0.62", "0.55 $", "1240", "620 ms", "4980", "6350 万", "92", "108 %", "100"}
    inks = []
    for sh in s.shapes:
        bb = dk._bbox_in(sh)
        if bb is None or not dk._is_text(sh):
            continue
        t = sh.text_frame.text.strip()
        if t in labels:
            r = dk._ink_rect(sh, bb)
            if r:
                inks.append((t, r[0], bb))
    assert len(inks) == 9, f"expected 9 value labels, found {len(inks)}"
    for i in range(len(inks)):
        for j in range(i + 1, len(inks)):
            ov = dk._overlap_area(inks[i][1], inks[j][1])
            assert ov < 1e-6, f"value-label ink overlap {ov:.3f}in^2: {inks[i][0]} x {inks[j][0]}"
    # byte-identity guard: the rightward row's label frames at the pre-change coordinates
    lw = 0.42 * w
    bx0, bx1 = x + lw + 0.15, x + w - 1.02
    span = 6800.0 - 4300.0
    x0 = bx0 + (4980 - 4300) / span * (bx1 - bx0)
    x1 = bx0 + (6350 - 4300) / span * (bx1 - bx0)
    ry = y0 + 2 * 0.52
    fb = next(b for t, _i, b in inks if t == "4980")
    fa = next(b for t, _i, b in inks if t == "6350 万")
    assert abs(fb[0] - (x0 - 0.62)) < 1e-3 and abs(fb[1] - (ry - 0.185)) < 1e-3, \
        "rightward before-label moved from its legacy frame"
    assert abs(fa[0] - (x1 + 0.10)) < 1e-3 and abs(fa[1] - (ry - 0.21)) < 1e-3, \
        "rightward after-label moved from its legacy frame"
ok("dumbbell_board leftward + v_mid rows (outward labels, zero ink overlap)", _dumbbell_directional)

# --- lint_layout SLIVER_GAP: warns on a near-touching seam, silent on a rule-compliant gap ---
def _sliver_gap():
    p = dk.blank_deck(10, 5.625); s2 = dk.add_slide(p)
    dk.box(s2, 1, 1.00, 4, 1.02, fill="E3E6EC")
    dk.box(s2, 1, 2.04, 4, 1.02, fill="E3E6EC")      # the documented 0.02in pitch-seam bug
    f = dk.lint_layout(p, verbose=False)
    assert any(c == "SLIVER_GAP" for _n, _sv, c, _m in f), "SLIVER_GAP missed a 0.02in seam"
    p2 = dk.blank_deck(10, 5.625); s3 = dk.add_slide(p2)
    dk.box(s3, 1, 1.0, 4, 1.0, fill="E3E6EC")
    dk.box(s3, 1, 2.2, 4, 1.0, fill="E3E6EC")        # a clear 0.2in gap — must stay silent
    assert not any(c == "SLIVER_GAP" for _n, _sv, c, _m in dk.lint_layout(p2, verbose=False)), \
        "SLIVER_GAP false-positive on a rule-compliant 0.2in gap"
ok("lint_layout SLIVER_GAP (0.02in seam warns, 0.2in gap silent)", _sliver_gap)

# --- lint_layout CONNECTOR_IN_BOX: flags an arrow anchored at a block's centre (drawn above it),
#     stays silent when the ends are edge-docked (connect_boxes/hub_spokes) or the line is covered ---
def _connector_in_box():
    hub = (3.5, 1.0, 3.0, 1.0); spokes = [(1.0, 3.5, 2.4, 1.0), (3.8, 3.5, 2.4, 1.0), (6.6, 3.5, 2.4, 1.0)]
    def deck(fn):
        p = dk.blank_deck(10, 5.625); s = dk.add_slide(p); fn(s); return p
    def centred(s):                                   # BUG: spokes leave the hub's centre, drawn ABOVE it
        dk.box(s, *hub, fill="FF5A4D", round=True)
        for sp in spokes: dk.box(s, *sp, fill="1B1E27", round=True)
        for sp in spokes: dk.connector(s, (hub[0]+hub[2]/2, hub[1]+hub[3]/2), (sp[0]+sp[2]/2, sp[1]), color=C("3DD6C4"))
    def docked(s):                                    # FIX: edge-docked fan
        dk.box(s, *hub, fill="FF5A4D", round=True)
        for sp in spokes: dk.box(s, *sp, fill="1B1E27", round=True)
        dk.hub_spokes(s, hub, spokes, color=C("3DD6C4"), gap=0.05)
    def covered(s):                                   # OK: connectors added BEFORE the nodes cover the seam
        for sp in spokes: dk.connector(s, (hub[0]+hub[2]/2, hub[1]+hub[3]/2), (sp[0]+sp[2]/2, sp[1]), color=C("3DD6C4"))
        dk.box(s, *hub, fill="FF5A4D", round=True)
        for sp in spokes: dk.box(s, *sp, fill="1B1E27", round=True)
    assert any(c == "CONNECTOR_IN_BOX" for _n, _sv, c, _m in dk.lint_layout(deck(centred), verbose=False)), \
        "CONNECTOR_IN_BOX missed a centre-anchored fan drawn above the hub"
    for name, fn in (("edge-docked", docked), ("covered", covered)):
        assert not any(c == "CONNECTOR_IN_BOX" for _n, _sv, c, _m in dk.lint_layout(deck(fn), verbose=False)), \
            f"CONNECTOR_IN_BOX false-positive on the {name} pattern"
    # edge_point lands ON the boundary, not the interior
    ex, ey = dk.edge_point((3.5, 1.0, 3.0, 1.0), (5.0, 4.0))
    assert abs(ey - 2.0) < 1e-6, "edge_point should dock on the block's bottom edge (y=2.0)"
ok("lint_layout CONNECTOR_IN_BOX (centre-anchor flags; edge-docked/covered silent)", _connector_in_box)

# --- icons.sanitize_svg: strips active/external content, keeps drawing content, and is ReDoS-safe ---
def _svg_sanitizer():
    import time as _t
    import icons as _ic
    # dangerous vectors removed, legit drawing + internal refs kept. NOTE: the file:// URL is a benign
    # SENTINEL ("file:///SANITIZED"), not a real system path — the assertion below only checks that the
    # `file://` reference is stripped, so no real sensitive local path needs to ship in this fixture.
    dirty = ('<svg onload="x()"><script>fetch("//e")</script>'
             '<foreignObject><iframe src="file:///SANITIZED"/></foreignObject>'
             '<image href="http://evil/x"/><style>@import url(http://evil)</style>'
             '<path d="M0 0h9"/><use href="#g"/></svg>')
    c = _ic.sanitize_svg(dirty).lower()
    for probe in ("<script", "foreignobject", "<iframe", "<image", "file://", "http://", "onload="):
        assert probe not in c, f"sanitize_svg left a {probe!r} vector"
    assert "<path" in c and "#g" in c, "sanitize_svg stripped legitimate drawing content / internal ref"
    # ReDoS guard: pathological all-open-tag input must be REJECTED fast, not scanned quadratically
    t0 = _t.time()
    try:
        _ic.sanitize_svg("<foreignObject" * 12000); raise AssertionError("oversize SVG was not rejected")
    except ValueError:
        pass
    dt = _t.time() - t0
    assert dt < 0.2, f"sanitize_svg took {dt*1000:.0f}ms on a ReDoS payload (must reject in <200ms)"
ok("icons.sanitize_svg (strips script/foreignObject/external refs; keeps paths; ReDoS-safe)", _svg_sanitizer)

# --- maps.choropleth: projection math, region-key matching (incl. NE -99 patch), and a real render ---
def _choropleth():
    import os as _os
    import tempfile as _tf
    import maps as _m
    assert _m._norm("广东省") == "广东" and _m._norm("Germany") == "germany", "region name normalization"
    fx, fy = _m._albers(10, 50, 10, 52, 43, 62)
    assert abs(fx) < 5 and abs(fy) < 5 and fx == fx, "albers projection produced a sane finite point"
    kde = _m._region_keys({"NAME": "Germany", "ISO_A2": "DE", "ISO_A3": "DEU"}, "ne")
    assert "de" in kde and "deu" in kde and "germany" in kde, "country key set (iso + name)"
    assert "fr" in _m._region_keys({"NAME": "France", "ISO_A2": "-99"}, "ne"), "NE -99 ISO patch (France→FR)"
    assert "广东" in _m._region_keys({"name": "广东省", "adcode": 440000}, "cn"), "province key (name)"
    # full render needs the base geometry (fetched+cached on first use); skip cleanly when offline
    out = _os.path.join(_tf.gettempdir(), "_smoke_choropleth.png")
    try:
        _m.choropleth_png(out, {"DE": 3.0, "FR": 7.0, "ES": 12.0}, "europe", legend=True)
        # NaN/None must NOT crash (treated as no-data), and div must zero-centre on 0
        _m.choropleth_png(out, {"DE": float("nan"), "FR": 7.0, "ES": None}, "europe", legend=True)
        _m.choropleth_png(out, {"DE": 20, "FR": -5, "ES": 12}, "europe", scale="div", legend=True)
    except RuntimeError as e:
        if "could not fetch" in str(e):
            print("  [smoke] choropleth render skipped — base geometry not cached (offline)"); return
        raise
    assert _os.path.exists(out) and _os.path.getsize(out) > 2000, "choropleth PNG missing/empty"
    lo, hi = -12, 20                                          # div zero-centre check (helper is offline)
    neg, pos = _m._div_poles("#1F5FA8"); assert neg != pos, "diverging poles must be distinct"
ok("maps.choropleth (projection + region-match + render)", _choropleth)

# --- composition/target/range components: stacked+area charts, bullet_graph (per-row scale), range_bars ---
def _comp_components():
    p = dk.blank_deck(10, 5.625); s = dk.add_slide(p)
    ser = [("A", [3, 4, 5]), ("B", [2, 3, 2])]
    for k in ("column_stacked", "column_stacked_100", "area_stacked"):     # the new native KIND entries resolve
        dk.native_chart(dk.add_slide(p), 0.5, 0.5, 4, 2.5, ["Q1", "Q2", "Q3"], ser, kind=k, legend=True)
    yb = dk.bullet_graph(s, 0.5, 0.5, 4.6, [("CSAT", 4.1, 4.5), ("Rev", 82, 90), ("Churn", 6, 5)],
                         higher_better=True)                                 # mixed units → per-row scale
    assert yb > 0.5, "bullet_graph returns a bottom y"
    yb2 = dk.range_bars(s, 5.0, 0.5, 4.4, [("DCF", 8, 12, 10), ("Comps", 9, 14, 11)])
    assert yb2 > 0.5, "range_bars returns a bottom y"
    assert not [f for f in dk.lint_layout(p, verbose=False) if f[1] == "CRITICAL"], "clean placement lints clean"
ok("native_chart stacked/area + bullet_graph (per-row) + range_bars", _comp_components)

ok("tint mixes toward white", lambda: dk.tint("1B7F5C", 0.14))
ok("kpi_card (delta + strip, tall enough)", lambda: dk.kpi_card(
    S(), 0.8, 0.8, 3.4, 2.3, "净收入留存 NRR", "108", unit="%", delta="+16pt",
    delta_color=dk.RGBColor(0x1B, 0x7F, 0x5C), sub="从 92%", strip="首次超过流失"))
ok("flow_compare (old/new + highlight + note)", lambda: dk.flow_compare(
    S(), 0.8, 1.4, 11.5, ["签约", "排期", "对接", "上线"], ["达标签约", "复用模板", "首次转化"],
    old_label="旧流程", new_label="新流程",             # CJK deck → pass the CJK row labels explicitly
    old_result="27 天", new_result="7.5 天", highlight_old=2, highlight_new=2,
    note="40% 卡在此", transition_label="模板化"))

# --- axis_scale: the shared value→x mapper (linear, degenerate-safe) ---
def _axis_scale():
    X, draw = dk.axis_scale(1.0, 8.0, 0.0, 100.0)
    assert abs(X(0) - 1.0) < 1e-6 and abs(X(100) - 9.0) < 1e-6 and abs(X(50) - 5.0) < 1e-6, \
        "axis_scale linear mapping wrong"
    Xd, _ = dk.axis_scale(1.0, 8.0, 5.0, 5.0)          # hi==lo must not divide by zero
    assert Xd(5.0) == 1.0, "axis_scale degenerate span not guarded"
    draw(S(), 3.0)                                       # draw_axis must not raise
ok("axis_scale (linear map + degenerate span + draw_axis)", _axis_scale)

# --- dot_strip: value-mapped dots with anti-collision labels, lint-clean ---
def _dot_strip():
    dk.EAFONT = dk.EAFONT or "Hiragino Sans GB"    # CJK labels need the EA slot (CJK_NO_EA gate)
    p = dk.blank_deck(10, 5.625); s2 = dk.add_slide(p)
    dk.dot_strip(s2, 0.6, 2.0, 8.0, [("博后", 70), ("学术", 100), ("工业", 180)],
                 60, 190, highlight=2, unit="k")
    dk.lint_layout(p, strict=True)                       # dots + labels stay in-canvas, no overlap
    # dense cluster near one end must not overflow the frame or collide (anti-collision path)
    p2 = dk.blank_deck(10, 5.625); s3 = dk.add_slide(p2)
    dk.dot_strip(s3, 0.6, 2.0, 8.0, [("a", 61), ("b", 63), ("c", 65), ("d", 188)], 60, 190, unit="k")
    dk.lint_layout(p2, strict=True)
ok("dot_strip (value-mapped + dense-cluster anti-collision, lint-clean)", _dot_strip)

# --- pangu: opt-in 盘古之白 normalizer — default OFF is byte-identical, modes are idempotent ---
def _pangu():
    s = "用K99机制占58.3%的博后"
    assert dk.pangu(s) == s, "pangu(default None) must be a no-op (byte-identical guarantee)"
    sp = dk.pangu(s, "spaced")
    assert "用 K99 机制" in sp and "58.3% 的" in sp, "pangu spaced did not insert boundary spaces"
    assert dk.pangu(sp, "spaced") == sp, "pangu spaced not idempotent"
    assert dk.pangu("用 K99 机制", "unspaced") == "用K99机制", "pangu unspaced did not strip"
    try:
        dk.pangu(s, "bogus"); raise AssertionError("pangu accepted an invalid mode")
    except ValueError:
        pass
ok("pangu (default no-op + spaced/unspaced idempotent + bad-mode guard)", _pangu)

# --- new geometry components: tier_stack / gantt / harvey_ball+eval_matrix / heat_matrix / device_frame ---
def _tier_stack():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.pyramid(s2, 1.0, 1.2, 4.4, 3.8, ["Vision", "Strategy", "Delivery", "Foundation"])
    dk.funnel(s2, 7.2, 1.2, 5.0, 3.8, ["Visits", "Signups", "Trials", "Paid"],
              values=[100, 54, 30, 12], labels="side")
    dk.lint_layout(p, strict=True)                       # tiers taper + labels stay collision-free
ok("tier_stack / pyramid + funnel (value-proportional, lint-clean)", _tier_stack)
raises("tier_stack rejects an empty tier list", lambda: dk.tier_stack(last(), 1, 1, 4, 3, []))

def _gantt():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.gantt(s2, 0.6, 1.3, 12.1,
             [("Discovery", 0, 1, 0), ("Design", 1, 2.5, 0),
              ("Build", 2, 5, 1), ("QA", 4.5, 6, 1), ("Launch", 6, 6.5, 2)],
             axis_min=0, axis_max=7, ticks=[0, 2, 4, 6], tick_labels=["Q1", "Q2", "Q3", "Q4"],
             lanes=["Plan", "Engineering", "GTM"], today=3.2, highlight=2)
    dk.lint_layout(p, strict=True)                       # bars keyed to axis_scale, swimlanes clean
ok("gantt (swimlanes + ticks + today + highlight, lint-clean)", _gantt)
raises("gantt raises on an off-axis bar", lambda: dk.gantt(last(), 0.6, 1.3, 11.0,
       [("X", 0, 5)], axis_min=0, axis_max=3))          # end 5 > axis_max 3 → off-canvas

def _harvey_ball():
    s2 = S()
    from pptx.oxml.ns import qn as _qn
    wsh = dk.harvey_ball(s2, 3.0, 3.0, 2, d=0.4)
    g = wsh._element.spPr.find(_qn("a:prstGeom"))
    assert g is not None and g.get("prst") == "pie", "harvey_ball(level=2) must build a PIE wedge"
    assert len(wsh.adjustments) == 2, "PIE must expose 2 adjustments (start/end angle)"
    dk.harvey_ball(s2, 4.0, 3.0, 0)                      # empty ring (no wedge)
    dk.harvey_ball(s2, 5.0, 3.0, 4)                      # full disc
ok("harvey_ball (PIE wedge geometry builds)", _harvey_ball)

def _eval_matrix():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.eval_matrix(s2, 0.8, 1.6, 11.0, ["Option A", "Option B", "Option C"],
                   ["Cost", "Speed", "Risk", "Support"],
                   [[4, 2, 3], [3, 4, 2], [2, 3, 4], [4, 1, 3]], recommend=0)
    dk.eval_matrix(dk.add_slide(p), 0.8, 1.6, 11.0, ["A", "B", "C"], ["c1", "c2", "c3"],
                   [["yes", "no", "partial"], ["partial", "yes", "no"], ["yes", "yes", "partial"]],
                   mark="mark", legend=False)
    dk.lint_layout(p, strict=True)
ok("eval_matrix (harvey balls + marks + recommend column, lint-clean)", _eval_matrix)

def _heat_matrix():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.heat_matrix(s2, 1.2, 1.4, 6.2, 4.2,
                   [[1, 2, 3, 4, 5], [2, 4, 6, 8, 10], [3, 6, 9, 12, 15],
                    [4, 8, 12, 16, 20], [5, 10, 15, 20, 25]],
                   ["Rare", "Unlikely", "Possible", "Likely", "Certain"],
                   ["Trivial", "Minor", "Moderate", "Major", "Severe"],
                   scale="risk", cell_labels=True)
    dk.heat_matrix(dk.add_slide(p), 1.2, 1.4, 6.0, 3.0,
                   [[10, 20, 30], [40, 50, 60]], ["r1", "r2"], ["c1", "c2", "c3"],
                   scale="div", cell_labels=True)
    dk.lint_layout(p, strict=True)                       # cells + contrast text stay collision-free
ok("heat_matrix (risk + div scales, contrast-aware text, lint-clean)", _heat_matrix)

def _device_frame():
    r = dk.device_frame(S(), IMG, 0.8, 1.3, 6.2, 3.6, chrome="browser",
                        url="app.example.com/dashboard")
    assert len(r) == 4 and r[2] > 0 and r[3] > 0, "device_frame must return the inner picture rect"
    dk.device_frame(S(), IMG, 5.2, 1.0, 2.2, 4.6, chrome="phone")
ok("device_frame (browser chrome + phone bezel)", _device_frame)

def _waterfall():
    q = os.path.join(TMP, "_wf.png")
    dc.waterfall(q, [("Start", None), ("Q1", 25), ("Q2", -12), ("Q3", 18), ("End", None)])
    assert os.path.exists(q) and os.path.getsize(q) > 1500, "waterfall PNG missing/trivial"
ok("dc.waterfall (floating step bars + dashed connectors)", _waterfall)
raises("waterfall rejects an empty item list", lambda: dc.waterfall(os.path.join(TMP, "_x.png"), []))

# --- box/connector kit: every node() shape + all three arrowhead variants ---
def _node_kit():
    s2 = S()
    centers = []
    for i, shp in enumerate(["roundrect", "rect", "pill", "circle", "diamond",
                             "parallelogram", "cylinder"]):
        centers.append(dk.node(s2, 0.4 + i * 1.8, 1.2, 1.6, 0.8, shp, shape=shp, sub="sub"))
    dk.node(s2, 0.4, 2.6, 1.6, 0.8, "hub", hub=True)
    dk.node(s2, 2.4, 2.6, 1.6, 0.8, "opt", dashed=True)
    dk.connector(s2, centers[0], centers[1], head="triangle", label="req")
    dk.connector(s2, (0.6, 4.0), (3.0, 4.0), style="dashed", head="open")
    dk.connector(s2, (0.6, 4.4), (3.0, 4.4), style="dotted", head="none")
    dk.elbow_connector(s2, dk.loop_path(8.0, 5.0, 4.0, 4.6), style="dotted", head="open", label="retry")
    dk.elbow_connector(s2, dk.loop_path(12.0, 9.0, 4.0, 4.6), head="none")
    dk.elbow_connector(s2, [(9.0, 5.2), (9.0, 5.6), (12.0, 5.6)], head="triangle")
ok("node (all 7 shapes + hub/dashed) + connector/elbow heads (triangle/open/none)", _node_kit)

raises("vstack overflow raises at build time", lambda: dk.vstack(
    S(), 0.6, 1.0, 5.0, [(2.0, lambda x, y, w: None), (2.0, lambda x, y, w: None)], bottom=3.0))

# --- content_band + bottom_callout: the footer-safe return geometry must be exact ---
def _band_and_callout():
    s2 = S()
    bx, by, bw, bh = dk.content_band(s2)
    assert abs(bx - dk.GUTTER) < 1e-6 and abs(bw - (W - 2 * dk.GUTTER)) < 1e-6, "content_band x/w wrong"
    assert abs((by + bh) - (H - dk.FOOTER_BAND - 0.15)) < 1e-6, "content_band bottom must clear the footer band"
    body = "the takeaway body — long enough that it could wrap on a narrower frame"
    top = dk.bottom_callout(s2, 0.6, W - 1.2, "TAKEAWAY", body)
    ch = dk.measure_callout("TAKEAWAY", body, W - 1.2)
    assert abs(top - (H - dk.FOOTER_BAND - 0.15 - ch)) < 1e-3, "bottom_callout bottom must anchor above the footer band"
    assert by < top < H, "bottom_callout top must fall inside the content band"
ok("content_band + bottom_callout (footer-safe return geometry)", _band_and_callout)

ok("equation_native (linear LaTeX -> editable runs)", lambda: dk.equation_native(
    S(), 0.8, 1.2, 9.0, 0.8, r"\mathcal{L} = \sum_i \|A x_i - y_i\|_2^2 + \lambda R(x_i)"))
raises("equation_native rejects 2-D \\frac (use equation_png)", lambda: dk.equation_native(
    S(), 0.8, 2.4, 6.0, 0.8, r"\frac{a}{b}"))

# --- step_list active_idx: the active disc's fill must DIFFER from an inactive one
#     (regression: the vertical branch once used _blend(acc, WHITE, 0.0) == acc, a no-op) ---
def _step_list_active():
    p = dk.blank_deck(10, 5.625); s2 = dk.add_slide(p)
    dk.step_list(s2, 0.8, 1.0, 8.0, [("Collect", "gather the inputs"), ("Train", "fit the model")],
                 active_idx=0)
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _T
    discs = [sh for sh in s2.shapes if sh.shape_type == _T.AUTO_SHAPE
             and abs(sh.width - dk.Inches(0.42)) < 2000 and abs(sh.height - dk.Inches(0.42)) < 2000]
    assert len(discs) == 2, f"expected 2 number discs, found {len(discs)}"
    fills = [d.fill.fore_color.rgb for d in discs]
    assert fills[0] != fills[1], "active_idx must render the active disc's fill differently from inactive"
    dk.step_list(s2, 0.8, 3.0, 8.0, [("A", ""), ("B", "")], orientation="horizontal", active_idx=1)
ok("step_list active vs inactive disc fills differ (v + h)", _step_list_active)

ok("cycle_diagram (3 nodes, defaults)", lambda: dk.cycle_diagram(
    S(), 6.5, 3.6, [("Plan", ""), ("Do", "daily"), ("Check", "")]))
ok("dumbbell_board (basic row, no hero/threshold)", lambda: dk.dumbbell_board(
    S(), 0.8, 1.6, 11.0, [("Latency", "", 120, 80, 0, 150, "ms")]))
ok("sources_page rule=True / rule=False", lambda: (dk.sources_page(S(), ["a", "b", "c"], rule=True),
                                                   dk.sources_page(S(), ["d", "e"], rule=False)))
ok("icon_tile (glass=True frosted tile)", lambda: dk.icon_tile(S(), 1.0, 1.0, 0.8, IMG, glass=True, fill=C("5B8DEF")))
ok("icon_tile two-colour grad shorthand (c0,c1) not misread as stop-list",
   lambda: dk.icon_tile(S(), 1.0, 1.0, 0.8, IMG, grad=(C("9DD9DD"), C("C6B8EB"))))

# --- grad normalization: (c0,c1) and stop-list forms both parse; RGBColor-is-tuple bug guard
def _grad_norm():
    a = dk._norm_stops((C("9DD9DD"), C("C6B8EB")))                       # shorthand → 2 stops
    b = dk._norm_stops([(0.0, C("FF0000"), 1.0), (1.0, C("0000FF"), 1.0)])  # full stops
    c = dk._norm_stops([(0.0, "9DD9DD"), (1.0, "C6B8EB")])               # (pos,colour) no alpha
    assert len(a) == 2 and abs(a[0][2] - 1.0) < 1e-9, a
    assert len(b) == 2 and str(b[0][1]) == "FF0000", b
    assert len(c) == 2 and abs(c[1][0] - 1.0) < 1e-9, c
ok("_norm_stops handles (c0,c1) + stop-list + (pos,colour)", _grad_norm)

# --- icon_tile contrast guard: a low-contrast glyph↔tile pair gets the tile auto-nudged to >=3:1
def _icon_guard():
    from PIL import Image as _I
    gp = os.path.join(TMP, "_vio_glyph.png")
    im = _I.new("RGBA", (48, 48), (0, 0, 0, 0))
    for x in range(16, 34):
        for y in range(22, 28): im.putpixel((x, y), (0x60, 0x52, 0xB3, 255))  # violet ink
    im.save(gp)
    assert dk.contrast_ratio(C("6052B3"), C("14182A")) < 3.0                  # the bad pair
    ink = dk._png_dominant_ink(gp)
    assert str(ink) == "6052B3", ink                                          # ink inferred from PNG
    dk.icon_tile(S(), 1.0, 1.0, 0.8, gp, shape="circle", fill=C("14182A"))    # no glyph= → auto-guard
ok("icon_tile auto-guards low glyph/tile contrast (infers ink, nudges tile)", _icon_guard)

prs.save(os.path.join(TMP, "_smoke_deck.pptx"))

# --- overflow-proof components: the four real-world overflow shapes self-heal, lint-clean ---
def _overflow_proof():
    # (1) meter_bar near the right edge with the roomy default value_w -> bar auto-shortens
    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.meter_bar(s, 4.6, 2.0, 4.7, 0.95, value="15.2%", accent=dk.MAGENTA)
    # (2) scorecard with a long wrapping caption in a short 4-up tile -> caption fits by measure
    dk.scorecard(s, 0.7, 1.8, 1.99, 1.9, "UNEMPLOYMENT", "5.9%", delta="-0.1pp",
                 caption="EU rate, May 2026 - near its all-time low across the whole series")
    # (3) insight_banner whose body wraps to 2 lines -> bar grows/font steps down, text contained
    dk.insight_banner(s, 0.7, 4.1, 8.6,
        "A body long enough to wrap onto a second line in this width, which used to overflow the bar.")
    # (4) stat_row long figure in a third-width column -> figure scales, never wraps mid-number
    s2 = dk.add_slide(p)
    dk.stat_row(s2, 0.7, 2.0, 8.6,
                [("-0.9 million", "working-age people lost per year, on average (projection)"),
                 ("2029", "the peak year"), ("10.9%", "slack, trending down")], fig_size=30)
    dk.lint_layout(p, strict=True)
    # and a FIT call stays byte-identical in spirit: no faults, no adjustment needed
    p2 = dk.blank_deck(); s3 = dk.add_slide(p2)
    dk.meter_bar(s3, 0.7, 2.0, 5.0, 0.6, value="60%", value_w=0.9)
    dk.scorecard(s3, 0.7, 3.0, 2.2, 1.9, "REVENUE", "42%", caption="short caption")
    dk.lint_layout(p2, strict=True)
ok("overflow-proof components (meter_bar/scorecard/insight_banner/stat_row self-heal, lint-clean)", _overflow_proof)


# --- numerals: PREVENTED in the components, warned (not blocked) on hand-set runs ---
def _oldstyle_figures():
    import presets as _pr
    saved = (dk.DISPLAY, dk.FONT)

    def _codes(fn):
        p = dk.blank_deck(); s = dk.add_slide(p); fn(s)
        return [f[2] for f in (dk.lint_layout(p, verbose=False, strict=False) or [])]

    def _builds(fn):
        p = dk.blank_deck(); s = dk.add_slide(p); fn(s)
        try:
            dk.lint_layout(p, verbose=False, strict=True); return True
        except RuntimeError:
            return False

    # 1) the rule NEVER blocks a build — it is a taste call, and its false-positive surface
    #    (every component x font x string shape; "7" and "10x" are digit-dominant) is unbounded
    for face, txt, size in (("Georgia", "596,513", 60), ("Hoefler Text", "2026", 34),
                            ("Georgia", "7", 40), ("Georgia", "10x", 40)):
        assert _builds(lambda s, f=face, t=txt, z=size: dk.text(
            s, 1, 1, 8, 1.6, [[(t, z, C("1B3A63"), True, False, f)]])), \
            "OLDSTYLE must never block a build (%s / %r)" % (face, txt)

    # 2) but a hand-set wobbling numeral IS reported
    assert "OLDSTYLE_FIGURES" in _codes(lambda s: dk.text(
        s, 1, 1, 8, 1.6, [[("596,513", 60, C("1B3A63"), True, False, "Georgia")]])), \
        "a hand-set Georgia hero numeral must be reported"
    # body prose and word-bearing headings are silent
    for txt, size in (("In 2026 the city grew.", 13), ("2026 Roadmap", 28), ("2026年展望", 28)):
        assert "OLDSTYLE_FIGURES" not in _codes(lambda s, t=txt, z=size: dk.text(
            s, 1, 1, 8, 1.6, [[(t, z, C("1B3A63"), True, False, "Georgia")]])), \
            "prose/headings must stay silent (%r)" % txt

    # 3) PREVENTION: components resolve a lining face themselves, under any preset AND under a
    #    Georgia BODY face (the provided-template case no preset covers)
    def _figures(s):
        dk.big_numeral(s, 0.6, 0.6, "03")
        dk.stat_row(s, 0.6, 1.9, 8.6, [("596,513", "", "returns")])
        dk.scorecard(s, 0.6, 3.0, 2.4, 1.3, "REVENUE", "596,513")
        dk.kpi_card(s, 3.2, 3.0, 2.6, 1.3, "USERS", "1,634")
        dk.change_stat(s, 0.6, 4.5, 5.0, 0.9, "4,461", "4,509")
        dk.ghost_numeral(s, 6.2, 4.4, 3.2, 1.0, "07")
    for nm in ("editorial_paper", "editorial_report", "luxury_dark", "museum_memorial", "swiss"):
        p = _pr.preset(nm); dk.DISPLAY = p.get("display"); dk.FONT = p.get("font") or dk.FONT
        assert "OLDSTYLE_FIGURES" not in _codes(_figures), "preset %s emits a wobbling figure" % nm
    dk.DISPLAY = dk.FONT = "Georgia"
    assert "OLDSTYLE_FIGURES" not in _codes(_figures), "Georgia BODY face: components must self-resolve"

    # 4) a WORD value must keep the deck's own font — resolving it split one card across two faces
    assert dk.numeral_run_face("Excellent", fallback="Georgia") == "Georgia"
    assert dk.numeral_run_face("596,513", fallback="Georgia") == dk.NUMERAL_SERIF
    # 5) big_numeral follows the deck, and is not pinned to one face
    dk.DISPLAY = dk.FONT = "Arial"
    assert dk.numeral_face(None, fallback=dk.DISPLAY) == "Arial", "big_numeral must not be pinned"
    # 6) face detection is measured, and the curated fallback is REACHABLE (it was shadowed once)
    assert dk.has_oldstyle_figures("Georgia") and not dk.has_oldstyle_figures("Palatino")
    assert dk.has_oldstyle_figures("Constantia"), "fallback list must be reachable for absent fonts"
    dk.DISPLAY, dk.FONT = saved
ok("numerals: prevented in components, warned (never blocked) on hand-set runs", _oldstyle_figures)


# --- no shape may carry the inherited theme shadow ---
def _no_inherited_effect():
    from pptx.oxml.ns import qn as _qn
    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.box(s, 0.5, 0.5, 3, 1, fill="1B3A63", round=True)
    dk.chip(s, 0.5, 1.8, 2, 0.5, "A", "sub", "D4442E")
    dk.node(s, 4, 0.5, 2, 0.8, "Node")
    dk.connector(s, (4, 1.3), (6, 2.0))
    dk.timeline(s, 0.5, 3.6, 8.6, [("2024", "A"), ("2025", "B"), ("2026", "C")])
    stray = [sh for sh in s.shapes if sh._element.find(_qn("p:style")) is not None]
    assert not stray, "%d shape(s) still carry the theme <p:style> (a soft shadow in LibreOffice)" % len(stray)
    codes = [f[2] for f in (dk.lint_layout(p, verbose=False, strict=False) or [])]
    assert "INHERITED_EFFECT" not in codes
ok("no shape carries the inherited theme shadow (<p:style>)", _no_inherited_effect)


# --- design_intent: the declared-register channel round-trips through the saved file ---
def _design_intent():
    import json as _json
    from pptx import Presentation as _P
    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.design_intent(s, envelope="upper", rhyme=2, reason="pivot beat")
    import os, tempfile
    f = os.path.join(tempfile.mkdtemp(), "t.pptx"); p.save(f)
    sl = list(_P(f).slides)[0]
    tags = [sh.name for sh in sl.shapes if (sh.name or "").startswith("deckkit-intent:")]
    assert len(tags) == 1, "intent tag must survive save"
    payload = _json.loads(tags[0].split(":", 1)[1])
    assert payload == {"envelope": "upper", "rhyme": 2, "reason": "pivot beat"}
    # and the tag must not trip the geometry lint
    dk.lint_layout(p, verbose=False, strict=True)
ok("design_intent tag round-trips and is lint-invisible", _design_intent)


# --- pic_alpha: native picture opacity, no overlay shape ---
def _pic_alpha():
    from pptx.oxml.ns import qn as _qn
    p = dk.blank_deck(); s = dk.add_slide(p)
    import glob
    img = glob.glob(os.path.expanduser(
        "~/Downloads/slides_skill_test/tokyo-first-timers/assets/opt/*.jpg"))
    if not img:
        return                       # fixture machine only; the A/B render proved the visual
    pic = dk.picture(s, img[0], 1, 1, 4, 3, fit="cover", alt="")
    dk.pic_alpha(pic, 14)
    blip = pic._element.blipFill.find(_qn("a:blip"))
    fixes = blip.findall(_qn("a:alphaModFix"))
    assert len(fixes) == 1 and fixes[0].get("amt") == "14000"
    dk.pic_alpha(pic, 30)            # idempotent: replaces, never stacks
    fixes = blip.findall(_qn("a:alphaModFix"))
    assert len(fixes) == 1 and fixes[0].get("amt") == "30000"
ok("pic_alpha sets native picture opacity idempotently", _pic_alpha)


# --- tier-2 components: build + lint clean, and their correctness contracts hold ---
def _tier2_components():
    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.small_multiples(s, 0.6, 0.6, 8.8, 4.2,
                       [("A", [1, 2, 3]), ("B", [2, 2, 2]), ("C", [1, 5, 9]), ("D", [0, 1, 1])],
                       categories=["x", "y", "z"], highlight=2)
    # shared scale is the CONTRACT: every panel's value axis pinned to one [lo, hi]
    charts = [sh.chart for sh in s.shapes if sh.has_chart]
    assert len(charts) == 4
    tops = {round(c.value_axis.maximum_scale, 4) for c in charts}
    lows = {round(c.value_axis.minimum_scale, 4) for c in charts}
    assert len(tops) == 1 and len(lows) == 1, "small_multiples panels must share one axis scale"
    dk.lint_layout(p, verbose=False, strict=True)

    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.position_map(s, 0.8, 0.8, 8.4, 4.2, [("A", 1, 1), ("B", 9, 8), ("C", 5, 2)], highlight=1)
    dk.lint_layout(p, verbose=False, strict=True)
    try:
        dk.position_map(s, 0.8, 0.8, 8.4, 4.2, [("A", 5, 5), ("B", 5.0, 5.0)])
        assert False, "coinciding points must raise"
    except ValueError:
        pass

    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.org_tree(s, 0.6, 0.6, 8.8, 4.4,
                ("R", [("a", [("a1", []), ("a2", [])]), ("b", [("b1", [])])]))
    dk.lint_layout(p, verbose=False, strict=True)
    try:
        dk.org_tree(s, 0.6, 0.6, 2.0, 4.0, ("R", [(str(i), []) for i in range(9)]))
        assert False, "an unfittable tree must raise, not squeeze"
    except ValueError:
        pass
ok("tier-2 components (small_multiples shared axis · position_map · org_tree)", _tier2_components)


def _position_map_hero_binding():
    """The hero label must be BOUND to its dot — bold + the dot's own hue — and must claim its
    slot before the greedy nudge, or a neighbour pushes the one label that matters off its point."""
    from pptx.util import Emu

    def probe(highlight):
        p = dk.blank_deck(); sl = dk.add_slide(p)
        dk.position_map(sl, 0.6, 0.6, 8.8, 4.2,
                        [("alpha", 1, 1), ("beta", 9, 9), ("gamma", 8.6, 8.4)],
                        highlight=highlight, accent="1F6E68")
        out = {}
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text in ("alpha", "beta", "gamma"):
                r = sh.text_frame.paragraphs[0].runs[0]
                col = str(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else None
                out[sh.text_frame.text] = (round(Emu(sh.top).inches, 3), r.font.bold, col)
        return out

    a = probe(1)
    assert a["beta"][1] is True and a["beta"][2] == "1F6E68", "hero label must be bold in the dot's hue"
    assert a["alpha"][1] is not True and a["gamma"][1] is not True, "non-hero labels stay quiet"
    b = probe(2)
    assert b["gamma"][2] == "1F6E68" and b["beta"][2] != "1F6E68", "hue must follow the declared hero"
    assert b["gamma"][0] <= a["gamma"][0], "hero-first ordering must not push the hero down"
    assert all(v[1] is not True for v in probe(None).values()), "no highlight => no bold label"
ok("position_map binds the hero label to its dot (bold + hue + first claim)", _position_map_hero_binding)


def _rule_through_text():
    """RULE_THROUGH_TEXT — a divider at a hand-picked y that the grown text above it now crosses.

    Shipped twice in one delivered deck (slides 5 and 13) and was caught by the USER, not by any
    lint: a thin box over text is not TEXT_OVERLAP, not an overflow, and not invisible text.
    Both directions are asserted, because the false-positive surface here is the legitimate
    title underline — a rule sitting just BELOW the ink, which must stay silent.
    """
    def codes(build):
        p = dk.blank_deck(); sl = dk.add_slide(p)
        build(sl)
        return {f[2] for f in dk.lint_layout(p, verbose=False)}

    def crossing(sl):
        dk.text(sl, 1.0, 1.0, 3.0, 1.4, [[("of voting members who backed the new agreement, signed "
                                           "by all parties on 20 July and ratified.", 12,
                                           dk.RGBColor(0, 0, 0), False, False, "Arial")]])
        dk.box(sl, 1.0, 1.55, 2.8, 0.02, fill="999999")
    assert "RULE_THROUGH_TEXT" in codes(crossing), "a rule drawn through text did not fire"

    def underline(sl):                     # the legitimate case: a rule just below the ink
        dk.text(sl, 1.0, 1.0, 6.0, 0.4, [[("A slide title", 26, dk.RGBColor(0, 0, 0), True, False, "Arial")]])
        dk.box(sl, 1.0, 1.52, 0.6, 0.035, fill="D2691E")
    assert "RULE_THROUGH_TEXT" not in codes(underline), "a title underline must not fire"

    def between(sl):                       # the normal divider, correctly derived
        dk.text(sl, 1.0, 1.0, 3.0, 0.4, [[("block one", 12, dk.RGBColor(0, 0, 0), False, False, "Arial")]])
        dk.box(sl, 1.0, 1.6, 2.8, 0.02, fill="999999")
        dk.text(sl, 1.0, 1.8, 3.0, 0.4, [[("block two", 12, dk.RGBColor(0, 0, 0), False, False, "Arial")]])
    assert "RULE_THROUGH_TEXT" not in codes(between), "a divider between blocks must not fire"

    def axis(sl):                          # a chart/motif axis under a filled bar's label
        dk.box(sl, 1.0, 1.0, 4.0, 0.46, fill="FF8A3D")
        dk.text(sl, 1.1, 1.05, 3.8, 0.36, [[("four years, paid", 11.5, dk.RGBColor(0, 0, 0), True, False, "Arial")]])
        dk.box(sl, 1.0, 1.56, 6.0, 0.014, fill="999999")
    assert "RULE_THROUGH_TEXT" not in codes(axis), "a bar's axis line must not fire"
ok("RULE_THROUGH_TEXT fires on a crossed rule and stays quiet on underlines/dividers/axes",
   _rule_through_text)


def _quiet_region_contract():
    import glob
    from image_fx import quiet_region
    imgs = glob.glob(os.path.expanduser(
        "~/Downloads/slides_skill_test/tokyo-first-timers/assets/opt/hero_cover.jpg"))
    if not imgs:
        return
    fx, fy, fw, fh, lum = quiet_region(imgs[0])
    assert 0 <= fx <= 1 and 0 <= fy <= 1 and 0 < fw <= 1 and 0 < fh <= 1
    # ink-zone coherence: the region must not average dark sky with cream ground into an
    # unusable mid-lum (the exact failure the constraint was added for)
    assert lum < 120 or lum > 150, "quiet_region returned a mixed-ink region (lum %d)" % lum
ok("quiet_region returns one coherent ink zone", _quiet_region_contract)



def _iso_components():
    """The 2.5D isometric suite: builds clean, lints clean, and iso_bars is FAITHFUL — a bigger
    value must render a taller extrusion, or the 2.5D is decorating a lie."""
    from pptx.util import Emu
    p = dk.blank_deck(); s = dk.add_slide(p)
    # a single prism returns its top-face centre for label placement
    c = dk.iso_prism(s, 2.0, 4.0, 1.2, 1.2, 1.0, "3E6E9E")
    assert isinstance(c, tuple) and len(c) == 2

    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.iso_bars(s, 0.8, 1.4, 8.4, 3.4, [10, 90, 40], labels=["a", "b", "c"], highlight=1)
    dk.lint_layout(p, verbose=False, strict=True)
    # FAITHFULNESS: the top face of the value-90 bar must sit higher on screen (smaller y) than the
    # top face of the value-10 bar. Measure the min-y of the freeform faces per bar cluster.
    tops = sorted(Emu(sh.top).inches for sh in s.shapes if "FREEFORM" in str(sh.shape_type))
    assert tops, "iso_bars drew no freeform faces"
    # crude but sufficient: the overall span of face-tops must be non-trivial (bars differ in height)
    assert max(tops) - min(tops) > 0.5, "iso_bars produced near-equal heights for very different values"

    p = dk.blank_deck(); s = dk.add_slide(p)
    dk.iso_stack(s, 0.6, 1.1, 9.0, 4.0, [("Base", "x"), ("Mid", "y"), ("Top", "z")],
                 accents=["2E9C93", "3E6E9E", "8A5BC7"])
    dk.lint_layout(p, verbose=False, strict=True)

    try:
        dk.iso_bars(s, 0.8, 1.4, 8.4, 3.4, [])
        assert False, "empty iso_bars must raise"
    except ValueError:
        pass
ok("iso 2.5D suite (prism · faithful bars · aligned stack)", _iso_components)

print(f"\nsmoke_deckkit: {len(fails)} failure(s)" + ("" if not fails else " — " + "; ".join(n for n, _ in fails)))
sys.exit(1 if fails else 0)

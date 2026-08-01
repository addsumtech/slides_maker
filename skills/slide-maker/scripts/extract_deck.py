#!/usr/bin/env python3
"""Pull the CONTENT out of an existing .pptx so a redesign can reuse it.

When the user hands you their own deck to *improve*, you need its real material —
the actual text per slide, the tables, the charts' numbers, the embedded figures, and
what they planned to SAY — not a guess at what's on it. This dumps all of it so the
rebuild is faithful to what they wrote and reuses their figures whole (see
design-principles: "use the source's own figures").

Pairs with render_deck.sh (which shows you what each slide LOOKS like) — run both:
render to see the design problems, extract to recover the content to carry forward.

**This tool is also the reconcile path.** `handoff-and-iteration.md` names it as the
preferred way to recover a user's hand-edits before a post-delivery rebuild, and that
procedure ends by declaring the build script canonical again. So anything this drops is
not merely missing from a report — it is deleted from the user's deck on the next
rebuild, by the procedure that exists to protect it. Two things it used to drop:

  * SPEAKER NOTES. There was no `notes_slide` reference in the file at all, so the
    user's entire narration — the thing a presented deck is actually *for* — was
    discarded by the tool whose stated job is to be faithful to what they wrote.
  * NATIVE CHARTS. A chart is a GraphicFrame: it matches neither PICTURE nor
    `has_table` nor `has_text_frame`, and the walk had no `else`, so its categories and
    series vanished without a word. Worse, a slide that was nothing BUT a chart came out
    as the literal string `_(empty / decorative only)_` — the tool did not just lose the
    data, it affirmatively reported that there had been none.

The rule that follows, and the one worth keeping: **never report absence you have not
verified.** A shape this script cannot read is now listed as UNEXTRACTED with its type,
so a lossy extraction announces itself instead of looking complete.

Usage:  python3 extract_deck.py /path/to/their_deck.pptx [out_dir]
Output: <out_dir>/content.md         — per-slide text, tables, chart data, notes, figures
        <out_dir>/assets/slideNN_*   — every embedded picture, saved whole
(default out_dir: ./extracted)
"""
import os
import sys

from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deckkit import open_presentation           # accepts .potx, which is how institutions ship

# Shape types with nothing to recover — reporting them as UNEXTRACTED would be noise, and noise
# is how a report teaches people to ignore it.
_DECORATIVE = {MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM}


def _chart_block(sh):
    """A chart's real numbers, as a readable table. Returns None when it cannot be read.

    python-pptx exposes categories and series on the plot; both are lazy, and a chart type it
    does not model raises rather than returning empty — which must surface as UNEXTRACTED, not
    as a chart with no data.
    """
    ch = sh.chart
    cats = [str(c) for c in ch.plots[0].categories]
    rows = []
    for ser in ch.series:
        vals = ["" if v is None else str(v) for v in ser.values]
        rows.append([str(ser.name or "")] + vals)
    head = [""] + cats
    widths = [max(len(r[i]) if i < len(r) else 0 for r in [head] + rows)
              for i in range(len(head))]
    def fmt(r):
        return " | ".join((r[i] if i < len(r) else "").ljust(widths[i]) for i in range(len(head)))
    try:
        kind = str(ch.chart_type).split(" ")[0]
    except Exception:
        kind = "chart"
    return "[{}]\n{}\n{}".format(kind, fmt(head), "\n".join(fmt(r) for r in rows))


def walk(shapes, slide_idx, found, assets, counter):
    """Recurse shapes (groups included), collecting text, tables, charts, pictures.

    Anything that matches no branch is recorded in found['unextracted'] rather than skipped —
    the silent `continue` at the end of this loop is what deleted every chart.
    """
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, slide_idx, found, assets, counter)
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = sh.image
                counter[0] += 1
                fn = "slide{:02d}_img{:02d}.{}".format(slide_idx, counter[0], img.ext)
                with open(os.path.join(assets, fn), "wb") as f:
                    f.write(img.blob)
                w = Emu(sh.width).inches
                h = Emu(sh.height).inches
                found["imgs"].append("{}  ({:.1f}x{:.1f} in)".format(fn, w, h))
            except Exception as e:
                found["imgs"].append("<unreadable picture: {}>".format(e))
            continue
        if sh.has_table:
            t = sh.table
            found["tables"].append(
                "\n".join(" | ".join(c.text.strip() for c in r.cells) for r in t.rows))
            continue
        if sh.has_chart:
            try:
                found["charts"].append(_chart_block(sh))
            except Exception as e:
                found["unextracted"].append("chart ({}) — could not read its data".format(e))
            continue
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                found["texts"].append(txt)
            continue
        if sh.shape_type not in _DECORATIVE:
            found["unextracted"].append(str(sh.shape_type))


def extract(path, out="./extracted"):
    assets = os.path.join(out, "assets")
    os.makedirs(assets, exist_ok=True)
    prs = open_presentation(path)
    n_slides = len(prs.slides._sldIdLst)
    W = Emu(prs.slide_width).inches
    H = Emu(prs.slide_height).inches
    lines = ["# Extracted content — {}".format(os.path.basename(path)),
             "\nSlide size: {:.2f} x {:.2f} in · {} slides\n".format(W, H, n_slides)]
    counter = [0]
    n_notes = n_charts = 0

    for i, slide in enumerate(prs.slides, 1):
        found = {"texts": [], "tables": [], "charts": [], "imgs": [], "unextracted": []}
        walk(slide.shapes, i, found, assets, counter)

        notes = ""
        # `has_notes_slide` is checked first on purpose: touching `.notes_slide` CREATES one, so
        # probing it directly would mutate the user's deck to answer a question about it.
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        lines.append("\n## Slide {}".format(i))
        if found["texts"]:
            lines.append("**Text:**")
            for t in found["texts"]:
                lines.append("- " + t.replace("\n", "  \n  "))
        if found["tables"]:
            lines.append("\n**Table(s):**")
            for tb in found["tables"]:
                lines.append("```\n" + tb + "\n```")
        if found["charts"]:
            n_charts += len(found["charts"])
            lines.append("\n**Chart data (rebuild with deckkit.native_chart — these are the real "
                         "numbers, not a reading of the picture):**")
            for cb in found["charts"]:
                lines.append("```\n" + cb + "\n```")
        if found["imgs"]:
            lines.append("\n**Figures (saved to assets/, reuse these whole):**")
            for im in found["imgs"]:
                lines.append("- " + im)
        if notes:
            n_notes += 1
            lines.append("\n**Speaker notes (what they planned to SAY — carry this into the "
                         "rebuild's notes, don't re-draft it):**")
            lines.append("> " + notes.replace("\n", "\n> "))
        if found["unextracted"]:
            lines.append("\n**⚠ UNEXTRACTED — this script could not read these; open the slide "
                         "and check what is on it before rebuilding:**")
            for u in found["unextracted"]:
                lines.append("- " + u)
        if not any(found[k] for k in ("texts", "tables", "charts", "imgs", "unextracted")) \
                and not notes:
            lines.append("_(empty / decorative only)_")

    with open(os.path.join(out, "content.md"), "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")
    return {"slides": n_slides, "images": counter[0], "charts": n_charts, "notes": n_notes,
            "out": out}


def main(argv):
    if len(argv) < 1:
        print("usage: python3 extract_deck.py /path/to/their_deck.pptx [out_dir]")
        return 2
    r = extract(argv[0], argv[1] if len(argv) > 1 else "./extracted")
    print("extracted {} slides -> {}/content.md  ({} images -> {}/assets/, "
          "{} chart(s), {} slide(s) with speaker notes)".format(
              r["slides"], r["out"], r["images"], r["out"], r["charts"], r["notes"]))
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))

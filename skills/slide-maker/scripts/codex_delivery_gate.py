#!/usr/bin/env python3
"""Strict, Codex-only delivery gate for the slide-maker workflow.

This script intentionally does not change deckkit, component_audit, or the
Claude Code workflow. It verifies the additional evidence Codex must collect
before claiming a deck is ready for delivery.
"""

from __future__ import annotations

import argparse
import ast
import hashlib
import importlib.util
import json
import os
import struct
import subprocess
import sys
from datetime import datetime, timezone
from collections import defaultdict
from pathlib import Path
from typing import Any


SCHEMA = "slide-maker-codex-evidence/v2"
RECEIPT_SCHEMA = "slide-maker-codex-delivery-receipt/v1"
BODY_FLOORS = {"presented": 13.5, "textheavy": 13.5, "selfread": 12.0}
# Each entry must be a code `lint_deck.py` can actually emit, normalised the way check_lint
# normalises one ("CJK TIGHT LEADING: slide 3 ..." -> "cjk_tight_leading"). Two of these were
# invented names -- `cjk_risk` and `color_envelope` -- that no linter output could ever match, so
# the gate believed it was enforcing them and could never fire. The suite now asserts every name
# here against the codes lint_deck really produces, because the existing test fed the gate its own
# expected shape (`lint["stats"]["warnings"] = ["card_dominance"]`) and so proved the logic while
# proving nothing about the vocabulary.
STRICT_STATS = {
    "card_dominance",
    "cjk_tight_leading",
    "envelope_monoculture",
    "flat_type",
    "size_sprawl",
    "small_type",
    "timid_cover",
}
# The gate had a blocking path for `severity == "error"` findings and a remediate-or-waive path
# for the `stats` stream, and NOTHING in between -- so a per-slide warning could not be held to
# any bar at all, however objective. That gap is where the accessibility floors live: measured on
# a delivered deck, an icon at 2.69:1 against its canvas (WCAG 1.4.11 asks 3:1) produced a warning
# on the `warnings` stream and passed every gate here.
#
# Deliberately TINY, and only floors with an arithmetic answer. A ratio either clears 3:1 or it
# does not; whether a component is over-used or a page is too dense is a judgment, and a judgment
# forced through a waiver form becomes a rubber stamp. Nothing that needs an opinion belongs here.
STRICT_WARNINGS = {
    "ICON CONTRAST",                                # WCAG 1.4.11, recolored monochrome icons
    "NON-TEXT CONTRAST",                            # WCAG 1.4.11, solid marks and connector lines
}
ICON_HELPERS = {"icon", "icon_card", "icon_tile", "icon_badge", "icon_ghost"}

sys.path.insert(0, str(Path(__file__).resolve().parent))
from written_reason import reason_width  # noqa: E402  (one shared definition, never a copy)
import re as _re  # noqa: E402

# The angle-bracket TEMPLATE placeholder — `<label-A>`, `<the question this room is asking>`. The
# raw `arc_divergence.py --template` / gate-skeleton output, pasted unedited, cleared every width
# floor because a placeholder is wide. This file already ran `startswith("<")` on fast_basis /
# none_opt_in for the same reason; this is that predicate, needing a closing `>` so a real "<50%"
# survives, shared across the arc / slides / checkpoint records.
_PLACEHOLDER = _re.compile(r"<[^<>]{2,}>")


def has_placeholder(value: Any) -> bool:
    """True when a string still carries an angle-bracket template placeholder (`<...>`)."""
    return bool(_PLACEHOLDER.search(str(value or "")))

# Kept in lockstep with render_deck._STRUCTURAL_ROLES and arc_divergence._STRUCTURAL: a cover is
# not a beat in the argument, and it is not a duplicate takeaway either.
STRUCTURAL_ROLES = frozenset((
    "cover", "agenda", "divider", "closing", "section", "thanks", "qa",
    "封面", "目录", "分隔", "收尾", "结尾", "致谢", "问答"))

TEMPLATE = {
    "schema": SCHEMA,
    "runtime": "codex",
    "delivery": "presented",
    # The tier is collected at the POST-BUILD review question (SKILL.md Step 5), with the rendered
    # deck visible — never at the interview. "fast" is the pre-selected default there.
    "review_effort": "fast",
    "fast_basis": "<only when review_effort is fast: >=12 chars — how fast was reached: the user's post-build choice, or 'post-build default — auto/not asked'>",
    # "none" (user declined review at the post-build question) additionally requires "none_opt_in":
    # >=12 chars quoting the user's decline, given AFTER they saw the rendered deck. Never a default.
    "deck": {
        "pptx": "deck.pptx",
        "sha256": "<sha256 of final deck.pptx>",
        "slide_count": 10,
    },
    "interview": {
        "mode": "answered",
        "record": "<user answers or auto-carved rationale>",
        # The one interview axis with a field of its own, because it is the one that goes missing
        # on a runtime with no choice UI and nothing downstream complains. A range, a time budget,
        # or a recorded derivation are all valid: "medium, 9-15" / "20 min, so ~18" / "user
        # declined — derived 11 from the ledger". What is NOT valid is silence, whose observed
        # result is a one-slide deck.
        "length": "<what the user said, or how the count was derived when they did not>",
    },
    "content": {
        "source_mode": "provided",
        "sources": [
            {
                "kind": "provided",
                "path": "README.md",
                "sha256": "<sha256>",
            }
        ],
        # WEB-RESEARCHED decks (source_mode == "web") MUST also fill these three — the floors from
        # content-planner.md §2(e), mirrored in the shared content checkpoint. Omit for a 'provided'
        # or 'none' deck.
        #   "coverage":  "<domain areas enumerated · covered · consciously cut (why)>",   # 全面
        #   "lifecycle": "<every featured product/version/entity confirmed live as of today · "
        #                "anything found discontinued/renamed + how the deck handles it>",  # (proactive)
        #   "provenance": {"summary": "checked N · confirmed N · fixed N · cut N",           # 准确
        #                  "method": "corroborated >=2 independent credible sources; MED labelled 'per public reporting'"},
        # and every claim_ledger row carries a "confidence": "HIGH"|"MED"|"LOW".
        # ONE ROW PER SLIDE — this list must cover every slide in the deck, and the gate
        # enforces that against deck.slide_count. Shown with three rows rather than one for a
        # blunt reason: a single-row example IS a one-slide deck, and a runtime filling this
        # template in copies its SHAPE. Measured: decks arriving at one page when the user named
        # no length. `slide_count: 10` above and a one-row list below contradicted each other,
        # and the example won, because an example is concrete and a number is not.
        "slides": [
            {
                "slide": 1,
                "role": "cover",
                "takeaway": "<one audience takeaway>",
                "evidence": ["README.md:1-10"],
            },
            {
                "slide": 2,
                "role": "problem",
                "takeaway": "<what the room is wrong about, or does not yet see>",
                "evidence": ["README.md:11-24"],
            },
            {
                "slide": 3,
                "role": "evidence",
                "takeaway": "<the claim this page makes, not its topic>",
                "evidence": ["README.md:25-40"],
            },
            # ... one row per remaining slide, through the closing page
        ],
        "claim_ledger": [
            {
                "claim": "<checkable claim>",
                "source": "README.md:1-10",
                "verified": True,
            }
        ],
        # THE ARC COMPETITION. This record bound the DESIGN competition (design.direction, with a
        # hashed directions.html) and recorded nothing about the CONTENT one — backwards, by the
        # skill's own reckoning: a wrong form costs one slide, a wrong arc costs the design plan and
        # the build underneath it. `arc_divergence.py` scores 2-3 candidates over one ledger; what
        # was missing was anywhere for its verdict to land on this path.
        "arc": {
            "chosen": "<the arc that won — must be one of candidates[].name>",
            "shape": "problem-turn-evidence",
            # The CANDIDATES themselves, as ARC OBJECTS — this gate runs arc_divergence.check() over
            # them. Shown as two dict skeletons (never a bare string: `--init` used to emit one
            # descriptive string and the gate then AttributeError'd on str.get). The `<…>` fields
            # are placeholders and the gate REFUSES them until filled — a filled competition is the
            # work, not a formality. `python3 scripts/arc_divergence.py --template` prints the shape.
            "candidates": [
                {"name": "<label-A>", "shape": "problem-turn-evidence",
                 "roles": ["problem", "evidence", "conclusion"],
                 "audience_question": "<the question this room is actually asking>",
                 "objection": "<the objection this arc pre-empts>",
                 "closing_ask": "<what the room should do or believe>",
                 "evidence": ["<claim-ledger id>", "<claim-ledger id>"]},
                {"name": "<label-B>", "shape": "recommendation-first",
                 "roles": ["conclusion", "evidence", "roadmap"],
                 "audience_question": "<a DIFFERENT question>",
                 "objection": "<a DIFFERENT objection>",
                 "closing_ask": "<a DIFFERENT ask>",
                 "evidence": ["<claim-ledger id>", "<claim-ledger id>"]},
            ],
            "rejected": [
                {"name": "<label-B>", "why_lost": "<one clause>"},
            ],
            "divergence_justified": "<only when the recomputed check flags: why this set is right anyway>",
        },
        "checkpoint": {"mode": "approved", "record": "<decision record>"},
    },
    "design": {
        "direction": {
            "branch": "clean",
            "artifact": "directions.html",
            "sha256": "<sha256>",
            "directions": [
                {
                    "id": "A",
                    "name": "A",
                    "bg": "#071820",
                    "accent": "#54D9D0",
                    "font_display": "Inter",
                    "font_body": "Inter",
                    "density": "minimal",
                    "cover": "low-left",
                    "skeleton": "rail",
                }
            ],
            "decision": "user-approved",
            "record": "<selection or auto-carve record>",
        },
        "type_scale": {"display": 34, "title": 24, "body": 14},
        # The governing picture, and the two it beat. `render_deck.py --gate-check` has required
        # this for a while and THIS record never picked it up — a bridged run could therefore
        # satisfy the Codex gate with a design nobody had chosen a concept for, and fail the shared
        # one. Same drift the anchor-proof and boldness comments below already record, caught a
        # third time. `via` carries the derivation rungs (topic -> core concepts -> visual language
        # -> motif); the MIDDLE rung is what separates a derived motif from an industry stereotype.
        "concept": {
            "chosen": "<what this deck's idea is a PICTURE of>",
            "via": "<core concepts> -> <visual language>",
            "rejected": [
                {"concept": "<runner-up>", "why_lost": "<one clause>"},
                {"concept": "<the other>", "why_lost": "<one clause>"},
            ],
        },
        "boldness": "balanced+",
        # Was the build fanned out (one author per section, fresh context each), and if not, why
        # not. Required from ~6 content slides up; "solo — <reason>" is always a legitimate
        # answer ("solo — this runtime has no subagent dispatch" is the normal one on Codex).
        # Mirrors render_deck.py --gate-check exactly — the two gates have drifted on duplicated
        # fields twice before.
        "build_shape": "fanout — <n> sections | solo — <reason>",
        # The resolved FILL-only vs TEXT-safe split. `render_deck.py --gate-check` has required
        # this since a deck shipped a chrome family at 2.4-3.3:1 — a hue that reads fine as a fill
        # measures 2-4:1 as small text on the same tint — and it even ships the hint
        # (`palette_audit.py --from-style <deck>/style.py`). This record never carried it, so the
        # one gate that could have caught that class on the Codex path could not see the palette
        # at all.
        "palette": "<FILL vs TEXT-safe split, per palette_audit.py>",
        # A motif that only RECURS is an ornament with a schedule. Three things it makes besides
        # itself; `page` takes `none - <reason>` because a deck with no page whose geometry the
        # idea could own must not invent one to fill this field.
        "motif_generates": {
            "background": "<what the motif makes the canvas do | flat by register - reason>",
            "markers": "<the numeral / icon / bullet system it implies>",
            "page": "<the slide whose GEOMETRY is the motif | none - reason>",
        },
        "signature_move": "<repeated, deliberate visual device>",
        "carried_by": [1, 5],
        # The ANCHOR PROOF — three rendered pages, three different failures. `signature` proves the
        # aesthetic risk survived the build; `complex` proves the design holds the deck's densest
        # page; `data` proves the charts speak the same visual language the type did.
        "signature_proof": [
            {
                "role": "signature",
                "slide": 1,
                "path": "render/slide-1.png",
                "sha256": "<sha256>",
                "pptx_sha256": "<must equal deck.sha256>",
            },
            {
                "role": "complex",
                "slide": 5,
                "path": "render/slide-5.png",
                "sha256": "<sha256>",
                "pptx_sha256": "<must equal deck.sha256>",
            },
            {
                "role": "data",
                "slide": 8,
                "path": "render/slide-8.png",
                "sha256": "<sha256>",
                "pptx_sha256": "<must equal deck.sha256>",
            },
        ],
        # ONE ROW PER SLIDE, as in content.slides above, and for the same reason: a one-row
        # example teaches a one-slide deck.
        "slides": [
            {
                "slide": 1,
                "function": "slide_01",
                "form": "cover",
                "runner_up": "editorial opener",
                "reason": "<why this form serves the takeaway>",
                "categorical": False,
                "components": [],
            },
            {
                "slide": 2,
                "function": "slide_02",
                "form": "<the form this content's SHAPE wants>",
                "runner_up": "<a form from a DIFFERENT family that it beat>",
                "reason": "<why this form serves the takeaway>",
                "categorical": False,
                "components": [],
            },
            # ... one row per remaining slide
        ],
        "checkpoint": {"mode": "approved", "record": "<decision record>"},
    },
    "build": {
        "script": "build_deck.py",
        "sha256": "<sha256>",
        "strict_layout": True,
    },
    "icons": [
        {
            "slide": 2,
            "family": "lucide",
            "asset": "assets/icons/feature.png",
            "sha256": "<sha256>",
            "rasterizer": "scripts/icons.py",
        }
    ],
    "visual_contract": {
        "manifest": "visual-contract.json",
        "sha256": "<sha256>",
        "result": "visual-contract-final.json",
        "result_sha256": "<sha256>",
        "pptx_sha256": "<must equal deck.sha256>",
    },
    "critics": [
        {
            "lens": "content",
            "review": "critic-content-round2.json",
            "sha256": "<sha256>",
            "pptx_sha256": "<must equal deck.sha256>",
        },
        {
            "lens": "design",
            "review": "critic-design-round2.json",
            "sha256": "<sha256>",
            "pptx_sha256": "<must equal deck.sha256>",
        },
    ],
    "thorough_panel": None,
    "arbiters": [],
    # THE ACTOR'S RENDER LOOK — one verdict per slide, covering every slide (Step 5). `ok` is a
    # valid verdict on a clean page; an empty/placeholder line is a slide nobody read.
    "render_selfcheck": {
        "slides": [
            {"n": 1, "verdict": "<what you saw on slide 1 — 'ok' or the defect + fix>"},
            # ... one verdict per remaining slide
        ],
    },
    "waivers": [],
}


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        value = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return value


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def write_receipt(path: Path, *, evidence_path: Path, build_script: Path, evidence: dict[str, Any]) -> None:
    """Write a final-file-bound receipt only after the strict gate has passed."""
    deck = evidence["deck"]
    receipt = {
        "schema": RECEIPT_SCHEMA,
        "status": "PASS",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "pptx": deck["pptx"],
        "pptx_sha256": deck["sha256"],
        "evidence_sha256": sha256_file(evidence_path),
        "build_script_sha256": sha256_file(build_script),
    }
    path.write_text(json.dumps(receipt, indent=2) + "\n", encoding="utf-8")


def resolve_path(root: Path, value: Any) -> Path | None:
    if not isinstance(value, str) or not value.strip():
        return None
    path = Path(value)
    return path if path.is_absolute() else root / path


def is_sha256(value: Any) -> bool:
    if not isinstance(value, str) or len(value) != 64:
        return False
    return all(char in "0123456789abcdef" for char in value.lower())


def require_string(value: Any, label: str, errors: list[str], minimum: int = 1) -> str | None:
    # Width, not codepoints — `written_reason.reason_width` counts an East-Asian wide character
    # as 2, so these floors mean "roughly this much information" rather than "roughly this much
    # Latin". Exactly a no-op for ASCII; the only records whose behaviour changes are CJK ones,
    # which were being refused for saying MORE in fewer characters.
    if not isinstance(value, str) or reason_width(value) < minimum:
        errors.append(f"{label} must be a non-empty string")
        return None
    return value.strip()


def check_hashed_file(
    root: Path,
    path_value: Any,
    expected_hash: Any,
    label: str,
    errors: list[str],
    minimum_bytes: int = 1,
) -> Path | None:
    path = resolve_path(root, path_value)
    if path is None:
        errors.append(f"{label}.path must be a non-empty path")
        return None
    if not path.is_file():
        errors.append(f"{label} file not found: {path}")
        return None
    if path.stat().st_size < minimum_bytes:
        errors.append(f"{label} is too small to be usable: {path}")
        return None
    if not is_sha256(expected_hash):
        errors.append(f"{label}.sha256 must be a 64-character SHA-256")
        return path
    actual_hash = sha256_file(path)
    if actual_hash != expected_hash:
        errors.append(f"{label} SHA-256 does not match: {path}")
    return path


def png_dimensions(path: Path) -> tuple[int, int] | None:
    info = png_info(path)
    return info[:2] if info is not None else None


def png_info(path: Path) -> tuple[int, int, int, int] | None:
    """Return PNG width, height, bit depth, and colour type without a Pillow dependency."""
    try:
        with path.open("rb") as handle:
            header = handle.read(26)
    except OSError:
        return None
    if len(header) != 26 or header[:8] != b"\x89PNG\r\n\x1a\n" or header[12:16] != b"IHDR":
        return None
    width, height = struct.unpack(">II", header[16:24])
    return width, height, header[24], header[25]


def slide_count_from_pptx(path: Path) -> int | None:
    try:
        from pptx import Presentation

        return len(Presentation(str(path)).slides)
    except Exception:
        return None


# Mirror of render_deck._ICON_NONE_CATEGORIES — the four high-bar reasons a categorical deck may
# carry no icon family, kept identical across the two gate paths so Codex and the shared path hold
# the same bar (they have drifted on icons before — this is why `categorical_slides` is one shared
# definition). Icons are the default on categorical content; skipping them names one of these.
ICON_NONE_CATEGORIES = ("motif-dominant", "editorial-register", "tiny-deck", "template-locked")


def _icon_waiver_ok(evidence: dict[str, Any]) -> bool:
    """True when the undeclared-categorical icon waiver carries a CLASSIFIED high-bar category.

    Symmetric with render_deck._icon_none_waived: naming a reason is no longer enough — the reason
    must classify from ICON_NONE_CATEGORIES, so a casual 'not category-rich' cannot clear a deck
    that reads as categorical and shipped zero icons.
    """
    for entry in evidence.get("waivers", []) or []:
        if not isinstance(entry, dict) or entry.get("kind") != "icon":
            continue
        if entry.get("scope") != "undeclared-categorical":
            continue
        if not require_string(entry.get("reason"), "waiver.reason", [], minimum=12):
            continue
        if str(entry.get("category") or "").strip().lower() in ICON_NONE_CATEGORIES:
            return True
    return False


def waived(evidence: dict[str, Any], kind: str, **matches: Any) -> bool:
    waivers = evidence.get("waivers", [])
    if not isinstance(waivers, list):
        return False
    for entry in waivers:
        if not isinstance(entry, dict) or entry.get("kind") != kind:
            continue
        if not require_string(entry.get("reason"), "waiver.reason", [], minimum=12):
            continue
        if all(entry.get(key) == value for key, value in matches.items()):
            return True
    return False


def parse_script(script_path: Path, errors: list[str]) -> tuple[dict[str, set[str]], bool]:
    try:
        tree = ast.parse(script_path.read_text(encoding="utf-8"), filename=str(script_path))
    except (OSError, SyntaxError) as exc:
        errors.append(f"cannot parse build script {script_path}: {exc}")
        return {}, False

    module_aliases: set[str] = {"deckkit"}
    direct_aliases: dict[str, str] = {}
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for item in node.names:
                if item.name == "deckkit":
                    module_aliases.add(item.asname or "deckkit")
        elif isinstance(node, ast.ImportFrom) and node.module == "deckkit":
            for item in node.names:
                direct_aliases[item.asname or item.name] = item.name

    def call_name(node: ast.expr) -> str | None:
        if isinstance(node, ast.Name):
            return direct_aliases.get(node.id, node.id)
        if isinstance(node, ast.Attribute) and isinstance(node.value, ast.Name):
            if node.value.id in module_aliases:
                return node.attr
        return None

    calls: dict[str, set[str]] = defaultdict(set)
    strict_layout = False

    class Visitor(ast.NodeVisitor):
        current_function: str | None = None

        def visit_FunctionDef(self, node: ast.FunctionDef) -> None:
            previous = self.current_function
            self.current_function = node.name
            self.generic_visit(node)
            self.current_function = previous

        visit_AsyncFunctionDef = visit_FunctionDef

        def visit_Call(self, node: ast.Call) -> None:
            nonlocal strict_layout
            name = call_name(node.func)
            if self.current_function and name:
                calls[self.current_function].add(name)
            if name == "lint_layout":
                strict_layout = any(
                    keyword.arg == "strict"
                    and isinstance(keyword.value, ast.Constant)
                    and keyword.value.value is True
                    for keyword in node.keywords
                ) or strict_layout
            self.generic_visit(node)

    Visitor().visit(tree)
    return dict(calls), strict_layout


def forbidden_icon_rasterizer_calls(script_path: Path, errors: list[str]) -> list[str]:
    """Find actual command invocations, not comments, that make preview thumbnails into icons."""
    try:
        tree = ast.parse(script_path.read_text(encoding="utf-8"), filename=str(script_path))
    except (OSError, SyntaxError) as exc:
        errors.append(f"cannot inspect build script for icon rasterization: {exc}")
        return []

    command_methods = {"run", "Popen", "call", "check_call", "check_output", "system", "popen"}
    blocked: list[str] = []
    for node in ast.walk(tree):
        if not isinstance(node, ast.Call):
            continue
        if isinstance(node.func, ast.Name):
            command_call = node.func.id in command_methods
        elif isinstance(node.func, ast.Attribute):
            command_call = node.func.attr in command_methods
        else:
            command_call = False
        if not command_call:
            continue
        strings = [
            value.value.lower()
            for value in ast.walk(node)
            if isinstance(value, ast.Constant) and isinstance(value.value, str)
        ]
        if any("qlmanage" in value for value in strings):
            blocked.append("qlmanage")
    return sorted(set(blocked))


def check_lint(lint: dict[str, Any], delivery: str, evidence: dict[str, Any], errors: list[str]) -> None:
    findings = lint.get("findings", [])
    if not isinstance(findings, list):
        errors.append("lint findings missing or malformed")
    else:
        blocking = [finding for finding in findings if isinstance(finding, dict) and finding.get("severity") == "error"]
        if blocking:
            errors.append(f"layout lint reports {len(blocking)} error(s)")

    # lint_deck.py's native JSON records pixel execution as a compact object.  Accept it directly
    # rather than forcing a hand-authored adapter that can accidentally misstate the check.
    pixels = lint.get("pixel_checks", [])
    if isinstance(pixels, dict):
        if pixels.get("ran") is not True or pixels.get("not_checked"):
            errors.append("pixel checks were skipped or incomplete in native lint output")
        pixels = [{"pass": True}]
    if not isinstance(pixels, list):
        errors.append("pixel_checks missing or malformed")
    else:
        failed = [item for item in pixels if isinstance(item, dict) and not item.get("pass", False)]
        if failed:
            errors.append(f"pixel checks report {len(failed)} failure(s)")

    floor = BODY_FLOORS.get(delivery, BODY_FLOORS["presented"])
    text_runs = lint.get("text_runs")
    if text_runs is None and isinstance(lint.get("deck"), dict):
        # Native lint reports a measured body median rather than role-by-role samples.  This is
        # conservative for the delivery floor and avoids treating generated source-note metadata
        # as body prose.
        median = lint["deck"].get("body_median_pt")
        text_runs = [{"role": "body", "size_pt": median, "exception": False}]
    if not isinstance(text_runs, list):
        errors.append("text_runs missing or malformed")
    else:
        undersized = [
            run
            for run in text_runs
            if isinstance(run, dict)
            and run.get("role") in {"body", "label", "footer"}
            and isinstance(run.get("size_pt"), (int, float))
            and run["size_pt"] < floor
            and not run.get("exception")
        ]
        if undersized:
            errors.append(f"{len(undersized)} body/label/footer run(s) are below {floor:g} pt")

    stats = lint.get("stats", {})
    warnings = stats.get("warnings", []) if isinstance(stats, dict) else []
    if not warnings and isinstance(lint.get("stats_warnings"), list):
        warnings = [
            "_".join(str(row).split(":", 1)[0].strip().lower().split())
            for row in lint["stats_warnings"]
        ]
    if not isinstance(warnings, list):
        errors.append("stats warnings missing or malformed")
    else:
        flagged = [
            warning
            for warning in warnings
            if warning in STRICT_STATS and not waived(evidence, "stats", warning=warning)
        ]
        if flagged:
            errors.append("stats warnings require remediation or an explicit waiver: " + ", ".join(flagged))

        # THE SAMENESS COMPOSITE, ported from the shared gate as a COMPOSITE — not as seven more
        # STRICT_STATS rows. That distinction is the whole calibration: `lint_deck` counts DISTINCT
        # monotony signals and blocks at >=4 with >=1 structural, precisely because any one of them
        # is legitimate on its own (a 小红书 carousel and a short status update each trip three
        # honestly). Adding them here per-warning would make this gate refuse decks the shared path
        # correctly ships — stricter, but wrong, which is worse than drifted.
        #
        # Measured before this: STRICT_STATS covered 2 of the 7 (card_dominance,
        # envelope_monoculture), so a Codex deck could fire LAYOUT SAMENESS + SKELETON VARIETY +
        # BOTTOM-STRIP MONOCULTURE + TITLE-RULE MONOCULTURE — four signals, a hard block on the
        # shared path — and hear nothing at all here.
        #
        # The code list is IMPORTED, never copied: this file already carries two comments about the
        # two gates drifting on a duplicated constant (`path` vs `png`, and the missing
        # `conservative` dial). One source, one rule.
        try:
            sys.path.insert(0, str(Path(__file__).resolve().parent))
            import lint_deck as _ld
            def _slug(codes):
                return {c.lower().replace(" ", "_").replace("-", "_") for c in codes}
            sameness_codes = _slug(_ld.SAMENESS_CODES)
            # IMPORTED, not retyped. This set used to be a hand-written literal sitting directly
            # under the comment above warning that copied constants drift — and it would have
            # silently disagreed the moment SAMENESS_STRUCTURAL changed.
            structural = _slug(_ld.SAMENESS_STRUCTURAL)
            timid_codes = _slug(_ld.TIMIDITY_CODES)
            timid_structural = _slug(_ld.TIMIDITY_STRUCTURAL)
        except Exception:                                     # pragma: no cover - import guard
            sameness_codes, structural = set(), set()
            timid_codes, timid_structural = set(), set()
        if sameness_codes:
            fired = {w for w in warnings if w in sameness_codes}
            if len(fired) >= 4 and (fired & structural) and not waived(evidence, "sameness"):
                errors.append(
                    "sameness: {} distinct deck-level monotony signals fired ({}) — the deck reads "
                    "as one template even where its forms vary. Redesign the repetition, or record "
                    '{{"kind": "sameness", "reason": "<why this deck repeats on purpose — name the '
                    'register>"}} in waivers.'.format(len(fired), ", ".join(sorted(fired))))
        # THE COUNTERWEIGHT, ported symmetrically. Without it the Codex profile would keep only the
        # half of the scale that punishes excess — the exact asymmetry that let a shared-path deck
        # be iterated flatter ten times with every gate reporting clean.
        if timid_codes:
            fired_t = {w for w in warnings if w in timid_codes}
            if len(fired_t) >= 2 and (fired_t & timid_structural) and not waived(evidence, "timidity"):
                errors.append(
                    "timidity: {} signals say this deck is measurably SAFE ({}) — no page carries a "
                    "protagonist that is not a sentence. Give the load-bearing pages a real one (a "
                    "figure, a chart, a form whose geometry IS the argument), or record "
                    '{{"kind": "timidity", "reason": "<the register that makes this restraint '
                    'deliberate>"}} in waivers.'.format(len(fired_t), ", ".join(sorted(fired_t))))

    # the accessibility floors on the per-slide `warnings` stream (see STRICT_WARNINGS)
    per_slide = lint.get("warnings", [])
    if isinstance(per_slide, list):
        hit: dict[str, set[int]] = {}
        for row in per_slide:
            text = row.get("text", "") if isinstance(row, dict) else str(row)
            slide = row.get("slide") if isinstance(row, dict) else None
            for code in STRICT_WARNINGS:
                if text.startswith(code + ":") and not waived(evidence, "a11y", warning=code):
                    hit.setdefault(code, set())
                    if isinstance(slide, int):
                        hit[code].add(slide)
        for code in sorted(hit):
            where = (" (slide%s %s)" % ("s" if len(hit[code]) > 1 else "",
                                        ", ".join(str(n) for n in sorted(hit[code])))
                     if hit[code] else "")
            errors.append(
                f"{code}{where} is below the WCAG 1.4.11 3:1 floor — remediate it, or record a "
                f"waiver {{\"kind\": \"a11y\", \"warning\": \"{code}\", \"reason\": \"…\"}} saying "
                f"why this mark is decorative")


def check_render_selfcheck(
    evidence: dict[str, Any], expected_slides: set[int], errors: list[str]
) -> None:
    """The actor's own render look, turned from prose into a trace — mirror of render_deck's gate.

    Step 5 asks the coordinator to read every slide PNG and record a one-line verdict per slide; that
    was prose with no backstop, so the cheap actor-side look (catching an overflow / cropped subject
    / wrong number BEFORE a critic round is spent) was the easiest step to skip silently. This makes
    it leave a mark: one non-placeholder verdict per slide, covering every slide. Same shape and
    honest limit as content.slides — it proves the trace exists, not that the eye judged well; the
    strong per-slide guarantee is the independent critic's coverage bind (check_critics).
    """
    rs = evidence.get("render_selfcheck")
    if isinstance(rs, dict) and require_string(rs.get("waived"), "render_selfcheck.waived",
                                               [], minimum=16):
        if has_placeholder(rs.get("waived")):
            errors.append("render_selfcheck.waived is still a <placeholder>")
        return
    slides = (rs or {}).get("slides") if isinstance(rs, dict) else None
    if not isinstance(slides, list) or not slides:
        errors.append("render_selfcheck.slides missing — Step 5's per-slide verdict, one line per "
                      "slide (a slide with no line was not looked at); or waive it in writing")
        return
    seen: dict[int, int] = {}
    for index, row in enumerate(slides):
        if not isinstance(row, dict):
            errors.append(f"render_selfcheck.slides[{index}] must be an object with n / verdict")
            continue
        n = row.get("n")
        if not isinstance(n, int) or isinstance(n, bool):
            errors.append(f"render_selfcheck.slides[{index}].n must be the slide number")
            continue
        if n in seen:
            errors.append(f"render_selfcheck.slides: two verdicts both claim slide {n}")
        seen[n] = index
        require_string(row.get("verdict"), f"render_selfcheck slide {n}.verdict", errors, minimum=4)
        if has_placeholder(row.get("verdict")):
            errors.append(f"render_selfcheck slide {n}.verdict is still a <placeholder>")
    if expected_slides and set(seen) != expected_slides:
        errors.append("render_selfcheck.slides must carry one verdict for every slide "
                      "(a slide with no verdict is a slide nobody looked at)")


def check_content(
    evidence: dict[str, Any], root: Path, expected_slides: set[int], errors: list[str]
) -> None:
    interview = evidence.get("interview")
    if not isinstance(interview, dict):
        errors.append("interview evidence missing")
    else:
        if interview.get("mode") not in {"answered", "auto"}:
            errors.append("interview.mode must be answered or auto")
        require_string(interview.get("record"), "interview.record", errors, minimum=12)
        # WHERE THE LENGTH ANSWER LANDS. `interview.record` is free text with a 12-char floor, so
        # it cannot tell an answered interview from a stub, and deck length is the axis that
        # actually goes missing on a runtime with no choice UI — nothing downstream demands it, so
        # nothing notices. Measured: decks arriving at ONE page when the user named no length.
        #
        # Deliberately not a number to validate against slide_count: the honest answers include a
        # range ("medium, 9-15"), a time budget ("20 min, so ~18"), and "user declined — derived 11
        # from the ledger's takeaway count". What is being checked is that the question was PUT and
        # its answer recorded, not that a particular integer was hit.
        require_string(interview.get("length"), "interview.length", errors, minimum=4)

    content = evidence.get("content")
    if not isinstance(content, dict):
        errors.append("content evidence missing")
        return
    source_mode = content.get("source_mode")
    if source_mode not in {"provided", "web", "none"}:
        errors.append("content.source_mode must be provided, web, or none")
        source_mode = "provided"
    sources = content.get("sources")
    if not isinstance(sources, list) or (source_mode != "none" and not sources):
        errors.append("content.sources must document the supplied source material")
    elif isinstance(sources, list):
        for index, source in enumerate(sources, start=1):
            label = f"content.sources[{index}]"
            if not isinstance(source, dict):
                errors.append(f"{label} must be an object")
                continue
            kind = source.get("kind")
            if kind == "provided":
                check_hashed_file(root, source.get("path"), source.get("sha256"), label, errors)
            elif kind == "web":
                require_string(source.get("locator"), f"{label}.locator", errors, minimum=12)
            else:
                errors.append(f"{label}.kind must be provided or web")

    slides = content.get("slides")
    if not isinstance(slides, list):
        errors.append("content.slides missing or malformed")
    else:
        rows = {row.get("slide"): row for row in slides if isinstance(row, dict) and isinstance(row.get("slide"), int)}
        if set(rows) != expected_slides:
            errors.append("content.slides must cover every final slide exactly once")
        # Two content slides carrying ONE memory sentence is the mechanical signature of a plan
        # written for the DECK and pasted down the column — and one of the two pages then has no
        # reason to exist. Structural rows are exempt: a cover and a closing legitimately restate
        # the deck's one sentence, and forbidding that would make authors invent a second.
        seen_takeaways: dict[str, int] = {}
        for number, row in rows.items():
            require_string(row.get("role"), f"content slide {number}.role", errors)
            require_string(row.get("takeaway"), f"content slide {number}.takeaway", errors, minimum=8)
            if has_placeholder(row.get("takeaway")):
                errors.append(f"content slide {number}.takeaway is still the `<placeholder>` from "
                              "the skeleton — the memory sentence has not been written")
            evidence_rows = row.get("evidence")
            if not isinstance(evidence_rows, list) or not all(isinstance(item, str) and item.strip() for item in evidence_rows):
                errors.append(f"content slide {number}.evidence must contain source references")
            elif any(has_placeholder(item) for item in evidence_rows):
                errors.append(f"content slide {number}.evidence is still a `<placeholder>` source trace")
            units = row.get("units")
            if units is not None and (not isinstance(units, int) or isinstance(units, bool) or units < 0):
                errors.append(f"content slide {number}.units must be a non-negative count")
            if str(row.get("role") or "").strip().lower() not in STRUCTURAL_ROLES:
                key = "".join(str(row.get("takeaway") or "").strip().lower().split())
                key = key.rstrip("。．.!！?？,，、;；")
                if key and key in seen_takeaways:
                    errors.append(f"content slides {seen_takeaways[key]} and {number} carry the "
                                  "SAME takeaway — give each its own memory sentence, or merge them")
                elif key:
                    seen_takeaways[key] = number

    # Mirrors what this gate already demands of `design.direction`: a competition is only a
    # competition if the losers are on the record. `picked contribution-first` alone is a sentence
    # the coordinator can write without any alternative having existed — the same reason the
    # content checkpoint's `arc gate:` line requires the losers and their clauses. Kept to ONE
    # rejected arc minimum, because arc_divergence.py accepts 2 candidates as a valid set.
    arc = content.get("arc")
    if not isinstance(arc, dict):
        errors.append("content.arc missing — the arc competition (agents/content-planner.md §3): "
                      "the arc that won, the ones it beat, and the divergence verdict")
    else:
        chosen = require_string(arc.get("chosen"), "content.arc.chosen", errors, minimum=4)
        candidates = arc.get("candidates")
        verdict = None
        if not isinstance(candidates, list) or len(candidates) < 2:
            errors.append("content.arc.candidates must carry the 2-3 candidate arcs THEMSELVES — "
                          "this gate scores them with arc_divergence.check(), so a pasted "
                          "`divergence` verdict is not evidence (a delivered deck passed with "
                          '"ok" while the script had never run for it). '
                          "`python3 scripts/arc_divergence.py --template` prints the shape")
        else:
            # Refuse the raw `--template` skeleton BEFORE scoring: its fields are angle-bracket
            # placeholders, which are wide enough to clear every width floor, so the unedited
            # skeleton passed the shared gate. `has_placeholder` is the same predicate this file
            # already applies to fast_basis/none_opt_in, tightened to need a closing `>`.
            ph = []
            for ci, cand in enumerate(candidates):
                if not isinstance(cand, dict):
                    ph.append(f"candidates[{ci}] is {type(cand).__name__}, not an arc object")
                    continue
                for field in ("name", "audience_question", "objection", "closing_ask"):
                    if has_placeholder(cand.get(field)):
                        ph.append(f"candidates[{ci}].{field}")
            if has_placeholder(arc.get("chosen")):
                ph.append("chosen")
            if ph:
                errors.append("content.arc still carries the raw --template placeholder(s): "
                              + ", ".join(ph[:8]) + " — the unedited skeleton clears every width "
                              "floor, so it would ship an arc competition that never happened")
            elif all(isinstance(c, dict) for c in candidates):
                try:
                    verdict = load_arc_checker().check(candidates)
                except Exception as exc:  # noqa: BLE001
                    errors.append(f"content.arc.candidates is not a readable candidate set: {exc} "
                                  "(same reader as the CLI, so arc_divergence.py <arcs>.json "
                                  "reproduces it)")
        # `if isinstance(c, dict)` — a candidate list of raw strings (the malformed skeleton) must
        # report a clean error, never crash on `str.get`. This line ran unconditionally before and
        # AttributeError'd on the exact file `--init` writes.
        names = [str(c.get("name") or "?") for c in candidates
                 if isinstance(c, dict)] if isinstance(candidates, list) else []
        if chosen and names and chosen not in names:
            errors.append(f"content.arc.chosen {chosen!r} is not one of the candidates "
                          f"({', '.join(names)}) — a winner from outside the field means the "
                          "competition scored a set the deck was not built from")
        rejected = arc.get("rejected")
        named: dict[str, Any] = {}
        if not isinstance(rejected, list) or not rejected:
            errors.append("content.arc.rejected must name every arc the winner beat — "
                          "a winner with no losers on the record is a derivation, not a choice")
        else:
            for index, row in enumerate(rejected):
                label = f"content.arc.rejected[{index}]"
                if not isinstance(row, dict):
                    errors.append(f"{label} must be an object")
                    continue
                nm = require_string(row.get("name"), f"{label}.name", errors, minimum=2)
                require_string(row.get("why_lost"), f"{label}.why_lost", errors, minimum=8)
                if nm:
                    named[nm] = row.get("why_lost")
        if names and chosen:
            missing = [x for x in names if x != chosen and x not in named]
            stray = [x for x in named if x not in names]
            if missing:
                errors.append("content.arc.rejected skips " + ", ".join(repr(m) for m in missing)
                              + " — every candidate that is not the winner needs its losing "
                                "clause; recording only the flattering half of the field is how a "
                                "competition quietly becomes a derivation again")
            if stray:
                errors.append("content.arc.rejected names " + ", ".join(repr(m) for m in stray)
                              + " — never in `candidates`, so it never competed")
        if verdict is not None:
            flagged = verdict.get("flagged") or []
            sketches = verdict.get("sketches") or []
            if flagged or sketches or verdict.get("no_ledger"):
                finding = "; ".join(filter(None, [
                    "{} pair(s) tell the same story".format(len(flagged)) if flagged else "",
                    "strawman candidate(s): " + ", ".join(x["name"] for x in sketches)
                    if sketches else "",
                    "no candidate names its ledger evidence" if verdict.get("no_ledger") else ""]))
                if reason_width(arc.get("divergence_justified")) < 16:
                    errors.append(f"the arc competition, recomputed here, does not hold: {finding}. "
                                  "Rediverge the candidates (move the CLAIM, not the wording) or "
                                  "record why the set is right anyway in "
                                  "content.arc.divergence_justified")

    ledger = content.get("claim_ledger")
    if not isinstance(ledger, list) or (source_mode != "none" and not ledger):
        errors.append("content.claim_ledger must contain verified claims")
    elif isinstance(ledger, list):
        for index, claim in enumerate(ledger, start=1):
            label = f"content.claim_ledger[{index}]"
            if not isinstance(claim, dict):
                errors.append(f"{label} must be an object")
                continue
            require_string(claim.get("claim"), f"{label}.claim", errors, minimum=8)
            require_string(claim.get("source"), f"{label}.source", errors, minimum=3)
            if claim.get("verified") is not True:
                errors.append(f"{label}.verified must be true")

    # A WEB-RESEARCHED deck ships on three floors (content-planner.md §2(e), mirrored in the shared
    # content checkpoint's coverage:/lifecycle:/provenance: lines): 全面 COMPREHENSIVE (a coverage map
    # + a proactive LIFECYCLE sweep so a discontinued/renamed product is never headlined), 充实
    # SUBSTANTIAL (concrete specifics), 准确 ACCURATE (each fact confidence-tagged + corroborated).
    # The Codex evidence carries them as structured fields so this gate checks what the shared
    # checkpoint states. Scoped to source_mode == "web": a 'provided' deck traces to its material and
    # 'none' is a stub. This closes the same gap the shared path had — a no-source deck that shipped
    # thin and headlined two discontinued products because none of the three was recorded anywhere.
    if source_mode == "web":
        require_string(content.get("coverage"),
                       "content.coverage — the domain areas enumerated · covered · cut (§2e 全面)",
                       errors, minimum=12)
        require_string(content.get("lifecycle"),
                       "content.lifecycle — every featured product/version/entity checked "
                       "live-vs-discontinued as of today (§2e; a headlined dead/renamed thing is a defect)",
                       errors, minimum=12)
        provenance = content.get("provenance")
        if not isinstance(provenance, dict):
            errors.append("content.provenance missing — the checked/confirmed/fixed/cut digest (§2e 准确)")
        else:
            require_string(provenance.get("summary"), "content.provenance.summary", errors, minimum=8)
        if isinstance(ledger, list):
            for index, claim in enumerate(ledger, start=1):
                if isinstance(claim, dict) and claim.get("confidence") not in {"HIGH", "MED", "LOW"}:
                    errors.append(f"content.claim_ledger[{index}].confidence must be HIGH, MED, or "
                                  "LOW so LOW/UNVERIFIED facts are visibly cut (§2e 准确)")

    checkpoint = content.get("checkpoint")
    if not isinstance(checkpoint, dict):
        errors.append("content.checkpoint missing")
    else:
        if checkpoint.get("mode") not in {"approved", "auto"}:
            errors.append("content.checkpoint.mode must be approved or auto")
        require_string(checkpoint.get("record"), "content.checkpoint.record", errors, minimum=12)
        if has_placeholder(checkpoint.get("record")):
            errors.append("content.checkpoint.record is still the `<placeholder>` from the skeleton")


def check_design(
    evidence: dict[str, Any],
    root: Path,
    expected_slides: set[int],
    deck_hash: str,
    errors: list[str],
) -> dict[int, dict[str, Any]]:
    design = evidence.get("design")
    if not isinstance(design, dict):
        errors.append("design evidence missing")
        return {}
    direction = design.get("direction")
    if not isinstance(direction, dict):
        errors.append("design.direction missing")
    else:
        branch = direction.get("branch")
        if branch not in {"clean", "provided-template", "generated-template", "mimic"}:
            errors.append("design.direction.branch is invalid")
        artifact_path = check_hashed_file(
            root,
            direction.get("artifact"),
            direction.get("sha256"),
            "design.direction.artifact",
            errors,
            minimum_bytes=512,
        )
        if artifact_path is not None and artifact_path.suffix.lower() in {".html", ".htm"}:
            markup = artifact_path.read_text(encoding="utf-8", errors="ignore").lower()
            if "<html" not in markup or "<body" not in markup:
                errors.append("design.direction.artifact must be a real HTML preview, not a placeholder file")
        choices = direction.get("directions")
        if not isinstance(choices, list) or not choices:
            errors.append("design.direction.directions missing")
        elif branch == "clean":
            names = [choice.get("name") for choice in choices if isinstance(choice, dict)]
            if len(choices) < 4 or len(set(names)) < 4 or any(not isinstance(name, str) or not name.strip() for name in names):
                errors.append("clean design direction needs four named preview directions")
            else:
                try:
                    validator = load_direction_validator()
                    result = validator.check(choices)
                    if result.get("flagged") and not require_string(
                        direction.get("diversity_waiver"), "design.direction.diversity_waiver", errors, minimum=12
                    ):
                        errors.append("direction preview has too-similar candidates without a named diversity waiver")
                except Exception as exc:
                    errors.append(f"design.direction candidates cannot pass the diversity check: {exc}")
            if artifact_path is not None and artifact_path.suffix.lower() in {".html", ".htm"}:
                markup = artifact_path.read_text(encoding="utf-8", errors="ignore")
                missing_names = [name for name in names if isinstance(name, str) and name not in markup]
                if missing_names:
                    errors.append("direction preview does not visibly include: " + ", ".join(missing_names))
        if direction.get("decision") not in {"user-approved", "auto-carve", "provided-template", "mode-a-mimic"}:
            errors.append("design.direction.decision is invalid")
        require_string(direction.get("record"), "design.direction.record", errors, minimum=12)

    scale = design.get("type_scale")
    if not isinstance(scale, dict):
        errors.append("design.type_scale missing")
    else:
        for key in ("display", "title", "body"):
            if not isinstance(scale.get(key), (int, float)):
                errors.append(f"design.type_scale.{key} must be numeric")
        if isinstance(scale.get("body"), (int, float)) and scale["body"] < BODY_FLOORS[evidence.get("delivery", "presented")]:
            errors.append("design.type_scale.body is below the delivery floor")

    # The documented dial is <conservative | balanced+ | bold | experimental> (SKILL.md,
    # agents/slide-design.md, review-rubrics.md). This set had `conservative` MISSING and
    # `deliberately-restrained` in its place — but that string is a value for the signature_move
    # FIELD ("deliberately restrained: <why>", the conservative dial's documented escape), never a
    # dial. So the one word an author is actually told to write failed this gate, and a field value
    # passed as a dial. The legacy string stays accepted so existing evidence files still load.
    if design.get("boldness") not in {"conservative", "balanced+", "bold", "experimental",
                                      "deliberately-restrained"}:
        errors.append("design.boldness must declare a supported direction "
                      "(conservative | balanced+ | bold | experimental)")
    # Mirrors render_deck.py --gate-check's `concept` contract EXACTLY — one picture with no
    # alternatives is not a choice, it is the first thing that came to mind. Kept deliberately at
    # the shared path's strictness (chosen + two rejected, each with the clause that lost it): the
    # `via` rungs are checked here only when present, because making them blocking on ONE path is
    # how the two gates drift, and this file already carries two comments about that costing the
    # repo. If the rungs become blocking, both paths change in the same commit.
    concept = design.get("concept")
    if not isinstance(concept, dict):
        errors.append("design.concept missing — name the governing picture and the two it beat")
    else:
        require_string(concept.get("chosen"), "design.concept.chosen", errors, minimum=8)
        rejected = concept.get("rejected")
        if not isinstance(rejected, list) or len(rejected) < 2:
            errors.append("design.concept.rejected must name TWO pictures the winner beat")
        else:
            for index, row in enumerate(rejected[:2]):
                if not isinstance(row, dict):
                    errors.append(f"design.concept.rejected[{index}] must be an object")
                    continue
                require_string(row.get("concept"), f"design.concept.rejected[{index}].concept", errors)
                require_string(row.get("why_lost"), f"design.concept.rejected[{index}].why_lost",
                               errors, minimum=8)

    require_string(design.get("palette"), "design.palette", errors, minimum=12)
    # Same threshold and same rule as render_deck.py --gate-check: from ~6 content slides the
    # build-shape decision must be recorded. Never blocks the CHOICE — solo is mandatory on a
    # runtime with no subagent dispatch — only the absence of a decision.
    if len(expected_slides) >= 7 and not str(design.get("build_shape", "")).strip():
        errors.append('design.build_shape missing on a {}-slide deck — "fanout — <n> sections" or '
                      '"solo — <reason>" ("solo — no subagent dispatch on this runtime" is the '
                      'normal Codex answer); the build step is 40-71% of a session, so the '
                      'decision is recorded, not assumed'.format(len(expected_slides)))
    # Same carve as signature_proof, and for the same reason: under a conservative dial with a
    # recorded "deliberately restrained" move there is no loud motif to be productive, and
    # demanding three products would push an author to invent a device so the field has an answer.
    _restrained = (str(design.get("boldness", "")).strip().lower() in {"conservative",
                                                                      "deliberately-restrained"}
                   and str(design.get("signature_move", "")).strip().lower()
                       .startswith("deliberately restrained"))
    generates = design.get("motif_generates")
    if not _restrained:
        if not isinstance(generates, dict):
            errors.append("design.motif_generates missing — name the three things the motif makes "
                          "besides itself (background, markers, page); a motif that only recurs is "
                          "an ornament with a schedule")
        else:
            for key in ("background", "markers", "page"):
                require_string(generates.get(key), f"design.motif_generates.{key}", errors, minimum=4)
    require_string(design.get("signature_move"), "design.signature_move", errors, minimum=12)
    carried_by = design.get("carried_by")
    if not isinstance(carried_by, list) or len(set(carried_by) & expected_slides) < min(2, len(expected_slides)):
        errors.append("design.carried_by must show the signature move on at least two slides")

    # Mirrors render_deck.py's carve with the SAME condition — the two gate paths disagreeing about
    # what an honest plan looks like has already cost this repo once (this key was spelled `path`
    # here and `png` there, so a bridged run wrote the field its own gate demanded and the other
    # rejected it).
    carved = (str(design.get("boldness", "")).strip().lower() == "conservative"
              and str(design.get("signature_move", "")).strip().lower()
                  .startswith("deliberately restrained"))
    # THE ANCHOR-PROOF CONTRACT LIVES IN scripts/anchor_proof.py, imported by BOTH gate paths.
    # This gate and render_deck.py --gate-check have already diverged once on this exact field (the
    # file key was spelled `path` here and `png` there, so a bridged run wrote what its own gate
    # demanded and the other rejected it). Shared module, one rule, no drift. What stays local is
    # the strictness: the Codex path binds every anchor PNG to a SHA-256 AND to the final PPTX
    # hash, which the shared path does not, and folding that into the shared module would have to
    # weaken it to the weaker of the two.
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    import anchor_proof as _ap

    proof = design.get("signature_proof")
    if carved and proof is None:
        pass                          # a conservative deck that declared it took no risk
    elif proof is None:
        errors.append("design.signature_proof missing")
    else:
        for line in _ap.faults(proof, n_slides=len(expected_slides),
                               expected_slides=expected_slides, carved=carved):
            errors.append("design.signature_proof: " + line)
        for index, anchor in enumerate(_ap.normalise(proof) or []):
            label = "design.signature_proof[%d] (%s)" % (index, anchor.get("role"))
            proof_path = check_hashed_file(
                root,
                _ap.anchor_file(anchor),
                anchor.get("sha256"),
                label,
                errors,
                minimum_bytes=512,
            )
            if proof_path is not None:
                dimensions = png_dimensions(proof_path)
                if dimensions is None or dimensions[0] < 640 or dimensions[1] < 360:
                    errors.append(label + " must be a rendered PNG of at least 640x360")
            if anchor.get("pptx_sha256") != deck_hash:
                errors.append(label + " is not bound to the final PPTX SHA-256")

    rows = design.get("slides")
    if not isinstance(rows, list):
        errors.append("design.slides missing or malformed")
        row_map: dict[int, dict[str, Any]] = {}
    else:
        row_map = {row.get("slide"): row for row in rows if isinstance(row, dict) and isinstance(row.get("slide"), int)}
        if set(row_map) != expected_slides:
            errors.append("design.slides must cover every final slide exactly once")
        for number, row in row_map.items():
            require_string(row.get("function"), f"design slide {number}.function", errors)
            form = require_string(row.get("form"), f"design slide {number}.form", errors)
            runner_up = require_string(row.get("runner_up"), f"design slide {number}.runner_up", errors)
            if form and runner_up and form.lower() == runner_up.lower():
                errors.append(f"design slide {number} must keep a distinct runner-up form")
            require_string(row.get("reason"), f"design slide {number}.reason", errors, minimum=12)
            if not isinstance(row.get("categorical"), bool):
                errors.append(f"design slide {number}.categorical must be boolean")
            components = row.get("components")
            if not isinstance(components, list) or not all(isinstance(component, str) and component for component in components):
                errors.append(f"design slide {number}.components must be a list of component names")

    checkpoint = design.get("checkpoint")
    if not isinstance(checkpoint, dict):
        errors.append("design.checkpoint missing")
    else:
        if checkpoint.get("mode") not in {"approved", "auto"}:
            errors.append("design.checkpoint.mode must be approved or auto")
        require_string(checkpoint.get("record"), "design.checkpoint.record", errors, minimum=12)
        if has_placeholder(checkpoint.get("record")):
            errors.append("design.checkpoint.record is still the `<placeholder>` from the skeleton")
    return row_map


def check_build(
    evidence: dict[str, Any], root: Path, supplied_script: Path, errors: list[str]
) -> tuple[Path | None, dict[str, set[str]]]:
    build = evidence.get("build")
    if not isinstance(build, dict):
        errors.append("build evidence missing")
        return None, {}
    script = check_hashed_file(root, build.get("script"), build.get("sha256"), "build.script", errors)
    if script is None:
        return None, {}
    if script.resolve() != supplied_script.resolve():
        errors.append("--build-script does not match evidence.build.script")
    if build.get("strict_layout") is not True:
        errors.append("build.strict_layout must be true")
    calls, has_strict_layout = parse_script(script, errors)
    for token in forbidden_icon_rasterizer_calls(script, errors):
        errors.append(
            f"build script invokes macOS Quick Look thumbnail generation ({token}); use scripts/icons.py / icon_png for transparent high-resolution icon assets"
        )
    if not has_strict_layout:
        errors.append("build script must call lint_layout(..., strict=True)")
    return script, calls


def check_components(
    evidence: dict[str, Any],
    components: dict[str, Any],
    design_rows: dict[int, dict[str, Any]],
    calls: dict[str, set[str]],
    errors: list[str],
) -> None:
    for slide, row in design_rows.items():
        function = row.get("function")
        if not isinstance(function, str) or function not in calls:
            errors.append(f"design slide {slide} function is absent from the build script")
            continue
        missing = [component for component in row.get("components", []) if component not in calls[function]]
        if missing:
            errors.append(f"design slide {slide} declares component(s) not called by {function}: {', '.join(missing)}")

    clusters = components.get("clusters", [])
    if not isinstance(clusters, list):
        errors.append("component audit clusters missing or malformed")
        return
    for cluster in clusters:
        if not isinstance(cluster, dict):
            continue
        slide = cluster.get("slide")
        suggestions = cluster.get("suggest", [])
        if slide not in design_rows or not isinstance(suggestions, list) or not suggestions:
            continue
        row = design_rows[slide]
        function = row.get("function")
        implemented = set(row.get("components", [])) & calls.get(function, set()) & set(suggestions)
        suppressed_by = {
            item for item in components.get("suppressed_by", []) if isinstance(item, str)
        }
        emitted_here = suppressed_by & calls.get(function, set()) & set(row.get("components", []))
        if not implemented and not emitted_here and not waived(
            evidence, "component", slide=slide, pattern=cluster.get("pattern")
        ):
            pattern = cluster.get("pattern", "unnamed cluster")
            errors.append(
                f"component audit cluster on slide {slide} ({pattern}) needs a suggested component or a documented waiver"
            )


def recompute_component_audit(script: Path, deck: Path, errors: list[str]) -> dict[str, Any] | None:
    command = [
        sys.executable,
        str(Path(__file__).with_name("component_audit.py")),
        str(script),
        str(deck),
        "--json",
    ]
    result = subprocess.run(command, capture_output=True, text=True)
    if result.returncode not in {0, 2}:
        detail = (result.stdout + result.stderr).strip().replace("\n", " ")
        errors.append("component audit could not inspect the final build/deck: " + detail[:300])
        return None
    try:
        value = json.loads(result.stdout)
    except json.JSONDecodeError as exc:
        errors.append(f"component audit returned invalid JSON: {exc}")
        return None
    if not isinstance(value, dict) or value.get("inspected") is not True:
        errors.append("component audit did not confirm inspection of the final deck")
        return None
    return value


def check_icons(
    evidence: dict[str, Any],
    root: Path,
    design_rows: dict[int, dict[str, Any]],
    calls: dict[str, set[str]],
    errors: list[str],
) -> None:
    icon_rows = evidence.get("icons", [])
    if not isinstance(icon_rows, list):
        errors.append("icons must be a list")
        return
    by_slide: dict[int, list[dict[str, Any]]] = defaultdict(list)
    families: set[str] = set()
    for index, row in enumerate(icon_rows, start=1):
        label = f"icons[{index}]"
        if not isinstance(row, dict) or not isinstance(row.get("slide"), int):
            errors.append(f"{label} must identify a slide")
            continue
        family = require_string(row.get("family"), f"{label}.family", errors)
        if family:
            families.add(family)
        asset = check_hashed_file(root, row.get("asset"), row.get("sha256"), label, errors, minimum_bytes=32)
        rasterizer = require_string(row.get("rasterizer"), f"{label}.rasterizer", errors, minimum=3)
        if rasterizer not in {"scripts/icons.py", "provided-hires"}:
            errors.append(f"{label}.rasterizer must be scripts/icons.py or provided-hires")
        if asset is not None:
            info = png_info(asset)
            if info is None:
                errors.append(f"{label} must be a readable PNG icon asset")
            else:
                width, height, _bit_depth, color_type = info
                if min(width, height) < 256:
                    errors.append(
                        f"{label} is only {width}x{height}px; Codex icon assets need a 256px minimum edge to avoid thumbnail blur"
                    )
                if color_type not in {4, 6}:
                    errors.append(
                        f"{label} must preserve transparent alpha (PNG colour type 4 or 6), not a matted thumbnail"
                    )
        by_slide[row["slide"]].append(row)
    if len(families) > 1:
        errors.append("all icon evidence must use one coherent icon family")

    # 🔴 THE OPT-IN HOLE. Everything below is keyed on a slide having been MARKED `categorical`,
    # which the run writes about itself. Mark nothing, and a deck with zero icons satisfies every
    # line of this function — measured on a real 12-slide Codex deck: 0 pictures in the file, and
    # this gate had nothing to say. A gate that only fires when the author opts in is not a
    # detector, and the shared path's equivalent is one (it hashes the icon-sized pictures in the
    # built file and compares them against the declared family).
    #
    # So: if NOTHING is categorical, that claim is itself checked against the deck. A deck with
    # repeated same-size label sets — the shape of a category row — and no icons anywhere is the
    # case this whole field exists for. Waivable per slide like the rest, because the detector
    # over-counts by construction (tables, timelines, stat rows) and an over-counting detector
    # must not hold a deck without an escape.
    # Read the BUILT FILE, not the claim. `lint_deck.categorical_slides` is the shared definition
    # the other gate path uses; importing it is what keeps the two from disagreeing again.
    looks_categorical: list[int] = []
    try:
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        import lint_deck as _ld
        from pptx import Presentation as _P
        deck_rel = (evidence.get("deck") or {}).get("pptx")
        if deck_rel:
            looks_categorical = _ld.categorical_slides(_P(str(root / deck_rel)))
    except Exception:                                     # pragma: no cover - detector is advisory
        looks_categorical = []

    declared = {s for s, d in design_rows.items() if d.get("categorical")}
    # The detector over-counts (tables, timelines, stat rows), so it never fails a deck by itself —
    # it fails the CONTRADICTION: pages that read as category sets, declared as none, and no icon
    # anywhere. Any one of those three being false leaves this silent.
    missed = sorted(set(looks_categorical) - declared)
    if missed and not icon_rows and not _icon_waiver_ok(evidence):
        errors.append(
            "slides {} carry parallel label sets (3+ short peers sharing a baseline across half the "
            "canvas — the shape of a category row) but are declared `categorical: false`, and the "
            "deck records no icon assets at all. Icons are the DEFAULT on categorical content — they "
            "aid the 1-second read and reinforce the system — so shipping zero here is a HIGH-bar "
            "choice. Mark them and build the family (scripts/icons.py), or record why an icon family "
            "would HURT, with a CLASSIFIED reason (one of: {}): "
            '{{"kind": "icon", "scope": "undeclared-categorical", "category": "motif-dominant", '
            '"reason": "<why an icon family would hurt here>"}}. A bare reason no longer clears it — '
            "name the class, the way the critic waiver names its own.".format(
                missed, " | ".join(ICON_NONE_CATEGORIES)))

    for slide, design in design_rows.items():
        if not design.get("categorical"):
            continue
        rows = by_slide.get(slide, [])
        if not rows:
            if not waived(evidence, "icon", slide=slide):
                errors.append(f"categorical slide {slide} needs rendered icon evidence or a documented waiver")
            continue
        function = design.get("function")
        helper_count = len(calls.get(function, set()) & ICON_HELPERS)
        if helper_count < 1:
            errors.append(f"categorical slide {slide} records icon assets but {function} has no deckkit icon helper call")


def load_review_validator() -> Any:
    validator_path = Path(__file__).with_name("validate_review.py")
    spec = importlib.util.spec_from_file_location("slide_maker_validate_review", validator_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load validate_review.py")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_direction_validator() -> Any:
    validator_path = Path(__file__).with_name("directions_diversity.py")
    spec = importlib.util.spec_from_file_location("slide_maker_directions_diversity", validator_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load directions_diversity.py")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_arc_checker() -> Any:
    # Same shape as the three loaders above, and for the same reason: the gate's verdict must come
    # from the code the CLI runs. `content.arc.divergence` used to be a STRING this record wrote
    # about itself, and a delivered deck passed with the two characters "ok" while the script had
    # never been invoked for it. A gate that reads a verdict is a gate the run can dictate to.
    validator_path = Path(__file__).with_name("arc_divergence.py")
    spec = importlib.util.spec_from_file_location("slide_maker_arc_divergence", validator_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load arc_divergence.py")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_visual_contract() -> Any:
    script_path = Path(__file__).with_name("codex_visual_contract.py")
    spec = importlib.util.spec_from_file_location("slide_maker_codex_visual_contract", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load Codex visual-contract checker")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def check_critics(
    evidence: dict[str, Any], root: Path, expected_slides: set[int], deck_hash: str, errors: list[str]
) -> dict[str, dict[str, Any]]:
    reviewed: dict[str, dict[str, Any]] = {}
    critics = evidence.get("critics")
    if not isinstance(critics, list):
        errors.append("critics must be a list")
        return reviewed
    effort = evidence.get("review_effort", "standard")
    if effort not in {"fast", "standard", "thorough", "none"}:
        errors.append("review_effort must be fast, standard, thorough, or none")
        return reviewed
    if effort == "none":
        # The user declined review at the POST-BUILD question, with the rendered deck visible.
        # Never a default, never derived: the gate wants the decline itself on record.
        basis = evidence.get("none_opt_in")
        if not require_string(basis, "none_opt_in", errors, minimum=12):
            return reviewed
        if isinstance(basis, str) and basis.lstrip().startswith("<"):
            errors.append("none_opt_in still contains the template placeholder — quote the user's post-build decline")
        if critics:
            errors.append("review_effort none must not carry critic reviews — record the tier that actually ran")
        return reviewed
    lenses = [row.get("lens") for row in critics if isinstance(row, dict)]
    thorough_scope = None
    if effort == "fast":
        basis = evidence.get("fast_basis")
        if not require_string(basis, "fast_basis", errors, minimum=12):
            return reviewed
        if isinstance(basis, str) and basis.lstrip().startswith("<"):
            errors.append("fast_basis still contains the template placeholder — record how fast was reached")
            return reviewed
        required_lenses = {"general"}
    else:
        required_lenses = {"content", "design"}
    if effort == "thorough":
        panel = evidence.get("thorough_panel")
        if not isinstance(panel, dict):
            errors.append("thorough review requires thorough_panel evidence")
        else:
            thorough_scope = panel.get("scope")
            if thorough_scope not in {"light", "full"}:
                errors.append("thorough_panel.scope must be light or full")
            require_string(panel.get("record"), "thorough_panel.record", errors, minimum=12)
    if not required_lenses.issubset(set(lenses)):
        errors.append("critic panel does not cover the required lenses")
    if len(set(row.get("review") for row in critics if isinstance(row, dict))) != len(critics):
        errors.append("each critic lens must use a distinct review artifact")

    try:
        validator = load_review_validator()
    except Exception as exc:
        errors.append(f"could not load critic validator: {exc}")
        return
    for index, critic in enumerate(critics, start=1):
        label = f"critics[{index}]"
        if not isinstance(critic, dict):
            errors.append(f"{label} must be an object")
            continue
        lens = critic.get("lens")
        if lens not in {"content", "design", "general"}:
            errors.append(f"{label}.lens is invalid")
            continue
        path = check_hashed_file(root, critic.get("review"), critic.get("sha256"), label, errors, minimum_bytes=32)
        if path is None:
            continue
        if critic.get("pptx_sha256") != deck_hash:
            errors.append(f"{label} is not bound to the final PPTX SHA-256")
        try:
            review = load_json(path)
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            errors.append(f"{label} cannot be read: {exc}")
            continue
        reviewer = review.get("reviewer")
        if not isinstance(reviewer, dict):
            errors.append(f"{label} must record reviewer provenance (isolated, human, or self-review)")
        else:
            origin = reviewer.get("origin")
            if origin not in {"isolated", "human", "self-review"}:
                errors.append(f"{label}.reviewer.origin must be isolated, human, or self-review")
            require_string(reviewer.get("identity"), f"{label}.reviewer.identity", errors, minimum=3)
            if reviewer.get("fresh_context") is not True:
                errors.append(f"{label}.reviewer.fresh_context must be true")
            if origin == "self-review" and not waived(evidence, "critic-independence", lens=lens):
                errors.append(f"{label} is self-review; record a critic-independence waiver instead of claiming independent consent")
        for issue in validator.validate_critic(review):
            errors.append(f"{label} invalid critic schema: {issue}")
        coverage = review.get("coverage", {})
        opened = coverage.get("slides_opened", []) if isinstance(coverage, dict) else []
        if not isinstance(opened, list) or not expected_slides.issubset(set(opened)):
            errors.append(f"{label} did not inspect every final slide")
        if not isinstance(coverage, dict) or coverage.get("stats_block_seen") is not True:
            errors.append(f"{label} did not confirm the statistics block")
        if not isinstance(coverage, dict) or coverage.get("contract_card_seen") is not True:
            errors.append(f"{label} did not confirm the design contract")
        passes = coverage.get("passes", []) if isinstance(coverage, dict) else []
        if not isinstance(passes, list) or not any(lens in str(item).lower() for item in passes):
            errors.append(f"{label} does not record its {lens} lens pass")
        if lens == "general":
            normalized = " ".join(str(item).lower() for item in passes)
            if "content" not in normalized or "design" not in normalized:
                errors.append(f"{label} generalist pass must explicitly cover content and design")
        probes = review.get("probes", {})
        if lens == "content" and not isinstance(probes, dict):
            errors.append(f"{label} lacks content probes")
        if lens == "content" and not require_string(probes.get("memory_sentence") if isinstance(probes, dict) else None, f"{label}.probes.memory_sentence", errors, minimum=12):
            pass
        if lens == "design":
            per_slide = probes.get("per_slide", []) if isinstance(probes, dict) else []
            inspected = {row.get("slide") for row in per_slide if isinstance(row, dict)}
            if not expected_slides.issubset(inspected):
                errors.append(f"{label} lacks a design probe for every final slide")
        reviewed[lens] = review

    if effort == "thorough" and thorough_scope == "full":
        arbiters = evidence.get("arbiters")
        if not isinstance(arbiters, list) or not arbiters:
            errors.append("full thorough review requires a final arbiter confirmation artifact")
            return reviewed
        for index, arbiter in enumerate(arbiters, start=1):
            label = f"arbiters[{index}]"
            if not isinstance(arbiter, dict):
                errors.append(f"{label} must be an object")
                continue
            path = check_hashed_file(root, arbiter.get("review"), arbiter.get("sha256"), label, errors, minimum_bytes=32)
            if path is None:
                continue
            if arbiter.get("pptx_sha256") != deck_hash:
                errors.append(f"{label} is not bound to the final PPTX SHA-256")
            try:
                review = load_json(path)
            except (OSError, ValueError, json.JSONDecodeError) as exc:
                errors.append(f"{label} cannot be read: {exc}")
                continue
            for issue in validator.validate_arbiter(review):
                errors.append(f"{label} invalid arbiter schema: {issue}")
            checks = review.get("checks")
            if not isinstance(checks, list):
                errors.append(f"{label} must contain final fix-confirmation checks")
            elif any(not row.get("resolved") or row.get("dulled") for row in checks if isinstance(row, dict)):
                errors.append(f"{label} reports an unresolved or dulled final fix")
    return reviewed


def check_visual_contract(
    evidence: dict[str, Any],
    root: Path,
    deck_path: Path | None,
    deck_hash: str,
    build_script: Path | None,
    reviews: dict[str, dict[str, Any]],
    errors: list[str],
) -> None:
    contract = evidence.get("visual_contract")
    if not isinstance(contract, dict):
        errors.append("visual_contract evidence is required for Codex delivery")
        return
    manifest_path = check_hashed_file(
        root, contract.get("manifest"), contract.get("sha256"), "visual_contract.manifest", errors, minimum_bytes=32
    )
    result_path = check_hashed_file(
        root, contract.get("result"), contract.get("result_sha256"), "visual_contract.result", errors, minimum_bytes=32
    )
    if contract.get("pptx_sha256") != deck_hash:
        errors.append("visual_contract is not bound to the final PPTX SHA-256")
    if manifest_path is None or result_path is None or deck_path is None or build_script is None:
        return
    try:
        manifest = load_json(manifest_path)
        result = load_json(result_path)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        errors.append(f"visual_contract cannot be read: {exc}")
        return
    if result.get("schema") != "slide-maker-codex-visual-contract-result/v1":
        errors.append("visual_contract result has the wrong schema")
    if result.get("pptx_sha256") != deck_hash:
        errors.append("visual_contract result was produced for a different PPTX")
    if result.get("manifest_sha256") != contract.get("sha256"):
        errors.append("visual_contract result is not bound to the supplied manifest")
    if result.get("passed") is not True:
        errors.append("visual_contract result contains failed local checks")
    try:
        checker = load_visual_contract()
        fresh = checker.evaluate(deck_path, manifest_path, build_script)
    except Exception as exc:
        errors.append(f"visual_contract could not be recomputed: {exc}")
        return
    if fresh.get("passed") is not True:
        errors.append("visual_contract fails when recomputed against the final build and PPTX")
    declared_zones = {
        row.get("id") for row in manifest.get("zones", []) if isinstance(row, dict) and isinstance(row.get("id"), str)
    }
    declared_icons = {
        row.get("id") for row in manifest.get("icons", []) if isinstance(row, dict) and isinstance(row.get("id"), str)
    }
    result_zones = {
        row.get("id") for row in result.get("zones", []) if isinstance(row, dict) and row.get("pass") is True
    }
    result_icons = {
        row.get("id") for row in result.get("icons", []) if isinstance(row, dict) and row.get("pass") is True
    }
    if result_zones != declared_zones or result_icons != declared_icons:
        errors.append("visual_contract result does not account for every declared zone and icon")
    if evidence.get("review_effort", "standard") == "none":
        # The user waived the review loop at the post-build question; the critic-attestation rows
        # are review artifacts and go with it. The deterministic recompute above still ran — the
        # visual-contract FLOOR holds at every tier, only the human-eye attestation is waived.
        return
    # At `fast` the single generalist carries both lenses, so its review holds the probe rows.
    design_review = reviews.get("design") or reviews.get("general")
    probes = design_review.get("probes", {}) if isinstance(design_review, dict) else {}
    hotspot_checks = probes.get("hotspot_checks", []) if isinstance(probes, dict) else []
    icon_checks = probes.get("icon_checks", []) if isinstance(probes, dict) else []
    checked_zones = {
        row.get("id") for row in hotspot_checks
        if isinstance(row, dict) and row.get("result") == "pass" and isinstance(row.get("observed"), str) and len(row["observed"].strip()) >= 12
    }
    checked_icons = {
        row.get("id") for row in icon_checks
        if isinstance(row, dict) and row.get("result") == "pass" and isinstance(row.get("observed"), str) and len(row["observed"].strip()) >= 12
    }
    if not declared_zones.issubset(checked_zones):
        errors.append("design critic did not visually attest every declared local hotspot")
    if not declared_icons.issubset(checked_icons):
        errors.append("design critic did not attest every declared icon's semantic fit")


def evaluate(
    lint: dict[str, Any], components: dict[str, Any], supplied_script: Path, evidence: dict[str, Any], root: Path
) -> list[str]:
    errors: list[str] = []
    if evidence.get("schema") != SCHEMA:
        errors.append(f"evidence.schema must be {SCHEMA}")
    if evidence.get("runtime") not in {"codex", "openai-gpt-bridged"}:
        errors.append("evidence.runtime must be codex or openai-gpt-bridged")
    delivery = evidence.get("delivery", "presented")
    if delivery not in BODY_FLOORS:
        errors.append("evidence.delivery must be presented, textheavy, or selfread")
        delivery = "presented"
    deck = evidence.get("deck")
    deck_hash = ""
    expected_slides: set[int] = set()
    deck_path: Path | None = None
    if not isinstance(deck, dict):
        errors.append("deck evidence missing")
    else:
        deck_path = check_hashed_file(root, deck.get("pptx"), deck.get("sha256"), "deck", errors, minimum_bytes=512)
        deck_hash = deck.get("sha256") if is_sha256(deck.get("sha256")) else ""
        count = deck.get("slide_count")
        if not isinstance(count, int) or count < 1:
            errors.append("deck.slide_count must be a positive integer")
        else:
            expected_slides = set(range(1, count + 1))
            if deck_path is not None:
                actual_count = slide_count_from_pptx(deck_path)
                if actual_count is None:
                    errors.append("could not read final PPTX to count slides")
                elif actual_count != count:
                    errors.append(f"deck.slide_count is {count}, but final PPTX contains {actual_count} slides")
    check_lint(lint, delivery, evidence, errors)
    if expected_slides:
        check_content(evidence, root, expected_slides, errors)
        check_render_selfcheck(evidence, expected_slides, errors)
        design_rows = check_design(evidence, root, expected_slides, deck_hash, errors)
        build_script, calls = check_build(evidence, root, supplied_script, errors)
        audited_components = components
        if build_script is not None and deck_path is not None:
            recomputed = recompute_component_audit(build_script, deck_path, errors)
            if recomputed is not None:
                audit_keys = ("used_forms", "clusters", "suppressed_by", "inspected")
                if any(components.get(key) != recomputed.get(key) for key in audit_keys):
                    errors.append("components JSON does not match a fresh audit of the final build and PPTX")
                audited_components = recomputed
        check_components(evidence, audited_components, design_rows, calls, errors)
        check_icons(evidence, root, design_rows, calls, errors)
        reviews = check_critics(evidence, root, expected_slides, deck_hash, errors)
        check_visual_contract(evidence, root, deck_path, deck_hash, build_script, reviews, errors)
    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description="Codex-only final delivery gate for slide-maker")
    parser.add_argument("--lint", type=Path, help="JSON from lint_layout")
    parser.add_argument("--components", type=Path, help="JSON from component_audit.py")
    parser.add_argument("--build-script", type=Path, help="final deck build script")
    parser.add_argument("--evidence", type=Path, help=".codex-deck-evidence.json")
    parser.add_argument("--receipt", type=Path, help="write a final-file-bound PASS receipt")
    parser.add_argument("--init", type=Path, help="write an evidence template and exit")
    args = parser.parse_args()

    if args.init:
        if args.init.exists():
            print(f"refusing to overwrite existing file: {args.init}", file=sys.stderr)
            return 2
        args.init.write_text(json.dumps(TEMPLATE, indent=2) + "\n", encoding="utf-8")
        print(f"wrote evidence template: {args.init}")
        return 0

    missing = [name for name in ("lint", "components", "build_script", "evidence") if getattr(args, name) is None]
    if missing:
        parser.error("required unless --init: " + ", ".join("--" + item.replace("_", "-") for item in missing))
    try:
        lint = load_json(args.lint)
        components = load_json(args.components)
        evidence = load_json(args.evidence)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"cannot read gate input: {exc}", file=sys.stderr)
        return 2
    errors = evaluate(lint, components, args.build_script, evidence, args.evidence.parent)
    if errors:
        print("CODEX DELIVERY GATE: BLOCKED")
        for error in errors:
            print(f"- {error}")
        return 1
    print("CODEX DELIVERY GATE: PASS")
    if args.receipt:
        try:
            write_receipt(args.receipt, evidence_path=args.evidence, build_script=args.build_script, evidence=evidence)
        except OSError as exc:
            print(f"cannot write delivery receipt: {exc}", file=sys.stderr)
            return 2
        print(f"CODEX DELIVERY RECEIPT: {args.receipt}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

#!/usr/bin/env python3
"""The direction the user PICKED must be the direction that shipped.

WHY. The branch-(c) gate renders four directions, the user clicks one, and the choice is recorded
as a sentence on the design checkpoint's `direction gate:` line. Nothing then compared the BUILT
deck to the direction that was picked. Measured on a real delivered deck: the chosen direction
declared a Georgia display face and a centred cover; the deck shipped with Helvetica Neue titles
and a low-left cover, because `style.py` set `display=` and every title passed `dk.FONT`, so the
DISPLAY slot was never read. The author reported it — "我选的 B 和实际最终模版并不一样" — and no
gate had anything to say, while two neighbouring checks in this same repo already existed for the
identical class of failure: `check_register_pixels` (a declared colour must reach the pixels) and
`check_style_applied` (a declared preset must actually be called).

WHAT IT COMPARES. Only what a machine can settle from the file:

    ground     the deck's own background colour vs the direction's `bg`
    accent     the direction's `accent` must APPEAR somewhere in the deck
    display    the typeface of the largest text run vs the direction's `display`
    body       the typeface of the body-size runs vs the direction's `body`
    cover      `centred` vs `low-left`, measured from the cover's largest text block

Composition beyond those two named values, motif and skeleton are NOT checked and say so: a
"skeleton: island" is a judgement, and a check that guesses at one is worse than a check that
names its own limit.

A DELIBERATE deviation is legitimate and common — the freshness rule can move a ground, a
contrast floor can move an accent. It is recorded in writing, per axis:

    "design_plan": {"direction_deviations": {"bg": "the freshness gate measured the picked
                     value as a repeat of a recent deck's"}}

An unrecorded one is what this exists to catch, because that is the version the user cannot see.

    python3 scripts/check_direction_applied.py <deck>.pptx [--json]
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Two thresholds, because the two questions are not the same one. GROUND is the first thing the
# eye lands on and the least varied thing in a run of decks, so it is held tight: #F2EFE6 and
# #E8DFC9 differ by 55 and the freshness gate treats them as different papers, which they are.
# An ACCENT only has to APPEAR somewhere, so a renderer's rounding or a shade picked for contrast
# should not read as absent.
NEAR_GROUND = 26
NEAR_ACCENT = 60
CENTRED = 0.09     # a block is centred when its centre is within this fraction of the canvas centre


def _hex(v):
    v = str(v or "").strip().lstrip("#").upper()
    return v if re.fullmatch(r"[0-9A-F]{6}", v) else None


def _dist(a, b):
    a, b = _hex(a), _hex(b)
    if not a or not b:
        return None
    return sum(abs(int(a[i:i + 2], 16) - int(b[i:i + 2], 16)) for i in (0, 2, 4))


def _face(stack):
    return (str(stack or "").split(",")[0].strip().strip("'\"").lower()) or None


def facts(pptx):
    """What the BUILT deck actually is: ground, colours present, display/body faces, cover block."""
    from pptx import Presentation
    prs = Presentation(str(pptx))
    W = prs.slide_width / 914400.0
    grounds, colours, runs = {}, {}, []
    cover_block = None
    for n, slide in enumerate(prs.slides, 1):
        x = slide._element
        bg = x.find(".//" + _P + "bg")
        if bg is not None:
            c = bg.find(".//" + _A + "srgbClr")
            if c is not None and c.get("val"):
                grounds[c.get("val").upper()] = grounds.get(c.get("val").upper(), 0) + 1
        for c in x.iter(_A + "srgbClr"):
            if c.get("val"):
                v = c.get("val").upper()
                colours[v] = colours.get(v, 0) + 1
        for r in x.iter(_A + "r"):
            rpr = r.find(_A + "rPr")
            t = r.find(_A + "t")
            if rpr is None or t is None or not (t.text or "").strip():
                continue
            try:
                sz = int(rpr.get("sz") or 0) / 100.0
            except ValueError:
                sz = 0.0
            lat = rpr.find(_A + "latin")
            runs.append((sz, _face(lat.get("typeface") if lat is not None else None),
                         (t.text or "").strip()))
        if n == 1:
            best = None
            for sh in slide.shapes:
                try:
                    if not (sh.has_text_frame and sh.text_frame.text.strip()):
                        continue
                    big = max((rr.font.size.pt if rr.font.size else 0)
                              for p in sh.text_frame.paragraphs for rr in p.runs) if any(
                        p.runs for p in sh.text_frame.paragraphs) else 0
                except Exception:
                    continue
                if best is None or big > best[0]:
                    best = (big, (sh.left or 0) / 914400.0, (sh.width or 0) / 914400.0,
                            (sh.top or 0) / 914400.0)
            if best:
                cover_block = {"size": best[0], "left": best[1], "w": best[2], "top": best[3]}
    ground = max(grounds, key=grounds.get) if grounds else None
    sized = [r for r in runs if r[0] > 0 and r[1]]
    display = max(sized, key=lambda r: r[0])[1] if sized else None
    # BODY is the most common face among runs BELOW the display size — not "the median size's
    # face". On a page with two runs the median IS the title, so a deck whose body face matched
    # perfectly was reported as diverging. Small samples are exactly when a check must not guess.
    if sized:
        top = max(r[0] for r in sized)
        below = [r for r in sized if r[0] < top] or sized
        body_faces = {}
        for _sz, face, _t in below:
            body_faces[face] = body_faces.get(face, 0) + 1
        body = max(body_faces, key=body_faces.get) if body_faces else None
    else:
        body = None
    return {"ground": ground, "colours": colours, "display": display, "body": body,
            "cover_block": cover_block, "canvas_w": W}


def picked(gates, directions):
    """Which of the rendered directions was chosen, from the recorded `direction_gate` line."""
    line = ""
    d = (gates or {}).get("design_plan")
    if isinstance(d, dict):
        line = str(d.get("direction_gate") or "")
        if isinstance(d.get("direction_pick"), str) and d["direction_pick"].strip():
            line = d["direction_pick"]
    m = re.search(r"picked\s+`?([^`·\n]+?)`?\s+(?:of\b|$)", line, re.I)
    want = (m.group(1) if m else line).strip().strip("`\"' ")
    if not want:
        return None, "no `direction_gate` line to read the pick from"
    names = [str(x.get("name") or "") for x in directions]
    for nm in names:
        if nm.lower() == want.lower():
            return nm, None
    for nm in names:                       # "B" or "B — Aperture" against "B — Aperture (bespoke)"
        if nm.lower().startswith(want.lower()) or want.lower().startswith(nm.lower()):
            return nm, None
    return None, "the recorded pick {!r} matches none of {}".format(want, names)


def check(pptx, gates=None, deck_dir=None):
    """Return (problems, facts). Never raises on a missing artifact — it says NOT CHECKED."""
    out = {}
    deck_dir = Path(deck_dir or Path(pptx).resolve().parent)
    dj = deck_dir / "directions.json"
    if not dj.is_file():
        out["note"] = ("no directions.json beside the deck — this deck did not go through the "
                       "direction gate, which is NOT the same as its direction being honoured")
        return [], out
    try:
        directions = json.loads(dj.read_text(encoding="utf-8"))
    except (OSError, ValueError) as exc:
        out["note"] = "cannot read {}: {}".format(dj, exc)
        return [], out
    name, why = picked(gates, directions)
    if not name:
        out["note"] = "cannot tell which direction was picked — {}".format(why)
        return [], out
    want = [x for x in directions if x.get("name") == name][0]
    got = facts(pptx)
    out.update({"picked": name, "measured": {k: got[k] for k in ("ground", "display", "body")}})

    dev = {}
    d = (gates or {}).get("design_plan")
    if isinstance(d, dict) and isinstance(d.get("direction_deviations"), dict):
        dev = {str(k).lower(): str(v) for k, v in d["direction_deviations"].items() if str(v).strip()}
    out["deviations_recorded"] = sorted(dev)

    problems = []

    def _flag(axis, msg):
        if axis in dev:
            out.setdefault("accepted", []).append("{}: {}".format(axis, dev[axis]))
        else:
            problems.append((axis, msg))

    dist = _dist(want.get("bg"), got["ground"])
    if got["ground"] is None:
        out.setdefault("unchecked", []).append("ground — the deck writes no slide background")
    elif dist is not None and dist > NEAR_GROUND:
        _flag("bg", "the picked direction's ground is #{} and the deck's is #{} — the first thing "
                    "the eye lands on is not the one that was chosen"
                    .format(_hex(want.get("bg")), got["ground"]))

    acc = _hex(want.get("accent"))
    if acc:
        # `or 999` here was a real bug: an EXACT match distance of 0 is falsy, so a perfectly
        # applied accent read as absent. Compare against None explicitly.
        near = [v for v in got["colours"]
                if (lambda d: d is not None and d <= NEAR_ACCENT)(_dist(acc, v))]
        if not near:
            _flag("accent", "the picked direction's accent #{} appears nowhere in the deck".format(acc))

    for axis in ("display", "body"):
        wf, gf = _face(want.get(axis)), got[axis]
        if wf and gf and wf != gf:
            _flag(axis, "the picked direction's {} face is {!r} and the deck sets {!r} — the slot "
                        "was declared and never read".format(axis, wf, gf))

    cb, cover = got["cover_block"], str(want.get("cover") or "").strip().lower()
    if cb and cover in ("centred", "low-left"):
        centre = (cb["left"] + cb["w"] / 2.0) / max(0.01, got["canvas_w"])
        is_centred = abs(centre - 0.5) <= CENTRED
        if cover == "centred" and not is_centred:
            _flag("cover", "the picked direction's cover is CENTRED and the deck's title block sits "
                           "at {:.0%} of the width".format(centre))
        if cover == "low-left" and is_centred:
            _flag("cover", "the picked direction's cover is LOW-LEFT and the deck's title block is "
                           "centred")
    elif cover:
        out.setdefault("unchecked", []).append(
            "cover — {!r} is a judgement this cannot settle from the file".format(cover))
    out.setdefault("unchecked", []).append("skeleton and motif — judgement, deliberately not guessed")
    return problems, out


def _selftest():
    import tempfile
    ok, bad = [], []
    tmp = Path(tempfile.mkdtemp(prefix="dirapply-"))

    # An EXACT match has distance 0, which is falsy — the first version wrote `(_dist(...) or 999)`
    # and therefore reported a perfectly applied accent as absent. These cases hold that shut.
    for a, b, near in (("#F2EFE6", "F2EFE6", True), ("#F2EFE6", "E8DFC9", False),
                       ("#2F6B5F", "2F6B5F", True), ("#2F6B5F", "306C60", True)):
        d = _dist(a, b)
        got = d is not None and d <= NEAR_GROUND
        (ok if got == near else bad).append(
            "#{} and #{} read as {}the same ground (distance {})".format(
                _hex(a), b, "" if near else "NOT ", d)
            if got == near else "{} vs {} gave {} (distance {})".format(a, b, got, d))

    dirs = [{"name": "A — Swiss"}, {"name": "B — Aperture (bespoke)"}]
    for line, want in (("picked `B — Aperture (bespoke)` of 4 rendered", "B — Aperture (bespoke)"),
                       ("picked B — Aperture of 4 rendered directions", "B — Aperture (bespoke)"),
                       ("picked Z of 4", None)):
        got, _why = picked({"design_plan": {"direction_gate": line}}, dirs)
        (ok if got == want else bad).append(
            "the pick reads out of {!r}".format(line[:34]) if got == want
            else "{!r} -> {!r}, wanted {!r}".format(line, got, want))

    (dirs_p := tmp / "directions.json").write_text(json.dumps(
        [{"name": "B", "bg": "#F2EFE6", "accent": "#2F6B5F", "display": "Georgia",
          "body": "Helvetica Neue", "cover": "centred"}]), encoding="utf-8")
    import warnings
    warnings.simplefilter("ignore")
    import deckkit as dk
    dk.set_palette(deep="1E1B18", magenta="2F6B5F", font="Helvetica Neue", display="Georgia")
    dk.set_ground("F2EFE6")
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    dk.text(s, 0.6, 2.0, 6.0, 0.8, [[("A left-set title", 40, dk.DEEP, True, False, dk.FONT)]])
    dk.text(s, 0.6, 3.0, 6.0, 0.4, [[("body", 13, dk.DEEP, False, False, dk.FONT)]])
    deck = tmp / "d.pptx"
    prs.save(str(deck))
    gates = {"design_plan": {"direction_gate": "picked `B` of 1 rendered"}}
    probs, f = check(deck, gates=gates, deck_dir=tmp)
    axes = {a for a, _m in probs}
    (ok if "display" in axes else bad).append(
        "a display face DECLARED and never read is caught — the exact drift a real deck shipped "
        "(Georgia picked, Helvetica Neue built)" if "display" in axes else str(axes))
    (ok if "cover" in axes else bad).append(
        "a centred cover that shipped left-set is caught" if "cover" in axes else str(axes))

    gates["design_plan"]["direction_deviations"] = {
        "display": "the presenter's brand type replaces the direction's serif",
        "cover": "the cover carries a logo lockup that wants the left rail"}
    probs2, f2 = check(deck, gates=gates, deck_dir=tmp)
    (ok if not [a for a, _m in probs2 if a in ("display", "cover")] else bad).append(
        "...and a deviation RECORDED IN WRITING is accepted rather than reported — moving a "
        "direction is legitimate, moving it silently is not")

    probs3, f3 = check(deck, gates={"design_plan": {}}, deck_dir=tmp)
    (ok if not probs3 and f3.get("note") else bad).append(
        "an unreadable pick says NOT CHECKED rather than clean — the two are different facts")

    for line in ok:
        print("  ok   " + line)
    for line in bad:
        print("  FAIL " + line)
    print("\n{} passed, {} failed".format(len(ok), len(bad)))
    return 1 if bad else 0


def main(argv=None):
    from _console import safe_stdio
    safe_stdio()
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("pptx", nargs="?")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--selftest", action="store_true")
    a = ap.parse_args(argv)
    if a.selftest:
        return _selftest()
    if not a.pptx:
        ap.print_help()
        return 2
    deck = Path(a.pptx)
    gates = {}
    gp = deck.resolve().parent / ".deck-gates.json"
    if gp.is_file():
        try:
            gates = json.loads(gp.read_text(encoding="utf-8"))
        except ValueError:
            gates = {}
    probs, f = check(deck, gates=gates)
    if a.json:
        print(json.dumps({"problems": [{"axis": x, "why": y} for x, y in probs], "facts": f},
                         indent=2, ensure_ascii=False))
        return 1 if probs else 0
    if f.get("note"):
        print("  [--] direction NOT CHECKED — {}".format(f["note"]))
        return 0
    print("direction picked: {}".format(f.get("picked")))
    for line in f.get("accepted", []):
        print("  [--] deviation recorded: {}".format(line))
    for line in f.get("unchecked", []):
        print("  [--] not checked: {}".format(line))
    if not probs:
        print("the deck is the direction that was picked, on every axis a file can settle.")
        return 0
    print("\n{} axis/axes diverge from the picked direction:\n".format(len(probs)))
    for axis, why in probs:
        print("  {}: {}".format(axis.upper(), why))
    print("\n  Bring the deck back to the direction, or record the deviation per axis:")
    print('    "design_plan": {"direction_deviations": {"<axis>": "<why>"}}')
    return 1


if __name__ == "__main__":
    sys.exit(main())

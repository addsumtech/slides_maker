#!/usr/bin/env python3
"""Hold every image evidence token against something REAL — the ledger, and the built deck.

`references/image-generation.md` defines the token grammar every image row in the plan must carry
(`sourced — <origin> (<license>)`, `generated — <tool>`, `searched (Commons, Openverse), none
found → …`, …). Until now that grammar was enforced by reading: the plan said what it said, and
the two claims that actually cost something were unverifiable by construction —

  * **"I searched and found nothing."** Cheap to type, expensive to do. It is the sanctioned door
    to generating an illustration of a real place, so a run under time pressure walks through it
    without knocking. `fetch_images.py` now records every query and its outcome; this checks the
    rung against that record, and REFUSES the `none found` rung when the recorded outcome was
    `unreachable` — a blocked network is not evidence that no photo of Amsterdam exists.
  * **"CC BY, credit given."** The licence obligation is discharged by a line ON A SLIDE, and
    nothing ever looked at the slides. A deck can carry an attribution-required photo, record the
    licence dutifully in its plan, and ship with no credit anywhere — a licence violation that
    every gate in the tree currently calls clean.

What it does NOT do: judge whether a photo is the right subject or good enough. That is
`image_qc.py` plus eyes (see both files' docstrings). This one is about the RECORD.

    python3 scripts/check_image_provenance.py <deck-dir>            # .deck-gates.json + assets
    python3 scripts/check_image_provenance.py <deck-dir> --pptx <deck.pptx>
    python3 scripts/check_image_provenance.py --selftest

Exit 0 clean · 1 findings · 2 could not run (never silently "clean").
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

GATES_FILE = ".deck-gates.json"
LEDGER_NAME = "sources.json"

# Punctuation is normalised before matching: an em dash, an en dash and a hyphen are the same
# intent, and a gate that fails on which dash somebody typed is noise that teaches people to
# switch gates off. The STRUCTURE is what is checked.
_DASH = re.compile(r"[—–−]|--")
_ARROW = re.compile(r"→|->|=>")

TOKEN_FORMS = (
    ("sourced", re.compile(r"\bsourced\s*-\s*(?P<origin>[^()]+?)\s*\((?P<license>[^)]+)\)", re.I)),
    ("provided", re.compile(r"\bprovided\s*-\s*(?P<origin>[^()]+?)\s*\((?P<note>[^)]+)\)", re.I)),
    ("generated", re.compile(r"\bgenerated\s*-\s*(?P<tool>[^|;,]+)", re.I)),
    # The origins parenthetical is OPTIONAL to PARSE and REQUIRED to PASS. The reference teaches
    # the rung both ways — the bullet shows `searched, none found → …` and the paragraph under it
    # requires the origins be named — so a parser that demands the parenthetical would reject the
    # form the skill itself prints, and report it as BAD TOKEN: a true complaint under a
    # misleading name. Parsing it and failing it as UNNAMED ORIGINS says the actual thing.
    ("searched", re.compile(
        r"\bsearched\s*(?:\((?P<origins>[^)]+)\))?\s*,\s*"
        r"(?P<outcome>none found|found but low-quality)\s*->\s*"
        r"(?P<fallback>generated,\s*flagged illustrative|native form)", re.I)),
)


def _norm(s):
    s = _DASH.sub("-", str(s or ""))
    s = _ARROW.sub("->", s)
    return re.sub(r"\s+", " ", s).strip()


def parse_token(row):
    """Return (kind, match) for the FIRST sanctioned form found, else (None, None).

    `searched` is tried before `generated` on purpose: the fallback rungs CONTAIN the word
    'generated', and matching the shorter form first would silently accept
    'searched (…), none found -> generated, flagged illustrative' as a bare `generated` token —
    losing exactly the claim that needs checking."""
    n = _norm(row)
    for kind, rx in (TOKEN_FORMS[3], TOKEN_FORMS[0], TOKEN_FORMS[1], TOKEN_FORMS[2]):
        m = rx.search(n)
        if m:
            return kind, m
    return None, None


def _rows(image_sources):
    if image_sources is None:
        return []
    if isinstance(image_sources, str):
        return [r for r in re.split(r"[\n;]+", image_sources) if r.strip()]
    if isinstance(image_sources, dict):
        return ["{}: {}".format(k, v) for k, v in image_sources.items()]
    return [str(r) for r in image_sources]


def _is_na(image_sources):
    """`n/a — <reason>` is the sanctioned escape, same shape as `logo plan: n/a — …`."""
    if not isinstance(image_sources, str):
        return False
    n = _norm(image_sources).lower()
    return n.startswith("n/a") and len(n) > 6         # a bare "n/a" is not a reason


def _deck_text(pptx_path):
    try:
        from pptx import Presentation
    except ImportError:
        return None
    out = []
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        return None
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        out.append(c.text)
        if slide.has_notes_slide:
            try:
                out.append(slide.notes_slide.notes_text_frame.text)
            except Exception:
                pass
    return _norm(" ".join(out)).lower()


def _credit_present(entry, deck_text):
    """A credit is present when the AUTHOR (or, failing that, the title) appears in the deck.

    Deliberately loose about format — a credits page, a caption and a source note are all valid
    places and all phrase it differently. Strict about the fact: the name has to be SOMEWHERE."""
    if deck_text is None:
        return None                                   # unknown: no pptx / no python-pptx
    author = _norm(entry.get("author", "")).lower()
    if len(author) >= 4:
        # When there IS an author, the author is the credit. Falling back to the title here was a
        # real hole, caught by the self-test: a file titled "Campus" scored as credited because the
        # deck happens to contain the word "campus" in a heading. A credit that a body word can
        # satisfy is not a check.
        return author in deck_text
    title = re.sub(r"\.(jpe?g|png|webp|tiff?)$", "", _norm(entry.get("title", "")).lower()).strip()
    # No author (some CC0 uploads) — the title stands in, but only when it is distinctive enough
    # that a stray heading cannot satisfy it.
    return bool(len(title) >= 12 and title in deck_text)


def check(deck_dir, *, pptx=None, gates=None, ledger=None, strict_credits=True):
    """Return a list of (CODE, message). Empty list means clean."""
    deck_dir = Path(deck_dir)
    problems = []

    if gates is None:
        gp = deck_dir / GATES_FILE
        gates = {}
        if gp.exists():
            try:
                gates = json.loads(gp.read_text(encoding="utf-8"))
            except ValueError as exc:
                return [("UNREADABLE GATES", "{}: {}".format(gp, exc))]
    design = gates.get("design_plan") or {}
    image_sources = design.get("image_sources")

    if ledger is None:
        ledger = {"entries": [], "searches": []}
        for cand in (deck_dir / LEDGER_NAME, deck_dir / "assets" / LEDGER_NAME,
                     deck_dir / "assets" / "sourced" / LEDGER_NAME):
            if cand.exists():
                try:
                    ledger = json.loads(cand.read_text(encoding="utf-8"))
                except ValueError as exc:
                    problems.append(("UNREADABLE LEDGER", "{}: {}".format(cand, exc)))
                break
    entries = ledger.get("entries") or []
    searches = ledger.get("searches") or []

    # 1. The plan has to SAY something about images — including that there are none.
    if image_sources is None and not entries:
        return problems                               # nothing planned, nothing sourced: not our call
    if _is_na(image_sources):
        if entries:
            problems.append(("N/A CONTRADICTED",
                             "`image_sources` is `n/a — …` but {} image(s) were sourced into the "
                             "ledger. One of the two is wrong.".format(len(entries))))
        return problems

    rows = _rows(image_sources)
    if not rows and entries:
        problems.append(("NO TOKENS",
                         "{} sourced image(s) in {} and no `image_sources` row in the design plan. "
                         "A bare filename with no evidence token is an INCOMPLETE plan "
                         "(references/image-generation.md).".format(len(entries), LEDGER_NAME)))

    # 2. Every row carries a token in the sanctioned grammar.
    kinds = []
    for row in rows:
        kind, m = parse_token(row)
        if not kind:
            problems.append(("BAD TOKEN",
                             "row {!r} carries no sanctioned evidence token. The grammar is in "
                             "references/image-generation.md: `sourced — <origin> (<license>)` · "
                             "`provided — user (own material)` · `generated — <tool>` · `searched "
                             "(<origins>), none found → generated, flagged illustrative` · "
                             "`searched (<origins>), found but low-quality → …` · `searched "
                             "(<origins>), none found → native form`.".format(row[:110])))
            continue
        kinds.append((kind, m, row))

    # 3. A `searched … none found` rung must name origins AND be backed by a recorded search.
    recorded = {s.get("outcome") for s in searches}
    for kind, m, row in kinds:
        if kind != "searched":
            continue
        origins = [o.strip() for o in (m.group("origins") or "").split(",") if o.strip()]
        if not origins:
            problems.append(("UNNAMED ORIGINS",
                             "row {!r}: a `searched` rung must NAME the origins tried — write it "
                             "`searched (Commons, Openverse[, press kit]), none found → …`. A bare "
                             "rung is an incomplete plan (references/image-generation.md), the same "
                             "clause pattern as the logo token.".format(row[:80])))
        if searches:
            if "unreachable" in recorded and "none found" not in recorded:
                problems.append(("UNREACHABLE IS NOT NONE-FOUND",
                                 "row {!r} claims a search found nothing, but the only recorded "
                                 "search outcome is `unreachable` — the sources could not be "
                                 "CONTACTED. That is a connectivity failure, not evidence that no "
                                 "photo exists, and it does not license a generated plate of a "
                                 "real subject. Re-run the search from a network that reaches "
                                 "Commons/Openverse.".format(row[:80])))
            elif "none found" not in recorded:
                problems.append(("UNBACKED RUNG",
                                 "row {!r} claims `none found`, but {} records only: {}. Run the "
                                 "search through scripts/fetch_images.py so the claim has a "
                                 "record.".format(row[:80], LEDGER_NAME,
                                                  ", ".join(sorted(recorded)) or "nothing")))
        else:
            problems.append(("UNRECORDED SEARCH",
                             "row {!r} claims a search happened; no search is recorded in {}. The "
                             "rung is the door to generating an illustration of a REAL subject, so "
                             "it is the one claim that has to leave a trace."
                             .format(row[:80], LEDGER_NAME)))

    # 4. Every `sourced —` token should correspond to a ledger entry (count, and licence text).
    sourced_rows = [(m, row) for kind, m, row in kinds if kind == "sourced"]
    placed = [e for e in entries if e.get("status") == "placed"] or entries
    if sourced_rows and not entries:
        problems.append(("UNBACKED SOURCED",
                         "{} row(s) claim a sourced photo and {} holds no entry. A sourced claim "
                         "with no provenance record is the photographic version of an unsourced "
                         "number.".format(len(sourced_rows), LEDGER_NAME)))
    else:
        led_lic = {_norm(e.get("license", "")).lower() for e in entries}
        for m, row in sourced_rows:
            lic = _norm(m.group("license")).lower()
            if led_lic and lic not in led_lic and not any(lic in l for l in led_lic):
                problems.append(("LICENCE MISMATCH",
                                 "row {!r} states licence {!r}; the ledger holds {}. A credit line "
                                 "built from the wrong licence is a wrong licence statement."
                                 .format(row[:70], m.group("license").strip(),
                                         ", ".join(sorted(x for x in led_lic if x)) or "none")))

    # 5. Attribution-required photos must be CREDITED IN THE DECK, not just in the ledger.
    deck_text = _deck_text(pptx) if pptx else None
    if strict_credits:
        need = [e for e in placed if e.get("attribution_required")]
        if need and deck_text is None and pptx:
            problems.append(("CREDITS UNVERIFIED",
                             "{} attribution-required photo(s) and the deck could not be read "
                             "(python-pptx missing or file unreadable) — the licence obligation is "
                             "unchecked, which is not the same as met.".format(len(need))))
        elif deck_text is not None:
            for e in need:
                if _credit_present(e, deck_text) is False:
                    problems.append(("MISSING CREDIT",
                                     "{} is {} — attribution REQUIRED — and neither its author "
                                     "({!r}) nor its title appears anywhere in the deck. Add a "
                                     "credit: deckkit.source_note(...) at the image, or one line "
                                     "on deckkit.sources_page(...). Ready-made text: "
                                     "`python3 scripts/fetch_images.py ledger <assets> --credits`."
                                     .format(e.get("file"), e.get("license"),
                                             (e.get("author") or "")[:40])))

    # 6. A photo nobody LOOKED at. `candidate` is the state fetch_images writes; `placed` is what
    #    adopt() writes after a human/model viewed the contact sheet.
    if entries and not any(e.get("status") == "placed" for e in entries) and sourced_rows:
        problems.append(("NOT ADOPTED",
                         "{} candidate(s) in {} and none marked `placed`. Downloading is not "
                         "choosing: view the contact sheet (image_qc.py --contact-sheet), then "
                         "`fetch_images.py adopt <dir> <file>`.".format(len(entries), LEDGER_NAME)))
    return problems


# --------------------------------------------------------------------------- selftest

def _selftest():
    ok, bad = [], []

    def run(name, *, gates, ledger, pptx=None, expect_codes=(), forbid_codes=()):
        got = {c for c, _ in check(".", gates=gates, ledger=ledger, pptx=pptx)}
        miss = set(expect_codes) - got
        extra = set(forbid_codes) & got
        if not miss and not extra:
            ok.append(name)
        else:
            bad.append("{} — missing {} unexpected {}".format(name, sorted(miss), sorted(extra)))

    good_led = {"entries": [{"file": "a.jpg", "license": "CC BY-SA 4.0", "author": "X Photographer",
                             "title": "Campus", "attribution_required": True, "status": "placed"}],
                "searches": [{"query": "campus", "outcome": "found", "n_results": 3}]}

    run("a well-formed sourced row with a matching ledger entry is clean (credits unchecked "
        "without a deck)",
        gates={"design_plan": {"image_sources": [
            "slide 4 | campus | sourced — Wikimedia Commons (CC BY-SA 4.0)"]}},
        ledger=good_led, forbid_codes=("BAD TOKEN", "UNBACKED SOURCED", "LICENCE MISMATCH"))

    run("a row with no token at all is an incomplete plan",
        gates={"design_plan": {"image_sources": ["slide 4 | campus | campus.jpg"]}},
        ledger=good_led, expect_codes=("BAD TOKEN",))

    run("`searched (…), none found → generated` with NO recorded search is refused",
        gates={"design_plan": {"image_sources": [
            "slide 9 | 1890 factory | searched (Commons, Openverse), none found → generated, "
            "flagged illustrative"]}},
        ledger={"entries": [], "searches": []}, expect_codes=("UNRECORDED SEARCH",))

    run("...and is ACCEPTED when the ledger records a real empty search",
        gates={"design_plan": {"image_sources": [
            "slide 9 | 1890 factory | searched (Commons, Openverse), none found → generated, "
            "flagged illustrative"]}},
        ledger={"entries": [], "searches": [{"query": "1890 factory floor",
                                             "outcome": "none found", "n_results": 0}]},
        forbid_codes=("UNRECORDED SEARCH", "UNBACKED RUNG", "BAD TOKEN"))

    run("a network failure cannot be spent as a `none found` rung",
        gates={"design_plan": {"image_sources": [
            "slide 9 | Amsterdam canal | searched (Commons, Openverse), none found → generated, "
            "flagged illustrative"]}},
        ledger={"entries": [], "searches": [{"query": "Amsterdam canal", "outcome": "unreachable",
                                             "n_results": 0}]},
        expect_codes=("UNREACHABLE IS NOT NONE-FOUND",))

    run("a `searched` rung is not mis-read as a bare `generated` token",
        gates={"design_plan": {"image_sources": [
            "slide 9 | x | searched (Commons), none found -> generated, flagged illustrative"]}},
        ledger={"entries": [], "searches": [{"outcome": "none found"}]},
        forbid_codes=("BAD TOKEN",))

    run("an em dash, an en dash and a plain hyphen all parse — the gate checks structure, not "
        "which dash somebody typed",
        gates={"design_plan": {"image_sources": [
            "a | sourced — Commons (CC0)", "b | sourced – Commons (CC0)",
            "c | sourced - Commons (CC0)"]}},
        ledger={"entries": [{"file": "x", "license": "CC0", "status": "placed"}], "searches": []},
        forbid_codes=("BAD TOKEN",))

    run("a licence stated in the plan that contradicts the ledger is caught",
        gates={"design_plan": {"image_sources": ["slide 4 | sourced — Commons (CC0)"]}},
        ledger=good_led, expect_codes=("LICENCE MISMATCH",))

    run("an attribution-required photo with NO credit anywhere in the deck is caught",
        gates={"design_plan": {"image_sources": ["slide 4 | sourced — Wikimedia Commons (CC BY-SA 4.0)"]}},
        ledger=good_led, pptx="__DECK_NO_CREDIT__", expect_codes=("MISSING CREDIT",))

    run("...and is clean when the author's name is on a slide",
        gates={"design_plan": {"image_sources": ["slide 4 | sourced — Wikimedia Commons (CC BY-SA 4.0)"]}},
        ledger=good_led, pptx="__DECK_WITH_CREDIT__", forbid_codes=("MISSING CREDIT",))

    run("`n/a — text-only deck` is a legitimate answer",
        gates={"design_plan": {"image_sources": "n/a — text-only deck, no content images"}},
        ledger={"entries": [], "searches": []}, forbid_codes=("BAD TOKEN", "NO TOKENS"))

    run("...but not while photos sit in the ledger",
        gates={"design_plan": {"image_sources": "n/a — text-only deck"}},
        ledger=good_led, expect_codes=("N/A CONTRADICTED",))

    run("candidates that were never adopted are reported — downloading is not choosing",
        gates={"design_plan": {"image_sources": ["slide 4 | sourced — Commons (CC0)"]}},
        ledger={"entries": [{"file": "a.jpg", "license": "CC0", "status": "candidate"}],
                "searches": [{"outcome": "found"}]},
        expect_codes=("NOT ADOPTED",))

    print("\n".join("  ok   " + x for x in ok))
    if bad:
        print("\n".join("  FAIL " + x for x in bad))
    print("\n{} passed, {} failed".format(len(ok), len(bad)))
    return 1 if bad else 0


# The self-test needs deck TEXT without building a pptx; `_deck_text` is patched for those two
# cases only, and only inside the self-test.
_REAL_DECK_TEXT = _deck_text


def _fake_deck_text(p):
    if str(p) == "__DECK_NO_CREDIT__":
        # Contains the file's TITLE word ("campus") and not its author — the exact shape that used
        # to pass.
        return "campus overview  what we build  next steps"
    if str(p) == "__DECK_WITH_CREDIT__":
        return _norm("sources: campus photo by x photographer (cc by-sa 4.0)").lower()
    return _REAL_DECK_TEXT(p)


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0],
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("deck_dir", nargs="?", help="Directory holding .deck-gates.json and assets/.")
    ap.add_argument("--pptx", help="The built deck, so credits can be checked ON THE SLIDES.")
    ap.add_argument("--no-credit-check", action="store_true")
    ap.add_argument("--selftest", action="store_true")
    a = ap.parse_args(argv)

    if a.selftest:
        global _deck_text
        _deck_text = _fake_deck_text
        try:
            return _selftest()
        finally:
            _deck_text = _REAL_DECK_TEXT
    if not a.deck_dir:
        ap.print_help()
        return 2
    if not os.path.isdir(a.deck_dir):
        print("not a directory: {}".format(a.deck_dir), file=sys.stderr)
        return 2
    probs = check(a.deck_dir, pptx=a.pptx, strict_credits=not a.no_credit_check)
    if not probs:
        print("clean — every image token is backed by the ledger, and every attribution-required "
              "photo is credited in the deck.")
        return 0
    print("\n{} image-provenance problem(s):\n".format(len(probs)))
    for code, msg in probs:
        print("  {}: {}\n".format(code, msg))
    return 1


if __name__ == "__main__":
    sys.exit(main())

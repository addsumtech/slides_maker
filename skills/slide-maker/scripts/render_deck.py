#!/usr/bin/env python3
"""Render a .pptx to one PNG per slide, so you can SEE each slide and catch
overflow / contrast / glyph problems before handing the deck back.

Cross-platform: works on macOS, Linux, WSL, and NATIVE Windows (PowerShell / cmd) —
no shell required. The .sh wrapper just delegates here.

Usage:
    python3 render_deck.py /path/to/deck.pptx [out_dir]
    # Windows:  python scripts\\render_deck.py C:\\path\\deck.pptx
Output: <out_dir>/slide01.png, slide02.png, ...   (default out_dir: ./render)

Requires: LibreOffice + pymupdf (python -m pip install pymupdf). One-time installs.
Override LibreOffice discovery with the SOFFICE env var (full path to the binary).
"""
import contextlib
import json
import os
import re
import sys
import shutil
import tempfile
import subprocess
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from written_reason import reason_width  # noqa: E402  (one shared definition, never a copy)


def find_soffice():
    """Locate the LibreOffice binary across macOS / Linux / WSL / native Windows.

    Order: $SOFFICE override -> anything on PATH -> known install locations.
    On Windows, prefer the sibling ``soffice.com`` (the console front-end that
    blocks until conversion finishes) over ``soffice.exe`` (which can detach and
    leave the PDF half-written).
    """
    def prefer_com(path):
        if path and path.lower().endswith("soffice.exe"):
            com = path[:-4] + ".com"
            if os.path.isfile(com):
                return com
        return path

    env = os.environ.get("SOFFICE")
    if env and os.path.isfile(env):
        return prefer_com(env)

    for cmd in ("soffice", "libreoffice", "soffice.com", "soffice.exe"):
        found = shutil.which(cmd)
        if found:
            return prefer_com(found)

    candidates = [
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        # Linux
        "/usr/bin/soffice", "/usr/bin/libreoffice", "/usr/local/bin/soffice",
        "/snap/bin/libreoffice", "/opt/libreoffice/program/soffice",
        # native Windows (default installer locations)
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # Windows install reached from WSL via the /mnt/c mount
        "/mnt/c/Program Files/LibreOffice/program/soffice.exe",
    ]
    # %ProgramFiles% may point somewhere non-default on Windows.
    for var in ("ProgramFiles", "ProgramFiles(x86)", "ProgramW6432"):
        base = os.environ.get(var)
        if base:
            candidates.append(os.path.join(base, "LibreOffice", "program", "soffice.exe"))

    for path in candidates:
        if os.path.isfile(path):
            return prefer_com(path)
    return None


# ── BATCHED GATE REPORTING ────────────────────────────────────────────────────────────────────
# `die()` used to be what it looks like: print one message, exit. On the hand-off gate that made
# every run report exactly ONE problem, and a thin `.deck-gates.json` then costs one fail → fix →
# re-run round-trip PER FIELD — measured shape: ~15 sequential stop classes (critic → build_shape
# → boldness → design_plan fields → type_scale → signature_proof → carried_by → concept →
# signature_move → form reach → arc → provenance → sameness → density), each round-trip re-sending
# the whole conversation at hand-off-time context, which is the most expensive point of the run.
# The repo already knew this: `codex_delivery_gate.py` accumulates an `errors` list and prints all
# of it, and `validate_review.py`'s docstring names "the ping-pong of one-error-at-a-time retries"
# as the anti-pattern it was built to avoid. This path was the only holdout.
#
# So in batching mode `die()` records the failure and raises `_GateStop`, which ends the CURRENT
# gate section (its remaining checks genuinely depend on what just failed) and lets every OTHER
# section still run. `check_handoff_gates` prints the whole list at the end. Nothing about the
# messages, the thresholds or the exit code changes — only how many runs it takes to see them.
#
# `_GateStop` derives from BaseException, like SystemExit and for the same reason: several gates
# call `die()` inside a `try/except Exception`, and an ordinary exception would be swallowed there
# — turning a hard gate into a silent pass. That is the one way this change could have made the
# file WORSE than the behaviour it replaces.
_COLLECTED = None       # list of (section, msg, code) while a batching run is in flight
_SECTION = None         # the gate section currently executing, for labelling


class _GateStop(BaseException):
    """One gate section failed; its remaining checks depend on what just failed."""


def _section(gates, key, what="an object"):
    """`gates[key]` as a dict, or a readable death — never an AttributeError.

    Every gate on this path read its record as `gates.get("x") or {}` and then called `.get()` on
    the result. That is correct for a missing key and wrong for a MIS-SHAPED one: a hand-written
    `.deck-gates.json` carrying `"a11y": "we don't need it"` — a plausible mistake, and the shape a
    model writing the file by hand actually produces — raised
    `AttributeError: 'str' object has no attribute 'get'`. Because AttributeError is not
    `_GateStop`, it escaped the section contract and took down the whole run with a traceback
    instead of naming the field. Measured on `sameness` and `a11y` alike before this existed.

    `scripts/deck_gates.py` writes the right shape; this is for the file someone edited by hand.
    """
    val = gates.get(key) if isinstance(gates, dict) else None
    if val is None or val == {} or val == "":
        return {}
    if not isinstance(val, dict):
        die("`.deck-gates.json` -> {!r} must be {}, not {} ({!r}). A waiver is a record, not a "
            "sentence:\n    {{\"{}\": {{\"waived\": \"<the reason>\", "
            "\"waived_category\": \"<kind>\"}}}}\n    `python3 scripts/deck_gates.py check "
            "<deck-dir>` reports every shape fault in one pass."
            .format(key, what, type(val).__name__, str(val)[:60], key))
    return val


@contextlib.contextmanager
def _gate_section(label):
    """Run one INDEPENDENT gate section: a failure inside it never suppresses the others."""
    global _SECTION
    prev, _SECTION = _SECTION, label
    try:
        yield
    except _GateStop:
        pass
    finally:
        _SECTION = prev


def die(msg, code=1):
    if _COLLECTED is not None:
        _COLLECTED.append((_SECTION, msg, code))
        raise _GateStop
    print(msg, file=sys.stderr)
    sys.exit(code)


def _png_is_flat(path):
    """True when the image is a single flat colour — i.e. not a render of anything.

    Measured: `signature_proof` required its PNGs to exist and be non-empty, and a 960x540
    rectangle of one grey satisfied the ANCHOR PROOF — the mechanism whose entire purpose is to
    put rendered evidence where a design decision is made. A proof that a blank file can satisfy
    proves nothing. Costs one PIL open; on any failure it returns False (unreadable is not a lie).
    """
    try:
        from PIL import Image
        with Image.open(path) as im:
            cols = im.convert("RGB").getcolors(maxcolors=64)
        return cols is not None and len(cols) <= 1
    except Exception:
        return False


def _report_carried_by(pptx_path, cb):
    """Say, as a number, how many `carried_by` slides are structurally distinct from the deck.

    WHY. `carried_by` was presence-checked only: a list of ≥2 slide numbers passed. Naming a slide
    is free, so the cheapest way to satisfy the field is to list three and put a token of the motif
    — a stripe, a matching hue — on two of them, while the actual structural work happens on one.
    Measured on a real build: carried_by=[1, 9, 12], and only slide 9 had a skeleton the deck did
    not otherwise use. The plan read bravely and the deck read safe, which is precisely the failure
    `boldness`/`signature_move` exist to prevent.

    NOT A FAILURE, deliberately. A signature move can legitimately live on colour, type or concept
    and touch no geometry at all — dying on that would push authors toward layout stunts. So this
    prints the count and warns; the author has to look at a number instead of at their own sentence.
    """
    try:
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        import lint_deck as _ld
        from pptx import Presentation as _P
        prs = _P(str(pptx_path))
        sw = prs.slide_width / 914400.0
        sh = prs.slide_height / 914400.0
        skels = [_ld._skeleton(_ld._boxes(s, sw, sh)) for s in prs.slides]
    except Exception as e:                      # never let a measurement break the delivery gate
        print(f"[gates] carried_by: not measurable on this deck ({type(e).__name__}) — "
              f"check by eye that each named slide is structurally distinct")
        return
    if not skels:
        return

    def _same(a, b):
        return len(a & b) / max(1, len(a | b)) >= 0.75

    # the deck's default page = the skeleton the most slides share
    modal, best = skels[0], 0
    for cand in skels:
        c = sum(1 for k in skels if _same(cand, k))
        if c > best:
            modal, best = cand, c
    distinct = [i for i in cb
                if isinstance(i, int) and 1 <= i <= len(skels) and not _same(skels[i - 1], modal)]
    print(f"[gates] carried_by structure: {len(distinct)}/{len(cb)} named slide(s) differ from the "
          f"deck's default skeleton (distinct: {distinct or 'none'})")

    # …and whether the DEVICE is drawn the same way on each, which is what the declaration claims.
    # The check above asks "is each named slide unlike the deck's default page" — three slides can
    # each be unusual in three unrelated ways and pass it.
    #
    # Scoped deliberately to the accent RULE, because that is the case that occurred and the only
    # one that is measurable without knowing what the move is. Measured on a delivered deck whose
    # plan declared carried_by=[5,10,11] as "the same line grammar": #1F3B2F at 41% of the canvas
    # and 3.2px · #8A8377 at 50% and 2.9px · #B4462A at 34%, 2.3px, and at a different y entirely.
    # Three colours, three spans, three weights.
    #
    # A generic structural test was tried first and DROPPED: mean pairwise skeleton likeness among
    # the named slides was 0.11 against a deck baseline of 0.12. That gap is noise, and reporting
    # it as evidence would be the confident-wrong-number failure this file keeps fixing.
    #
    # NOT A FAILURE, same carve as above: a signature move may live on colour, type or concept. If
    # the named slides do not all carry a rule, the device is not a rule and nothing is said.
    idx = [i - 1 for i in cb if isinstance(i, int) and 1 <= i <= len(skels)]
    if len(idx) >= 2:
        try:
            slides = list(prs.slides)
            rules = []
            for i in idx:
                bx = _ld._boxes(slides[i], sw, sh)
                cand = [b for b in bx if b["solid"] and not b["text"] and not b["bg"]
                        and b["fill"] and b["w"] > 0.15 * sw and b["h"] < 0.02 * sh
                        and b["t"] < 0.80 * sh]          # exclude the footer rule
                rules.append(max(cand, key=lambda z: z["w"]) if cand else None)
        except Exception:
            rules = []
        if rules and all(r is not None for r in rules):
            cols = {r["fill"] for r in rules}
            spans = [r["w"] / sw for r in rules]
            thick = [r["h"] * 144 for r in rules]
            span_off = (max(spans) - min(spans)) > 0.10
            thick_off = (max(thick) / max(min(thick), 0.1)) > 1.25
            if len(cols) > 1 or span_off or thick_off:
                bits = []
                for i, r in zip(idx, rules):
                    bits.append(f"p{i + 1} #{r['fill']} {r['w'] / sw:.0%}w {r['h'] * 144:.1f}px")
                print(f"        ⚠ the carried_by slides draw their rule THREE different ways: "
                      + " · ".join(bits) + ". `carried_by` says they carry ONE move; a device "
                      f"redrawn per page is not one move, it is three. Match the colour, the span "
                      f"and the weight, or say in the hand-off what the shared move actually is.")
            else:
                print(f"        carried_by rule is consistent across the named slides "
                      f"(#{list(cols)[0]}, {min(spans):.0%}-{max(spans):.0%} width)")
    if len(distinct) < 2:
        print("        ⚠ the signature move is doing structural work on fewer than 2 of the slides "
              "that claim it. If it lives on colour/type/concept instead, fine — say which in the "
              "hand-off. If it does not, the deck is safer than its plan and this is the moment to "
              "fix that, not after the user says '不够大胆'.")


def _report_form_reach(pptx_path):
    """Say how much of the CATALOGUE the build actually reached for, at the moment of hand-off.

    WHY, measured on this skill's own last deck. The library exposes ~174 helpers. The build called
    15, of which exactly one (`stat_row`) was a composed form; the other 59 shapes were raw
    `box`/`text`. `component_audit.py` reports precisely this and SKILL.md's pre-flight points at
    it — but the pointer is prose, and the author (me) never ran it. The cost was concrete: slide 3
    hand-built a track + fill + label out of two boxes, and `meter_bar`'s docstring opens by calling
    itself "the correct, reusable form of the hand-built 'track + fill + number' row".

    Bespoke composition is NOT the failure — a Mondrian page and a tessellation cannot come from a
    catalogue, and forcing components onto them would make the deck worse. The failure is never
    having LOOKED, which is invisible unless somebody prints the number. So this prints it and
    names the one command that answers it, and never blocks.
    """
    try:
        here = Path(__file__).resolve().parent
        sys.path.insert(0, str(here))
        import component_audit as _ca
        # the build script sits next to the deck by convention (build_deck.py / build_<name>.py)
        cands = sorted(Path(pptx_path).parent.glob("build*.py"))
        if not cands:
            return
        called = _ca._script_calls(str(cands[0]))
    except Exception:
        return
    forms = sorted(set(called) & set(_ca.FORM_GUARANTEE))
    prims = sum(1 for n in ("box", "text") if n in called)
    if len(forms) >= 3 or not prims:
        print(f"[gates] form reach: {len(forms)} named component(s) — {', '.join(forms) or 'none'}")
        return None
    print(f"[gates] form reach: {len(forms)} of {len(_ca.FORM_GUARANTEE)} named components "
          f"({', '.join(forms) or 'none'}); the rest of {cands[0].name} is raw box/text.")
    print("        Bespoke composition is legitimate and often right — but confirm it was CHOSEN, "
          "not defaulted to: `python3 scripts/sigs.py --list` (or --search <shape>) is one call and "
          "answers it. Measured cost of skipping it once: a meter_bar rebuilt out of two boxes.")
    # Returning the count is what turns this from a line nobody reads into a question somebody
    # answers. It still never blocks on the NUMBER — only on the absence of a decision. See the
    # `form_reach` waiver in the gate below.
    return {"forms": forms, "total": len(_ca.FORM_GUARANTEE), "script": cands[0].name}


def _report_plan_files(pptx_path):
    """Plan files left in the deliverable folder — the checkpoint that became a file.

    `checkpoint-convention.md` says it plainly: the checkpoint artifact is a compact table PASTED
    INTO THE CONVERSATION, and `content-plan.md` / `design-plan.md` must not be written into the
    deliverable folder — the conversation is the record, and a folder the user opens should hold
    the deck, not the working notes that produced it.

    That rule lived only as one sentence inside a reference file, which is the position most
    likely to be skipped: nothing anywhere reported a plan file, so writing one cost nothing and
    the user found the clutter instead of a gate. Advisory rather than blocking — a user may have
    ASKED for plan files, and that carve is in the same sentence — but it is now said out loud at
    the moment the folder is handed over.
    """
    try:
        d = Path(pptx_path).parent
        stray = sorted(p.name for p in d.glob("*.md")
                       if p.name.lower() in ("content-plan.md", "design-plan.md",
                                             "content_plan.md", "design_plan.md"))
    except Exception:
        return
    if stray:
        print("[gates] plan files in the deliverable folder: %s — the checkpoint artifact belongs "
              "in the conversation, not on disk (checkpoint-convention.md). Delete them unless the "
              "user asked for plan files." % ", ".join(stray))


def _report_palette_drift(pptx_path, declared):
    """Name the slides where deckkit's OWN default accents survived into a deck with its own palette.

    WHY. `set_palette()` exists because a component's colour defaults BIND AT IMPORT — SKILL.md
    says so ("a bare `deckkit.MAGENTA = ...` does NOT re-theme components whose signature default
    binds at import"). The natural build script does not call it: it defines its palette as local
    constants and passes them where the API takes a colour. Every component that takes one
    IMPLICITLY then keeps deckkit blue or magenta.

    Measured on a real 10-slide build whose declared palette was a single bound teal: `title_bar`
    drew its accent rule in deckkit BLUE on all eight interior slides and `bottom_callout` set its
    label in deckkit MAGENTA — two foreign hues on a deck whose entire argument was one colour
    meaning one thing. Zero hard findings, zero warnings, every hand-off gate passed. The
    semantic-colour contract is a required plan FIELD and nothing compared it to the pixels.

    Deliberately narrow: it looks only for deckkit's own shipped constants, not for "any hue not in
    the palette". A generic scan would flag every photo, chart theme and template colour; this one
    answers a bounded question — did the LIBRARY's defaults leak past the deck's own choices? —
    which is the failure that actually happens. Prints, never dies: a deck may legitimately keep
    one deckkit hue if it declared it.
    """
    try:
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        import deckkit as dk
        import lint_deck as _ld
        from pptx import Presentation
    except Exception:
        return
    want = set(re.findall(r"[0-9A-Fa-f]{6}", json.dumps(declared) if not isinstance(declared, str)
                          else declared))
    want = {w.upper() for w in want}
    # deckkit's shipped accent constants — the ones a component reaches for when the caller passes
    # nothing. Neutrals (DEEP/SLATE/MUTE/WHITE/TINT/LIGHT) are excluded on purpose: they are the
    # library's ink and paper, and a deck that keeps them has not lost its identity.
    lib = {}
    for name in ("BLUE", "MAGENTA", "TEAL", "GOLD", "STEEL", "VIOLET", "GREEN"):
        v = getattr(dk, name, None)
        if v is not None:
            lib[str(v).upper()] = name
    leaked = {h: n for h, n in lib.items() if h not in want}
    if not want or not leaked:
        return
    try:
        prs = Presentation(pptx_path)
        sw = prs.slide_width / 914400.0
        sh = prs.slide_height / 914400.0
    except Exception:
        return
    hits = {}
    for i, slide in enumerate(prs.slides, 1):
        try:
            boxes = _ld._boxes(slide, sw, sh)
        except Exception:
            continue
        for b in boxes:
            for key in ("fill", "line", "ink"):
                v = b.get(key)
                h = str(v).upper() if v else ""
                if len(h) == 6 and h in leaked:
                    hits.setdefault(leaked[h], set()).add(i)
    if not hits:
        return
    print("[gates] palette drift: deckkit's own default accent(s) survived into a deck that "
          "declared its own palette —")
    for name, pages in sorted(hits.items()):
        print("        {} ({}) on slide(s) {}".format(
            name, getattr(dk, name), ", ".join(str(p) for p in sorted(pages))))
    print("        These are component defaults bound at IMPORT, not choices: title_bar's accent "
          "rule, a callout's label, a marker. Call deckkit.set_palette(...) once after import so "
          "the whole component set follows the deck, or pass the colour explicitly at each call.")
    print("        (A hue you DID choose is not drift — declare it in the design plan's `palette` "
          "field and this goes quiet.)")


def _report_file_observations(pptx_path):
    """What the DECK shows, for the case where there is no record to read.

    Every hand-off gate keys off `.deck-gates.json`, so a deck with no record gets one message
    about the record and nothing about itself — which teaches that the failure is paperwork. The
    deck that motivated this had five slides of category content and zero icons; the record being
    absent is why nothing said so, but writing the record would not have fixed it either.

    File-derived only, no declaration involved, and it never decides anything: it prints and
    returns, and the caller still dies on the missing record.
    """
    try:
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        import lint_deck as _ld
        from pptx import Presentation as _P
        prs = _P(str(pptx_path))
        ev = _ld.icon_evidence(prs)
    except Exception:
        return
    if ev["icons"] == 0 and len(ev["categorical"]) >= 2:
        print("[gates] the deck itself, before the record: slides {} read as category sets and it "
              "carries NO icon at all.".format(", ".join(map(str, ev["categorical"]))))
        print("        On category-rich content an icon family is a design must "
              "(slide-design self-verify (g)) — this is a real miss, not a bookkeeping one, and "
              "it survives writing the record.")
        # FLUSH. `die()` writes to stderr unbuffered while this is stdout, block-buffered when
        # piped — so without the flush the one line about the actual deck lands BELOW forty lines
        # of JSON-shape instructions, which is where a reader has already stopped.
        sys.stdout.flush()


# The FOUR high-bar reasons a deck may carry no icon family. Icons are the DEFAULT on any deck with
# categorical / multi-item / conceptual content — they aid the 1-second read and reinforce the
# system, and skipping them is a real but HIGH-bar choice (user directive, 2026-08). A bare
# "not category-rich" no longer clears it — that was the casual skip this set closes. Each names a
# case where an icon family would genuinely HURT, not just be optional:
#   motif-dominant     a strong constructed motif icons would dilute (the ≤3-loud motif budget)
#   editorial-register a data / editorial register (FT/Economist) icons would cheapen into corporate
#   tiny-deck          a 1-2 slide ask
#   template-locked    a registered/provided template that carries its own marks or forbids them
# 🔴 They are NOT a licence to drop icons from the CATEGORICAL slides that do have them — a deck may
# be motif-dominant AND still put icons on its one roster page. The category explains the REST.
_ICON_NONE_CATEGORIES = ("motif-dominant", "editorial-register", "tiny-deck", "template-locked")


def _icon_none_category_holds(cat, pptx_path, flagged):
    """Is the CLASSIFIED reason true of the built file? Returns (ok, why_not).

    The category was checked against a fixed list of strings and nothing else, so any of the four
    words cleared the gate. Measured on this repo's own deck: `motif-dominant` was written for a
    deck, accepted, and the first human reader's first note was "icons should be here". The word
    was doing the work, not the fact. Each class makes a CHECKABLE claim about the artifact, so
    check it:

      motif-dominant     the deck must actually carry a LOUD motif — the thing icons would dilute
      tiny-deck          it must actually be small (<=2 content slides)
      template-locked    it must actually be built on a template, not a blank deck
      editorial-register left as declared: "this register would cheapen with icons" is a taste
                         claim about a look, and forcing a measurement onto it would invent one.

    🔴 And on every class, the waiver must name EVERY slide the gate flagged. Naming one of five
    was enough before, which let a re-decision cover the page the author had already thought about
    and skip the four they had not.
    """
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
    except Exception:
        return True, ""                       # never fail the gate on the reader itself
    if cat == "motif-dominant":
        try:
            import deckkit as _dk
            loud = [n for n, sl in enumerate(prs.slides, 1)
                    for sh in sl.shapes if _dk._is_motif(sh, loud=True)]
        except Exception:
            return True, ""
        if not loud:
            return False, ("category 'motif-dominant' claims icons would dilute a strong motif, but "
                           "the built deck carries NO loud motif at all (nothing tagged with "
                           "deckkit.tag_motif(..., loud=True)). Either the motif is not built yet, "
                           "or this is the casual skip wearing the strongest word on the list")
    if cat == "tiny-deck":
        n = len(prs.slides._sldIdLst)
        if n > 4:
            return False, ("category 'tiny-deck' on a {}-slide deck — the carve is for a 1-2 slide "
                           "ask, not for any deck whose author would rather not draw icons"
                           .format(n))
    if cat == "template-locked":
        # python-pptx's DEFAULT template already ships eleven named layouts, so "has layout names"
        # proves nothing — `blank_deck()` cleared this on the first try. Compare against that known
        # set instead: a deck whose layouts are exactly the stock ones is not carrying anyone's
        # template.
        STOCK = {"Title Slide", "Title and Content", "Section Header", "Two Content", "Comparison",
                 "Title Only", "Blank", "Content with Caption", "Picture with Caption",
                 "Title and Vertical Text", "Vertical Title and Text"}
        try:
            names = {str(getattr(l, "name", "")) for l in prs.slide_layouts}
        except Exception:
            names = set()
        if names and not (names - STOCK):
            return False, ("category 'template-locked' but the deck carries only python-pptx's "
                           "stock layouts — it is a blank deck, and there is no template to be "
                           "locked by. `deckkit.open_template()` is what makes this class true")
    return True, ""


def _icon_none_waived(gates):
    """True when `icon_family: none` was re-decided against the BUILT slides AND classifies WHY.

    Two things, not one. (1) It must name the slides it was re-decided against — proof someone
    looked at pixels, not a plan sentence written before any slide existed. (2) It must classify the
    reason from `_ICON_NONE_CATEGORIES` — because icons are now the default on categorical content,
    so the bar for skipping is a genuinely strong reason (a motif they'd dilute, a register they'd
    cheapen), not 'the content isn't category-rich'. Symmetric with the critic waiver's
    `waived_category`: a waiver that cannot name its own class is the casual skip in disguise.
    """
    d = _section(gates, "design_plan")
    named = d.get("icon_none_checked")
    cat = str(d.get("icon_none_category") or "").strip().lower()
    return (isinstance(named, list) and len(named) >= 1) and cat in _ICON_NONE_CATEGORIES


def _report_icon_waiver(pptx_path, fam, gates=None, delivery=None):
    """Name the slides that contradict an `icon_family: none - <reason>` waiver.

    WHY. The waiver is meant to stay satisfiable — a deliberately icon-free deck is a real choice.
    But the reason is free text written at PLAN time, before any slide exists, and nothing ever
    revisits it. Measured on a real build: the plan said 'none — 概念性内容,图标会变装饰', and the
    deck then shipped three category slides (church vs market · pronk vs vanitas · four genres)
    that are exactly the entity-rich case SKILL.md calls a design must. The waiver was not a lie
    when written; it was never re-tested against what got built. So test it against the built file.

    A peer group = 3+ text boxes at the same size, aligned in a row or a column: the shape of a
    category set. It over-counts (tables, timelines), which is why this prints slides and does not
    die — but naming '2, 6, 7' is much harder to wave past than a general nag.
    """
    if not isinstance(fam, str):
        return
    if not fam.strip().lower().startswith("none"):
        # THE OTHER DIRECTION, and it is the one the incident actually took. This guard was
        # installed because "a deck shipped with ZERO icons through every automated gate" — and it
        # then checked only `icon_family: none` against a deck that HAS icons, i.e. a stale record
        # in the harmless direction. A deck declaring `icon_family: lucide-outline, 24 icons` and
        # shipping none passed silently: the guard installed for the incident could not fire on it.
        #
        # Counted by IMAGE IDENTITY, not by position+size. The geometry signature below is a chrome
        # detector — it asks "does this rectangle recur?" — and a real icon SET is many different
        # pictures at similar sizes, so counting distinct geometries miscounts it. Hashing the blob
        # answers the question actually being asked: how many distinct small pictures are in here?
        try:
            from pptx import Presentation as _P1
            import hashlib as _h1
            _prs1 = _P1(str(pptx_path))
            _blobs, _per = {}, {}
            for _i1, _s1 in enumerate(_prs1.slides):
                for _sh1 in _s1.shapes:
                    try:
                        if _sh1.width / 914400.0 >= 1.2 or _sh1.height / 914400.0 >= 1.2:
                            continue
                        # `.image` raises on any non-picture shape, which is the test — no
                        # MSO_SHAPE_TYPE import needed, and nothing silently swallows a NameError.
                        _k1 = _h1.sha256(_sh1.image.blob).hexdigest()[:16]
                    except Exception:
                        continue
                    _blobs[_k1] = _blobs.get(_k1, 0) + 1
                    _per.setdefault(_k1, set()).add(_i1)
            _n_sl1 = len(_prs1.slides) or 1
            # A logo is ONE picture stamped on most slides; an icon set is many pictures on a
            # minority. Drop any blob appearing on more than half the deck.
            _icons = {k: v for k, v in _blobs.items() if len(_per[k]) <= max(2, 0.5 * _n_sl1)}
            if not _icons:
                # 🔴 die, not print. This guard was installed because "a deck shipped with ZERO
                # icons through every automated gate" — and for its whole life it answered that
                # incident with a printed line and exit 0, which is the same outcome the incident
                # had. The Codex path already BLOCKS this exact case (`check_icons`: "categorical
                # slide N needs rendered icon evidence or a documented waiver", plus a check that
                # the slide's build function calls a deckkit icon helper), so the shared path was
                # the weaker of the two on the one class of failure this field exists for.
                #
                # Only THIS direction is fatal. The `icon_family: none` arm below stays a print
                # because its peer-group detector over-counts by design (tables, timelines) — its
                # own docstring says so, and a detector that over-counts must not hold a deck.
                die("`design_plan.icon_family` declares {!r} but the deck contains NO icon-sized "
                    "picture that is not deck chrome.\n"
                    "  The record and the file disagree: either the icons were never built, or the "
                    "field is stale.\n\n"
                    "  Build them (`scripts/icons.py <family>:<name> <out>.png`, placed with "
                    "deckkit.icon / icon_tile),\n"
                    "  or record the truth — `\"icon_family\": \"none — <why this deck carries no "
                    "icons>\"` — which is a\n"
                    "  legitimate choice on a figure-dominated or sober-register deck and is "
                    "checked in the other direction.\n\n"
                    "  (This is the direction the guard was written for. It printed and exited 0 "
                    "until now, which is\n"
                    "  exactly what happened the first time.)".format(fam.strip()[:60]))
        except Exception:
            pass
        return
    # Caught on its first run: a rebuild ADDED icons and the plan record still said `none`, because
    # nothing ever compares the record to the file. Say so before looking for peer groups — a record
    # that disagrees with the deck makes every other field in it suspect too.
    try:
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        import lint_deck as _ld0
        from pptx import Presentation as _P0
        _prs0 = _P0(str(pptx_path))
        _sw0, _sh0 = _prs0.slide_width / 914400.0, _prs0.slide_height / 914400.0
        # A repeated LOGO is icon-SIZED and would be miscounted as an icon family — measured: a
        # deck stamping one 0.6in mark on 8 slides counted 8 "icons" and would have been told its
        # record was stale. Chrome repeats at the same geometry; content does not. Drop any picture
        # whose position+size recurs on 3+ slides before counting.
        _sig, _n_sl = {}, len(_prs0.slides) or 1
        for _s0 in _prs0.slides:
            for _b0 in _ld0._boxes(_s0, _sw0, _sh0):
                if _b0.get("pic") and not _b0.get("bg") and _b0["w"] < 1.2 and _b0["h"] < 1.2:
                    _k = (round(_b0["l"], 1), round(_b0["t"], 1), round(_b0["w"], 1), round(_b0["h"], 1))
                    _sig[_k] = _sig.get(_k, 0) + 1
        # Chrome repeats on MOST slides; a content icon set appears on a minority of them. A bare
        # ">=3 occurrences = chrome" rule was wrong in both directions (measured): it cleared a
        # 3-icon row that recurred on 3 of 8 slides, which is content, while a logo on 8 of 8 is not.
        n_ic = sum(c for c in _sig.values() if c <= max(2, 0.5 * _n_sl))
        if n_ic >= 2:
            print(f"[gates] icon waiver: the plan says `icon_family: none` but the deck contains "
                  f"{n_ic} icon-sized image(s) — the record is stale. Update it to the family you "
                  f"actually shipped; a plan that disagrees with the file is not a record of anything.")
            return
    except Exception:
        pass
    try:
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        import lint_deck as _ld
        from pptx import Presentation as _P
        prs = _P(str(pptx_path))
        sw, sh = prs.slide_width / 914400.0, prs.slide_height / 914400.0
    except Exception:
        return
    # The peer-group rule now lives in lint_deck.categorical_slides — ONE definition, imported by
    # this path and by codex_delivery_gate. It was inline here and absent there, which is exactly
    # how the two gates came to disagree about what "checked" means for icons.
    hits = _ld.categorical_slides(prs, sw, sh)
    if len(hits) >= 2:
        print(f"[gates] icon waiver: `icon_family: none` but slides {hits} carry parallel label sets "
              f"(3+ peers in a row/column) — the entity-rich case SKILL.md calls a design must.")
        print("        ⚠ Re-decide NOW against the built slides, not against the plan sentence: either "
              "give these a family, or restate in the hand-off why THESE specific slides are better "
              "without one. A plan-time waiver that was never re-tested is how a deck ships zero icons.")
        # 🔴 And at deck scale it stops being a prompt. The detector OVER-COUNTS by construction
        # (tables, timelines, stat rows share the peer shape), so one or two slides can never hold a
        # deck — that caveat is what keeps it honest. But a deck with ZERO icon-sized pictures
        # anywhere AND three or more pages of parallel label sets is not an over-count; it is the
        # shape of a plan sentence that was written once and never re-tested. Measured twice in one
        # session: a Kimi-built deck and a deck built here, both with `icon_family: none - <reason>`,
        # both shipping zero icons past every gate. Same composite shape as sameness/timidity —
        # several weak signals before a hold, and a written escape that must name the SLIDES.
        try:
            _tot = sum(c for c in _sig.values() if c <= max(2, 0.5 * _n_sl))
        except Exception:
            _tot = 0
        # The SAME floors sameness and timidity carry, for the same reason: a 小红书 carousel, a
        # poster, or a five-slide ask is not a deck that owes an icon family. Without these the
        # block fired on a 3-slide portrait fixture — caught by the suite, which is what it is for.
        try:
            _asp = float(prs.slide_width) / float(prs.slide_height)
        except Exception:
            _asp = 1.78
        # Threshold lowered 3 → 2 (user directive: icons are a near-universal must). The detector
        # over-counts (tables/timelines/stat-rows share the peer shape), so 2 will sometimes fire on
        # a deck that legitimately wants none — but the cost of a false fire is now ONE line
        # (classify the reason), while the cost of a false PASS is a categorical deck that shipped
        # zero icons. Asymmetric, so the threshold sits on the side that asks the question. The
        # floors (aspect, non-surface, >=8 body slides below) still exclude carousels/posters/tiny
        # asks, and the waiver stays satisfiable for the real motif-dominant / editorial cases.
        # …and the CLASSIFIED reason must be true of the built file, not merely be one of four
        # accepted words. `motif-dominant` on a deck with no loud motif, `tiny-deck` on twelve
        # slides, `template-locked` on a blank deck — each was accepted, because the check compared
        # a string to a list. Measured here: `motif-dominant` cleared this gate on a deck whose
        # first human reader's first note was "icons should be here".
        _waived = _icon_none_waived(gates)
        if _waived:
            _cat = str(_section(gates, "design_plan").get("icon_none_category") or "").strip().lower()
            _ok, _why = _icon_none_category_holds(_cat, pptx_path, hits)
            if not _ok:
                _waived = False
                print("  [gates] icon waiver REJECTED — {}".format(_why))
            else:
                _named = {str(x).strip().lower().replace("slide", "").strip()
                          for x in (_section(gates, "design_plan").get("icon_none_checked") or [])}
                # Naming only SOME of the flagged slides is REPORTED, not refused. It was a hold
                # for one run and it turned out to reject a legitimate `editorial-register` waiver
                # that had named two of a dozen flagged pages — and the problem this batch actually
                # had was a category that was FALSE of the built file, which the check above now
                # catches precisely. A second, blunter rule that fires on lawful use is how a gate
                # earns the reflex to waive it.
                _missing = [h for h in hits if str(h) not in _named]
                if _missing:
                    print("  [gates] icon waiver names {} of {} flagged slide(s) — not re-decided: "
                          "{}. Look at those before the hand-off; the category covers the REST, not "
                          "the pages you have not opened."
                          .format(len(hits) - len(_missing), len(hits),
                                  ", ".join(str(m) for m in _missing)))
        _cheap = (len(hits) >= 2 and _tot == 0 and _asp >= 1.2
                  and delivery != "surface" and not _waived)
        # Only NOW pay for the real content-slide count — the same `body_n` sameness uses, which
        # excludes the cover, the closer and any declared appendix run. Raw slide count is the
        # wrong number: a 12-slide deck whose slides 5+ are declared reference material has three
        # content pages, and holding it for lacking an icon family is nonsense. Measured by the
        # suite on exactly that fixture.
        _body = 99
        if _cheap:
            try:
                _st, _ = _sameness_stats(str(pptx_path), delivery or "presented")
                _body = int(_st.get("body_n") or 0)
            except Exception:
                _body = 99
        if _cheap and _body >= 8:
            die("`icon_family: none` on a deck with ZERO icon-sized pictures and {} pages of "
                "parallel label sets ({}) — icons are the DEFAULT on categorical content: they aid "
                "the 1-second read and reinforce the system, so skipping them is a HIGH-bar choice, "
                "not a plan sentence.\n"
                "    Build the family (scripts/icons.py <lib>:<name> <out>.png, placed with "
                "deckkit.icon / icon_tile) — this is almost always the right move —\n"
                "    or record the re-decision AGAINST THE BUILT SLIDES with a CLASSIFIED reason "
                "(one of: {}):\n"
                '    {{"design_plan": {{"icon_family": "none — <why an icon family would HURT here>", '
                '"icon_none_checked": ["slide {}", "..."], "icon_none_category": "motif-dominant"}}}}\n'
                "    A bare 'not category-rich' no longer clears it; name the class of reason, the "
                "way the critic waiver names its own.".format(
                    len(hits), hits, " | ".join(_ICON_NONE_CATEGORIES), hits[0]))


def _tail(text, limit=4000):
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return "...<truncated>...\n" + text[-limit:]


# Exactly what a render writes: slide01.png … slideNN.png and the two bookend thumbnails. Used to
# decide what may be deleted when the out dir is shared with the user's own files.
_RENDER_PNG = re.compile(r"^(slide\d{2,}|thumb_first|thumb_last|contact(_\d{2})?)\.png$")


def _rels_targets(xml_bytes, base_dir):
    """Parse a .rels part -> {rId: normalized package path}.

    Parsed as XML, not by regex: OOXML allows the Id/Target attributes in either order, and a
    Target may be absolute ("/ppt/slides/slide1.xml"). Getting this wrong is not a cosmetic bug —
    a rId that fails to resolve silently drops a slide from the deck order, which shifts every
    index after it and writes real slides into the wrong PNG filenames.
    """
    import posixpath
    import xml.etree.ElementTree as ET
    out = {}
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return out
    for rel in root:
        rid = rel.get("Id")
        tgt = rel.get("Target")
        if not rid or not tgt or (rel.get("TargetMode") or "") == "External":
            continue
        if tgt.startswith("/"):
            out[rid] = tgt.lstrip("/")
        else:
            out[rid] = posixpath.normpath(posixpath.join(base_dir, tgt))
    return out


def _slide_fingerprints(pptx):
    """One content hash per slide, in deck order, plus any reason the deck cannot be diffed.

    Each slide's hash covers its own XML, its rels, the BYTES of every media part it references,
    and a DECK-GLOBAL digest (presentation.xml, theme, masters, layouts and the media those
    reference) — so a re-cropped photo, a swapped plate, a new theme or a different canvas size all
    count as changes.

    Returns (hashes, blockers). `blockers` is a list of human-readable reasons the caller must fall
    back to a full render; it is never a soft signal. Anything that makes the slide->page mapping
    uncertain belongs here, because a wrong mapping writes a real slide into another slide's PNG.
    """
    import hashlib
    import zipfile
    blockers = []
    with zipfile.ZipFile(pptx) as z:
        names = set(z.namelist())
        _part_cache = {}

        def part_digest(name):
            """Hash a package part ONCE. A background plate shared by every slide was previously
            re-read and re-hashed per referencing slide, making cost O(slides x media-bytes)."""
            d = _part_cache.get(name)
            if d is None:
                d = hashlib.sha256(z.read(name)).digest()
                _part_cache[name] = d
            return d
        pres_rels = _rels_targets(z.read("ppt/_rels/presentation.xml.rels"), "ppt") \
            if "ppt/_rels/presentation.xml.rels" in names else {}

        # ---- deck-global digest -------------------------------------------------------
        # docProps/* is excluded on purpose: core.xml carries a modified-timestamp that changes on
        # every save, which would force a full render every time and delete the feature.
        # ppt/media and ppt/notesSlides are excluded here — slide media is covered per-slide below,
        # and notes never reach a rendered pixel. But media referenced by a LAYOUT, MASTER or THEME
        # is deck-global and IS folded in: swapping a master's background image in place changes no
        # slide's XML, so without this the whole deck would render stale under "no slide changed".
        gh = hashlib.sha256()
        global_media = set()
        for n in sorted(names):
            if not n.startswith("ppt/"):
                continue
            if n.startswith(("ppt/slides/", "ppt/media/", "ppt/notesSlides/")):
                continue
            gh.update(n.encode())
            gh.update(z.read(n))
            if n.startswith(("ppt/slideLayouts/_rels/", "ppt/slideMasters/_rels/", "ppt/theme/_rels/")):
                base = n.rsplit("/_rels/", 1)[0]
                for tgt in _rels_targets(z.read(n), base).values():
                    if "/media/" in "/" + tgt:
                        global_media.add(tgt)
        for m in sorted(global_media):
            if m in names:
                gh.update(m.encode())
                gh.update(part_digest(m))
        global_digest = gh.digest()

        # ---- deck order ---------------------------------------------------------------
        import re
        pres = z.read("ppt/presentation.xml").decode("utf-8", "replace")
        rids = re.findall(r'<p:sldId[^>]*r:id="([^"]+)"', pres)
        order = [pres_rels.get(r) for r in rids]
        if any(t is None for t in order):
            blockers.append("a slide reference could not be resolved in presentation.xml.rels")
            order = [t for t in order if t]

        fps = []
        for sname in order:
            if sname not in names:
                # Never hash a missing part as b"" — every affected slide would collapse to the
                # same constant and real edits would stop registering, forever.
                blockers.append("slide part missing from the package: {}".format(sname))
                continue
            h = hashlib.sha256()
            h.update(global_digest)
            body = z.read(sname)
            h.update(body)
            if b'type="slidenum"' in body:
                # An auto slide-number field renumbers itself inside a subset.
                blockers.append("deck uses auto slide-number fields")
            if re.search(rb'<p:sld\b[^>]*\bshow="0"', body):
                # LibreOffice DROPS hidden slides from the PDF, so PDF page N is no longer deck
                # slide N — on the full path too. Verified with a 4-slide deck: one hidden slide
                # produced a 3-page PDF.
                blockers.append("deck has hidden slides")
            rname = sname.replace("slides/", "slides/_rels/") + ".rels"
            rel_bytes = z.read(rname) if rname in names else b""
            h.update(rel_bytes)
            for tgt in sorted(_rels_targets(rel_bytes, "ppt/slides").values()):
                if "/media/" in "/" + tgt and tgt in names:
                    h.update(tgt.encode())
                    h.update(part_digest(tgt))
            fps.append(h.hexdigest())

        if len(fps) != len(rids):
            blockers.append("deck order could not be fully resolved ({} of {} slides)".format(
                len(fps), len(rids)))
    # de-dup, order-stable
    seen, uniq = set(), []
    for b in blockers:
        if b not in seen:
            seen.add(b); uniq.append(b)
    return fps, uniq


def _subset_pptx(src, keep_idx, dest):
    """Copy `src` keeping ONLY the 0-indexed slides in `keep_idx` (order preserved).

    Deleting sldId entries is safe and cheap; the orphaned parts stay in the package and
    LibreOffice ignores them. Verified pixel-identical to the same slide rendered from the
    full deck, which is what makes the fast path trustworthy rather than merely fast.
    """
    import shutil
    from pptx import Presentation
    shutil.copy(src, dest)
    prs = Presentation(dest)
    lst = prs.slides._sldIdLst
    for i, sid in enumerate(list(lst)):
        if i not in keep_idx:
            lst.remove(sid)
    prs.save(dest)
    return dest


def _profile_pool_dir():
    """Where the reusable LibreOffice profiles live. Same convention as the icon cache."""
    env = os.environ.get("SLIDE_MAKER_CACHE")
    if env:
        return os.path.join(env, "lo-profiles")
    if sys.platform == "win32":
        base = os.environ.get("LOCALAPPDATA") or os.path.expanduser(r"~\AppData\Local")
    elif sys.platform == "darwin":
        base = os.path.expanduser("~/Library/Caches")
    else:
        base = os.environ.get("XDG_CACHE_HOME") or os.path.expanduser("~/.cache")
    return os.path.join(base, "slide-maker", "lo-profiles")


PROFILE_POOL_SLOTS = 8       # more concurrent renders than that, and the extras go throwaway


class _Profile(object):
    """A LibreOffice user profile to render with, plus how to let it go.

    Building a user profile from scratch is most of what a headless `--convert-to` run spends
    its time on: measured 3.71s per export with a fresh profile against 2.16s once the profile
    already exists — ~42% of the export, on every render, spent recreating the same directory.

    The old code paid that every time by design, and the design had two real reasons, both of
    which this pool has to keep:

      1. **Parallel renders must not fight.** The large-deck section fan-out renders several
         decks at once, and two `soffice` processes pointed at ONE profile collide: measured,
         one of the pair exits 1 having written no PDF at all.
      2. **The user's own LibreOffice may be open.** A render must not touch the GUI's profile.

    So this is a POOL, not a shared singleton: N private slots, each held under an exclusive
    `flock` for the duration of one render. Concurrent renders land on different slots; a slot
    is only ever used by one process at a time, which is exactly the invariant a throwaway
    profile gave for free. Nothing here goes near LibreOffice's default profile.

    When every slot is busy, or the platform has no `flock` (Windows), or anything at all goes
    wrong, we hand back a throwaway profile — i.e. precisely the old behaviour. The pool is an
    optimization; it is never the reason a deck fails to render.
    """

    def __init__(self, path, pooled, lock_fh=None):
        self.path = path
        self.pooled = pooled
        self._lock_fh = lock_fh

    def release(self, discard=False):
        if not self.pooled:
            shutil.rmtree(self.path, ignore_errors=True)
            return
        if discard:
            # This slot was in hand for a render that failed. It may be the cause (a profile
            # torn by a killed process, or one written by a different LibreOffice version), so
            # do not leave it for the next run to trip over.
            shutil.rmtree(self.path, ignore_errors=True)
        if self._lock_fh is not None:
            try:
                import fcntl
                fcntl.flock(self._lock_fh.fileno(), fcntl.LOCK_UN)
            except Exception:
                pass
            try:
                self._lock_fh.close()
            except Exception:
                pass


def _acquire_profile():
    """Take a pooled profile slot, or fall back to a throwaway one."""
    if os.environ.get("SLIDE_MAKER_NO_PROFILE_POOL"):
        return _Profile(tempfile.mkdtemp(prefix="lo_render_"), pooled=False)
    try:
        import fcntl                                    # POSIX only; Windows keeps the old path
    except ImportError:
        return _Profile(tempfile.mkdtemp(prefix="lo_render_"), pooled=False)
    try:
        root = _profile_pool_dir()
        os.makedirs(root, exist_ok=True)
        for i in range(PROFILE_POOL_SLOTS):
            slot = os.path.join(root, "slot{:02d}".format(i))
            os.makedirs(slot, exist_ok=True)
            # The lock file lives BESIDE the slot, not inside it: `release(discard=True)` deletes
            # the slot tree, and unlinking the file we hold the lock on would drop the lock.
            fh = open(os.path.join(root, "slot{:02d}.lock".format(i)), "a+")
            try:
                fcntl.flock(fh.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
            except OSError:
                fh.close()
                continue                                # another render holds this slot
            return _Profile(slot, pooled=True, lock_fh=fh)
    except Exception:
        pass
    return _Profile(tempfile.mkdtemp(prefix="lo_render_"), pooled=False)


def _render_pdf(soffice, src, outdir):
    """pptx -> pdf via a private LibreOffice profile, into an EMPTY private directory.

    `outdir` must not already contain `<src-stem>.pdf`. Rendering into a directory that may hold a
    previous PDF is how a FAILED conversion gets read back as success: the caller checks only that
    the file exists, and a stale one satisfies that. Reproduced with real LibreOffice on a deck it
    refuses to convert — the run printed "rendered N slides" and exit 0 over untouched output.
    Returns (pdf_path, result, cmd); the caller must check `result.returncode`.

    The profile comes from `_acquire_profile()` — a pooled, warm one when a slot is free, a
    throwaway otherwise. On a non-zero exit with a POOLED profile we throw that slot away and
    retry ONCE on a throwaway, so a poisoned slot degrades into the old behaviour instead of
    into a failed render. A genuine bad deck fails both times and is reported as before.
    """
    pdf = os.path.join(outdir, os.path.splitext(os.path.basename(src))[0] + ".pdf")

    def _once(profile):
        cmd = [soffice, "-env:UserInstallation=" + Path(profile.path).as_uri(),
               "--headless", "--convert-to", "pdf", "--outdir", outdir, src]
        try:
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                                    text=True, timeout=300)
        except subprocess.TimeoutExpired:
            profile.release(discard=True)
            die("LibreOffice render exceeded 300s and was killed — the .pptx may be malformed or "
                "hostile (e.g. a decompression bomb). Nothing was produced from {}.".format(src))
        return result, cmd

    profile = _acquire_profile()
    result, cmd = _once(profile)
    if result.returncode != 0 and profile.pooled:
        profile.release(discard=True)
        print("note: render failed on a pooled LibreOffice profile — discarding it and retrying "
              "with a fresh one", file=sys.stderr)
        retry = _Profile(tempfile.mkdtemp(prefix="lo_render_"), pooled=False)
        result, cmd = _once(retry)
        retry.release()
        return (pdf, result, cmd)
    profile.release()
    return (pdf, result, cmd)


# --- parallel rasterization -------------------------------------------------------------
# Rasterizing is the one render stage that is embarrassingly parallel: each page is an
# independent draw into its own file. It is also, on an image-heavy deck, the stage that
# actually costs something.
#
# Two things about the shape below are load-bearing, both measured rather than reasoned:
#
#  1. **Chunk contiguously; open the PDF ONCE per worker.** The obvious version — one task
#     per page — was measured at 18.5s against 2.9s serial on a 16-page image deck: 6.5x
#     SLOWER, because every task re-opens and re-parses the whole PDF. Chunked, the same
#     deck rasterizes in 1.27s (4 workers).
#  2. **Never let this fail a render.** It is an optimization on a step that already worked.
#     Any failure — no process pool in a locked-down sandbox, a worker dying, pickling
#     trouble — returns False and the caller runs the original serial loop. A deck that
#     renders slower is a nuisance; a deck that does not render is a broken run.
#
# Verified byte-identical to the serial loop (sha256 over every PNG, text and image decks).
# `SLIDE_MAKER_RENDER_WORKERS=1` forces the serial path.
RASTER_MIN_PAGES = 8          # below this the pool costs more than the pages save
RASTER_MAX_WORKERS = 4        # measured: 4 workers ~= 8 workers (1.27s vs 1.25s), half the RAM
# Page COUNT was the wrong gate on its own: what a worker pool saves is DECODE work, and a page's
# decode cost is set by its weight, not by how many neighbours it has. Measured serial rasterize:
#
#     18-page text/CJK deck   2.5-5.5 KB/page,  0 images   0.05-0.08s   (2.9-4.7 ms/page)
#     16-page image deck       84.5 KB/page,   16 images   0.47s        (29.3 ms/page)
#
# Fanning out costs ~70 ms in process starts and `import fitz` per worker, so on the text decks the
# pool spent more than the whole job — measured +2.1% wall clock, i.e. the optimization ran and made
# the render slower. On the image deck it is a clear -11.4%. PDF bytes-per-page separates the two by
# ~15x, which is a wide enough margin that a threshold in the middle is not a tuned constant: below
# it there is not enough decode work to pay a single process start, whatever the page count.
# This changes only WHICH path runs; both are already asserted byte-identical by
# tests/test_render_parallel.py, so a mis-set threshold costs milliseconds and never a wrong PNG.
RASTER_MIN_BYTES_PER_PAGE = 20 * 1024


def _rasterize_chunk(args):
    """Render a contiguous run of pages in one worker. Module-level so it is picklable.

    `thumbs` maps a page index this chunk owns to a thumbnail name. The bookend thumbnails
    are produced HERE, by the worker that already rasterized that page at 2x, and only after
    the whole chunk is done — because a page's thumbnail is not independent of whether that
    page was rendered at 2x first. Measured: `thumb_first` of an image-heavy deck hashes
    differently when rendered from a fresh document than when rendered after the 2x pass
    (both individually deterministic). MuPDF is scaling a cached decode in one case and
    doing a scaled decode in the other. So the thumbnail follows its page into the worker,
    which keeps the ordering — 2x pass, then thumbnails — byte-for-byte what it has always
    been. Rendering them in the parent instead silently changes two PNGs the critic looks at.
    """
    pdf, idxs, out, thumbs = args
    import fitz
    doc = fitz.open(pdf)
    try:
        for i in idxs:
            doc[i].get_pixmap(matrix=fitz.Matrix(2, 2)).save(
                os.path.join(out, "slide{:02d}.png".format(i + 1)))
        for i, name in sorted(thumbs.items()):
            page = doc[i]
            zoom = 240.0 / max(1.0, page.rect.width)
            page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(
                os.path.join(out, name + ".png"))
    finally:
        doc.close()
    return len(idxs)


def _raster_workers(n_pages, pdf=None):
    """How many workers to rasterize `n_pages` with — 1 means 'use the serial loop'.

    `pdf` is optional so an explicit `SLIDE_MAKER_RENDER_WORKERS` and the page-count floor still
    answer without it; when it is given, page WEIGHT gets the final say (see
    RASTER_MIN_BYTES_PER_PAGE). An unreadable size is treated as heavy — the pool is the safe
    guess there, since being wrong costs ~70 ms and the serial path is the fallback anyway.
    """
    env = os.environ.get("SLIDE_MAKER_RENDER_WORKERS")
    if env:
        try:
            return max(1, min(int(env), n_pages))
        except ValueError:
            pass                                  # a typo'd override must not break the render
    if n_pages < RASTER_MIN_PAGES:
        return 1
    if pdf is not None:
        try:
            if os.path.getsize(pdf) / float(n_pages) < RASTER_MIN_BYTES_PER_PAGE:
                return 1                          # not enough decode work to pay a process start
        except Exception:
            pass
    try:
        cores = os.cpu_count() or 1
    except Exception:
        cores = 1
    return max(1, min(RASTER_MAX_WORKERS, cores - 1, n_pages))


def _rasterize_parallel(pdf, n_pages, out):
    """Try to rasterize all pages across worker processes. Returns True if it wrote them all.

    Returns False (writing nothing, or leaving a partial set the serial loop then overwrites)
    whenever the fast path is unavailable or fails, so the caller can fall back.
    """
    workers = _raster_workers(n_pages, pdf)
    if workers < 2:
        return False
    per = -(-n_pages // workers)
    parts = [list(range(k * per, min(n_pages, (k + 1) * per))) for k in range(workers)]
    parts = [p for p in parts if p]
    # Hand each bookend thumbnail to whichever chunk owns its page (see _rasterize_chunk).
    want = {0: "thumb_first", n_pages - 1: "thumb_last"} if n_pages else {}
    jobs = [(pdf, p, out, {i: nm for i, nm in want.items() if i in p}) for p in parts]
    try:
        from concurrent.futures import ProcessPoolExecutor
        with ProcessPoolExecutor(max_workers=len(parts)) as ex:
            done = list(ex.map(_rasterize_chunk, jobs))
    except Exception as exc:
        print("note: parallel rasterize unavailable ({}) — using the serial path".format(exc),
              file=sys.stderr)
        return False
    if sum(done) != n_pages:
        # A worker returned but wrote fewer pages than it was given. Do not trust a partial
        # set: say so and let the serial loop rewrite every page.
        print("note: parallel rasterize covered {}/{} pages — redoing serially".format(
            sum(done), n_pages), file=sys.stderr)
        return False
    return True


def _sha256(path):
    import hashlib
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


GATES_FILE = ".deck-gates.json"


# Structural rows are excluded from the per-slide takeaway uniqueness check: a cover and a closing
# legitimately restate the deck's one sentence, and blocking that would push authors to invent a
# second sentence for a page whose whole job is to repeat the first. Same vocabulary
# `arc_divergence.py` excludes from its `order` axis, for the same reason.
_STRUCTURAL_ROLES = frozenset((
    "cover", "agenda", "divider", "closing", "section", "thanks", "qa",
    "封面", "目录", "分隔", "收尾", "结尾", "致谢", "问答"))

# Values that look filled and say nothing. Kept SHORT on purpose: a long blocklist starts rejecting
# honest short answers, and the width floor beside it already catches most emptiness. These are the
# strings a run reaches for when it is filling a field rather than answering it.
_STUBS = frozenset((
    "n/a", "na", "none", "tbd", "todo", "-", "--", "—", "...", "…", "?", "x",
    "source", "sources", "see source", "see sources", "same", "ditto",
    "无", "暂无", "待定", "见来源", "同上", "略", "无来源", "见上"))


# The angle-bracket TEMPLATE placeholder. `python3 scripts/arc_divergence.py --template` and this
# gate's own die() messages print fillable skeletons whose fields read `<label-A>`,
# `<the question this room is actually asking>`, `<the one sentence this slide leaves behind>`.
# Measured: the raw skeleton, pasted with ZERO edits, passed the arc / slides / checkpoint gates —
# every text floor here checks WIDTH, and a placeholder is wide. The template is SUPPOSED to pass
# its own divergence checker (test_arc_divergence pins that — a skeleton that modelled a strawman
# would teach one), so the guard lives in the GATE, not in arc_divergence's rules: the skeleton may
# pass as a template, it may not pass as a delivered deck. Same predicate codex_delivery_gate
# already applies to `fast_basis`/`none_opt_in` (startswith "<"), tightened to REQUIRE a closing `>`
# so a real "<50%" or "reduced latency <10ms" survives while "<the ask>" does not.
_PLACEHOLDER = re.compile(r"<[^<>]{2,}>")


def _has_placeholder(v):
    """True when a string still carries an angle-bracket template placeholder (`<...>`)."""
    return bool(_PLACEHOLDER.search(str(v or "")))


def _is_stub(v):
    """True when a field is filled with a token that answers nothing — or a raw template placeholder."""
    t = re.sub(r"[\s。．.、,，;；:：]+$", "", str(v or "").strip().lower())
    return (not t) or t in _STUBS or _has_placeholder(v)


def _norm_takeaway(v):
    """Normalise a memory sentence for equality: whitespace and terminal punctuation only.

    Deliberately NOT fuzzy. A near-duplicate pair is a judgement call the author may have made on
    purpose; an exact duplicate is a plan that was written once and pasted down the column, and
    that is the only case worth failing a build over.
    """
    t = str(v or "").strip().lower()
    t = re.sub(r"\s+", "", t)
    return re.sub(r"[。．.!！?？,，、;；]+$", "", t)


def _arc_verdict(candidates):
    """Score the candidate arcs with `arc_divergence.check` AT GATE TIME.

    Imported lazily and by name so this file states the dependency where it is used: the gate's
    verdict has to come from the same code the CLI runs, or the two drift and the gate starts
    certifying a rule the script no longer applies. Any exception propagates to the caller, which
    reports it with the CLI command that reproduces it.
    """
    import arc_divergence                                          # noqa: PLC0415
    return arc_divergence.check(candidates)


def _deck_slide_count(pptx):
    """How many slides the deck ACTUALLY has — the number a coverage claim is checked against.

    Deliberately fails loud rather than returning 0: a count of zero would make every coverage
    claim vacuously complete, which is the failure mode this number exists to close.
    """
    from pptx import Presentation
    try:
        return len(Presentation(pptx).slides)
    except Exception as exc:
        die("cannot read {} to count its slides ({}). The coverage gate compares the critic's "
            "`slides_opened` against the real deck, so an unreadable deck is not a pass."
            .format(pptx, exc))


CRITIC_WAIVER_KINDS = {
    "no-dispatch-on-host":
        "the runtime cannot dispatch a subagent (record it in the capability ledger too)",
    "already-reviewed-minor-edit":
        "a 1-2 slide edit to a deck that already passed its loop",
    "user-waived":
        "the user was asked and chose to ship over it",
    "external-deck":
        "a deck this skill did not author (redesign diagnosis / critique-only run)",
    # The loop RAN and did not converge. Every kind above describes a loop that was SKIPPED, so a
    # deck that reached its round cap with majors still open had no honest route: `user-waived`
    # asserts "the user was asked and chose to ship over it", which is a claim about a conversation
    # that did not happen. Measured on a real build — the builder's words were "the four options
    # force either a lie or a red gate", and it correctly left the gate red rather than fabricate
    # one. A record with no state for the commonest non-consent ending is a record that teaches
    # people to lie to it.
    "cap-reached-majors-open":
        "the loop RAN to its round cap and majors remain open (needs `open` + `surfaced_to_user`)",
}
# Every message that offers the categories renders them from this dict. They were hand-copied into
# four places and the fifth kind reached none of them, so the file that ACCEPTED
# `cap-reached-majors-open` also told every reader it did not exist — which is the same lie the
# fifth kind was added to prevent.
CRITIC_WAIVER_LIST = " | ".join(sorted(CRITIC_WAIVER_KINDS))


def _critic_waiver_menu(indent="    "):
    return "\n".join("{}{:28s} {}".format(indent, k, v)
                     for k, v in sorted(CRITIC_WAIVER_KINDS.items()))


def _critic_effort(critic):
    """How much reviewing happened, said in the unit the record can actually back.

    `rounds` is the only hand-typed field in the critic block — `validate_review.py --record`
    preserves one but never derives it, because nothing inside a review file says which round it
    belongs to. `reviews_seen` is the machine-counted one. These lines used to print
    `rounds` with a bare `"?"` fallback, which reported "consented after ? round(s)" on every
    tool-written record: the number the tool DOES know was sitting in the same dict, unread.
    """
    if isinstance(critic.get("rounds"), int):
        return "{} round(s)".format(critic["rounds"])
    if isinstance(critic.get("reviews_seen"), int):
        return "{} review(s) [rounds not recorded]".format(critic["reviews_seen"])
    return "an unrecorded number of rounds"


def _direction_gate(design, deck_dir):
    """The direction competition, RE-SCORED here — the design side's twin of the arc gate.

    `directions_diversity.py` is a real detector: given the four directions of a real build it
    reported `TOO SIMILAR Brutalist vs Swiss (palette 37.9, matched palette+type+density)` and
    `NO BESPOKE DIRECTION: every candidate is a preset (or a motif-less colourway)`. Both were
    true. NEITHER reached the design checkpoint, because nothing in the flow required the script
    to be run — the author caught the first by eye and never noticed the second.

    That is the same asymmetry the arc gate was written to close: the CONTENT competition is
    re-scored from its candidates at hand-off, while the DESIGN competition's verdict was a
    sentence somebody typed. So this reads the candidates and scores them here. A deck whose look
    was not chosen from alternatives records the named carve instead — the `logo plan: n/a — …`
    shape — because "no competition ran" and "a competition ran and I did not write it down" must
    not look identical."""
    dg = design.get("direction_gate")
    if isinstance(dg, str) and _norm_na(dg):
        return                                        # a recorded carve: locked / mimic / tiny ask
    if not isinstance(dg, dict):
        die('`design_plan.direction_gate` is missing. The look was either CHOSEN from rendered\n'
            '    alternatives (branch c) or it was not, and both are recordable:\n'
            '      "direction_gate": {"candidates": "directions.json" | [ {...}, ... ],\n'
            '                         "picked": "<the direction the user chose>"}\n'
            '      "direction_gate": "n/a - <locked template | mimic | user supplied the look | '
            'tiny ask>"\n'
            '    It is re-scored here with scripts/directions_diversity.py, the way the arc\n'
            '    competition is re-scored with arc_divergence — a verdict somebody typed is not\n'
            '    evidence that the check ran.')
    cands = dg.get("candidates")
    if isinstance(cands, str):
        cp = os.path.join(deck_dir, cands)
        if not os.path.exists(cp):
            die("`design_plan.direction_gate.candidates` points at {!r}, which does not exist "
                "beside the deck.".format(cands))
        try:
            with open(cp, encoding="utf-8") as fh:
                cands = json.load(fh)
        except ValueError as exc:
            die("`design_plan.direction_gate.candidates` ({}) is not valid JSON: {}"
                .format(cands, exc))
    if not isinstance(cands, list) or len(cands) < 2:
        die("`design_plan.direction_gate.candidates` needs the LIST of directions that were "
            "shown (2-4 of them). One direction is not a competition.")
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        import directions_diversity                                  # noqa: PLC0415
        r = directions_diversity.check(cands)
    except Exception as exc:                                          # never silently
        print("  [--] DIRECTION GATE: NOT RE-SCORED — {}: {}".format(type(exc).__name__, exc))
        return
    faults = []
    if r["flagged"]:
        faults.append("{} pair(s) read as skins of one idea: {}".format(
            len(r["flagged"]),
            "; ".join("{} vs {} (palette {}, matched {})".format(
                p["a"], p["b"], p["palette_distance"], ", ".join(p["matched_axes"]) or "none")
                for p in r["flagged"])))
    if r["no_bespoke"]:
        faults.append("NO BESPOKE DIRECTION — every candidate is a preset or a motif-less "
                      "colourway; at least one must be a register invented for THIS topic "
                      "(a dict carrying its own `cover_motif` + `ambient_motif`)")
    if r["colourway_excess"]:
        faults.append("more than one motif-less colourway: {}".format(
            ", ".join(r["colourway_excess"])))
    if faults and not str(dg.get("waived", "")).strip():
        die("the direction competition does not hold up when re-scored:\n    - "
            + "\n    - ".join(faults)
            + "\n    REDIVERGE them, or record why the set stands: "
              '"direction_gate": {..., "waived": "<the reason>"}.')
    if faults:
        print("  [gates] direction gate: {} finding(s) WAIVED — {}".format(
            len(faults), str(dg["waived"])[:90]))
    else:
        print("  [gates] direction gate: {} direction(s) re-scored — none is a skin of another, "
              "{} bespoke".format(len(cands), len(r["bespoke"])))


def _norm_na(v):
    t = str(v or "").strip().lower().replace("—", "-").replace("–", "-")
    return t.startswith("n/a") and len(t) > 6


def _image_provenance_gate(design, pptx):
    """Evidence tokens vs the provenance ledger vs the deck's own text.

    Kept thin on purpose: the rules live in scripts/check_image_provenance.py so that the Codex
    delivery gate can call the SAME code. Two gate paths re-implementing one contract is how this
    repo previously ended up with a key spelled `path` on one side and `png` on the other."""
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        import check_image_provenance as cip
    except Exception as exc:                       # never silently: a missing checker is a finding
        print("  [gates] image provenance NOT checked — {}: {}".format(type(exc).__name__, exc))
        return
    deck_dir = os.path.dirname(os.path.abspath(pptx)) or "."
    gates = {"design_plan": design}
    probs = cip.check(deck_dir, pptx=pptx, gates=gates)
    if probs:
        die("image provenance:\n\n" + "\n\n".join(
            "  {}: {}".format(c, m) for c, m in probs))


def check_handoff_gates(pptx, mode="presented", gate_check=False):
    """Run EVERY hand-off gate, then report EVERY failure — one run, one fix pass.

    This is the batching wrapper described at `die()`. It exists because the gate used to stop at
    the first problem, so a thin record cost one round-trip per field at the most expensive moment
    of the run. Now each independent section reports for itself and this prints the whole list.

    What it deliberately does NOT change: the checks, their thresholds, their messages, or the
    exit code. A batched report is a cheaper way to see the same refusals — never a softer one.
    Two properties keep that true, and both are asserted by `tests/test_gate_batching.py`:

      * a section that fails still fails — the exit code is the max of the collected codes, and a
        deck with any problem never prints a pass line;
      * failures do not MASK each other, in either direction. Within one section the first stop
        still wins (its later checks read values the failed one was supposed to establish, so
        continuing would report invented follow-on faults), and across sections nothing is
        suppressed.

    A structural failure — no `.deck-gates.json`, unreadable JSON, an unknown recorded delivery —
    is still terminal: it is collected and reported alone, because every later gate reads what it
    could not produce. That case never had a ping-pong problem, since it is one message either way.
    """
    global _COLLECTED
    outer, _COLLECTED = _COLLECTED, []
    try:
        try:
            _handoff_gate_checks(pptx, mode, gate_check)
        except _GateStop:
            pass                    # structural stop: report what we have, which is that one
        problems = _COLLECTED
    finally:
        _COLLECTED = outer
    if not problems:
        return
    n = len(problems)
    head = ("render_deck: {} hand-off gate(s) failed. ALL of them are listed below — fix them in "
            "ONE pass, then re-run.".format(n) if n > 1 else
            "render_deck: 1 hand-off gate failed.")
    print(head, file=sys.stderr)
    for i, (section, msg, _code) in enumerate(problems, start=1):
        label = " {}".format(section) if section else ""
        print("\n[{}/{}]{}\n{}".format(i, n, label, msg), file=sys.stderr)
    sys.exit(max(code for _s, _m, code in problems))


def _find_build_script(pptx, design):
    """The build script for this deck: recorded, else the `build_<stem>.py` convention beside it.

    Discovery can legitimately fail — a build script is not a deliverable and the plan files are
    deliberately never written into the deck folder — so a miss is reported LOUDLY and does not
    block. Only a resolved script that contradicts the record does.
    """
    rec = design.get("build_script")
    if rec:
        p = Path(rec).expanduser()
        return p if p.is_file() else None
    deck = Path(pptx)
    for cand in (deck.parent / f"build_{deck.stem}.py",
                 Path.cwd() / f"build_{deck.stem}.py",
                 deck.parent / "build.py"):
        if cand.is_file():
            return cand
    return None


def _style_applied_gate(design, pptx):
    """`design_plan.style_pick` names a register; the build must actually apply it.

    Both delivery gates required this field as a STRING and neither verified it — measured by
    grep, `presets.apply` / `set_geometry` / `set_ground` appeared in no gate script at all. So a
    deck recording `style_pick: "brutalist for engineering - beat blueprint"` and built with
    deckkit's stock defaults passed here and on the Codex path alike: the competition ran, the
    winner was written down, and nothing carried it into the build. `check_style_applied.py` owns
    the logic so both paths ask the same question of the same code.
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    try:
        import check_style_applied as csa
        names = csa.preset_names()
    except Exception as exc:                       # never fail a render on the checker itself
        print(f"  [--] STYLE APPLIED: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    script = _find_build_script(pptx, design)
    if script is None:
        print("  [--] STYLE APPLIED: NOT CHECKED — no build script found beside the deck "
              "(looked for build_<stem>.py). NOT the same as clean: record it as "
              "`design_plan.build_script` and the register a deck DECLARES gets verified against "
              "the one it APPLIES.")
        return
    code, msg = csa.evaluate(design.get("style_pick"), script.read_text(encoding="utf-8"),
                             names, design.get(csa.WAIVER_KEY),
                             design.get(csa.LOOK_SOURCE_KEY))
    if code == 1:
        die("`design_plan.style_pick` " + msg.split("—", 1)[-1].strip()
            + f"\n    (build script: {script})")
    if code == 2:
        print(f"  [--] STYLE APPLIED: NOT CHECKED — {msg}")


def _LD_A11Y_CODES():
    import lint_deck as _ld
    return _ld.A11Y_BLOCKING


def _LD_A11Y_WCAG():
    import lint_deck as _ld
    return _ld.A11Y_WCAG


def _LD_A11Y_ALL():
    import lint_deck as _ld
    return _ld.A11Y_CODES


def _check_a11y(pptx, delivery, gates):
    """The accessibility floors lint already measured and no gate on this path ever read.

    Measured by grep before this existed: MISSING ALT-TEXT, NO SLIDE TITLE, DUPLICATE SLIDE TITLES
    and READING ORDER were emitted as advisory `[warn]`s and consumed by NOTHING on either runtime,
    while the codex path held only the two WCAG contrast codes. A deck could therefore ship with no
    image described, no slide titled and the reading order scrambled, and both delivery gates
    called it clean. This repo has run that experiment already — the deck-level sameness signals
    were warns nobody read, which is why they became a gate.

    None of these needs an opinion, which is why they can hold a deck when the critic's taste calls
    cannot: a shape either carries a description or it does not, a title is either first in z-order
    or it is not, a ratio either clears 3:1 or it does not.

    A deck may still be handed over with them open — a throwaway internal review, a deck for one
    sighted person — but in WRITING, with the count, so the user learns what they are accepting:

        {"a11y": {"waived": "internal 10-minute review, no distribution",
                  "waived_category": "not-distributed"}}
    """
    if not isinstance(gates, dict):
        gates = {}
    stats, _aspect = _sameness_stats(pptx, delivery)      # one lint run, already cached above
    warns = stats.get("slide_warns")
    if warns is None:
        print("  [--] A11Y: NOT CHECKED — this lint build does not surface the per-slide warn "
              "stream. NOT the same as clean.")
        return
    codes = _LD_A11Y_CODES()
    advisory = [c for c in _LD_A11Y_ALL() if c not in codes]
    hits, noted = {}, {}
    for w in warns:
        text = str(w.get("text") if isinstance(w, dict) else w)
        code = text.split(":", 1)[0].strip()
        if code in codes:
            hits.setdefault(code, []).append(w.get("slide") if isinstance(w, dict) else None)
        elif code in advisory:
            noted.setdefault(code, []).append(w.get("slide") if isinstance(w, dict) else None)
    if noted:
        print("  [--] a11y advisory: {} — not held here (a statement slide may carry an "
              "off-canvas title on purpose), but check them if this deck is distributed."
              .format(" · ".join("{} (slide {})".format(c, ", ".join(str(x) for x in v))
                                 for c, v in sorted(noted.items()))))
    waiver = _section(gates, "a11y")
    if not hits:
        print("[gates] a11y: 0 of {} floor(s) fired (alt-text · slide titles · reading order · "
              "non-text contrast)".format(len(codes)))
        return
    fired = " · ".join("{} ({} slide{})".format(c, len(v), "" if len(v) == 1 else "s")
                       for c, v in sorted(hits.items()))
    if waiver.get("waived"):
        why = str(waiver.get("waived") or "")
        if len(why.strip()) < 24:
            die("`a11y.waived` needs a real reason, not '{}'. Say what makes these acceptable on "
                "THIS deck — who reads it, and how.".format(why.strip()))
        print("[gates] a11y WAIVED [{}] — {} — NOT ACCESSIBLE\n        {}\n        Say this in the "
              "hand-off note too; a waiver the user never sees is a silence."
              .format(waiver.get("waived_category") or "uncategorised", fired, why))
        return
    wcag = [c for c in hits if c in _LD_A11Y_WCAG()]
    die("this deck does not clear the accessibility floors:\n    - " + fired
        + ("\n    {} of those are WCAG ratios, which are arithmetic rather than judgement."
           .format(len(wcag)) if wcag else "")
        + "\n    Fix: deckkit.alt_text(shape, '<one line>') on informative images (alt='' for "
          "purely decorative); give every slide a title (an off-canvas title is a sanctioned trick "
          "for statement slides) and add it FIRST so z-order matches reading order; raise "
          "icon/mark contrast to 3:1.\n    Or waive in writing: "
          '{"a11y": {"waived": "<who reads this deck, and how>", "waived_category": "<kind>"}}')


def _register_kit_note(pptx, gates):
    """The declared register HAS a buildable surface — say so if the deck used none of it.

    `presets.apply()` sets a palette, a corner radius, a ground and a font. Rendered side by side,
    one page through all 18 presets came out as 18 colourways of one page: no memphis bands, no
    bauhaus primitive, no glass, no overprint, no scanlines. `register_surface.py` builds those,
    for the six registers that have a kit so far — and a capability nobody is told about at the
    moment it applies is the same as one that does not exist. This is the telling. A NOTE, never a
    hold: a deliberate quiet treatment of a loud register is a real design choice, and holding a
    finished deck over an unused helper would teach the waiver reflex.
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    try:
        import register_surface as rs
        import check_style_applied as csa
        d = (gates or {}).get('design_plan') or {}
        if not isinstance(d, dict):
            return
        reg, conf = csa.declared_preset(d.get('style_pick'), csa.preset_names(),
                                        d.get(csa.LOOK_SOURCE_KEY))
        if conf != 'sure' or not reg or not rs.has(reg):
            return
        deck_dir = Path(pptx).resolve().parent
        used = any('register_surface' in f.read_text(encoding='utf-8', errors='ignore')
                   for f in list(deck_dir.glob('*.py')) + list(deck_dir.glob('**/*.py'))[:60])
    except Exception:
        return
    if used:
        print("[gates] `{}` was built with its surface kit, not just its palette".format(reg))
        return
    print("[gates] `{}` HAS a buildable surface — `register_surface.ground()` paints its own "
          "furniture and hands back the content rect, and `.card()` gives its card FORM. This deck "
          "used none of it, so what shipped is the register's colourway. A deliberate quiet "
          "treatment is a legitimate choice; forgetting the kit exists is not.".format(reg))


def _register_keep_note(pptx, gates):
    """Say, at hand-off, that an INVENTED register is about to be forgotten.

    The skill's own example library holds four bespoke registers; this user's look history holds
    nine that were designed, shipped and lost. Measured by grep, no script anywhere wrote one: the
    mechanism was "remember to edit the markdown", which is not a mechanism. So the skill invents
    good registers and forgets every one, and every deck starts its design from zero.

    A NOTE, never a hold. Keeping a register is the user's call about their own collection, and a
    gate that blocked a finished deck over a library entry would be the waive-reflex failure this
    repo keeps re-learning. But silence is how it stayed at four.
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    try:
        import save_register as sr
        name, _body = sr.entry_for(Path(pptx).resolve().parent)
    except Exception:
        return
    if not name:
        return
    try:
        already = sr.is_kept(name)          # identity, not string equality — `save_register` will
    except Exception:                       # answer "already kept" for a glossed name, and a
        already = False                     # reminder that disagrees with the tool it points at
                                            # gets ignored on the second hand-off and every one after
    if already:
        print("[gates] register `{}` is already kept in your collection".format(name))
        return
    print("[gates] this deck invented `{}` — keep it, or it is gone when the folder is:\n"
          "        python3 scripts/save_register.py {}\n"
          "        (the library has 4 invented registers; a look history usually has more that "
          "were shipped and lost)".format(name, Path(pptx).resolve().parent))


def _register_guard_gate(pptx, gates):
    """A declared register must be OBEYED, not merely paletted.

    `check_style_applied.py` verifies the CALL; `check_register_pixels.py` verifies the COLOUR and
    says in its own docstring that colour is all it judges. Neither can see that a deck declared
    `brutalist` and shipped rounded cards with soft shadows — measured, 18 registers applied to one
    page produced 18 pages differing only in ground, radius and rule weight. This checks the
    shape-level prohibitions each preset's own `guard` states, from `presets.FORBIDS`.
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    try:
        import check_register_guard as crg
    except Exception as exc:
        print(f"  [--] REGISTER GUARD: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    try:
        violations, facts = crg.check(pptx, None, gates if isinstance(gates, dict) else {})
    except Exception as exc:
        print(f"  [--] REGISTER GUARD: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    if facts.get("note"):
        print("  [--] register guard: " + facts["note"])
        return
    if not violations:
        print("[gates] register guard: {} obeys {}".format(
            facts.get("register"), " · ".join(facts.get("rules") or [])))
        return
    die("the deck declares `{}` and breaks its own guard:\n    - ".format(facts.get("register"))
        + "\n    - ".join("{}: {}".format(c.upper(), m) for c, m in violations)
        + "\n    These are the register's OWN prohibitions (presets.PRESETS[…]['guard']), and only "
          "the ones a machine can settle — the rest stay prose."
        + "\n    Re-run alone: python3 scripts/check_register_guard.py {}".format(pptx))


def _register_pixels_gate(pptx):
    """The declared register must reach the RENDERED PIXELS — and must not be the last deck's.

    SKILL.md names two rules that "survive no matter what": never ship deckkit's default blue, and
    never reuse the last deck's scheme. Both were prose, and the nearest gate checked something
    weaker: `_style_applied_gate` verifies the CALL — that `presets.apply("brutalist")` appears in
    the build script. A deck that calls it and then sets the tokens back by hand passes, and a
    BESPOKE register (the case this skill actively encourages, with no preset call to find) is
    skipped by definition. Measured on the deck built in this repo's own session: its whole
    terminal register was set by hand, and nothing verified any of it landed.

    So this asks the render. `check_register_pixels.py` owns the measurement; both delivery paths
    call the same code so they cannot drift.
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    try:
        import check_register_pixels as crp
    except Exception as exc:                       # never fail a render on the checker itself
        print(f"  [--] REGISTER PIXELS: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    deck_dir = Path(pptx).resolve().parent
    taste = None
    try:
        import registry
        t = registry.taste_file()
        taste = str(t) if t else None
    except Exception:
        pass                                       # no registry footprint: the freshness half sits out
    try:
        probs, facts = crp.check(deck_dir, taste=taste)
    except Exception as exc:
        print(f"  [--] REGISTER PIXELS: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    if facts.get("waived"):
        print("[gates] register pixels: waived in writing (design_plan.register_pixels_waived).")
        return
    codes = [c for c, _ in probs]
    if "NO RENDERS" in codes:
        print("  [--] REGISTER PIXELS: NOT CHECKED — no renders beside the deck. NOT the same as "
              "clean: render first, and the register a deck DECLARES gets checked against the one "
              "it SHOWS.")
        return
    if facts.get("band"):
        print("[gates] look history: {}".format(facts["band"]))
    if not probs:
        print("[gates] register pixels: {} of {} declared colour(s) reached the pages{}".format(
            len(facts.get("present") or []), len(facts.get("declared") or []),
            "" if not facts.get("note") else " — " + facts["note"]))
        return
    die("the register this deck DECLARES did not reach its PIXELS:\n    - "
        + "\n    - ".join("{}: {}".format(c, m) for c, m in probs)
        + "\n    Re-run alone: python3 scripts/check_register_pixels.py {}".format(deck_dir))


def _surface_gate(pptx, gates):
    """A canvas format's contract, checked against the built deck instead of trusted.

    `formats.py` registers each design surface — margins, platform-UI safe zones, whether columns
    work, whether the surface carries chrome, and now the ABSOLUTE type floors and fill range a
    PRINTED board needs. Measured by grep, nothing downstream ever consumed it: `import formats`
    appeared in two files, both producers. So every per-surface rule in `references/canvas-formats.md`
    was advisory by construction, on exactly the surfaces where the mistake is least recoverable —
    a story caption under the swipe bar is invisible to whoever built it on a desktop, and a poster
    is wrong only once it has been printed a metre wide.
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    try:
        import check_surface as cs
    except Exception as exc:                       # never fail a render on the checker itself
        print(f"  [--] SURFACE: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    design = (gates or {}).get("design_plan") or {} if isinstance(gates, dict) else {}
    waive = design.get("surface_sections_waived")
    try:
        probs, facts = cs.check(pptx, design.get("format"), waive,
                                design.get("surface_section_terms"))
    except Exception as exc:
        print(f"  [--] SURFACE: NOT CHECKED — {exc.__class__.__name__}: {exc}")
        return
    if facts.get("note"):
        print("  [--] SURFACE: " + facts["note"])
        return
    head = "[gates] surface: {}".format(facts.get("format"))
    for extra in ("floors", "fill", "ink"):
        if facts.get(extra):
            head += " · " + facts[extra]
    print(head)
    if probs:
        die("this deck breaks the contract of the canvas it is built on:\n    - "
            + "\n    - ".join("{}: {}".format(c, m) for c, m in probs)
            + "\n    Re-run alone: python3 scripts/check_surface.py {}".format(pptx))


def _handoff_gate_checks(pptx, mode="presented", gate_check=False):
    """Refuse --deliverables until the quality gates have actually run.

    The gates that guard a deck's quality — the design plan, the independent critic, the
    primary-source spot-check — are all *self-reported prose artifacts*. The model that skipped
    them is the same model that writes the hand-off note claiming they happened, so skipping is
    invisible: the note reads identically either way. This skill's own enforcement invariant
    ranks prose last for exactly that reason (deterministic lint > required-field > checklist >
    prose), yet these three sat at the bottom of it.

    So put the check where the incentive is. `--deliverables` produces the PDF and the preview
    page — the artifacts the user actually asked for — which makes it the one step nobody skips.
    Requiring evidence *here* costs nothing on a deck that ran its loop and blocks the shortcut on
    one that didn't.

    Evidence lives in `<deck-dir>/.deck-gates.json`:

        {"critic": {"verdict": "consent", "rounds": 2, "review": "reviews/r2.json"},
         "provenance": {"claims": [{"claim": "<the claim>", "verdict": "CONFIRMED",
                                    "url": "https://<primary source>"}]}}

    (A summary tally like {"checked": 87} is REJECTED on purpose — a tally is written by the same
    pass that would have skipped the refutation.)

    A gate may be waived — quick decks exist — but a waiver is a written reason that travels with
    the deck AND names which kind of skip it is, not a silence:

        {"critic": {"waived": "1-slide fix to an already-reviewed deck",
                    "waived_category": "already-reviewed-minor-edit"}}

    Set SLIDE_MAKER_SKIP_GATES=1 to bypass entirely (CI smoke tests, throwaway renders).
    """
    if os.environ.get("SLIDE_MAKER_SKIP_GATES"):
        # Say so, and never let the caller print a pass. Measured: with this variable set on a deck
        # carrying NO .deck-gates.json at all, `--gate-check` exited 0 under the line "all hand-off
        # gates pass — the deck may be handed over". A bypass is legitimate for CI smoke renders and
        # throwaway previews; asserting that skipped gates PASSED is the one thing it must not do,
        # because that sentence is the artifact a hand-off note is written from.
        print("[gates] SKIPPED — SLIDE_MAKER_SKIP_GATES=1 is set. NOTHING was checked: not the "
              "critic record, not the design plan, not provenance, not density.")
        if gate_check:
            die("`--gate-check` is the hand-off gate itself, so bypassing it has no honest meaning. "
                "Unset SLIDE_MAKER_SKIP_GATES (it exists for CI smoke renders), or record a written, "
                "classified waiver in .deck-gates.json — a waiver is visible, an env var is not.")
        return
    path = os.path.join(os.path.dirname(os.path.abspath(pptx)) or ".", GATES_FILE)
    if not os.path.isfile(path):
        # SAY WHAT THE FILE ITSELF SHOWS, before saying a word about missing paperwork. Every gate
        # below reads `.deck-gates.json`, so with no record they all report the same thing — write
        # the record — and a reader learns the problem is bookkeeping. It is not. Measured on a
        # delivered 13-slide deck that had no record at all: the FILE showed five slides of
        # category content and zero icons, a real design miss that no amount of paperwork fixes,
        # and every gate stayed silent about it because the gates read the record, not the deck.
        _report_file_observations(pptx)
        die("--deliverables is the hand-off run, and this deck has no record that its quality "
            "gates ran.\n"
            "  Missing: {}\n\n"
            "  The independent critic (Step 5) is the one that cannot be self-certified — you are "
            "not the final judge of your own deck.\n"
            "  Let the loop WRITE it, so the record is evidence and not a claim — one flag on a\n"
            "  step you already run before acting on any review:\n\n"
            "    python3 scripts/validate_review.py critic <review.json> --record {}\n\n"
            "  (that fills `critic` from the review itself, with its path + sha256, and the gate\n"
            "  re-reads it). Then add the two blocks only you can supply. Full shape:\n\n"
            '    {{"critic": {{"verdict": "consent", "rounds": 2}},\n'
            '     "design_plan": {{"boldness": "balanced+", "signature_move": "<the one risk>",\n'
            '                     "carried_by": [4, 6, 8], "form_ledger": "<family tally>",\n'
            '                     "icon_family": "<family | none - reason>",\n'
            '                     "motif_generates": {{"background": "<what the motif makes the canvas do>",\n'
            '                                        "markers": "<the numeral/icon system it implies>",\n'
            '                                        "page": "<the slide whose GEOMETRY is the motif, or none - reason>"}},\n'
            '                     "palette": "<FILL vs TEXT-safe split, per palette_audit.py>",\n'
            '                     "type_scale": {{"display": 34, "title": 24, "body": 14}},\n'
            '                     "signature_proof": [{{"role": "signature", "slide": 6, "png": "render/slide06.png"}},\n'
            '                                         {{"role": "complex", "slide": 9, "png": "render/slide09.png"}},\n'
            '                                         {{"role": "data", "slide": 11, "png": "render/slide11.png"}}]}},\n'
            '     "provenance": {{"claims": [{{"claim": "<the claim>", "verdict": "CONFIRMED",\n'
            '                                "url": "https://<primary source>"}}]}}}}\n\n'
            "  (A summary tally like {{\"checked\": 87}} is REJECTED on purpose — a tally is "
            "written by the same pass that would have skipped the refutation.)\n\n"
            "  A gate you deliberately skipped is waived in writing AND CLASSIFIED, not omitted:\n\n"
            '    {{"critic": {{"waived": "why this deck does not need a critic round",\n'
            '                 "waived_category": "{}",\n'
            '                 "inline_ran": true}}}}\n\n'
            "  (`inline_ran` is required only for `no-dispatch-on-host`; "
            "`cap-reached-majors-open` — the loop RAN and did not converge — additionally needs\n"
            "  `open` and `surfaced_to_user`.)\n\n"
            "  Renders without --deliverables are unaffected; iterate freely."
            .format(path, os.path.dirname(path) or ".", CRITIC_WAIVER_LIST))
    # --- resolve the DELIVERY once, above every gate that keys off it -------------------------
    # Measured divergence: the type-scale floor read a `delivery` key out of this JSON while the
    # density gate three blocks below read the raw `mode` the CLI flags set. One run could
    # therefore enforce two different deliveries, and a NOTE claiming "using the recorded value"
    # was true of only one of them. Resolve it here, once, and hand the SAME value to both.
    #
    # `mode` is what --selfread / --surface / --textheavy set. A recorded key wins over it, since
    # it travels with the deck; an unrecognised recorded value dies rather than silently falling
    # back to the presented floor, which would be a legibility floor quietly not applied.
    # BODY_FLOORS is deliberately the SAME three keys the Codex delivery gate validates, so the
    # two paths cannot disagree; `surface` is a CLI mode rather than a recorded delivery, and it
    # maps to `selfread` because a poster is read up close — the same reason --surface already
    # exempts it from the presented density budget.
    _KNOWN_DELIVERY = ("presented", "textheavy", "selfread", "surface")
    _DELIVERY_ALIAS = {"surface": "selfread"}
    try:
        gates = json.load(open(path, encoding="utf-8"))
    except Exception as e:
        die("{} is not readable JSON: {}".format(path, e))

    _recorded = str(gates.get("delivery") or "").strip()
    if _recorded and _recorded not in _KNOWN_DELIVERY:
        die("`delivery` in {} is {!r}, which is not a delivery mode. One of: {}. An unrecognised "
            "value used to fall through to the presented floor, i.e. a legibility floor silently "
            "not applied to the deck it was recorded for."
            .format(path, _recorded, ", ".join(_KNOWN_DELIVERY)))
    if _recorded and mode != "presented" and _recorded != mode:
        print("[gates] NOTE: {} records delivery={!r} but you passed {!r} — using the RECORDED "
              "value for every gate in this run. Fix whichever is wrong; the two must not "
              "disagree about which deck this is.".format(path, _recorded, mode))
    delivery = _recorded or mode          # one of _KNOWN_DELIVERY, un-aliased

    with _gate_section('critic'):
        critic = _section(gates, "critic")
        if critic.get("waived"):
            # A waiver is legitimate — quick decks and hosts without subagent dispatch are real —
            # but an UNCLASSIFIED one is just a sentence, and the model that skipped the loop writes
            # the same sentence as the model that ran it. The Codex delivery gate has required a
            # distinct schema-valid review artifact per lens for a while; this path accepted any
            # string. Measured: a hand-typed waiver carried a whole deck through `all hand-off gates
            # pass` without an independent critic ever seeing it. So: name the category, and say
            # whether the lenses ran at all.
            WAIVER_KINDS = CRITIC_WAIVER_KINDS
            reason = critic["waived"]
            kind = critic.get("waived_category")
            if not isinstance(reason, str) or reason_width(reason) < 24:
                die("`critic.waived` must be a written reason that travels with the deck, not a "
                    "placeholder. Say what was skipped and why, in a sentence someone can disagree "
                    "with later.")
            if kind not in WAIVER_KINDS:
                die("`critic.waived_category` must name WHICH kind of skip this is — an unclassified "
                    "waiver is indistinguishable from never having run the loop. One of:\n"
                    + _critic_waiver_menu()
                    + "\n\n  If none of these fit, the honest move is to run the critic.")
            if kind == "cap-reached-majors-open":
                # This kind asserts the OPPOSITE of the others — that the work happened — so it carries
                # the evidence: which findings survived, and whether the user was actually told. Without
                # `open` it degrades into "we tried", which is what every other waiver already says.
                _open = critic.get("open")
                if not isinstance(_open, list) or not _open or not all(
                        isinstance(x, str) and x.strip() for x in _open):
                    die("`waived_category: cap-reached-majors-open` must list the surviving findings in "
                        "`\"open\": [\"...\"]` — one short string each. This category claims the loop "
                        "RAN, so the deck owes the reader what it ran INTO; an empty list is a consent "
                        "verdict wearing a waiver's name.")
                if not isinstance(critic.get("surfaced_to_user"), bool):
                    die("`waived_category: cap-reached-majors-open` must record "
                        "`\"surfaced_to_user\": true|false` — whether these open majors were put in "
                        "front of the user. Shipping over a known major is a decision; the record has "
                        "to say whose it was.")
            if kind == "no-dispatch-on-host" and "inline_ran" not in critic:
                die("`waived_category: no-dispatch-on-host` must also record `\"inline_ran\": true|false` "
                    "— whether the content and design lenses were at least run inline. 'Ran inline in "
                    "the author's own context' and 'was never reviewed' are different claims, and the "
                    "hand-off note reads identically for both unless this file separates them.")
            print("[gates] critic WAIVED [{}] — NOT INDEPENDENTLY REVIEWED".format(kind))
            print("        {}".format(reason))
            if kind == "no-dispatch-on-host":
                print("        lenses run inline: {} — inline review is the author grading "
                      "themselves.".format("yes" if critic.get("inline_ran") else "NO"))
            print("        Say this in the hand-off note too; a waiver the user never sees is a "
                  "silence.")
        elif critic.get("verdict") == "consent":
            # A record the model TYPED at hand-off is self-certification: the model that skipped the
            # loop writes the same JSON as the model that ran it. So when the record points at the
            # review artifact (validate_review.py --record puts it there), re-read that artifact and
            # verify it still says what the record claims. A hand-written record still passes — but it
            # is LABELLED self-reported, so the two are distinguishable instead of identical.
            src = critic.get("source")
            if src:
                if not os.path.isfile(src):
                    die("`critic.source` points at {} — which does not exist. The gate re-reads the "
                        "review artifact rather than trusting the summary; restore the file, re-run "
                        "`validate_review.py critic <review.json> --record <deck-dir>`, or waive in "
                        "writing.".format(src))
                digest = critic.get("sha256")
                if digest and _sha256(src) != digest:
                    die("`critic.source` ({}) has changed since it was recorded — the sha256 no longer "
                        "matches. Re-validate it with `validate_review.py critic <review.json> "
                        "--record <deck-dir>` so the record and the evidence agree.".format(src))
                try:
                    review = json.load(open(src, encoding="utf-8"))
                except Exception as exc:
                    die("`critic.source` ({}) is not readable JSON: {}".format(src, exc))
                if review.get("verdict") != "consent":
                    die("the recorded verdict says consent, but the review at {} says {!r}. The "
                        "artifact wins.".format(src, review.get("verdict")))
                hard = [f for f in (review.get("findings") or [])
                        if isinstance(f, dict) and f.get("severity") in ("blocker", "major")]
                if hard:
                    die("the review at {} consents while still carrying {} blocker/major finding(s) — "
                        "that is a contract violation (agents/critic.md: any blocker/major -> revise). "
                        "Fix them and re-review, or waive in writing.".format(src, len(hard)))
                # --- bind the coverage claim to the DECK, not just to the file --------------------
                # Until this ran, "verified" meant only "the artifact exists and still hashes to what
                # was recorded". Measured: a schema-valid review of a 15-slide deck declaring
                # slides_opened=[1] was accepted, recorded with a sha256, and printed as verified, and
                # every hand-off gate passed. `slides_opened` is the anti-skim field; nothing compared
                # it to the deck it claims to have read. SKILL.md Step 5 already tells the coordinator
                # to make this comparison by hand ("lists every slide in the critic's ASSIGNED scope —
                # whole deck for a sole critic; its section's range for a per-section critic"); the
                # Codex delivery gate already mechanises it. This is the shared path catching up.
                _cov = review.get("coverage") or {}
                _opened = {v for v in (_cov.get("slides_opened") or []) if isinstance(v, int)}
                _scope = _cov.get("scope")
                if isinstance(_scope, (list, tuple)) and len(_scope) == 2 \
                        and all(isinstance(v, int) for v in _scope):
                    _expect = set(range(_scope[0], _scope[1] + 1))
                    _what = "its declared scope (slides {}-{})".format(*_scope)
                else:
                    _expect = set(range(1, _deck_slide_count(pptx) + 1))
                    _what = "the whole deck ({} slides)".format(len(_expect))
                _missing = sorted(_expect - _opened)
                if _missing:
                    die("the review at {} consents for {}, but `coverage.slides_opened` never lists "
                        "slide(s) {}{}.\n"
                        "  A critic can only judge what it opened, and consent on an unopened slide is "
                        "not a verdict — it is a gap the record renders as a pass.\n"
                        "  Re-dispatch the critic over the missing slides, or — for a per-section "
                        "critic — declare the range it was assigned:\n"
                        '      "coverage": {{"scope": [4, 9], "slides_opened": [4,5,6,7,8,9], ...}}'
                        .format(src, _what,
                                ", ".join(str(m) for m in _missing[:12]),
                                " (+{} more)".format(len(_missing) - 12) if len(_missing) > 12 else ""))
                print("[gates] critic consented after {} — verified against {} "
                      "(opened {}/{} slides)".format(
                          _critic_effort(critic), os.path.basename(src),
                          len(_opened & _expect), len(_expect)))
            else:
                # 🔴 A consent with NO review artifact is REFUSED, not labelled. It used to print
                # "SELF-REPORTED" and pass — the last self-certification hole in the gate set: the
                # model that skipped the loop writes the same `{"verdict":"consent","rounds":N}` as
                # the model that ran it, and a label the two share is not a distinction. Every other
                # Step-1/2 verdict became a re-readable artifact this way (arc recomputed, slides
                # table, provenance per-claim); the critic is the last one, and it is the costliest
                # to fake because a real review must OPEN every slide (the coverage bind above).
                #
                # The escape is not a weaker consent — it is the honest WAIVER path above
                # (`waived` + `waived_category`), which prints NOT INDEPENDENTLY REVIEWED. A host
                # that genuinely cannot dispatch a subagent records `no-dispatch-on-host` +
                # `inline_ran`; a host that CAN dispatch has no reason to self-report, because
                # producing the artifact is `validate_review.py --record` on a review it already
                # ran. So a consent MUST carry `source`; there is no legitimate source-less consent.
                die("critic records `verdict: consent` but no `source` — the review artifact the "
                    "gate re-reads.\n"
                    "  A typed consent is self-certification: a skipped loop writes the identical "
                    "JSON to a real one, so consent now REQUIRES the recorded review (path + "
                    "sha256), not a summary.\n"
                    "  Ran the critic? Dispatch it against the schema and record it — the shape it "
                    "is judged by IS the shape you ask for:\n"
                    "      python3 scripts/validate_review.py --schema critic   # -> subagent "
                    "structured-output schema\n"
                    "      python3 scripts/validate_review.py critic <review.json> --record {dir}\n"
                    "  Could NOT run an independent critic (no subagent dispatch, or a 1-2 slide "
                    "edit)? That is the WAIVER, and it is honest about not being independent:\n"
                    '      {{"critic": {{"waived": "<why, >=24 chars>", "waived_category": '
                    '"no-dispatch-on-host", "inline_ran": true}}}}'
                    .format(dir=os.path.dirname(path) or "."))
            if critic.get("corroborated_by"):
                # An arbiter pass is only corroboration when it CORROBORATES. Read what it actually
                # said: a Job-2 payload reporting an unresolved finding, a dulled strength, or a
                # regressed neighbour is the opposite of a confirmation, and printing it as one is how
                # a failed verification round became a hand-off credential.
                _open = critic.get("arbiter_open") or []
                if _open:
                    _lines = []
                    for c in _open[:6]:
                        bits = []
                        if not c.get("resolved"):
                            bits.append("NOT resolved" + (": " + c["still_wrong"] if c.get("still_wrong") else ""))
                        if c.get("dulled"):
                            bits.append("dulled a named strength")
                        if c.get("regressions"):
                            bits.append("regressed " + "; ".join(map(str, c["regressions"])))
                        _lines.append("    - {}: {}".format(c.get("finding_ref") or "?", " · ".join(bits)))
                    die("the arbiter pass recorded against this deck reports {} item(s) that are still "
                        "open, so it is not a corroboration:\n{}\n\n"
                        "  Fix them and re-run the round, or — if you are shipping over it — say so in "
                        "writing where the user can see it. The loop RAN here, so the honest category "
                        "is the one that says so:\n"
                        '    {{"critic": {{"waived": "shipping over: <the open item and why>",\n'
                        '                 "waived_category": "cap-reached-majors-open",\n'
                        '                 "open": ["<each surviving finding>"],\n'
                        '                 "surfaced_to_user": true|false}}}}\n\n'
                        "  Use `user-waived` ONLY if you actually asked the user and they chose to "
                        "ship — it is a claim about a conversation, and writing it for a conversation "
                        "that did not happen is the failure this classification exists to catch."
                        .format(len(_open), "\n".join(_lines)))
                print("[gates] consent corroborated by {} arbiter pass(es), no open items".format(
                    len(critic["corroborated_by"])))
        elif critic.get("verdict") == "revise":
            die("the last critic review returned verdict=revise. Fix the blockers and re-run the "
                "loop, or record a waiver with the reason you are shipping over it. The loop RAN "
                "here, so the honest category is the one that says so:\n"
                '    {"critic": {"waived": "shipping over: <the surviving finding and why>",\n'
                '                "waived_category": "cap-reached-majors-open",\n'
                '                "open": ["<each surviving finding>"],\n'
                '                "surfaced_to_user": true|false}}\n\n'
                "  Use `user-waived` ONLY if you actually asked the user and they chose to ship — it "
                "is a claim about a conversation, and writing it for a conversation that did not "
                "happen is the failure this classification exists to catch.")
        else:
            die("{} has no usable `critic` record — needs {{\"verdict\": \"consent\"|\"revise\"}} or "
                "{{\"waived\": \"<reason>\", \"waived_category\": \"<one of these>\"}}:\n{}"
                .format(path, _critic_waiver_menu()))

    with _gate_section('design_plan'):
        # The design plan is the art director's output (Step 2). Self-authoring one is indistinguishable
        # from dispatching for it — unless the record has to carry the fields the dispatch produces.
        # `icon_family` joins the four originals for one measured reason: a deck shipped with ZERO
        # icons through every automated gate, while a missing LOGO was caught by a required
        # checkpoint token. Icons are called a design must on every branch and had no field, no
        # column and no check anywhere. The token grammar mirrors `logo plan:` — a family name or
        # an explicit `none — <reason>` — so a deliberately icon-free deck is always satisfiable.
        # `palette` joins them for the same measured reason as `icon_family`. SKILL.md already
        # states the two-token rule -- a hue used as TEXT must itself clear 4.5:1, so keep a bright
        # FILL token and a darker TEXT twin. The rule is correct and still easy to break, because
        # the check is per-PAIR and a build touches dozens of them. Measured on one deck: the author
        # declared the rule in the design plan and then broke it FOUR times -- a vivid ochre set as a
        # label on its own pale slab (2.34:1), coral emphasis text on a coral tint (4.19:1), a table
        # highlight (4.19:1), a muted grey carrying real content on cream (4.26:1). None were
        # reckless; each was a pair nobody was thinking about while computing contrast for a
        # different pair, and each surfaced at render time or in review, a round later.
        # `scripts/palette_audit.py` resolves the whole matrix in one call, so the field is cheap to
        # fill honestly and cannot be filled at all without having run something.
        # type_scale and signature_proof were gated on the CODEX delivery path only, so on the shared
        # path typography was the one pillar of the visual language nobody had to resolve (palette,
        # icons and forms all had required fields), and the signature move was accepted as a SENTENCE
        # with nothing showing it survived into the render. That asymmetry is the exact shape of a bug
        # this repo already fixed once: the critic waiver was schema-checked for Codex and a hand-typed
        # string everywhere else, and it carried a real deck through "all gates pass".
        DESIGN_FIELDS = ("concept", "boldness", "signature_move", "carried_by", "form_ledger",
                         "icon_family", "palette", "type_scale", "signature_proof",
                         "motif_generates", "style_pick", "image_sources")
        # `image_sources` — the per-image EVIDENCE TOKENS (references/image-generation.md owns the
        # grammar). It joins the required set for the reason `logo plan:` did: the decision that
        # matters is not "which photo" but "where did it come from, and what licence rides with
        # it", and a field nobody is required to fill is a decision nobody is required to make.
        # A deck with no content images writes `"image_sources": "n/a — <why>"`, exactly like a
        # `logo plan: n/a — …` line. The token grammar and the ledger cross-check live in
        # scripts/check_image_provenance.py, which runs right after this block.
        # `style_pick` — the TOPIC-adapted look choice (references/design-by-topic.md): the preset or
        # bespoke register chosen for the SUBJECT's domain, the nearest rival it beat + the one clause,
        # and the domain cliché it avoided. It exists because the look was keyed on PURPOSE only and a
        # tech deck defaulted to the dark_tech/synthwave cliché the domain map warns against. On a
        # LOCKED look (provided/registered template or a Mode-A mimic) the look is not domain-picked —
        # write `"style_pick": "n/a — <locked: template | mimic | provided>"` (a non-empty string
        # passes, exactly like a `logo plan: n/a — …` line).
        # THE RESTRAINT CARVE — built on the escape the skill ALREADY documents, not a new one.
        # agents/slide-design.md: under a *conservative* dial (user-requested or purpose-defaulted for a
        # sober defense / regulatory / status deck) "the risk is OPTIONAL: take a modest, restrained
        # signature move if one fits, OR — if none does — fill the field with the one-clause
        # `deliberately restrained: <why>` so the field is never blank either way."
        #
        # That escape existed in prose and nowhere in the gate, so an honest 5-minute lab-meeting plan
        # was rejected for lacking a rendered `signature_proof` — a PNG proving an aesthetic risk
        # survived the build, for a deck that deliberately took none. The only way out was
        # {"waived": …}, which also switches off palette, type_scale and icon_family: the real choice
        # was "invent a risk" or "abandon all design gating", and both corrupt the record.
        #
        # What it deliberately does NOT do: it does not relax `signature_move` (the field is still never
        # blank — you must WRITE why restraint is the position) and it does not relax `carried_by`,
        # `palette`, `type_scale`, `icon_family` or `form_ledger`. And it cannot be claimed above the
        # conservative dial: at balanced+ and above a real signature move is required, not optional.
        design = _section(gates, "design_plan")

        # BUILD SHAPE — was the build fanned out, and if not, why not. The build step is 40-71% of all
        # model-active minutes (SKILL.md, five measured sessions), and the fan-out rule that addresses
        # it lived only in prose: a 13-slide deck was built solo at 241 round-trips against a ~125
        # budget, with the batching rules in context and unfollowed, and nothing anywhere asked why.
        # Same pattern as form_reach: the gate never blocks the CHOICE — solo is legitimate everywhere
        # and mandatory on hosts without subagent dispatch — only the absence of a decision.
        if design and not design.get("waived"):
            _n_slides = _deck_slide_count(pptx)
            _shape = str(design.get("build_shape", "")).strip()
            if _n_slides >= 7 and not _shape:
                die("`design_plan.build_shape` is missing on a {}-slide deck. From ~6 content slides "
                    "up the build FANS OUT (one author per section, fresh context each — SKILL.md "
                    "'Scaling up'), because the build step is 40-71% of a session's model-active "
                    "minutes and a saturated context is where the batching rules stop being followed.\n"
                    '    "build_shape": "fanout — <n> sections"\n'
                    '    "build_shape": "solo — <reason>"  (e.g. "solo — host has no subagent '
                    'dispatch", "solo — one tightly-coupled argument")\n'
                    "  Solo is a legitimate answer on every host; what is not legitimate is nobody "
                    "having decided.".format(_n_slides))
            elif _shape:
                print("[gates] build shape: {}".format(_shape[:100]))

        _dial = str(design.get("boldness", "")).strip().lower()
        _move = str(design.get("signature_move", "")).strip().lower()
        # Validate the enum. Every dial-keyed branch in this file and in codex_delivery_gate.py is an
        # equality test against "conservative", so an unrecognised value is not a loud error — it is a
        # SILENT demotion to "not carved, not conservative". Verified: `"boldness": "BANANA"` printed
        # "all hand-off gates pass". A typo must not be a way out of a dial-keyed rule.
        _DIALS = ("conservative", "balanced+", "bold", "experimental")
        if design and not design.get("waived") and _dial and _dial not in _DIALS:
            die("`design_plan.boldness` is {!r}, which is not a dial. One of: {}.\n"
                "  (Every dial-keyed rule tests for `conservative`, so an unrecognised value silently "
                "reads as 'not conservative' rather than failing — which makes a typo an escape.)"
                .format(design.get("boldness"), " | ".join(_DIALS)))
        _carved = _dial == "conservative" and _move.startswith("deliberately restrained")
        if design.get("waived"):
            print("[gates] design plan WAIVED — {}".format(design["waived"]))
        elif design:
            # `motif_generates` takes the SAME carve as `signature_proof`: under a conservative dial
            # with a recorded "deliberately restrained" move there is no loud motif to be productive,
            # and demanding three things it makes would push an author to invent a device so the field
            # has an answer — the exact failure agents/slide-design.md names ("never invent an artifact
            # so this field has an answer"). A tiny 1-2 slide ask carves out for the same reason the
            # anchor proof does; that one is not detectable here, so it rides the design_plan waiver.
            _skip = {"signature_proof", "motif_generates"} if _carved else set()
            required = [f for f in DESIGN_FIELDS if f not in _skip]
            missing = [f for f in required if not design.get(f)]
            if missing:
                hint = ("\n    palette: the resolved FILL-only vs TEXT-safe split — run\n"
                        "      python3 scripts/palette_audit.py --from-style <deck>/style.py\n"
                        "    and paste what it hands back. A hue that works as a fill can read at\n"
                        "    2-4:1 as text on the same tint; the matrix is what stops that being\n"
                        "    found a render later." if "palette" in missing else "")
                die("`design_plan` is missing {}. These are the art director's outputs "
                    "(agents/slide-design.md, Step 2) — a plan without them was not designed, it was "
                    "defaulted. Fill them, or waive with a reason.{}".format(", ".join(missing), hint))
            # DECLARED -> APPLIED. `style_pick` was required as a STRING and never verified: measured
            # by grep, `presets.apply` / `set_geometry` / `set_ground` appear in no gate script at
            # all, so a deck recording "brutalist for engineering - beat blueprint" and built with
            # deckkit's stock defaults passed here and on the Codex path alike. The competition ran,
            # the winner was written down, and nothing carried it into the build.
            _style_applied_gate(design, pptx)
            # The tokens are checked against the RECORD (assets/sources.json) and against the built
            # deck, not merely against the grammar: a `searched, none found` rung with no recorded
            # search, and a CC BY photo with no credit on any slide, are both invisible to a
            # field-presence check and both were shipping.
            _image_provenance_gate(design, pptx)
            _direction_gate(design, os.path.dirname(os.path.abspath(pptx)) or ".")
            scale = design["type_scale"]
            if not isinstance(scale, dict) or not all(
                    isinstance(scale.get(k), (int, float)) for k in ("display", "title", "body")):
                die('`type_scale` must resolve the three tiers as numbers, e.g. '
                    '{"display": 34, "title": 24, "body": 14}. SIZE SPRAWL tells authors to draw sizes '
                    '"from the deck\'s declared type-scale tokens" — this is where they get declared. '
                    'A deck with no scale does not have restrained typography, it has whatever each '
                    'slide happened to pick.')
            # same floors the Codex delivery gate uses, so the two paths cannot disagree about what
            # counts as legible body type
            BODY_FLOORS = {"presented": 13.5, "textheavy": 13.5, "selfread": 12.0}
            # `delivery` is resolved ONCE at the top of this function (recorded key > CLI mode, with
            # surface aliased to selfread) and is the same value the density gate uses. Hardcoding
            # "presented" here made --selfread INERT for this floor: a self-read deck with body 12pt
            # died citing the *presented* floor while the same flag correctly drove density.
            floor = BODY_FLOORS[_DELIVERY_ALIAS.get(delivery, delivery)]
            if scale["body"] < floor:
                die(f'`type_scale.body` is {scale["body"]}pt, under the {floor}pt floor for a '
                    f'{delivery} deck — that is a legibility floor, not a style choice.')
            if not (scale["display"] > scale["title"] > scale["body"]):
                die(f'`type_scale` is not a scale: display {scale["display"]} > title {scale["title"]} '
                    f'> body {scale["body"]} must hold, or the tiers do not rank and the hierarchy is '
                    f'decorative rather than structural.')
            if _carved:
                print("[gates] boldness=conservative with a 'deliberately restrained:' signature move — "
                      "signature_proof not required (there is no risk to prove); every other field is")
            proof = None if _carved else design["signature_proof"]
            if proof is not None:
                # THE CONTRACT LIVES IN scripts/anchor_proof.py, imported by BOTH gate paths.
                # `png` or `path`: the Codex delivery gate spells this key `path` in its own evidence file,
                # and a Codex run keeps BOTH records (references/codex-runtime.md). Demanding one spelling
                # here would reject the field an OpenAI-bridged run naturally writes — the same evidence,
                # rejected for its key name. That divergence is exactly why the rule is no longer written
                # out twice: this file and codex_delivery_gate.py import the same module, so they cannot
                # drift apart again. Only the FILE checking differs by design (Codex binds a SHA-256 to
                # the final PPTX; here it is existence plus a size floor), and flattening the two would
                # weaken the stricter one.
                sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
                import anchor_proof as _ap
                _n = _deck_slide_count(pptx)
                _bad = _ap.faults(proof, n_slides=_n,
                                  expected_slides=set(range(1, _n + 1)) if _n else None)
                if _bad:
                    die("`signature_proof` (the ANCHOR PROOF, Step 4):\n    - "
                        + "\n    - ".join(_bad)
                        + "\n  A move that exists only as a sentence in the plan is the documented "
                          "failure: it gets sanded back to the safe catalogue during the build and "
                          "nobody notices, because the plan still reads bravely. The two anchors "
                          "beside it catch the other two: a design approved on a spacious page that "
                          "cannot hold the deck's densest one, and charts that obey none of the "
                          "palette/type decisions made against text.")
                for _a in _ap.normalise(proof):
                    proof_file = _ap.anchor_file(_a)
                    png = Path(proof_file)
                    if not png.is_absolute():
                        png = Path(pptx).parent / png
                    if not png.exists() or png.stat().st_size < 512:
                        die('`signature_proof` {} anchor ({}) does not exist or is empty — render '
                            'the slide first (render_deck.py <deck> <dir>) and point at the real PNG.'
                            .format(_a.get("role"), proof_file))
                    if _png_is_flat(png):
                        die('`signature_proof` {} anchor ({}) is a single flat colour — that is not '
                            'a render of a slide. Measured: a 960x540 rectangle of one grey used to '
                            'satisfy the ANCHOR PROOF, the mechanism whose whole purpose is to put '
                            'rendered evidence where the design decision is made.'
                            .format(_a.get("role"), proof_file))
            # ── THE MATERIAL PROBE (Step 2's opening block) ──────────────────────────────────
            # The one required declaration on the design plan that cannot be written without
            # having MADE something. Every other field is a sentence, and a deck shipped six grey
            # rectangles under the sentence `signature move: 封面自己演示论点` — true of nothing on
            # the page. The concept gates all passed; no step ever asked what the device is made of.
            probe = design.get("material_probe")
            if not isinstance(probe, dict):
                die('`design_plan.material_probe` is missing. Step 2 opens by BUILDING one real '
                    'slide — the signature page in the register you invented — rendering it, and '
                    'looking at it, before any of this plan is written.\n'
                    '    {"material_probe": {"png": "render/slideNN.png",\n'
                    '                        "safe_version": "<what the DEFAULT version of this '
                    'page would have been — if it is about the same thing, the register is a look, '
                    'not a move>"}}')
            ppng = Path(str(probe.get("png") or ""))
            if not ppng.is_absolute():
                ppng = Path(pptx).parent / ppng
            if not str(probe.get("png") or "") or not ppng.exists() or ppng.stat().st_size < 512:
                die('`design_plan.material_probe.png` must point at a REAL rendered slide '
                    '({!r} is missing or empty). The probe is the artifact; a plan describing it '
                    'is not.'.format(str(probe.get("png") or "")))
            if _png_is_flat(ppng):
                die('`design_plan.material_probe.png` is a single flat colour — not a render.')
            if reason_width(probe.get("safe_version")) < 20:
                die('`design_plan.material_probe.safe_version` must say in one line what the SAFE '
                    'version of that page would have been. It is the whole test: if the honest '
                    'answer is "about the same thing", the register is a look rather than a move, '
                    'and that is worth discovering before twenty declarations are written on it.')
            print("[gates] material probe: {} · safe version would have been: {}".format(
                probe.get("png"), str(probe.get("safe_version"))[:90]))

            cb = design["carried_by"]
            if not isinstance(cb, list) or len(cb) < 2:
                die("`carried_by` must name at least 2 slides where the signature move does structural "
                    "work. One brave slide among eleven safe ones is a tonal break, not a position.")
            # SKILL.md, slide-design.md and review-rubrics.md all name the same three answers as the
            # SAFE CATALOGUE rather than a signature move. Nothing checked, so the literal example
            # passed: a plan whose signature_move was the string "a big number" was accepted and
            # printed approvingly by this gate. A denylist is trivially evaded by paraphrase and that
            # is fine — what it closes is the case that actually happens, the example copied verbatim
            # because it was the nearest words to hand. Judging whether a REAL move is bold stays the
            # critic's distinctiveness axis; this only refuses the three the skill already disowned.
            # CONCEPT — what the deck's idea is a picture of, and the two pictures it beat.
            # The pipeline diverged on STYLE (the direction gate: "the same four slide types … only the
            # style differs") and on LAYOUT (form-selection's per-slide runner-up) and never on the
            # IDEA. A motif does not fill that hole: it is chosen as an attribute of a preset picked
            # first and capped at <=3 appearances, so a governing image is structurally forbidden from
            # governing. The field is cheap by design — three sentences at plan time, no extra dispatch
            # and no extra round trip — so the only thing worth checking is that three genuinely
            # different pictures were considered, not one relabelled twice.
            _con = design.get("concept")
            if _con is not None and not design.get("waived"):
                _rej = (_con.get("rejected") if isinstance(_con, dict) else None) or []
                _win = str((_con.get("chosen") if isinstance(_con, dict) else _con) or "").strip()
                if not isinstance(_con, dict) or not _win or len(_rej) < 2:
                    die('`design_plan.concept` must name the governing image AND the two it beat:\n'
                        '    "concept": {"chosen": "<what this deck is a picture of>",\n'
                        '                "rejected": [{"concept": "<the runner-up>", "why_lost": "<one clause>"},\n'
                        '                             {"concept": "<the other>", "why_lost": "<one clause>"}]}\n'
                        "  One picture with no alternatives is not a choice, it is the first thing that "
                        "came to mind — which is the default this field exists to interrupt.")
                _names = [_win.lower()] + [" ".join(str((r or {}).get("concept", "")).lower().split())
                                           for r in _rej]
                if len(set(n for n in _names if n)) < 3:
                    die("`design_plan.concept` lists the same picture more than once ({}). Three "
                        "governing images for one argument, not one relabelled — a network, an "
                        "organism and a pair of hands want different motifs, different colour logic "
                        "and different covers, which is the whole reason to choose between them."
                        .format(" · ".join(_names)))
                for r in _rej:
                    if not str((r or {}).get("why_lost", "")).strip():
                        die("`design_plan.concept.rejected` needs `why_lost` on each entry — a rejected "
                            "concept with no reason is a decoration on the record, not a decision.")
                print("[gates] concept: {} · beat {}".format(
                    _win[:60], " · ".join(str((r or {}).get("concept", "?"))[:28] for r in _rej)))
            _sm = " ".join(str(design["signature_move"]).lower().split())
            _CATALOGUE = ("a big number", "a nice gradient", "a full-bleed photo", "a full bleed photo")
            if any(_sm == c or _sm.rstrip(".") == c for c in _CATALOGUE):
                die('`signature_move` is {!r} — the skill names that (with "a nice gradient" and "a '
                    'full-bleed photo") as the SAFE CATALOGUE, explicitly NOT a signature move '
                    '(SKILL.md Step 2 · agents/slide-design.md self-verify (h) · '
                    'references/review-rubrics.md distinctiveness).\n'
                    "  The field wants the ONE aesthetic RISK a template would not take, scoped to "
                    "where it lands, and doing structural work on the carried_by slides — the motif "
                    "becoming the shape of the content, not a decoration repeated.\n"
                    "  If this deck genuinely takes no risk, that is a legitimate answer with its own "
                    'arm: set `boldness: conservative` and write `signature_move: "deliberately '
                    'restrained: <why>"`.'.format(design["signature_move"]))
            print("[gates] design plan: boldness={} · signature={} · carried_by={}".format(
                design["boldness"], str(design["signature_move"])[:48], cb))
            _report_carried_by(pptx, cb)
            # the declared palette, re-tested against the BUILT file — same reason as the icon
            # waiver above: a plan field written before any slide exists proves nothing about
            # the slides.
            _report_palette_drift(pptx, design.get("palette"))
            _report_plan_files(pptx)
            _report_icon_waiver(pptx, design.get("icon_family"), gates, delivery)
            _low_reach = _report_form_reach(pptx)
            # 🔴 A report that never asks for an answer is a line people learn to scroll past. This one
            # printed `1 of 23 named components; the rest is raw box/text` on a delivered deck and let
            # it through — and three of that deck's review findings were defects the unused components
            # prevent by construction (a label grazing its bar, a value floating off a track's
            # centreline, a reference line drawn three different ways).
            #
            # It still does NOT block on the number. Bespoke composition is legitimate and often the
            # signature move itself; a Mondrian page cannot come from a catalogue. What it blocks on is
            # the absence of a DECISION — the failure is never having looked, and `sigs.py --list` is
            # one call. Same shape as the sameness waiver: a written reason, in the record.
            if _low_reach is not None:
                _fr = design.get("form_reach")
                _why = str((_fr or {}).get("waived") or "").strip() if isinstance(_fr, dict) else ""
                if reason_width(_why) < 24:
                    die("form reach is {n} of {t} named components and the rest of {s} is raw "
                        "box/text, with no recorded reason.\n"
                        "  Look once — `python3 scripts/sigs.py --list` (or --search <shape>) — then "
                        "either build the component, or record WHY the hand-rolled form is the right "
                        "one for this deck:\n"
                        '    {{"design_plan": {{..., "form_reach": {{"waived": "<why bespoke here — '
                        'name the form and what a component would have cost it>"}}}}}}\n'
                        "  This never blocks on the NUMBER; it blocks on there being no decision. "
                        "Measured: a deck shipped at 1 of 23, and three of its review findings were "
                        "defects the unused components prevent by construction."
                        .format(n=len(_low_reach["forms"]), t=_low_reach["total"],
                                s=_low_reach["script"]))
                print("[gates] form reach WAIVED — %s" % _why[:150])
        else:
            die("no `design_plan` record. Step 2 dispatches agents/slide-design.md as the deck's art "
                "director; nothing else decides deck rhythm or the signature move.\n"
                # Built from DESIGN_FIELDS so this template cannot drift behind the gate again — it
                # listed SIX of the eight for a while, so copying it verbatim produced a record that
                # died on "missing type_scale, signature_proof".
                + "    {\"design_plan\": {" + ", ".join('"%s": ...' % f for f in DESIGN_FIELDS) + "}}\n"
                '    (type_scale is {"display": 34, "title": 24, "body": 14}; signature_proof is the\n'
                '     THREE-anchor list [{"role": "signature"|"complex"|"data", "slide": N,\n'
                '     "png": "render/slideNN.png"}, ...] — rendered evidence that the move survived, that\n'
                '     the design holds the densest page, and that the charts speak the same language)\n'
                '    or {"design_plan": {"waived": "<reason>"}}')

    with _gate_section('content.arc'):
        # THE ARC COMPETITION, ported from the Codex record so both paths hold the same bar. The design
        # side has had a rendered competition for years (the direction gate) and the content side got
        # one in 3e4eddb — but its verdict landed only on the content CHECKPOINT, which is prose in a
        # conversation, while every other Step-1/2 decision reaches this file. That asymmetry put the
        # cheaper decision under a gate and left the costlier one out: a wrong form costs one slide, a
        # wrong arc costs the design plan and the build underneath it.
        #
        # Shape is IDENTICAL to codex_delivery_gate's `content.arc`, deliberately — the two gates have
        # already drifted on a duplicated field twice (`path` vs `png`, the missing `conservative`
        # dial), and a third spelling would be the same mistake with a new name.
        #
        # 🔴 THE VERDICT IS RECOMPUTED HERE, NOT READ. The first version required `chosen` + one
        # `rejected` with a clause + a `divergence` string, and its own comment said the quiet part:
        # "`picked contribution-first` on its own is a sentence anyone can write without a competition
        # having happened". The fix it chose — demand the losers and their clauses — raised the price
        # of fabricating the record without making it impossible, because every one of those fields is
        # still prose the run authors about itself.
        #
        # Measured on a delivered deck: `content.arc.divergence` was the two-character string "ok",
        # `chosen`/`rejected` were filled in plausibly at build time, the gate passed, and
        # `arc_divergence.py` had never been invoked for that deck. The session transcript shows three
        # invocations on the one deck that had a real interview and zero on the two that were
        # delegated. A gate that reads a verdict is a gate the run can dictate to.
        #
        # So the gate takes the CANDIDATES and runs `arc_divergence.check()` itself. Faking it now
        # requires authoring 2-3 arcs that mechanically diverge on four axes and each carry comparable
        # ledger evidence — which is not a cheaper way to pass, it is the work. Same move `e7336e6`
        # made for icons ("ask the FILE, so the check survives a run that writes no record"), applied
        # to the costliest Step-1 decision in the pipeline.
        content = _section(gates, "content")
        if content.get("waived"):
            print("[gates] arc competition WAIVED — {}".format(content["waived"]))
        else:
            arc = content.get("arc")
            if not isinstance(arc, dict):
                die("`content.arc` is missing — the arc competition (Step 1, "
                    "agents/content-planner.md §3).\n"
                    "  2-3 candidate arcs over ONE ledger. Put the CANDIDATES here; this gate scores "
                    "them itself:\n\n"
                    '    "content": {"arc": {"chosen": "<the name of the arc that won>",\n'
                    '                        "candidates": [<2-3 arc objects — '
                    'python3 scripts/arc_divergence.py --template>],\n'
                    '                        "rejected": [{"name": "<runner-up>", '
                    '"why_lost": "<one clause>"}]}}\n\n'
                    '  Or waive it: {"content": {"waived": "<why this deck had one possible arc>"}}.')
            cands = arc.get("candidates")
            if not isinstance(cands, list) or len(cands) < 2:
                die("`content.arc.candidates` must carry the 2-3 candidate arcs THEMSELVES, not a "
                    "report about them.\n"
                    "  This gate runs scripts/arc_divergence.py over them right here, so a pasted "
                    "verdict is no longer evidence:\n"
                    "  a delivered deck passed with `\"divergence\": \"ok\"` while the script had "
                    "never run for it.\n"
                    "  `python3 scripts/arc_divergence.py --template` prints the fillable shape "
                    "(name · shape · roles ·\n"
                    "  audience_question · objection · closing_ask · evidence[ledger ids]).")
            # The raw `--template` skeleton passed this gate: its fields are angle-bracket
            # placeholders, and every check below measured WIDTH, which a placeholder has. Refuse a
            # candidate set that still carries them BEFORE scoring — arc_divergence's own template is
            # supposed to pass its divergence checker, so the skeleton is caught here, at delivery,
            # not in the divergence axes. `<label-A>` for a `chosen` was literally what shipped.
            _ph = []
            for _ci, _c in enumerate(cands):
                if not isinstance(_c, dict):
                    _ph.append("candidates[{}] is {} not an arc object".format(
                        _ci, type(_c).__name__))
                    continue
                for _f in ("name", "audience_question", "objection", "closing_ask"):
                    if _has_placeholder(_c.get(_f)):
                        _ph.append("candidates[{}].{}".format(_ci, _f))
            if _has_placeholder(arc.get("chosen")):
                _ph.append("chosen")
            for _ri, _r in enumerate(arc.get("rejected") or []):
                if isinstance(_r, dict):
                    for _f in ("name", "why_lost"):
                        if _has_placeholder(_r.get(_f)):
                            _ph.append("rejected[{}].{}".format(_ri, _f))
            if _ph:
                die("`content.arc` still carries the raw `--template` placeholder(s): {}.\n"
                    "  The unedited skeleton passes every width floor — `<label-A>` is wide — so it "
                    "would ship an arc\n"
                    "  competition that never happened. Fill the fields, or waive the arc with a real "
                    "reason.".format(", ".join(_ph[:8]) + (" …" if len(_ph) > 8 else "")))
            try:
                _verdict = _arc_verdict(cands)
            except Exception as exc:                              # noqa: BLE001
                die("`content.arc.candidates` is not a readable candidate set: {}\n"
                    "  Same reader as the CLI, so `python3 scripts/arc_divergence.py <arcs>.json` "
                    "reproduces this exactly.".format(exc))
            _names = [str(c.get("name") or "?") for c in cands if isinstance(c, dict)]
            _chosen = str(arc.get("chosen") or "").strip()
            if not _chosen:
                die("`content.arc.chosen` is empty — name the arc that won.")
            if _chosen not in _names:
                die("`content.arc.chosen` is {!r}, which is not one of the candidates ({}).\n"
                    "  A winner from outside the field means the competition scored a set the deck "
                    "was not built from.".format(_chosen, ", ".join(_names)))
            _rejected = arc.get("rejected")
            if not isinstance(_rejected, list) or not _rejected:
                die("`content.arc.rejected` must name every arc the winner beat, with the clause "
                    "that lost it.\n"
                    "  A winner with no losers on the record is a derivation wearing a "
                    "competition's clothes.")
            _named = {}
            for _i, _row in enumerate(_rejected):
                if not isinstance(_row, dict) or not str(_row.get("name") or "").strip() \
                        or reason_width(_row.get("why_lost")) < 8:
                    die("`content.arc.rejected[{}]` needs a `name` and a `why_lost` clause — the "
                        "reason is the whole point.".format(_i))
                _named[str(_row["name"]).strip()] = _row["why_lost"]
            _missing = [n for n in _names if n != _chosen and n not in _named]
            if _missing:
                die("`content.arc.rejected` skips {}, which competed and lost.\n"
                    "  Every candidate that is not the winner needs its losing clause; recording "
                    "only the flattering half of\n"
                    "  the field is how a competition quietly becomes a derivation again."
                    .format(", ".join(repr(m) for m in _missing)))
            _stray = [n for n in _named if n not in _names]
            if _stray:
                die("`content.arc.rejected` names {}, which never competed — it is not in "
                    "`candidates`.".format(", ".join(repr(m) for m in _stray)))
            _flagged = _verdict.get("flagged") or []
            _sketch = _verdict.get("sketches") or []
            _noledger = bool(_verdict.get("no_ledger"))
            if _flagged or _sketch or _noledger:
                _finding = "; ".join(filter(None, [
                    "{} pair(s) tell the same story: {}".format(
                        len(_flagged), ", ".join("{}~{}".format(p["a"], p["b"]) for p in _flagged))
                    if _flagged else "",
                    "strawman candidate(s): {}".format(", ".join(
                        "{} carries {}/{} evidence units".format(x["name"], x["evidence"], x["top"])
                        for x in _sketch)) if _sketch else "",
                    "no candidate names its ledger evidence" if _noledger else ""]))
                if reason_width(arc.get("divergence_justified")) < 16 \
                        or _has_placeholder(arc.get("divergence_justified")):
                    die("the arc competition, RECOMPUTED here, does not hold: {}.\n"
                        "  Rediverge the candidates — move the CLAIM, not the wording: put the "
                        "evidence before the claim instead\n"
                        "  of after it, argue to a different ask, pre-empt a different objection. "
                        "Or keep the set and record why in\n"
                        "  `content.arc.divergence_justified`. This is the REDIVERGE-or-justify "
                        "rule the script prints; it is\n"
                        "  never an auto-kill.".format(_finding))
                print("[gates] arc competition: {} candidates scored HERE — flagged ({}) · "
                      "justified: {}".format(len(cands), _finding, arc["divergence_justified"]))
            else:
                print("[gates] arc competition: {} candidates scored HERE — none is a rewording of "
                      "another, none is a sketch; won: {}".format(len(cands), _chosen))

    with _gate_section('content.slides'):
        # THE CONTENT CHECKPOINT BECOMES AN ARTIFACT ON THIS PATH TOO. Step 1 ends in a table —
        # `# | 角色 | 记忆句 | 承载证据 | units` (references/checkpoint-convention.md) — and on the
        # shared path that table lived ONLY as prose in the conversation. Under the per-deck AUTO
        # WAIVER it is posted as an FYI rather than a stop, so nothing waits for it and nothing
        # reports its absence.
        #
        # Measured across one session: the table was posted for the ONE deck that had a real
        # interview and for NEITHER of the two that opened with "you decide the rest" — and those two
        # are the decks whose design was judged flat and whose direction was judged wrong. A
        # per-slide takeaway table is exactly the artifact that catches a wrong direction before a
        # single slide is built: the reader sees twelve memory sentences and none is about the thing
        # they asked for.
        #
        # 🔴 THE FIELD NAME IS `content.slides`, NOT A NEW ONE, because codex_delivery_gate has
        # required this exact record all along ("content.slides must cover every final slide exactly
        # once", with role/takeaway/evidence per row). The asymmetry was never that the artifact had
        # not been invented — it was that the CODEX path demanded it and the shared path did not, so
        # a Claude Code or Kimi run could deliver with no per-slide plan at all. A third spelling
        # would have been the drift those two gates have already suffered twice.
        #
        # This gate cannot check that a table was pasted into a chat, and does not pretend to. It
        # checks that the per-slide plan EXISTS, covers every slide exactly once, and that no two
        # content slides carry the same memory sentence — the cheap mechanical signature of a plan
        # written for the deck rather than per slide.
        _content = _section(gates, "content")
        _sw = _content.get("slides_waived")
        if _sw:
            if reason_width(_sw) < 16 or _has_placeholder(_sw):
                die("`content.slides_waived` needs a real reason (>=16 wide), not a token or a "
                    "`<placeholder>` — say what this deck is that a per-slide plan does not apply to.")
            print("[gates] content checkpoint rows WAIVED — {}".format(_sw))
        else:
            _rows = _content.get("slides")
            _n_slides = _deck_slide_count(pptx)
            if not isinstance(_rows, list) or not _rows:
                die("`content.slides` is missing — the content checkpoint's table "
                    "(references/checkpoint-convention.md),\n"
                    "  one row per slide. Same shape codex_delivery_gate has always required:\n\n"
                    '    "content": {"slides": [\n'
                    '        {"slide": 1, "role": "cover",\n'
                    '         "takeaway": "<the one sentence this slide leaves behind>",\n'
                    '         "evidence": ["<SOURCE TRACE — a locator, not a label>"],\n'
                    '         "units": 1},\n'
                    '        ...]}\n\n'
                    "  Read top to bottom, the `takeaway` column IS the takeaway spine. Under the "
                    "auto waiver the table is still\n"
                    "  posted in chat as an FYI — the waiver removes the stop, never the record; "
                    "this field IS that record.\n"
                    '  Or waive it: {"content": {"slides_waived": "<why this deck has no per-slide '
                    'plan>"}}.')
            _seen, _takeaways = {}, {}
            for _i, _row in enumerate(_rows):
                if not isinstance(_row, dict):
                    die("`content.slides[{}]` must be an object with slide / role / takeaway / "
                        "evidence.".format(_i))
                _n = _row.get("slide")
                if not isinstance(_n, int) or isinstance(_n, bool):
                    die("`content.slides[{}].slide` must be the slide NUMBER, got {!r}."
                        .format(_i, _row.get("slide")))
                if _n in _seen:
                    die("`content.slides`: rows {} and {} both claim slide {} — one slide, one row."
                        .format(_seen[_n], _i, _n))
                _seen[_n] = _i
                if not str(_row.get("role") or "").strip():
                    die("`content.slides[{}].role` is empty — name the beat this slide plays "
                        "(cover / problem / evidence / conclusion …).".format(_i))
                if reason_width(_row.get("takeaway")) < 8 or _is_stub(_row.get("takeaway")):
                    die("`content.slides[{}]` (slide {}) has no real 记忆句 / takeaway.\n"
                        "  It is the sentence the room keeps after the slide is gone. A label like "
                        "{!r} is a topic, and a topic\n"
                        "  is what a slide is ABOUT, not what it SAYS."
                        .format(_i, _n, _row.get("takeaway")))
                _ev = _row.get("evidence")
                if not isinstance(_ev, list) or not _ev or \
                        any(not isinstance(x, str) or _is_stub(x) for x in _ev) or \
                        reason_width(" ".join(str(x) for x in _ev)) < 4:
                    die("`content.slides[{}]` (slide {}) has no 承载证据.\n"
                        "  A LIST of SOURCE TRACEs — a locator (\"Fig 3 / p.4 ¶2\", a table cell, a "
                        "short verbatim span), or an\n"
                        "  honest statement of what kind of claim it is. It is the cheapest "
                        "per-slide grounding catch on the\n"
                        "  path delegation uses most.".format(_i, _n))
                _u = _row.get("units")
                if _u is not None and (not isinstance(_u, int) or isinstance(_u, bool) or _u < 0):
                    die("`content.slides[{}].units` must be a non-negative count, got {!r}."
                        .format(_i, _u))
                if str(_row.get("role")).strip().lower() not in _STRUCTURAL_ROLES:
                    _key = _norm_takeaway(_row.get("takeaway"))
                    if _key in _takeaways:
                        die("`content.slides`: slides {} and {} carry the SAME takeaway.\n"
                            "  Two content slides with one memory sentence is the mechanical "
                            "signature of a plan written for the\n"
                            "  DECK rather than per slide — and one of the two pages has no reason "
                            "to exist. Give each its own\n"
                            "  sentence, or merge them.".format(_takeaways[_key], _n))
                    _takeaways[_key] = _n
            if _n_slides and set(_seen) != set(range(1, _n_slides + 1)):
                _miss = sorted(set(range(1, _n_slides + 1)) - set(_seen))
                _extra = sorted(set(_seen) - set(range(1, _n_slides + 1)))
                die("`content.slides` must cover every final slide exactly once — {}{}\n"
                    "  Every slide is a beat someone chose; a plan that stops short of the deck is "
                    "the half that got built\n"
                    "  without one.".format(
                        "no row for slide(s) {}. ".format(_miss) if _miss else "",
                        "row(s) for slide(s) {} that the deck does not have. ".format(_extra)
                        if _extra else ""))
            print("[gates] content plan: {} row(s), one per slide, {} distinct takeaway(s)"
                  .format(len(_rows), len(_takeaways)))
            # The units column exists to make an about-to-be-empty or about-to-be-dense page visible
            # AT THE CHECKPOINT rather than at the render. Printed, never fatal: a 0 on a pure-image
            # beat and a 6 on a spoken one are both legitimate, and only the author knows which.
            _thin = sorted(r["slide"] for r in _rows if isinstance(r.get("units"), int)
                           and not isinstance(r.get("units"), bool) and r["units"] == 0
                           and str(r.get("role", "")).strip().lower() not in _STRUCTURAL_ROLES)
            _dense = sorted(r["slide"] for r in _rows if isinstance(r.get("units"), int)
                            and not isinstance(r.get("units"), bool) and r["units"] >= 6)
            if _thin:
                print("[gates]   units=0 on content slide(s) {} — an empty beat, or a count nobody "
                      "filled".format(_thin))
            if _dense:
                print("[gates]   units>=6 on slide(s) {} — check these against the density dial "
                      "before the build".format(_dense))

    with _gate_section('checkpoints'):
        # DELEGATION MUST BE VISIBLE AT DELIVERY. The per-deck auto waiver turns the two 🔴 stops into
        # FYIs, and SKILL.md is explicit that it "removes the stop, never the record". But a run that
        # posted both FYIs and a run that posted neither hand over identical decks, so the user finds
        # out which one they got by reading the slides.
        #
        # Field names and the mode vocabulary are codex_delivery_gate's, unchanged
        # (`content.checkpoint` / `design.checkpoint`, mode `approved` | `auto`, plus a `record`):
        # that gate has required both since it was written, and this is the second artifact the
        # shared path was missing rather than a new idea.
        #
        # This section does not — cannot — verify that text appeared in a conversation. Its teeth are
        # the sections above it: `content.slides` and `design_plan` ARE the checkpoints' artifacts,
        # and this ledger prints each checkpoint's mode next to whether its artifact exists. An
        # `approved` claimed over a waived artifact is then a contradiction on a single line.
        _cn = _section(gates, "content")
        _MODES = ("approved", "auto")
        _have = {
            "content": ("slides x{}".format(len(_cn.get("slides") or []))
                        if isinstance(_cn.get("slides"), list) and _cn.get("slides")
                        else ("WAIVED" if _cn.get("slides_waived") else "NO ARTIFACT")),
            "design": ("design_plan"
                       if isinstance(gates.get("design_plan"), dict)
                       and not (_section(gates, "design_plan")).get("waived") else "WAIVED")}
        _ck = {"content": (_section(gates, "content")).get("checkpoint"),
               "design": (_section(gates, "design")).get("checkpoint")
               or (_section(gates, "design_plan")).get("checkpoint")}
        _bad = [k for k, v in _ck.items()
                if not isinstance(v, dict) or v.get("mode") not in _MODES
                or reason_width(v.get("record")) < 12 or _has_placeholder(v.get("record"))]
        if _bad:
            die("the 🔴 checkpoint record is missing or malformed for: {}.\n\n"
                '    "content": {{"checkpoint": {{"mode": "approved", "record": "<what was shown '
                'and what came back>"}}}},\n'
                '    "design":  {{"checkpoint": {{"mode": "auto", "record": "<what was shown and '
                'what came back>"}}}}\n\n'
                "  `approved` = presented, and the user approved before the build continued.\n"
                "  `auto`     = presented under the per-deck AUTO WAIVER, and the build continued "
                "(the waiver removes\n"
                "               the stop, never the record).\n"
                "  Same two values codex_delivery_gate has always accepted. There is no third one: a "
                "run that SKIPPED the\n"
                "  checkpoint has nothing true to write here, and that is the point — delegation "
                "changes WHO approves,\n"
                "  not WHETHER the step happened.".format(", ".join(sorted(_bad))))
        print("[gates] checkpoint ledger — content: {} ({}) · design: {} ({})".format(
            _ck["content"]["mode"], _have["content"],
            _ck["design"]["mode"], _have["design"]))
        # NOTED, not fatal, and the distinction is the lesson this whole file is about. `approved`
        # over a waived artifact LOOKS like a contradiction and is not one: a content checkpoint
        # carries a deck memory sentence, a ledger digest and an emotional-curve line besides the
        # per-slide table, so a user can genuinely approve one whose table was waived for a written
        # reason. A rule that blocks a legitimate state is worse than no rule — it teaches authors
        # to write whichever value gets them through. The ledger line above is the deliverable; this
        # is the sentence under it.
        _odd = [k for k in ("content", "design")
                if _ck[k]["mode"] == "approved" and "WAIVED" in _have[k]]
        if _odd:
            print("[gates]   note: {} says `approved` over a WAIVED artifact — legitimate if the "
                  "user approved the rest of that checkpoint, worth a second look if not"
                  .format(", ".join(sorted(_odd))))

    with _gate_section('render_selfcheck'):
        # THE ACTOR'S OWN LOOK, TURNED FROM PROSE INTO A TRACE. Step 5 tells the coordinator to read
        # every slide PNG and "record a one-line verdict for EVERY slide — a slide with no line was
        # not checked." That was prose with no backstop, so the cheap actor-side look (the one that
        # catches an overflow, a cropped subject, a wrong number BEFORE a critic round is spent) was
        # the easiest step to skip silently — the "passed because nothing looked" failure this skill
        # keeps re-finding.
        #
        # This makes the look leave a mark: one verdict per slide, covering every slide, no
        # placeholders. Same SHAPE and same honest LIMIT as content.slides — it proves the trace
        # exists, not that the eye judged well (a lazy "ok" on a bad slide still passes). The strong
        # guarantee that every slide was JUDGED is the independent critic's coverage bind above; this
        # is the cheaper pre-filter in front of it, and its value is that a MISSING slide is now
        # visible instead of silent. Genuinely no render to look at (a --static/no-render edge, or a
        # deck this gate should not touch)? Waive it in writing.
        _rs = _section(gates, "render_selfcheck")
        _rw = _rs.get("waived")
        if _rw:
            if reason_width(_rw) < 16 or _has_placeholder(_rw):
                die("`render_selfcheck.waived` needs a real reason (>=16 wide), not a token/"
                    "placeholder — say why there is no rendered deck to look at.")
            print("[gates] render self-check WAIVED — {}".format(_rw))
        else:
            _sl = _rs.get("slides")
            _n = _deck_slide_count(pptx)
            if not isinstance(_sl, list) or not _sl:
                die("`render_selfcheck.slides` is missing — Step 5's per-slide verdict, one line per "
                    "slide (a slide with no line was not looked at):\n\n"
                    '    "render_selfcheck": {"slides": [\n'
                    '        {"n": 1, "verdict": "ok — signature move lands, 4.9% buried below the '
                    'line"},\n'
                    '        {"n": 2, "verdict": "teal glyph on aqua tile <3:1 — recoloured"},\n'
                    '        ...]}\n\n'
                    "  Read render/contact.png, then EVERY slide PNG in one message, and write the "
                    "verdict as you go.\n"
                    '  No rendered deck to look at? Waive it: {"render_selfcheck": {"waived": '
                    '"<why>"}}.')
            if _n and len(_sl) != _n:
                die("`render_selfcheck.slides` has {} verdict(s) for a {}-slide deck — one per "
                    "slide, covering every slide. A slide with no verdict is a slide nobody looked "
                    "at.".format(len(_sl), _n))
            _seen_rs = {}
            _limit_rs = _n or len(_sl)
            for _i, _row in enumerate(_sl):
                if not isinstance(_row, dict):
                    die("`render_selfcheck.slides[{}]` must be an object with n / verdict."
                        .format(_i))
                _rn = _row.get("n")
                if not isinstance(_rn, int) or isinstance(_rn, bool) or not 1 <= _rn <= _limit_rs:
                    die("`render_selfcheck.slides[{}].n` must be a slide number in 1..{}, got {!r}."
                        .format(_i, _limit_rs, _row.get("n")))
                if _rn in _seen_rs:
                    die("`render_selfcheck.slides`: two verdicts both claim slide {} — one per "
                        "slide.".format(_rn))
                _seen_rs[_rn] = _i
                if reason_width(_row.get("verdict")) < 4 or _is_stub(_row.get("verdict")):
                    die("`render_selfcheck.slides[{}]` (slide {}) has no real verdict. It is the "
                        "one line you write having LOOKED at the slide — `ok` is fine on a clean "
                        "page, but an empty/placeholder line is a slide that was not read."
                        .format(_i, _rn))
            print("[gates] render self-check: {} verdict(s), one per slide".format(len(_sl)))

    with _gate_section('provenance'):
        # Provenance: a self-filled tally proves nothing — the refutation pass is what the gate is FOR.
        # Require per-claim verdicts, so "confirmed" means someone tried to break it and could not.
        prov = _section(gates, "provenance")
        if prov.get("waived"):
            print("[gates] provenance WAIVED — {}".format(prov["waived"]))
        elif prov:
            claims = prov.get("claims")
            if not isinstance(claims, list) or not claims:
                die("`provenance` needs a per-claim `claims` list, not a summary tally. A tally is "
                    "written by the same pass that would have skipped the check.\n"
                    '    "claims": [{"claim": "...", "verdict": "CONFIRMED|WRONG|PARTLY-WRONG|'
                    'UNVERIFIABLE", "url": "https://..."}]')
            ALLOWED = {"CONFIRMED", "WRONG", "PARTLY-WRONG", "UNVERIFIABLE"}
            bad = [c for c in claims if c.get("verdict") not in ALLOWED]
            if bad:
                die("{} claim row(s) carry no valid verdict (one of {}).".format(
                    len(bad), " / ".join(sorted(ALLOWED))))
            nourl = [c for c in claims if c.get("verdict") == "CONFIRMED" and not c.get("url")]
            if nourl:
                die("{} claim(s) are CONFIRMED with no primary-source URL. Confirmed-without-a-source "
                    "is the exact failure this gate exists to catch.".format(len(nourl)))
            unresolved = [c for c in claims if c["verdict"] in ("WRONG", "PARTLY-WRONG")]
            if unresolved:
                die("{} claim(s) still verdict WRONG / PARTLY-WRONG. Fix or cut them before hand-off:\n"
                    "  - {}".format(len(unresolved),
                                    "\n  - ".join(str(c.get("claim"))[:90] for c in unresolved[:5])))
            tally = {}
            for c in claims:
                tally[c["verdict"]] = tally.get(c["verdict"], 0) + 1
            print("[gates] provenance: {} claim(s) adversarially checked — {}".format(
                len(claims), " · ".join("{} {}".format(v, k) for k, v in sorted(tally.items()))))
        else:
            print("[gates] no provenance record — fine for a deck built from the user's own material; "
                  "a research-sourced deck should carry one.")

    with _gate_section('sameness'):
        # ── SAMENESS: the deck-level monotony the [stats] block measured and nobody read ──────────
        # Every deck-level "this deck is one page repeated" signal the linter computes is a
        # `warns.append` printed under a line that says the stats are advisory; `lint()` returns only
        # the hard-finding count, and nothing on this path ever read the warns. So blandness was
        # DETECTED deterministically and blocked nothing. This turns the measurement into a gate —
        # and only the measurement: whether a deck is TIMID stays the critic's taste call (which the
        # skill deliberately caps as non-blocking), whether it is REPETITIVE is a share of slides
        # agreeing with each other, which is a defect with a concrete fix. `agents/critic.md` states
        # that exact test, and it is why this can hold a deck while the distinctiveness axis cannot.
        #
        # 🔴 THAT PARAGRAPH USED TO END HERE, and the asymmetry it describes was the bug: "whether a
        # deck is TIMID stays the critic's taste call" left the safe side of the scale with no
        # deterministic voice at all. Some of timidity IS measurable — whether any page has a
        # protagonist that is not a sentence, whether the palette does any work — and that half now
        # has a gate of its own, below.
        _check_sameness(pptx, delivery, gates)

    with _gate_section('a11y'):
        _check_a11y(pptx, delivery, gates)

    with _gate_section('surface'):
        _surface_gate(pptx, gates)

    with _gate_section('register_pixels'):
        _register_pixels_gate(pptx)

    with _gate_section('register_guard'):
        _register_guard_gate(pptx, gates)
        _register_kit_note(pptx, gates)
        _register_keep_note(pptx, gates)

    with _gate_section('timidity'):
        _check_timidity(pptx, delivery, gates)

    with _gate_section('density'):
        # ── DENSITY: a slide is a visual aid, not a document ────────────────────────
        # This one is a gate rather than a warning because the warning already existed and was
        # already ignored — twice, by the same author, on two consecutive decks. Measured: one
        # deck shipped with 8 of 12 slides over the presented text budget, the next with 12 of 12
        # (loads of 81-144 words against a budget of ~40), and both times the per-slide TEXT WALL
        # line was read and dismissed as advisory. The skill's OWN reference deck — the file it
        # tells every builder to copy — runs at a median of 27 words a slide, so the budget is not
        # unrealistic; what failed was that nothing made ignoring it cost anything.
        # A deck may legitimately be denser (a self-read leave-behind, a spec sheet). That is what
        # the waiver is for: not a rule against text, a rule against text arriving by default.
        #
        # DELIVERY MODE. The budget here is lint_deck's budget, taken from lint_deck: 70 words for a
        # presented deck, 120 for a self-read one, and no budget at all on a poster (`surface`) or a
        # deck whose density the user CHOSE at Q4 (`textheavy`). A gate that fires on a mode the
        # interview offers, the rubric protects and the lint deliberately passes does not enforce
        # anything for long — it teaches the author to paste a waiver, and after that it is decoration.
        txt = gates.get("density")
        # the SAME resolved delivery the type floor used — reading raw `mode` here is how one run
        # enforced two different deliveries
        if delivery in ("surface", "textheavy"):
            print("[gates] density: not applied — %s deck (the user chose this density, or the "
                  "surface has no per-slide budget)" % delivery)
            return
        over, total, median = _density_stats(pptx, budget=70 if delivery == "presented" else 120)
        if total:
            if isinstance(txt, dict) and txt.get("waived"):
                print("[gates] density: {}/{} slide(s) over the presented text budget, median {} "
                      "words — WAIVED: {}".format(over, total, median, str(txt["waived"])[:110]))
            elif over * 3 > total:
                die("{} of {} slides are over the {} text budget (median {} words a slide; "
                    "aim ~40, warn >{}). The skill's own reference deck runs at 27.\n"
                    "    A slide is a visual aid for a speaker — the sentences belong in the speaker "
                    "notes, which this deck already has.\n"
                    "    Cut the on-slide prose, or record the deliberate choice:\n"
                    '    "density": {{"waived": "why this deck is meant to be read, not presented"}}'
                    .format(over, total, delivery, median, 70 if delivery == "presented" else 120))
            else:
                print("[gates] density: {}/{} slide(s) over the text budget, median {} words a slide"
                      .format(over, total, median))


SAMENESS_WAIVER_KINDS = {
    "series-frame":
        "a carousel / board / poster series where the REPEATED FRAME is the artifact",
    "register-uniform":
        "a deliberately single-register deck (留白 / ink-wash, a uniformly dark editorial briefing)",
    "template-locked":
        "a registered or provided template whose grid this deck may not break",
    "reference-run":
        "the repeated pages are reference material (first try design_intent(role='appendix') — "
        "that is the real fix, and it takes those pages out of the count)",
    "user-waived":
        "the user was shown the finding and chose to ship over it",
}


def _sameness_stats(pptx, delivery):
    """The deck-level sameness codes, taken from lint_deck's OWN measurement.

    Same reason `_density_stats` calls `reading_load`: the number in the [stats] line and the
    number in this gate have to be ONE number by construction. Nothing is re-implemented here —
    a second implementation is how lint came to say 136 words a slide while the gate said 4.

    `renders_dir=None` keeps lint's normal auto-discovery of ./render beside the deck, INCLUDING
    its staleness guard, so the gate inherits the same freshness rule the CLI has.
    """
    import contextlib
    import io
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    import lint_deck as _ld
    from pptx import Presentation
    try:
        prs = Presentation(pptx)
        aspect = (prs.slide_width / float(prs.slide_height)) if prs.slide_height else 0.0
    except Exception as exc:
        die("the sameness gate could not open {} ({}). An unreadable deck is not a pass."
            .format(pptx, exc))
    stats, buf = {}, io.StringIO()
    gates_path = os.path.join(os.path.dirname(os.path.abspath(pptx)) or ".", GATES_FILE)
    try:
        with contextlib.redirect_stdout(buf):        # lint prints its whole report; we want the data
            _ld.lint(pptx, mode=delivery, renders_dir=None, static_ok=True,
                     gates_path=gates_path, stats_out=stats)
    except SystemExit as exc:
        die("the sameness gate could not lint {} (lint_deck exited {}). An unreadable deck is not "
            "a pass.".format(pptx, exc.code))
    return stats, aspect


def _check_sameness(pptx, delivery, gates):
    """Block a hand-off when the deck measurably repeats itself, unless the repetition is declared.

    Scope and threshold are calibrated against 11 decks BUILT AND LINTED in the registers this
    skill itself prescribes, not against intuition. The raw signal count is a bad gate: a 6-slide
    template-locked status update with ZERO hard findings already reaches 3 families, a 9-card
    小红书 carousel built exactly to canvas-formats.md's DNA reaches 7, and an appendix-heavy
    defense deck reaches 5. Each of those is a deck the skill tells you to build.

    So the scope is three DETERMINISTIC properties of the deck rather than a taxonomy of registers
    (the registers have no working wire today: --briefing / --selfread / --textheavy relieve
    nothing on the sameness side, and --surface relieves three signals only at n=1):
        body_n >= 8    — lint's own body run, cover/closer/appendix excluded. Kills the whole
                         status-update class, whose measured cliff was SIX total slides.
        landscape      — kills the carousel class, whose repeated frame IS the artifact.
        not surface    — a single canvas has no deck-level rhythm to vary.
    and the threshold is >= 4 distinct codes WITH at least one structural code. Three is where a
    competent consistent system lands (repeated layout + one card + one footer strip); the
    structural requirement excludes the one 4-code set that means "same frame, varied bodies".
    """
    if not isinstance(gates, dict):
        gates = {}
    waiver = _section(gates, "sameness")
    stats, aspect = _sameness_stats(pptx, delivery)
    fired = tuple(stats.get("sameness_codes") or ())
    body_n = int(stats.get("body_n") or 0)
    ran = bool(stats.get("render_signals_ran"))
    skip_reason = stats.get("render_skip_reason")
    total = len(_LD_SAMENESS_CODES())
    could_run = total if ran else total - len(_LD_SAMENESS_RENDER_DEPENDENT())
    tail = "" if ran else "  · NOT CHECKED: {} ({})".format(
        " · ".join(_LD_SAMENESS_RENDER_DEPENDENT()), skip_reason or "no renders beside the deck")

    if delivery == "surface" or aspect < 1.2 or body_n < 8:
        why = ("a single-canvas surface" if delivery == "surface"
               else "a portrait/square canvas — a series' repeated frame is the artifact"
               if aspect < 1.2 else "%d content slide(s), under the 8 this is calibrated for" % body_n)
        print("[gates] sameness: not applied — {} (the per-signal [stats] warnings still print)"
              .format(why))
        return

    structural = [c for c in fired if c in _LD_SAMENESS_STRUCTURAL()]
    blocks = len(fired) >= 4 and bool(structural)
    listed = " · ".join(fired) if fired else "none"

    if waiver:
        reason = waiver.get("waived")
        kind = waiver.get("waived_category")
        if not isinstance(reason, str) or reason_width(reason) < 40:
            die("`sameness.waived` must be a written reason that names the REGISTER this deck is "
                "in, not a circumstance — a sentence someone can disagree with later (>=40 chars).")
        if kind not in SAMENESS_WAIVER_KINDS:
            die("`sameness.waived_category` must name WHICH kind of deliberate repetition this "
                "is. One of:\n" + "\n".join("    {:18s} {}".format(k, v)
                                            for k, v in sorted(SAMENESS_WAIVER_KINDS.items())))
        recorded = waiver.get("codes")
        if not isinstance(recorded, list) or set(recorded) != set(fired):
            die("`sameness.codes` records {} but this deck now fires {}. A waiver written for a "
                "different state of the deck does not certify this one — re-read the measurement "
                "and rewrite the reason against it.\n    codes: {}"
                .format(sorted(recorded) if isinstance(recorded, list) else recorded,
                        sorted(fired) or "nothing", json.dumps(sorted(fired))))
        if not blocks:
            print("[gates] sameness: recorded waiver is NOT needed — {} of {} signal(s) fired ({})"
                  .format(len(fired), could_run, listed) + tail)
        else:
            print("[gates] sameness: WAIVED [{}] — {} of {} fired: {}".format(
                kind, len(fired), could_run, listed) + tail)
            print("        {}".format(reason.strip()))
        return

    if blocks:
        die("SAMENESS: {} of {} deck-level signals fired across {} content slides —\n"
            "        {}\n\n"
            "  Each is a MEASURED share, not a taste call: this deck repeats one page structure, "
            "one content vehicle, one piece of chrome or one vertical envelope across most of its "
            "pages.\n"
            "  Fix the deck — rotate the canvas architecture, break the card grid on the WOW "
            "slide, move the takeaway off the bottom strip on some pages, let one page bleed —\n"
            "  or record why the repetition IS the design:\n\n"
            '    "sameness": {{"waived": "<why this deck repeats on purpose>",\n'
            '                 "waived_category": "{}",\n'
            '                 "codes": {}}}\n\n'
            "  Categories: {}"
            .format(len(fired), could_run, body_n, " · ".join(fired),
                    " | ".join(sorted(SAMENESS_WAIVER_KINDS)), json.dumps(sorted(fired)),
                    "; ".join("{} = {}".format(k, v)
                              for k, v in sorted(SAMENESS_WAIVER_KINDS.items()))))

    # The boundary the count cannot decide: one more signal would have crossed the line, and one
    # signal did not run. Passing here would be a verdict about a measurement that was never taken.
    if len(fired) == 3 and structural and not ran:
        die("SAMENESS: 3 of the {} signals that could run fired ({}), and {} did not run ({}).\n"
            "  At three fired the verdict turns on the signal that is missing, so this deck cannot "
            "be decided as it stands.\n"
            "  Render it —  python3 scripts/render_deck.py {} <out>  — and re-run --gate-check, "
            "or record the waiver."
            .format(could_run, listed, " · ".join(_LD_SAMENESS_RENDER_DEPENDENT()),
                    skip_reason or "no renders beside the deck", pptx))

    print("[gates] sameness: {} of {} deck-level signal(s) fired across {} content slides{}"
          .format(len(fired), could_run, body_n,
                  " — " + listed if fired else "") + tail)


def _LD_SAMENESS_CODES():
    import lint_deck as _ld
    return _ld.SAMENESS_CODES


def _LD_SAMENESS_STRUCTURAL():
    import lint_deck as _ld
    return _ld.SAMENESS_STRUCTURAL


def _LD_SAMENESS_RENDER_DEPENDENT():
    import lint_deck as _ld
    return _ld.SAMENESS_RENDER_DEPENDENT


def _LD_TIMIDITY_CODES():
    import lint_deck as _ld
    return _ld.TIMIDITY_CODES


def _LD_TIMIDITY_STRUCTURAL():
    import lint_deck as _ld
    return _ld.TIMIDITY_STRUCTURAL


TIMIDITY_WAIVER_KINDS = {
    "register-restrained":
        "a deliberately quiet register (ink-wash, mono spec sheet, memorial) where restraint IS the design",
    "text-is-the-artifact":
        "the words themselves are the deliverable (a quote deck, a legal/《条款》 read-through)",
    "template-locked":
        "a registered or provided template this deck may not push past",
    "user-waived":
        "the user was shown the finding and chose to ship over it",
}


def _check_timidity(pptx, delivery, gates):
    """🔴 THE COUNTERWEIGHT. Hold a deck that is measurably SAFE, the way sameness holds one that
    is measurably repetitive.

    Every other blocking signal in this pipeline punishes excess. Measured on a real 12-page build:
    ten iterations, each one driven by an advisory, each one making the deck flatter — the dark
    pivot page deleted for ONE-OFF CANVAS FLIP, content cut for TEXT WALL, the type scale collapsed
    for SIZE SPRAWL. Every one of those advisories names the ambitious repair FIRST ("enrich with a
    second column of substance", "repeat the treatment as a divider family"); subtraction is simply
    the cheaper way to make the number go away, and with feedback on one side only, the cheap way
    always wins. The user's verdict on that deck was "设计能力变弱了" — and NOTHING in the pipeline
    had said so, because the one force that could (the critic's distinctiveness axis) is
    non-blocking at the default dial and lives inside a review that can be declined.

    Shaped exactly like the sameness gate, for the same reason: several weak signals, a structural
    requirement, a size floor, and a waiver that must NAME the register. It blocks at >= 2 codes
    with >= 1 structural — TIMID COVER and FLAT TYPE are one fact counted twice, so type drama
    alone can never hold a deck.

    It stands down entirely under `boldness: conservative` with a recorded `deliberately
    restrained:` move, because that is the dial saying "restraint is the position" — the same carve
    signature_proof already honours.
    """
    if not isinstance(gates, dict):
        gates = {}
    stats, aspect = _sameness_stats(pptx, delivery)
    fired = tuple(stats.get("timidity_codes") or ())
    body_n = int(stats.get("body_n") or 0)

    design = _section(gates, "design_plan")
    dial = str(design.get("boldness", "")).strip().lower()
    move = str(design.get("signature_move", "")).strip().lower()
    if dial == "conservative" and move.startswith("deliberately restrained"):
        print("[gates] timidity: not applied — boldness=conservative with a recorded "
              "`deliberately restrained` move; restraint IS the position here")
        return
    if delivery == "surface" or aspect < 1.2 or body_n < 8:
        why = ("a single-canvas surface" if delivery == "surface"
               else "a portrait/square canvas" if aspect < 1.2
               else "%d content slide(s), under the 8 this is calibrated for" % body_n)
        print("[gates] timidity: not applied — {} (the per-signal [stats] warnings still print)"
              .format(why))
        return

    structural = [c for c in fired if c in _LD_TIMIDITY_STRUCTURAL()]
    blocks = len(fired) >= 2 and bool(structural)
    listed = " · ".join(fired) if fired else "none"
    # A signal that could not run is not a signal that passed — say which, exactly as sameness does.
    ran = bool(stats.get("render_signals_ran"))
    import lint_deck as _ld
    tail = "" if ran else "  · NOT CHECKED: {} ({})".format(
        " · ".join(_ld.TIMIDITY_RENDER_DEPENDENT),
        stats.get("render_skip_reason") or "no renders beside the deck")
    waiver = _section(gates, "timidity")

    if waiver:
        reason = waiver.get("waived")
        kind = waiver.get("waived_category")
        if not isinstance(reason, str) or reason_width(reason) < 40:
            die("`timidity.waived` must NAME the register that makes this restraint deliberate — "
                "a sentence someone can disagree with later (>=40 wide).")
        if kind not in TIMIDITY_WAIVER_KINDS:
            die("`timidity.waived_category` must say WHICH kind of deliberate restraint this is. "
                "One of:\n" + "\n".join("    {:22s} {}".format(k, v)
                                          for k, v in sorted(TIMIDITY_WAIVER_KINDS.items())))
        recorded = waiver.get("codes")
        if not isinstance(recorded, list) or set(recorded) != set(fired):
            die("`timidity.codes` records {} but this deck now fires {}. A waiver written for a "
                "different state of the deck does not certify this one.".format(
                    sorted(recorded or []), sorted(fired)))
        print("[gates] timidity: {} of {} signal(s) fired ({}) — WAIVED [{}]: {}{}".format(
            len(fired), len(_LD_TIMIDITY_CODES()), listed, kind, str(reason)[:110], tail))
        return

    if blocks:
        die("this deck measures SAFE on {} of {} signals ({}) — and unlike every other gate here, "
            "that is a finding about ambition rather than about error.\n"
            "    Nothing is broken; the deck simply takes no position a template would not have "
            "taken. The repairs the advisories name FIRST are the ones that fix this — enrich a "
            "page with a real protagonist (a figure, a chart, a form whose geometry IS the "
            "argument), give the deck a rhythm event, let the palette mean something.\n"
            "    If the restraint is the design, say which register:\n"
            '    "timidity": {{"waived": "<the register, named>", "waived_category": "<{}>", '
            '"codes": {}}}'.format(len(fired), len(_LD_TIMIDITY_CODES()), listed,
                                   " | ".join(sorted(TIMIDITY_WAIVER_KINDS)), sorted(fired)))
    print("[gates] timidity: {} of {} signal(s) fired ({}){}".format(
        len(fired), len(_LD_TIMIDITY_CODES()), listed, tail))


def _density_stats(pptx, budget=70):
    """(slides over the text budget, total slides, median load).

    Calls `lint_deck.reading_load` — the SAME function the per-slide TEXT WALL warning calls —
    so the number in the warning and the number in this gate are one number by construction.
    An earlier version of this only shared the string-level word counter and re-implemented the
    per-slide accumulation and the chrome filter, which is where all the drift actually lives:
    its `sz <= 10.5 and len(t) < 40` skip had no POSITION test, so it was not a footer filter
    but an amnesty for small type anywhere on the slide. Measured on one deck: lint said 136
    words a slide, the gate said 4, and a wall of 10.5pt prose passed. Sharing a helper is not
    the same as sharing the measurement — share the measurement.
    """
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from lint_deck import _boxes, reading_load
        from pptx import Presentation
        prs = Presentation(pptx)
        sw = prs.slide_width / 914400.0
        sh = prs.slide_height / 914400.0
    except Exception:
        return 0, 0, 0
    loads = [reading_load(sl, _boxes(sl, sw, sh), sh) for sl in prs.slides]
    if not loads:
        return 0, 0, 0
    loads.sort()
    return sum(1 for x in loads if x > budget), len(loads), loads[len(loads) // 2]



def _contact_sheet(out, n_pages):
    """One image of the WHOLE deck, written beside the per-slide PNGs as `contact.png`.

    🔴 ADDITIVE, and the distinction is not a nicety. Measured on a real 12-page deck, a page in
    this sheet is 460px wide — 46 px/inch against the render's 144. At that size a 23pt title is
    15px and readable, and **13.5pt body text is 8.6px and a 10pt source line is 6.4px: neither
    can be read at all**. So this image answers DECK-LEVEL questions only — the light/dark rhythm,
    whether the bookends bookend, form variety, whether one chrome treatment is stamped on every
    page, a canvas flip that lands on exactly one slide. It is the view the deck-level checks
    (ONE-OFF CANVAS FLIP · TITLE-RULE MONOCULTURE · FLAT RHYTHM · BOTTOM-STRIP MONOCULTURE)
    measure, finally visible to a reader.

    It CANNOT answer a per-slide question — typography, contrast, a label grazing its bar, an
    overlap, whether a number is right — and it must never be used as if it had. The render
    self-check still reads slides page by page; this sheet is what you look at BEFORE that, to
    know which pages deserve the attention. Substituting it for the per-slide read would be the
    same failure this skill keeps finding elsewhere: passing a check by not looking.

    Silent no-op without Pillow, and never fatal — a missing contact sheet must not fail a render.
    """
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return None
    import glob as _glob
    pngs = sorted(_glob.glob(os.path.join(out, "slide*.png")),
                  key=lambda p: int(re.sub(r"\D", "", os.path.basename(p)) or 0))
    if len(pngs) < 2:
        return None                                      # a 1-page deck has no deck-level shape
    # SPLIT rather than shrink past legibility. A single sheet of 40 pages measures 1890x2860, and
    # an image that tall is downscaled to fit before anyone looks at it — measured, a 0.55 factor,
    # which takes a tile from 460px to 252px and a 23pt title from 15px to 8px. The sheet would
    # still be produced, still be called the deck-level view, and no longer show anything: a
    # silent degradation, which is the failure mode this skill treats as worse than a loud one.
    # 12 per sheet is the size that was verified readable on a real deck, so a long deck gets
    # several readable sheets instead of one useless one.
    PER = 12
    chunks = [pngs[i:i + PER] for i in range(0, len(pngs), PER)]
    made = []
    try:
        w0, h0 = Image.open(pngs[0]).size
        tw = 460
        th = max(1, int(round(tw * h0 / float(w0))))
        pad, lab = 10, 16
        for ci, chunk in enumerate(chunks):
            cols = 4 if len(chunk) > 6 else (3 if len(chunk) > 2 else 2)
            rows = -(-len(chunk) // cols)
            sheet = Image.new("RGB", (cols * tw + (cols + 1) * pad,
                                      rows * (th + lab) + (rows + 1) * pad), (235, 235, 238))
            d = ImageDraw.Draw(sheet)
            for i, p in enumerate(chunk):
                im = Image.open(p).convert("RGB").resize((tw, th), Image.LANCZOS)
                x = pad + (i % cols) * (tw + pad)
                y = pad + (i // cols) * (th + lab + pad)
                sheet.paste(im, (x, y))
                d.rectangle([x, y, x + tw, y + th], outline=(180, 180, 186))
                d.text((x + 3, y + th + 2), "%02d" % (ci * PER + i + 1), fill=(60, 60, 66))
            name = "contact.png" if len(chunks) == 1 else "contact_%02d.png" % (ci + 1)
            dest = os.path.join(out, name)
            sheet.save(dest)
            made.append(dest)
        return made
    except Exception:
        return None                                      # never fail a render over a preview

def _read_json(path):
    """Read a JSON file, returning {} on any error (missing / unreadable / malformed) — the
    design-checkpoint gate treats an unreadable evidence file as 'nothing recorded', never a crash."""
    try:
        with open(path, encoding="utf-8") as f:
            d = json.load(f)
        return d if isinstance(d, dict) else {}
    except Exception:
        return {}


def _pptx_slide_count(pptx):
    """Count slides in a .pptx by its package parts — cheap, and it runs BEFORE LibreOffice is
    located, so the design-checkpoint gate can fail early without a rasterizer present."""
    try:
        import re
        import zipfile
        with zipfile.ZipFile(pptx) as z:
            return sum(1 for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n))
    except Exception:
        return 0


def _content_plan_slide_count(deck_dir):
    """How many slides the APPROVED CONTENT PLAN (Step 1) names — read from either evidence file
    (the shared `.deck-gates.json` → `content.slides`, or the Codex `.codex-deck-evidence.json`).
    0 means no content plan is recorded yet (a test fixture or an ad-hoc render), so the Step-2
    design gate does not apply."""
    for fn in (GATES_FILE, ".codex-deck-evidence.json"):
        d = _read_json(os.path.join(deck_dir, fn))
        content = d.get("content")
        sl = content.get("slides") if isinstance(content, dict) else None
        if isinstance(sl, list) and sl:
            return len(sl)
    return 0


def _design_plan_and_checkpoint_present(deck_dir):
    """True iff a Step-2 design plan AND its design checkpoint are recorded — in the shared
    `.deck-gates.json` (`design_plan` + `design_plan.checkpoint`, or a top-level `design.checkpoint`)
    OR the Codex `.codex-deck-evidence.json` (`design` + `design.checkpoint`). This is the PRESENCE
    gate that forces Step 2 to have HAPPENED before the first full render; the DEPTH of the plan is
    validated separately by the hand-off gate (`--gate-check`) and `codex_delivery_gate.py`. Kept
    lenient on purpose — its job is to stop a design plan being reconstructed post-hoc at hand-off,
    not to re-judge the plan the deeper gates already own."""
    def _ck(obj):
        return isinstance(obj, dict) and obj.get("mode") in ("approved", "auto")
    g = _read_json(os.path.join(deck_dir, GATES_FILE))
    dp = g.get("design_plan")
    if isinstance(dp, dict) and dp:
        design = g.get("design")
        if _ck(dp.get("checkpoint")) or (isinstance(design, dict) and _ck(design.get("checkpoint"))):
            return True
    c = _read_json(os.path.join(deck_dir, ".codex-deck-evidence.json"))
    de = c.get("design")
    if isinstance(de, dict) and de and _ck(de.get("checkpoint")):
        return True
    return False


def main(argv):
    # --deliverables (alias --final): ALSO park the PDF beside the .pptx and write viewer.html.
    # OFF by default: while a deck is still being iterated, those two are pure churn — they are
    # regenerated every round, clutter the deck root, and go stale the moment the user hand-edits
    # the .pptx. They are produced once, at hand-off, when the user says the deck is final.
    deliverables = False
    # --fast: re-render ONLY the slides whose content changed since the last render.
    # Measured on a real 18-slide deck: full = 9.1s LibreOffice + 24.4s rasterize; one changed
    # slide = 2.5s + 0.7s. Rasterization, not the PDF export, is the dominant cost — so the win
    # comes from rasterizing one page, and subsetting the pptx makes the export cheap too.
    fast = False
    # --slides N[,M]: render ONLY these 1-indexed slides. Unlike --fast (which DIFFS against a
    # cache and therefore renders everything on a first run), this is an explicit "show me page N"
    # — the SIGNATURE PROOF that opens SKILL.md Step 4, and any "re-render the page I edited" loop.
    only = None
    argv = list(argv)
    while "--fast" in argv:
        argv.remove("--fast")
        fast = True
    for i, a in enumerate(list(argv)):
        if a == "--slides" or a.startswith("--slides="):
            raw = a.split("=", 1)[1] if "=" in a else (argv[i + 1] if i + 1 < len(argv) else "")
            try:
                only = sorted({int(t) for t in raw.replace(" ", "").split(",") if t})
            except ValueError:
                die("--slides wants 1-indexed slide numbers, e.g. --slides 1,6")
            if not only:
                die("--slides wants at least one slide number, e.g. --slides 6")
            argv.remove(a)
            if "=" not in a and raw in argv:
                argv.remove(raw)
            break
    # --gate-check runs the hand-off gates and exits, rendering nothing. The gates were reachable
    # only through --deliverables, which Step 6 deliberately makes a decline-able OFFER ("want a PDF
    # and a browser preview?"), so on every deck where the user said no, the strongest gate in the
    # skill never ran at all. Step 6 now calls this unconditionally, whatever the user answers.
    # Delivery mode — the SAME flags lint_deck.py takes, spelled the same way, because the gate
    # below enforces the lint's budget and the two must never disagree about which deck this is.
    mode = "presented"
    for flag, name in (("--selfread", "selfread"), ("--surface", "surface"),
                       ("--textheavy", "textheavy")):
        while flag in argv or ("--mode=" + name) in argv:
            argv = [a for a in argv if a not in (flag, "--mode=" + name)]
            mode = name
    # --static: accepted and intentionally inert. lint_deck.py takes it to silence NO BUILDS on a
    # deck whose user opted OUT of appear-builds; this tool already lints with static_ok=True, so
    # there is nothing for it to switch. It is consumed rather than rejected because callers pass
    # it by symmetry with lint_deck.py (the test suite does, on every gate call) — and before the
    # unknown-flag guard below existed it was silently absorbed as the OUTPUT DIRECTORY.
    while "--static" in argv or "--mode=static" in argv:
        argv = [a for a in argv if a not in ("--static", "--mode=static")]
    gate_only = False
    if "--gate-check" in argv:
        argv = [a for a in argv if a != "--gate-check"]
        gate_only = True
    for flag in ("--deliverables", "--final"):
        while flag in argv:
            argv.remove(flag)
            deliverables = True
    if fast and deliverables:
        # Decided here, not after LibreOffice has already run: --deliverables needs a whole-deck
        # PDF, which a subset render cannot produce. Failing late meant either an exit-1 after a
        # successful render, or (with nothing changed) a silent exit-0 that produced no PDF and no
        # viewer.html at the exact moment the hand-off contract required them.
        die("--deliverables needs a full-deck render — drop --fast for the hand-off run")
    if only and deliverables:
        # Same reason as --fast: a subset cannot produce the whole-deck PDF the hand-off promises.
        die("--deliverables needs a full-deck render — drop --slides for the hand-off run")
    if only and fast:
        # Contradictory intents: --fast decides WHICH slides to render, --slides declares them.
        die("--slides and --fast both choose the slide set — pass one")
    # Every flag this tool knows has been consumed by now, so anything still starting with `--`
    # is one it does not take — and until this guard existed it was not an error. `argv[0]` is
    # the deck and `argv[1]` is the output dir, so `render_deck.py deck.pptx --gate-check
    # --briefing` resolved to out="--briefing" and ran the hand-off gate at the `presented`
    # floor: a briefing deck (lint budget ~150 words) silently held to the ~40-word budget, with
    # nothing printed either way. That is the same silent-fallthrough class the `delivery` key
    # below already dies on ("a legibility floor silently not applied"); a mode flag typed at the
    # CLI deserves the same treatment. NB `--briefing` is a real lint_deck.py mode that this tool
    # genuinely does not implement — the message says so rather than pretending it is a typo.
    _stray = [a for a in argv if a.startswith("--")]
    if _stray:
        _hint = ""
        if any(a in ("--briefing", "--mode=briefing") for a in _stray):
            _hint += ("\n  `--briefing` is a lint_deck.py mode; the hand-off gate has no briefing "
                      "floor yet, so it would have run at `presented`. Lint the deck with "
                      "`lint_deck.py <deck> --briefing` and record `delivery` in .deck-gates.json.")
        die("unrecognised option(s): " + " ".join(_stray) + "\n  render_deck.py takes: --slides "
            "N[,M] · --fast · --deliverables/--final · --gate-check · --selfread · --textheavy · "
            "--surface (each mode also spelled --mode=NAME). The output dir is POSITIONAL: "
            "render_deck.py <deck>.pptx [out_dir]." + _hint)
    if not argv:
        die("usage: python3 render_deck.py /path/to/deck.pptx [out_dir] "
            "[--fast | --slides N[,M]] [--deliverables] [--gate-check]")
    pptx = argv[0]
    out = argv[1] if len(argv) > 1 else "./render"

    if not os.path.isfile(pptx):
        die("no such file: " + pptx)

    if gate_only:
        check_handoff_gates(pptx, mode, gate_check=True)
        print("[gates] all hand-off gates pass — the deck may be handed over")
        return 0

    # Checked here, before LibreOffice runs — same reason the --fast/--slides conflicts are:
    # failing after a successful render wastes the render and reads as a late surprise.
    if deliverables:
        check_handoff_gates(pptx, mode)

    # ── Step-2 DESIGN CHECKPOINT — a BRANCH-INVARIANT pre-build gate, enforced before the first
    # FULL render (and before LibreOffice is even located: an unplanned deck should not render, with
    # or without a rasterizer). A recorded content plan (Step 1) with NO design plan + design
    # checkpoint (Step 2) means the deck is about to be rendered/reviewed on a design that was never
    # planned or approved — the exact failure where `design_plan` gets reconstructed post-hoc to pass
    # the hand-off gate. It runs on EVERY Q1 branch: a generated hero / style gate confirms the LOOK,
    # it is NOT the design checkpoint. Reads BOTH evidence files (shared + Codex). EXEMPT: a `--slides`
    # probe (the Step-2 material probe / hero-checkpoint sample / signature proof all run this way,
    # before the plan is final); a deck with fewer slides than the plan (a probe/sample build); a deck
    # with no content plan or a 1–3 slide tiny ask (`_cp >= 4`). --gate-check / --deliverables already
    # run the deeper hand-off design_plan gate, so they are past this by construction.
    if only is None and not gate_only and not deliverables:
        _dd = os.path.dirname(os.path.abspath(pptx)) or "."
        _cp = _content_plan_slide_count(_dd)
        if _cp >= 4 and _pptx_slide_count(pptx) >= _cp and not _design_plan_and_checkpoint_present(_dd):
            die("STEP 2 NOT DONE — this deck has an approved content plan ({0} slides) but no design "
                "plan + design checkpoint recorded, and it is about to be FULL-rendered.\n"
                "  Step 2 (design plan + \U0001f534 design checkpoint) is BRANCH-INVARIANT — it runs on "
                "every Q1 choice: design-a-clean-one, a provided template, AND a generated visual "
                "identity. A generated hero / style gate confirms the LOOK; it is NOT the design "
                "checkpoint.\n"
                "  Produce the design plan — per-slide form ledger · deck rhythm · signature "
                "move under a boldness dial · the 3 design musts (builds/icons/formats) · "
                "semantic colour · density · logo/motif — post the \U0001f534 DESIGN "
                "CHECKPOINT, and record it: in `.deck-gates.json` as `design_plan` + "
                "`design_plan.checkpoint` ({{\"mode\": \"approved\"|\"auto\", \"record\": \"…\"}}), "
                "or on the Codex path in `.codex-deck-evidence.json` as `design` + `design.checkpoint`.\n"
                "  Rendering ONE probe slide first is expected and exempt — use `--slides N`.".format(_cp))

    soffice = find_soffice()
    if not soffice:
        die(
            "LibreOffice not found — needed to render slides for the verify + critic loop.\n"
            "  macOS:         brew install --cask libreoffice\n"
            "  Debian/Ubuntu: sudo apt install libreoffice\n"
            "  Windows:       winget install TheDocumentFoundation.LibreOffice\n"
            "                 (or choco install libreoffice-fresh)\n"
            "  other:         https://www.libreoffice.org/download\n"
            "  (or set the SOFFICE env var to the full path of the soffice binary)"
        )

    # Decide full vs incremental BEFORE spending anything on LibreOffice.
    # When out == the deck's own folder (the code below explicitly tolerates the user passing "."),
    # <deck>.pdf and viewer.html there are the HAND-OFF deliverables, not render products. Nothing
    # in this run may delete or overwrite them.
    _deck_dir = os.path.dirname(os.path.abspath(pptx)) or "."
    _out_is_deck_dir = os.path.abspath(out) == os.path.abspath(_deck_dir)

    cache_path = os.path.join(out, ".render-cache.json")
    # Cheap, but not free on a big deck — and a full render still needs them, both to validate the
    # page count and to leave a cache the NEXT --fast can diff against. Guarded so an unreadable
    # package falls back to LibreOffice's own diagnostic instead of a zipfile traceback (v3.5.1
    # behaviour, which the friendly error in troubleshooting-faq.md documents).
    try:
        fps, blockers = _slide_fingerprints(pptx)
    except Exception as exc:
        fps, blockers = [], ["could not read the .pptx package: {}".format(exc)]
    changed, skip_reason = None, None
    if fast:
        prev = None
        try:
            with open(cache_path, encoding="utf-8") as f:
                cached = json.load(f)
            prev = cached.get("fingerprints") if isinstance(cached, dict) else None
            if not (isinstance(prev, list) and all(isinstance(x, str) for x in prev)):
                prev = None                     # a JSON list/garbage must not crash the tool
        except Exception:
            prev = None
        if blockers:
            # Correctness beats speed: anything that makes the slide->page mapping uncertain
            # forces a full render.
            skip_reason = blockers[0]
        elif not prev:
            skip_reason = "no previous render cache"
        elif len(prev) != len(fps):
            # Slides inserted or deleted shifts every index after the edit; PNG filenames would
            # silently point at the wrong slides.
            skip_reason = "slide count changed ({} -> {})".format(len(prev), len(fps))
        else:
            changed = [i for i, h in enumerate(fps) if h != prev[i]]

            def _usable(name):                  # a 0-byte file from a killed rasterize is NOT a render
                p = os.path.join(out, name)
                try:
                    return os.path.getsize(p) > 0
                except OSError:
                    return False

            missing = [i for i in range(len(fps)) if not _usable("slide{:02d}.png".format(i + 1))]
            # thumbnails are regenerated only when a bookend slide is re-rendered, so a missing
            # thumb must pull its bookend into the changed set or the critic's poster test runs
            # on a stale or absent image while the run prints success
            if not _usable("thumb_first.png"):
                missing.append(0)
            if not _usable("thumb_last.png"):
                missing.append(len(fps) - 1)
            changed = sorted(set(changed) | set(i for i in missing if 0 <= i < len(fps)))

    try:
        _st = os.stat(pptx)
        pptx_stat = (_st.st_mtime_ns, _st.st_size)
    except OSError:
        pptx_stat = None
    if only is not None:
        if blockers:
            # The SAME correctness bar as --fast: anything that makes the slide->page mapping
            # uncertain (hidden slides, an unresolvable part, auto slide-number fields that
            # renumber inside a subset) must not be papered over just because the user named a page.
            die("cannot render a subset of this deck: {}\n"
                "  drop --slides and render the whole deck".format(blockers[0]))
        bad = [n for n in only if not (1 <= n <= len(fps))]
        if bad:
            die("--slides {} out of range — this deck has {} slide(s)".format(
                ",".join(str(b) for b in bad), len(fps)))
        changed = [n - 1 for n in only]
    incremental = (fast or only is not None) and changed is not None and 0 < len(changed) < len(fps)
    if fast and changed is not None and len(changed) == len(fps) and len(fps):
        # Every slide changed (a rebuild that touched everything, or a deck-global edit such as a
        # theme/canvas change). A subset of "all slides" is just a full render with extra steps.
        skip_reason = "every slide changed"
    if fast and changed is not None and not changed:
        print("no slide changed since the last render — nothing to re-render")
        print("next: python3 {} {} --renders {}".format(
            os.path.join(os.path.dirname(os.path.abspath(__file__)), "lint_deck.py"), pptx, out))
        return 0



    # Give this invocation a LibreOffice profile NO OTHER PROCESS IS USING: lets parallel
    # renders (the large-deck section fan-out) run at once without fighting a shared profile
    # lock, and lets the render work even while the user has the LibreOffice GUI open.
    # Without this, concurrent/coexisting soffice calls silently produce no PDF.
    # `_acquire_profile()` keeps that invariant with an exclusively-flocked slot from a
    # persistent pool (warm profile, ~1.5s cheaper per export) and falls back to the original
    # throwaway `mkdtemp` whenever a slot is unavailable — see `_Profile` for the whole story.
    src_pptx, keep = pptx, None
    tmp_subset = tmp_dir = None
    if incremental:
        keep = changed
        tmp_dir = tempfile.mkdtemp(prefix="lo_subset_")
        tmp_subset = os.path.join(tmp_dir, "subset.pptx")
        _subset_pptx(pptx, set(keep), tmp_subset)
        src_pptx = tmp_subset

    # A subset renders into its OWN temp dir. Writing out/subset.pdf would collide between two
    # concurrent runs sharing a render dir (which the per-invocation LibreOffice profile above
    # exists to allow), and a crashed run would leave a file that breaks the render-only cleanup.
    # ALWAYS convert into a private empty directory, never straight into `out`: `out` may already
    # hold a <deck>.pdf (a previous render, or a --deliverables artefact when out is the deck
    # folder), and a stale file there would make a failed conversion look like a successful one.
    if tmp_dir is None:
        tmp_dir = tempfile.mkdtemp(prefix="lo_render_out_")
    import atexit
    atexit.register(shutil.rmtree, tmp_dir, True)       # a die() mid-render must not leak a deck copy
    pdf, result, cmd = _render_pdf(soffice, src_pptx, tmp_dir)
    try:
        import fitz  # pymupdf
    except ImportError:
        die("pymupdf not installed — run: {} -m pip install pymupdf".format(
            os.path.basename(sys.executable) or "python"))

    if result.returncode != 0 or not os.path.isfile(pdf):
        detail = [
            "LibreOffice failed to convert {} (exit {}).".format(pptx, result.returncode)
            if result.returncode != 0 else "LibreOffice produced no PDF from {}.".format(pptx),
            "Command: " + " ".join(cmd),
            "Exit code: {}".format(result.returncode),
        ]
        stdout = _tail(result.stdout)
        stderr = _tail(result.stderr)
        if stdout:
            detail.append("stdout:\n" + stdout)
        if stderr:
            detail.append("stderr:\n" + stderr)
        detail.append(
            "Check that the file opens, close any open copy, and in sandboxed runtimes "
            "rerun the render with the permissions needed for LibreOffice."
        )
        die("\n".join(detail))

    # Open the PDF BEFORE the cleanup: a non-zero exit and a missing file were already fatal, but a
    # zero-exit run that writes a truncated/garbage PDF passed both checks, and the cleanup had
    # already deleted the previous render by the time fitz raised.
    try:
        _probe = fitz.open(pdf)
        _npages = _probe.page_count
        _probe.close()
    except Exception as exc:
        die("LibreOffice produced an unreadable PDF from {} ({}). The previous render in {} was "
            "left untouched.".format(pptx, exc, out))
    if _npages < 1:
        die("LibreOffice produced an empty PDF (0 pages) from {}. The previous render in {} was "
            "left untouched.".format(pptx, out))

    # Clear the previous render only NOW, after LibreOffice has actually produced a readable PDF. Doing it
    # earlier meant a failed conversion destroyed the last good render and left nothing at all —
    # the user lost working output to a run that produced none.
    # Wiping the render dir BEFORE deciding is what made --fast a no-op: every PNG looked
    # "missing", so every slide looked changed. An incremental run keeps the existing PNGs and
    # overwrites only the ones it re-renders.
    if os.path.isdir(out) and not incremental:
        entries = os.listdir(out)
        own_pdf = os.path.splitext(os.path.basename(pptx))[0] + ".pdf"
        render_only = all(
            (e.startswith(("slide", "thumb_")) and e.endswith(".png"))
            or e == "viewer.html" or e == own_pdf          # only THIS deck's fallback pdf is ours
            or e in (".DS_Store", "Thumbs.db", ".render-cache.json")   # OS junk + our own cache only
            for e in entries)                              # make the dir look deletable
        if render_only:
            shutil.rmtree(out, ignore_errors=True)
        else:
            # out holds files that are NOT ours (worst case: the user passed "." — the pptx's own
            # directory). NEVER rmtree it; clear only our previous render products.
            for e in entries:
                if _out_is_deck_dir and (e == "viewer.html" or not _RENDER_PNG.match(e)):
                    # out IS the deck folder: only files matching our own strict output pattern
                    # (slideNN.png / thumb_first|last.png) are ours. A user's slide_background.png
                    # or thumb_hero.png sitting beside the deck was being deleted silently.
                    continue
                if ((e.startswith(("slide", "thumb_")) and e.endswith(("png",)))
                        or e == "contact.png" or e.startswith("contact_")
                        or e == "viewer.html"):
                    try:
                        os.remove(os.path.join(out, e))
                    except OSError:
                        pass
    os.makedirs(out, exist_ok=True)


    doc = fitz.open(pdf)
    skip_cache = False
    if only is not None:
        # A --slides run rasterized SOME pages; the fingerprints in hand describe ALL of them.
        # Caching them would tell the next --fast "nothing changed" while most PNGs are stale —
        # precisely the lie the incremental path exists to avoid. Drop the cache instead, so the
        # next --fast honestly does a full render.
        skip_cache = True
    if incremental:
        # PDF page k corresponds to deck slide keep[k] — write ONLY those PNGs and leave the
        # rest of render/ untouched.
        if doc.page_count != len(keep):
            die("incremental render produced {} page(s) for {} changed slide(s) — refusing to "
                "write PNGs that may be mismatched; re-run without --fast".format(
                    doc.page_count, len(keep)))
        for k, page in enumerate(doc):
            page.get_pixmap(matrix=fitz.Matrix(2, 2)).save(
                os.path.join(out, "slide{:02d}.png".format(keep[k] + 1)))
        # thumbnails only matter when a bookend slide moved
        for name, idx in (("thumb_first", 0), ("thumb_last", len(fps) - 1)):
            if idx in keep:
                page = doc[keep.index(idx)]
                zoom = 240.0 / max(1.0, page.rect.width)
                page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(os.path.join(out, name + ".png"))
        n_pages = len(fps)
    else:
        pages = list(enumerate(doc, 1))
        # Fast path first; it returns False on a small deck or any trouble, and then both
        # loops below run exactly as they always did. When it succeeds it has ALSO written
        # the two bookend thumbnails, from the workers that own those pages — see
        # _rasterize_chunk for why they cannot be re-rendered here instead.
        drawn_in_parallel = _rasterize_parallel(pdf, doc.page_count, out)
        if not drawn_in_parallel:
            for i, page in pages:
                page.get_pixmap(matrix=fitz.Matrix(2, 2)).save(
                    os.path.join(out, "slide{:02d}.png".format(i)))
        # bookend thumbnails (~240px wide) for the critic's poster test — first + last slide small,
        # the scale at which a cover either survives or dies. Same PyMuPDF path, no new deps.
        if pages and not drawn_in_parallel:
            for name, (_, page) in (("thumb_first", pages[0]), ("thumb_last", pages[-1])):
                zoom = 240.0 / max(1.0, page.rect.width)
                page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(os.path.join(out, name + ".png"))
        n_pages = doc.page_count
        if n_pages != len(fps) and fps:
            # LibreOffice emitted a different number of pages than the deck has slides (hidden
            # slides are the known cause). slideNN.png no longer means deck slide NN, so no cache
            # may claim this render — and the user must be told rather than quietly trusting it.
            print("WARNING: {} slides in the deck but {} PDF page(s) rendered — slideNN.png may "
                  "not correspond to deck slide NN. Not caching fingerprints; --fast is disabled "
                  "until this is resolved.".format(len(fps), n_pages), file=sys.stderr)
            skip_cache = True
    doc.close()
    if incremental:
        pdf = None                              # a subset PDF is not the deck's PDF
        # Drop the previous FULL render's intermediate PDF, which the new PNGs have outdated —
        # but never when `out` IS the deck folder, where that same filename is the user's
        # hand-off deliverable. Deleting it there silently destroyed a --deliverables artefact.
        if not _out_is_deck_dir:
            try:
                os.remove(os.path.join(out, os.path.splitext(os.path.basename(pptx))[0] + ".pdf"))
            except OSError:
                pass

    # Record fingerprints so the NEXT run can diff against them. Written only after the PNGs
    # actually landed, so a crashed render never leaves a cache claiming work that did not happen.
    # Two guards: (1) if the .pptx changed WHILE we rendered, the fingerprints we hold describe
    # state that was never rasterized — caching them would freeze that slide stale forever;
    # (2) a failed write must DELETE the cache, because a cache older than the PNGs is exactly the
    # "no slide changed" lie this whole path must not tell.
    if not skip_cache and fps:
        # A stat() beats re-hashing the whole package: we only need to know whether the file moved
        # under us during the ~10s render, and mtime+size answers that for a fraction of the cost.
        try:
            st = os.stat(pptx)
            moved = (st.st_mtime_ns, st.st_size) != pptx_stat
        except OSError:
            moved = True
        if moved:
            skip_cache = True
            print("note: the .pptx changed during the render — not caching fingerprints",
                  file=sys.stderr)
    if skip_cache:
        try:
            os.remove(cache_path)
        except OSError:
            pass
    else:
        try:
            tmp_cache = cache_path + ".tmp"
            with open(tmp_cache, "w", encoding="utf-8") as f:
                json.dump({"fingerprints": fps}, f)
            os.replace(tmp_cache, cache_path)
        except OSError:
            try:
                os.remove(cache_path)
            except OSError:
                pass
            print("note: could not write the render cache — the next --fast will do a full render",
                  file=sys.stderr)

    # The PDF is an INTERMEDIATE of this render (pptx -> PDF -> PNG), so it always exists. Whether
    # it is promoted to a deliverable beside the .pptx is the user's call at hand-off.
    pdf_dest = pdf
    if deliverables and pdf is None:
        die("--deliverables needs a full-deck render; re-run without --fast")
    if deliverables:
        pdf_dest = os.path.join(os.path.dirname(os.path.abspath(pptx)) or ".",
                                os.path.splitext(os.path.basename(pptx))[0] + ".pdf")
        try:
            if os.path.abspath(pdf_dest) != os.path.abspath(pdf):
                shutil.move(pdf, pdf_dest)   # move, not replace: the source is a temp dir that may
        except OSError:                      # sit on a different filesystem
            pdf_dest = pdf                   # couldn't move (odd mount/permissions)

    # Self-contained flip-through viewer — parked BESIDE the .pptx (deck root), same as the PDF, so
    # the user finds it without digging into render/. It references the PNGs through the render subdir
    # (relative to the viewer's own location), so a plain double-click works. One file:// link, any
    # browser, any OS (arrow keys / click / thumbnail strip). Zero dependencies, zero network.
    deck_dir = os.path.dirname(os.path.abspath(pptx)) or "."
    viewer = None
    if deliverables:
        rel = os.path.relpath(os.path.abspath(out), deck_dir)
        pref = "" if rel in (".", "") else rel.replace(os.sep, "/").rstrip("/") + "/"
        slides = [pref + "slide{:02d}.png".format(i) for i in range(1, n_pages + 1)]
        viewer = os.path.join(deck_dir, "viewer.html")
        try:
            with open(viewer, "w", encoding="utf-8") as f:
                f.write(_viewer_html(os.path.splitext(os.path.basename(pptx))[0], slides))
        except OSError:
            viewer = None
    # sweep a stale viewer.html left inside the render dir by an older build (it now lives at root)
    stale = os.path.join(out, "viewer.html")
    if viewer and os.path.abspath(stale) != os.path.abspath(viewer) and os.path.exists(stale):
        try: os.remove(stale)
        except OSError: pass
    if incremental:
        print("{}: {} of {} slides re-rendered ({}) -> {}".format(
            "slides render" if only is not None else "fast render",
            len(keep), len(fps), ", ".join(str(i + 1) for i in keep), out))
        # If a hand-off already produced the deck-root pair, they now lag the deck. Say so loudly:
        # a stale PDF someone opens and reviews is the failure this whole path is built to avoid.

    else:
        print("rendered {} slides -> {}".format(n_pages, out))
        if fast and skip_reason:
            print("(--fast fell back to a full render: {})".format(skip_reason))
    _cs = _contact_sheet(out, n_pages)
    if _cs:
        print("contact sheet -> {}  (DECK-LEVEL view only: rhythm, bookends, form variety, "
              "chrome repetition. Body text is ~8px here and unreadable — the per-slide visual "
              "read is still required and this does not replace it.)".format(", ".join(_cs)))
    if not deliverables:
        print("pdf/viewer: not generated (deck still in progress) — at hand-off, once the user "
              "confirms the deck is final, re-run with --deliverables")
    else:
        print("pdf: {}".format(pdf_dest))
    if viewer:
        print("preview: {}  (open in a browser; arrow keys flip)".format(
            Path(viewer).resolve().as_uri()))
    # Any render that is NOT the hand-off run leaves an already-delivered pair behind the deck.
    # This must fire on the full path too: re-rendering a delivered deck is the likeliest way to
    # end up with a PDF someone opens and reviews after it stopped matching the .pptx.
    if not deliverables:
        _stale = [f for f in (os.path.splitext(os.path.basename(pptx))[0] + ".pdf", "viewer.html")
                  if os.path.isfile(os.path.join(_deck_dir, f))]
        if _stale:
            print("WARNING: {} at the deck root {} now STALE — re-run with --deliverables (and "
                  "without --fast) before handing the deck over".format(
                      " and ".join(_stale), "is" if len(_stale) == 1 else "are"), file=sys.stderr)
    print("next: python3 {} {} --renders {}  # render-time lint, then the actor-critic loop".format(
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "lint_deck.py"), pptx, out))
    # The render self-check is the single largest round-trip sink in the pipeline: one image Read
    # per slide, one message each, every message re-sending the whole conversation. SKILL.md Step 5
    # says to batch those reads, and a measured run showed the prose alone did not move it — the
    # build chained its shell commands 98% of the time and still read every PNG in its own message.
    # So the instruction is repeated HERE, in tool output the reader cannot skip, with the actual
    # paths already assembled. Cheap to print, and it names the files so there is nothing to look up.
    _pngs = sorted(f for f in os.listdir(out)
                   if f.startswith("slide") and f.endswith(".png")) if os.path.isdir(out) else []
    if len(_pngs) > 1:
        _sheets = sorted(f for f in os.listdir(out)
                         if f == "contact.png" or f.startswith("contact_"))
        _cs = ", ".join(os.path.join(out, f) for f in _sheets)
        if _sheets:
            # Deck-level FIRST, then every page. The order matters: the sheet tells you which
            # pages deserve scrutiny, and it is worthless for judging any of them — body text is
            # ~8px there. It is a way to arrive at the per-slide read informed, never a way to
            # skip it, and the line below still asks for ALL of them.
            print("      first read {}  — the whole deck at once: rhythm, bookends, form "
                  "variety, chrome repetition. It CANNOT settle a per-slide question (body text "
                  "is unreadable at that size), so it never replaces the reads below.".format(_cs))
        print("      then read ALL {} slide PNGs in ONE message (one tool block per slide, same "
              "message), and record a one-line verdict per slide:".format(len(_pngs)))
        print("      " + "  ".join(os.path.join(out, f) for f in _pngs))


def _viewer_html(title, slides):
    """Single-file dark flip-through viewer: big slide + thumbnail rail + keyboard/click nav."""
    import json
    return """<!DOCTYPE html>
<html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title} — preview</title>
<style>
  :root {{ color-scheme: dark; }}
  * {{ margin: 0; box-sizing: border-box; }}
  body {{ background: #17191f; color: #cfd4de; font: 13px/1.4 system-ui, sans-serif;
         height: 100vh; display: flex; flex-direction: column; user-select: none; }}
  header {{ padding: 8px 14px; display: flex; align-items: center; gap: 12px; }}
  header b {{ color: #fff; font-weight: 600; }}
  #stage {{ flex: 1; display: flex; align-items: center; justify-content: center;
            min-height: 0; padding: 0 12px; cursor: pointer; }}
  #main {{ max-width: 100%; max-height: 100%; box-shadow: 0 6px 30px rgba(0,0,0,.5);
           border-radius: 4px; }}
  #rail {{ display: flex; gap: 6px; overflow-x: auto; padding: 10px 14px; flex: none; }}
  #rail img {{ height: 62px; border-radius: 3px; opacity: .45; cursor: pointer;
               border: 2px solid transparent; }}
  #rail img.on {{ opacity: 1; border-color: #5b8def; }}
  #num {{ margin-left: auto; font-variant-numeric: tabular-nums; color: #8a93a6; }}
  kbd {{ background:#2a2e38; border-radius:3px; padding:1px 5px; font-size:11px; color:#9aa3b2; }}
</style></head><body>
<header><b>{title}</b><span>&larr;/&rarr; or click to flip &nbsp;<kbd>F</kbd> fullscreen</span><span id="num"></span></header>
<div id="stage"><img id="main" alt="slide"></div>
<div id="rail"></div>
<script>
const S = {slides}; let i = 0;
const main = document.getElementById('main'), rail = document.getElementById('rail'),
      num = document.getElementById('num');
S.forEach((src, k) => {{ const t = document.createElement('img'); t.src = src; t.loading = 'lazy';
  t.onclick = () => go(k); rail.appendChild(t); }});
function go(k) {{ i = (k + S.length) % S.length; main.src = S[i];
  num.textContent = (i + 1) + ' / ' + S.length;
  [...rail.children].forEach((t, k2) => t.classList.toggle('on', k2 === i));
  rail.children[i].scrollIntoView({{ inline: 'center', block: 'nearest', behavior: 'smooth' }});
  if (i + 1 < S.length) (new Image()).src = S[i + 1]; }}
document.getElementById('stage').onclick = () => go(i + 1);
addEventListener('keydown', e => {{
  if (e.key === 'ArrowRight' || e.key === ' ' || e.key === 'PageDown') go(i + 1);
  else if (e.key === 'ArrowLeft' || e.key === 'PageUp') go(i - 1);
  else if (e.key === 'Home') go(0); else if (e.key === 'End') go(S.length - 1);
  else if (e.key.toLowerCase() === 'f') document.documentElement.requestFullscreen?.(); }});
go(0);
</script></body></html>
""".format(title=title.replace("&", "&amp;").replace("<", "&lt;"), slides=json.dumps(slides))


try:                                            # console safety: a legacy code page must
    from _console import safe_stdio             # degrade a tick, never kill the report
    safe_stdio()
except Exception:
    pass


if __name__ == "__main__":
    main(sys.argv[1:])

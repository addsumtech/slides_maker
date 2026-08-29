#!/usr/bin/env python3
"""export_notes — pull the speaker notes out of a .pptx into a plain-text rehearsal script.

A presenter often rehearses from the notes alone, away from the slides. This reads each
slide's notes (deckkit.speaker_notes writes them) and emits a clean, numbered text file.

Usage:
    python3 export_notes.py deck.pptx [notes.txt]
    # default out: <deck>.notes.txt  beside the deck
"""
import sys
import os
from pptx import Presentation


def _largest_run_text(slide, _max_len=70):
    """The slide's visually dominant line — its title in everything but the XML.

    Ranks by point size, then by reading order (top, then left), so a big number in the corner
    does not beat the headline it annotates. Returns "" when the slide carries no sized text.
    """
    best = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            txt = "".join(r.text for r in para.runs).strip()
            if not txt:
                continue
            size = max((r.font.size.pt for r in para.runs if r.font.size), default=0)
            key = (size, -(sh.top or 0), -(sh.left or 0))
            if best is None or key > best[0]:
                best = (key, txt)
    if not best or best[0][0] == 0:
        return ""
    txt = " ".join(best[1].split())
    return txt if len(txt) <= _max_len else txt[:_max_len - 1] + "…"


def export_notes(pptx_path, out_path=None):
    prs = Presentation(pptx_path)
    if out_path is None:
        out_path = os.path.splitext(pptx_path)[0] + ".notes.txt"
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        # Slide title, for orientation — the whole point of the header line.
        # `shapes.title` only exists when the layout HAS a title placeholder, and deckkit builds
        # every slide on slide_layouts[6] (blank) with headings drawn by dk.text()/title_bar().
        # That is the same fact the NO SLIDE TITLE lint reports. So on the skill's own primary
        # output this read None on every slide and the rehearsal script came out as a bare
        # "--- Slide 1 ---" list: the one presenter-facing artifact, with no orientation at all,
        # exactly where it was needed most. Fall back to the visually dominant line — the largest
        # run on the slide, which is what a reader would call its title anyway.
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
        if not title:
            title = _largest_run_text(slide)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        header = f"--- Slide {i}" + (f": {title}" if title else "") + " ---"
        lines.append(header)
        lines.append(notes if notes else "(no notes)")
        lines.append("")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    return out_path, len(prs.slides._sldIdLst)


try:                                            # console safety: a legacy code page must
    from _console import safe_stdio             # degrade a tick, never kill the report
    safe_stdio()
except Exception:
    pass


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: python3 export_notes.py deck.pptx [notes.txt]")
        raise SystemExit(2)
    out, n = export_notes(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
    print(f"exported notes for {n} slides -> {out}")

#!/usr/bin/env python3
"""A declared register must be OBEYED, not just paletted — the machine-checkable half.

WHY. Measured on this toolchain: the same content run through all 18 `presets.apply()` calls
produced 18 pages that differed only in ground colour, corner radius and rule weight. memphis had
none of its coloured header bands, bauhaus none of its primitives, glassmorphism no glass card.
Every register's real look lives in its `surface` and `guard` fields — precise, executable
sentences — and `apply()` sets only palette, geometry tokens and ground, so the look is executed
by the author reading the prose or not at all.

Nothing measured the difference. `check_register_pixels.py` says so in its own docstring: it
judges COLOUR IDENTITY only. So "declare brutalist, ship its palette on rounded pastel cards"
cleared every gate — the 「只是一些颜色的搭配就说使用了这个模板」 failure, verifiable and
unverified.

This checks the SHAPE-LEVEL prohibitions, from `presets.FORBIDS`:

    rounded            prstGeom in a rounded family — swiss/brutalist/ink_wash/blueprint set
                       radius=0 precisely because their own guards forbid rounded cards
    gradient           a gradFill anywhere on the shape
    soft-shadow        shadow.inherit left True: the theme's soft shadow, never switched off
    proportional-face  a run in a face outside the mono whitelist (terminal is mono for EVERY run)
    confetti           more than ONE oversized primitive — bauhaus's guard says one hero shape,
                       "never a confetti of shapes (that is memphis)"

🔴 NARROW ON PURPOSE, in two directions. Only 7 of 18 registers declare prohibitions, and only
properties readable straight off the OOXML are listed: "photography carries ALL the colour", "no
diagrams drawn straight on the navy", "no title that isn't a full-sentence conclusion" are real
rules in `guard` and none is decidable without judging meaning. They stay prose. A register with
no `FORBIDS` entry is REPORTED as unchecked, never as clean — and a check that fired on lawful
composition would teach the reflex to waive it, which is the failure this whole area is about.

    python3 scripts/check_register_guard.py <deck.pptx> [--register <name>]
    python3 scripts/check_register_guard.py --selftest

Exit 0 clean or unchecked · 1 violations · 2 could not run.
"""
from __future__ import annotations

import argparse
import json
import os
import sys
import warnings
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
MONO_FACES = {"menlo", "consolas", "monaco", "courier", "courier new", "sf mono", "dejavu sans mono",
              "roboto mono", "ibm plex mono", "jetbrains mono", "source code pro", "fira code",
              "cascadia mono", "cascadia code", "liberation mono", "andale mono", "pt mono"}
PRIMITIVE_GEOMS = {"ellipse", "rect", "triangle", "isoscelesTriangle", "diamond", "roundRect"}
BIG_PRIMITIVE = 0.06          # share of the canvas above which a primitive is "oversized"


def _is_mono(face):
    f = str(face or "").strip().lower()
    return any(f == m or f.startswith(m) for m in MONO_FACES)


def _shape_facts(sh, canvas):
    """The five machine-readable properties, per shape."""
    x = sh._element
    # A TEXT BOX is not a primitive. Its prstGeom is `rect` like any square, so counting geometry
    # alone made a 5x1.5in title block an "oversized primitive" and reported confetti on a page
    # carrying exactly one hero circle. Found on a real page, not in the fixtures — every fixture
    # shape here was text-free, which is precisely why they missed it.
    has_text = False
    try:
        has_text = bool(sh.has_text_frame and (sh.text_frame.text or "").strip())
    except Exception:
        has_text = False
    prst = x.find(".//" + _A + "prstGeom")
    geom = prst.get("prst") if prst is not None else None
    try:
        area = ((sh.width or 0) / 914400.0) * ((sh.height or 0) / 914400.0)
    except TypeError:
        area = 0.0
    # `shadow.inherit` is True whenever a shape carries no explicit <a:effectLst> — which is the
    # DEFAULT state of every text box, and a text box has nothing to inherit FROM: python-pptx's
    # add_textbox writes no <p:style>, so there is no effectRef pointing at a theme effect. Reading
    # inherit alone therefore flagged four shapes on a page built entirely from deckkit helpers,
    # under a message that blamed "a raw add_shape()", and it fired that way under brutalist,
    # swiss, risograph AND bauhaus — 4 of the 7 registers that declare prohibitions, on a clean
    # build. A gate that fires on a deck built exactly to spec is not a floor; it is training for
    # the waiver reflex. The real condition is a live effect REFERENCE with no explicit effect of
    # its own, which is exactly what `add_shape()` leaves behind and what a theme shadow rides in on.
    eff = x.find(".//" + _P + "style/" + _A + "effectRef")
    try:
        ref_idx = int(eff.get("idx") or 0) if eff is not None else 0
    except (TypeError, ValueError):
        ref_idx = 0
    try:
        inherit = bool(sh.shadow.inherit)
    except Exception:
        inherit = False
    return {
        "geom": geom,
        "rounded": bool(geom and "round" in geom.lower()),
        "gradient": x.find(".//" + _A + "gradFill") is not None,
        "soft-shadow": bool(inherit and ref_idx > 0),
        "faces": {r.get("typeface") for r in x.iter(_A + "latin") if r.get("typeface")},
        "big_primitive": bool(geom in PRIMITIVE_GEOMS and canvas and not has_text
                              and area >= canvas * BIG_PRIMITIVE),
        "area": area,
        "rect": _rect(sh),
    }


def _has_text(sh):
    try:
        return bool(sh.has_text_frame and (sh.text_frame.text or "").strip())
    except Exception:
        return False


def _rect(sh):
    try:
        return (sh.left or 0, sh.top or 0, (sh.left or 0) + (sh.width or 0),
                (sh.top or 0) + (sh.height or 0))
    except TypeError:
        return None


BACKING_INSIDE = 0.6      # how much of the other thing sits within the primitive
BACKING_WEIGHT = 0.08     # ...and how much of the primitive it takes up, so a speck cannot exempt


def _is_backing(rect, others):
    """A primitive with something sitting ON it is a CARD or a PANEL, not a hero form.

    bauhaus's guard is one hero primitive per slide against a confetti of shapes. deckkit draws a
    card as a shape PLUS a separate text box, so a `has_text` test on the shape itself never sees
    the pairing, and two ordinary side-by-side cards read as two oversized primitives — measured on
    a page with one box and one callout, both flagged. The same blindness covers the panel behind a
    chart, an image, or an icon cluster: it carries no text either, and it is just as obviously
    furniture. So the test is CONTENT ON IT, whatever that content is — at least
    `BACKING_WEIGHT` of the primitive's own area, so a stray dot cannot launder a real hero shape.
    """
    if not rect:
        return False
    x0, y0, x1, y1 = rect
    own = max(1, (x1 - x0) * (y1 - y0))
    for other in others:
        if not other or other == rect:
            continue
        tx0, ty0, tx1, ty1 = other
        o_area = max(1, (tx1 - tx0) * (ty1 - ty0))
        if o_area >= own:                      # a bigger thing is not sitting ON this one
            continue
        ox = max(0, min(x1, tx1) - max(x0, tx0))
        oy = max(0, min(y1, ty1) - max(y0, ty0))
        inter = ox * oy
        if inter >= BACKING_INSIDE * o_area and inter >= BACKING_WEIGHT * own:
            return True
    return False


def _iter(slide):
    for sh in slide.shapes:
        if sh.shape_type == 6 and hasattr(sh, "shapes"):
            for inner in sh.shapes:
                yield inner
        else:
            yield sh


def check(pptx, register=None, gates=None):
    """Return (violations, facts)."""
    import presets
    from pptx import Presentation

    facts = {}
    if register is None:
        d = (gates or {}).get("design_plan") or {}
        # `check_style_applied.declared_preset` already owns this parse, INCLUDING the three-way
        # confidence it needed for the generated-template branch. Writing a second parser here cost
        # a real false positive immediately: a bespoke deck whose pick read "…beat blueprint-the-
        # preset because…" was checked as `blueprint`, because a substring search cannot tell the
        # register a deck DECLARES from the rival it says it BEAT. One parser, one answer.
        try:
            import check_style_applied as csa
            register, conf = csa.declared_preset(d.get("style_pick"), csa.preset_names(),
                                                 d.get(csa.LOOK_SOURCE_KEY))
            if conf != "sure":
                facts["confidence"] = conf
                register = None
        except Exception:
            register = None
    facts["register"] = register
    if not register:
        if facts.get("confidence") == "unsure":
            facts["note"] = ("`style_pick` names a preset but also carries a non-preset qualifier, "
                             "so the look may not be preset-based — NOT CHECKED rather than checked "
                             "against a register this deck may not have declared")
            return [], facts
        facts["note"] = ("no preset register declared (design_plan.style_pick names none) — a "
                         "bespoke or generated look has no FORBIDS to check, which is not the same "
                         "as clean")
        return [], facts
    rules = presets.FORBIDS.get(register)
    if not rules:
        facts["note"] = ("{!r} declares no machine-checkable prohibitions — its `guard` is real but "
                         "needs judgement of meaning, so it stays prose and is NOT CHECKED here"
                         .format(register))
        return [], facts
    facts["rules"] = list(rules)

    prs = Presentation(str(pptx))
    canvas = (prs.slide_width / 914400.0) * (prs.slide_height / 914400.0)
    hits = {}
    for n, slide in enumerate(prs.slides, 1):
        bigs = 0
        shapes = list(_iter(slide))
        facts_by_shape = [(sh, _shape_facts(sh, canvas)) for sh in shapes]
        on_page = [f["rect"] for _sh, f in facts_by_shape if f["rect"]]
        for sh, f in facts_by_shape:
            if f["big_primitive"] and not _is_backing(f["rect"], on_page):
                bigs += 1
            for rule in rules:
                if rule == "confetti" or rule == "proportional-face":
                    continue
                if f.get(rule):
                    hits.setdefault(rule, []).append(n)
            if "proportional-face" in rules:
                for face in f["faces"]:
                    if not _is_mono(face):
                        hits.setdefault("proportional-face", []).append((n, face))
        if "confetti" in rules and bigs > 1:
            hits.setdefault("confetti", []).append((n, bigs))

    WHY = {
        "rounded": ("{r}'s guard forbids rounded cards, and its preset sets radius=0 for exactly "
                    "that reason — a rounded corner here is the register being overridden after "
                    "apply() set it"),
        "gradient": "{r}'s guard forbids gradients — the fills are flat in this register",
        "soft-shadow": ("{r}'s guard forbids soft shadows. These shapes left shadow.inherit True, "
                        "so they carry the theme's default shadow: deckkit's own helpers switch it "
                        "off, so this is usually a raw add_shape() that skipped it"),
        "proportional-face": ("terminal is mono for EVERY run including headings — a proportional "
                              "face breaks the register"),
        "confetti": ("bauhaus's guard is one hero primitive per slide, \"never a confetti of shapes "
                     "(that is memphis)\""),
    }
    out = []
    for rule, where in sorted(hits.items()):
        if rule == "proportional-face":
            faces = sorted({f for _n, f in where})
            pages = sorted({n for n, _f in where})
            out.append((rule, "{} — found {} on slide(s) {}".format(
                WHY[rule].format(r=register), ", ".join(repr(f) for f in faces[:4]),
                ", ".join(str(p) for p in pages[:8]))))
        elif rule == "confetti":
            out.append((rule, "{} — slide(s) {} carry more than one oversized primitive ({})".format(
                WHY[rule].format(r=register),
                ", ".join(str(n) for n, _c in where[:6]),
                ", ".join("{}:{}".format(n, c) for n, c in where[:6]))))
        else:
            pages = sorted(set(where))
            out.append((rule, "{} — {} shape(s) on slide(s) {}".format(
                WHY[rule].format(r=register), len(where),
                ", ".join(str(p) for p in pages[:8]))))
    return out, facts


# --------------------------------------------------------------------------- selftest

def _selftest():
    import tempfile
    import deckkit as dk
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    ok, bad = [], []
    tmp = Path(tempfile.mkdtemp(prefix="regguard-"))
    _n = [0]

    def build(fn):
        prs = dk.blank_deck()
        s = dk.add_slide(prs)
        fn(s, prs)
        _n[0] += 1
        p = tmp / ("d%d.pptx" % _n[0])
        prs.save(str(p))
        return p

    def codes(path, reg):
        return {c for c, _m in check(path, register=reg)[0]}

    # THE FALSE-POSITIVE FLOOR, first because it is the failure that makes a gate worthless. An
    # ordinary page built from nothing but deckkit helpers must be clean under EVERY register that
    # declares prohibitions. Measured before this test existed: 4 of the 7 fired on exactly such a
    # page — brutalist, swiss, risograph and bauhaus each reported a soft shadow that was four TEXT
    # BOXES (nothing to inherit from), and bauhaus additionally called two ordinary side-by-side
    # cards a confetti of oversized primitives. A gate that fires on a deck built exactly to spec
    # is not a floor; it is training for the waiver reflex.
    import presets as _pre

    def _ordinary(s, p):
        y = dk.title_bar(s, "A page built the ordinary way", kicker="probe")
        dk.box(s, 0.8, y + 0.2, 4.0, 1.6)
        dk.text(s, 0.9, y + 0.4, 3.8, 0.5, [[("body copy", 14, dk.DEEP, False, False)]])
        dk.callout(s, 5.2, y + 0.2, 4.0, 1.6, "NOTE", "a callout")

    # The register must be APPLIED first — that is what a real deck does, and radius=0 / the mono
    # face are half of what makes the page lawful. apply() mutates deckkit's module globals, so the
    # snapshot below keeps this test from changing the answer of every test after it.
    _snap = {k: getattr(dk, k) for k in dir(dk) if k.isupper()}
    dirty = {}
    try:
        for r in sorted(_pre.FORBIDS):
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                _pre.apply(r)
            got = sorted(codes(build(_ordinary), r))
            if got:
                dirty[r] = got
    finally:
        for k, v in _snap.items():
            setattr(dk, k, v)
    ok.append("a page built from nothing but deckkit helpers is clean under all {} registers that "
              "declare prohibitions — the false-positive floor".format(len(_pre.FORBIDS))) \
        if not dirty else bad.append("clean deckkit page flagged: {}".format(dirty))

    def _card_with_type(s, p):
        c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(4))
        c.fill.solid(); c.fill.fore_color.rgb = dk.TINT
        c.line.fill.background(); c.shadow.inherit = False
        d = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1), Inches(5), Inches(4))
        d.fill.solid(); d.fill.fore_color.rgb = dk.TINT
        d.line.fill.background(); d.shadow.inherit = False
        dk.text(s, 0.8, 1.3, 4.4, 1.0, [[("type set on the form", 14, dk.DEEP, False, False)]])
        dk.text(s, 6.3, 1.3, 4.4, 1.0, [[("and on the other", 14, dk.DEEP, False, False)]])

    got = codes(build(_card_with_type), "bauhaus")
    ok.append("two big forms with TYPE SET ON THEM are cards, not a confetti of hero primitives — "
              "deckkit draws a card as a shape plus a separate text box, so counting geometry "
              "alone made every two-card page a bauhaus violation") \
        if "confetti" not in got else bad.append(str(got))

    # rounded — the case radius=0 exists to prevent
    got = codes(build(lambda s, p: dk.box(s, 1, 1, 3, 2, fill=dk.DEEP, round=True)), "swiss")
    ok.append("a rounded card under `swiss` is caught — its own guard forbids them and its preset "
              "sets radius=0 for that reason") if "rounded" in got else bad.append(str(got))
    got = codes(build(lambda s, p: dk.box(s, 1, 1, 3, 2, fill=dk.DEEP)), "swiss")
    ok.append("...and a square one passes") if "rounded" not in got else bad.append(str(got))

    got = codes(build(lambda s, p: dk.box(s, 1, 1, 3, 2,
                                          grad=[(0, dk.DEEP, 1.0), (1, dk.TINT, 1.0)])), "risograph")
    ok.append("a gradient under `risograph` is caught — its shadows are HARD and its fills flat") \
        if "gradient" in got else bad.append(str(got))

    def _raw_shadow(s, p):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2))
        c.fill.solid()
        c.fill.fore_color.rgb = dk.DEEP
        c.line.fill.background()          # deliberately NOT switching shadow.inherit off

    got = codes(build(_raw_shadow), "bauhaus")
    ok.append("a raw add_shape() that never switched shadow.inherit off is caught under `bauhaus` "
              "— this is the mistake the author of this checker made while building the very "
              "comparison that motivated it") if "soft-shadow" in got else bad.append(str(got))

    def _one_hero(s, p):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(0), Inches(4.6), Inches(4.6))
        c.fill.solid()
        c.fill.fore_color.rgb = dk.DEEP
        c.line.fill.background()
        c.shadow.inherit = False

    def _confetti(s, p):
        for i in range(3):
            c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4 + i * 3.0), Inches(1),
                                   Inches(2.4), Inches(2.4))
            c.fill.solid()
            c.fill.fore_color.rgb = dk.DEEP
            c.line.fill.background()
            c.shadow.inherit = False

    ok.append("ONE oversized primitive under `bauhaus` passes — that IS the register") \
        if "confetti" not in codes(build(_one_hero), "bauhaus") else bad.append("hero flagged")
    ok.append("...and three of them are caught: its guard says one hero, never a confetti") \
        if "confetti" in codes(build(_confetti), "bauhaus") else bad.append("confetti missed")

    def _prop(s, p):
        dk.text(s, 1, 1, 4, 0.4, [[("heading", 20, dk.DEEP, True, False, "Helvetica Neue")]])

    def _mono(s, p):
        dk.text(s, 1, 1, 4, 0.4, [[("heading", 20, dk.DEEP, True, False, "Menlo")]])

    ok.append("a proportional face under `terminal` is caught — mono for EVERY run, headings too") \
        if "proportional-face" in codes(build(_prop), "terminal") else bad.append("prop missed")
    ok.append("...and an all-mono page passes") \
        if "proportional-face" not in codes(build(_mono), "terminal") else bad.append("mono flagged")

    # a register with no machine-checkable guard must say so, not imply clean
    _p = build(lambda s, p: dk.box(s, 1, 1, 3, 2, fill=dk.DEEP, round=True))
    v, f = check(_p, register="consulting")
    ok.append("a register whose guard needs judgement is REPORTED as unchecked, not passed — "
              "\"titles must be full-sentence conclusions\" is real and undecidable from the file") \
        if not v and f.get("note") else bad.append("consulting: {} {}".format(v, f))
    v, f = check(_p, register=None)
    ok.append("a bespoke look with no declared preset is reported as unchecked too") \
        if not v and f.get("note") else bad.append("bespoke: {} {}".format(v, f))

    print("\n".join("  ok   " + x for x in ok))
    if bad:
        print("\n".join("  FAIL " + x for x in bad))
    print("\n{} passed, {} failed".format(len(ok), len(bad)))
    return 1 if bad else 0


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0],
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", nargs="?")
    ap.add_argument("--register", help="force a register instead of reading design_plan.style_pick")
    ap.add_argument("--selftest", action="store_true")
    a = ap.parse_args(argv)
    if a.selftest:
        return _selftest()
    if not a.pptx:
        ap.print_help()
        return 2
    if not os.path.exists(a.pptx):
        print("no such file: {}".format(a.pptx), file=sys.stderr)
        return 2
    gates = {}
    gp = Path(a.pptx).resolve().parent / ".deck-gates.json"
    if gp.exists():
        try:
            gates = json.loads(gp.read_text(encoding="utf-8"))
        except (ValueError, OSError):
            gates = {}
    try:
        violations, facts = check(a.pptx, a.register, gates)
    except Exception as exc:
        print("cannot run: {}: {}".format(exc.__class__.__name__, exc), file=sys.stderr)
        return 2
    print("register: {}".format(facts.get("register") or "none declared"))
    if facts.get("rules"):
        print("  checked prohibitions: " + " · ".join(facts["rules"]))
    if facts.get("note"):
        print("  [--] " + facts["note"])
    if not violations:
        print("the declared register is obeyed, on every property a machine can settle.")
        return 0
    print("\n{} violation(s):\n".format(len(violations)))
    for code, msg in violations:
        print("  {}: {}\n".format(code.upper(), msg))
    return 1


try:
    from _console import safe_stdio
    safe_stdio()
except Exception:
    pass


if __name__ == "__main__":
    sys.exit(main())

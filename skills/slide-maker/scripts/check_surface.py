#!/usr/bin/env python3
"""Every canvas format's contract, checked against the BUILT deck instead of trusted.

`formats.py` is a registry of design surfaces: margins, platform-UI safe zones, whether columns
work, whether the surface carries deck chrome, how dense it may be. Measured by grep, NOTHING
consumed it downstream — `import formats` appears in exactly two files, both of which are
producers. A build script opts into the registry voluntarily, and no check afterwards can tell
whether it did. So every per-surface rule in `references/canvas-formats.md` was advisory by
construction, on the surfaces where getting it wrong is least recoverable: a story caption placed
under the swipe bar is invisible to the person who built it on a desktop, and a poster is wrong
only once it has been printed at a metre wide.

What this checks, per format, recovering the format from the canvas SIZE (a built PPTX carries
nothing else):

  TYPE FLOOR       a PRINTED surface is read at a fixed distance, so its point sizes are absolute,
                   not a share of the canvas. lint's canvas-relative floor is right for a projected
                   deck and meaningless for a poster — on a 33in-wide A0 it would demand ~45pt body,
                   while deckkit's own cover cap of 46pt prints a title that cannot be read across a
                   hall. This applies the three-distance floors the format declares.
  MISSING SECTION  a surface may declare content it is not finished without. A poster declares
                   methods and limitations: the billboard style that reads best is also the style
                   that tends to drop exactly the two things a passer-by cannot reconstruct.
  SAFE ZONE        text inside the platform UI zones a vertical/social format reserves.
  COLUMNS          side-by-side splits on a canvas whose registry entry says `columns_ok=False`.
  DECK CHROME      a footer strip on a `social` surface, which carries no deck furniture.

    python3 scripts/check_surface.py <deck.pptx> [--format NAME] [--waive-sections "<reason>"]
    python3 scripts/check_surface.py --selftest

Deliberately NOT here: margins and full-bleed. `lint_deck.py` already owns off-canvas and bleed,
including the `design_intent(envelope="bleed")` / `bleed_intent()` declarations a hero legitimately
makes — a second opinion about the same pixels would disagree with the first on every cover.

🔴 It checks the format CONTRACT, not the design. A poster can pass every line here and still be
an unreadable wall — density, hierarchy and figure quality belong to the lint, the sameness gate
and the critic. Exit 0 clean · 1 findings · 2 could not run.
"""
from __future__ import annotations

import argparse
import os
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

EMU = 914400.0


def _fmt_for(prs, name=None):
    import formats
    if name:
        return formats.get(name)
    return formats.match(prs.slide_width / EMU, prs.slide_height / EMU)


def _is_ground(sh, W, H):
    """A full-bleed background is the canvas, not content — counting it would score every deck 100%."""
    try:
        return ((sh.width or 0) / EMU >= W * 0.97) and ((sh.height or 0) / EMU >= H * 0.97)
    except TypeError:
        return False


def _runs(shape):
    """[(text, pt, top_in, height_in)] for one shape's paragraphs, sizes resolved where set."""
    out = []
    if not getattr(shape, "has_text_frame", False):
        return out
    for para in shape.text_frame.paragraphs:
        txt = "".join(r.text or "" for r in para.runs) or (para.text or "")
        if not txt.strip():
            continue
        sizes = [r.font.size.pt for r in para.runs if r.font.size is not None]
        if not sizes and para.font.size is not None:
            sizes = [para.font.size.pt]
        out.append((txt.strip(), max(sizes) if sizes else None))
    return out


def _shapes(slide):
    for sh in slide.shapes:
        if sh.shape_type == 6 and hasattr(sh, "shapes"):        # GROUP
            for inner in sh.shapes:
                yield inner
        else:
            yield sh


# Deliberately NOT a size->role classifier. The first version guessed the role from the size and
# then judged the size against that role's floor, which is circular: a 27pt run on an A0 board was
# classified "section" because it cleared 75% of the 36pt section floor, and then failed the 36pt
# floor it had just been assigned to. Measured on a real poster built for this check. The three
# distances are asserted directly instead — every run readable at 1m, SOME run readable at 5m, and
# a middle tier existing at all — which needs no classification.


def check(pptx, format_name=None, waive_sections=None, extra_terms=None):
    """Return (problems, facts).

    `extra_terms` — {section_label: [word, …]} — extends the keyword list a required section is
    recognised by. The built-in lists are English and Chinese, which is where this skill's decks
    have been built; a poster written in Dutch, German or Japanese would otherwise be told it has
    no methods section while its Methode/Methoden/手法 heading sits at the top of it. Terms travel
    with the deck as `design_plan.surface_section_terms`, so the fix is a data line rather than an
    edit to this file — and rather than a waiver, which would switch the check OFF instead of
    teaching it the language.
    """
    from pptx import Presentation
    prs = Presentation(str(pptx))
    facts = {"canvas": "%.2fx%.2f in" % (prs.slide_width / EMU, prs.slide_height / EMU)}
    fmt = _fmt_for(prs, format_name)
    if fmt is None:
        facts["note"] = ("this canvas ({}) matches no registered format, so no surface contract "
                         "could be applied. NOT the same as clean — pass --format <name> if it is "
                         "a registered surface built at a nonstandard size."
                         .format(facts["canvas"]))
        return [], facts
    facts["format"] = "{} ({})".format(fmt.name, fmt.label)
    problems = []

    import formats
    fl = formats.floors(fmt)
    W, H = prs.slide_width / EMU, prs.slide_height / EMU

    sized_runs, text_all, rects = [], [], []
    for idx, slide in enumerate(prs.slides, 1):
        for sh in _shapes(slide):
            try:
                top = (sh.top or 0) / EMU
                left = (sh.left or 0) / EMU
                wid = (sh.width or 0) / EMU
                hei = (sh.height or 0) / EMU
            except TypeError:
                top = left = wid = hei = 0.0
            if wid > 0.02 and hei > 0.02 and not _is_ground(sh, W, H):
                rects.append((left, top, wid, hei))
            for txt, pt in _runs(sh):
                text_all.append(txt)
                if pt is not None:
                    sized_runs.append((idx, txt, pt))
                # ── SAFE ZONE ──────────────────────────────────────────────────────────────
                if fmt.safe_top and top < fmt.safe_top - 0.02:
                    problems.append(("SAFE ZONE",
                                     "slide {}: text sits {:.2f}in from the top, inside the {:.2f}in "
                                     "the {} surface reserves for platform UI. On a desktop preview "
                                     "it looks fine; in the app a profile bar covers it: {!r}"
                                     .format(idx, top, fmt.safe_top, fmt.name, txt[:50])))
                if fmt.safe_bottom and top + hei > H - fmt.safe_bottom + 0.02:
                    problems.append(("SAFE ZONE",
                                     "slide {}: text reaches {:.2f}in from the bottom, inside the "
                                     "{:.2f}in the {} surface reserves for the swipe/caption bar: "
                                     "{!r}".format(idx, H - (top + hei), fmt.safe_bottom,
                                                   fmt.name, txt[:50])))
            # ── DECK CHROME on a social surface ────────────────────────────────────────────
            if (fmt.chrome == "social" and getattr(sh, "has_text_frame", False)
                    and hei and hei <= 0.5 and wid >= W * 0.55 and top >= H * 0.82
                    and (sh.text_frame.text or "").strip()):
                problems.append(("DECK CHROME",
                                 "slide {}: a full-width strip of text sits in the bottom band of a "
                                 "SOCIAL surface ({!r}). Feed formats carry no deck footer — the "
                                 "platform's own UI lives there, and a footer reads as a stray "
                                 "caption.".format(idx, (sh.text_frame.text or "").strip()[:40])))

        # ── COLUMNS on a canvas the registry says cannot take them ─────────────────────────
        if not fmt.columns_ok:
            texts = [sh for sh in _shapes(slide)
                     if getattr(sh, "has_text_frame", False)
                     and (sh.text_frame.text or "").strip()
                     and (sh.width or 0) / EMU < W * 0.46]
            for i, a in enumerate(texts):
                for b in texts[i + 1:]:
                    at, ah = (a.top or 0) / EMU, (a.height or 0) / EMU
                    bt, bh = (b.top or 0) / EMU, (b.height or 0) / EMU
                    overlap = min(at + ah, bt + bh) - max(at, bt)
                    if overlap > 0.45 * min(ah, bh) and min(ah, bh) > 0.2:
                        problems.append(("COLUMNS",
                                         "slide {}: two text blocks share a horizontal band, each "
                                         "under half the canvas width. `{}` is registered "
                                         "columns_ok=False — a hand-held portrait surface reads as "
                                         "one column; side-by-side halves the measure and the type "
                                         "with it. Stack them.".format(idx, fmt.name)))
                        break
                else:
                    continue
                break

    # ── TYPE FLOOR — the printed-at-actual-size case ───────────────────────────────────────
    if fl:
        facts["floors"] = ", ".join("{} >={}pt".format(r, fl[r])
                                    for r in ("display", "section", "body") if r in fl)
        sized = sized_runs
        if not sized:
            problems.append(("TYPE FLOOR",
                             "no run on this deck carries an explicit point size, so the printed "
                             "floors could not be applied to a single one. On a projected deck "
                             "inherited sizes are fine; on a board printed at {:.0f}in wide they "
                             "are the difference between readable and not.".format(W)))
        else:
            body_floor = fl.get("body")
            if body_floor:
                under = [r for r in sized if r[2] < body_floor]
                if under:
                    worst = min(under, key=lambda t: t[2])
                    problems.append(("TYPE FLOOR",
                                     "{} run(s) fall under the {}pt BODY floor for a printed {} — "
                                     "smallest {:.0f}pt on slide {} ({!r}). This floor is ABSOLUTE: "
                                     "the board is read at a fixed distance, so a share of the "
                                     "canvas says nothing about whether it can be read. Cut words "
                                     "rather than shrinking them."
                                     .format(len(under), body_floor, fmt.label, worst[2],
                                             worst[0], worst[1][:40])))
            disp = fl.get("display")
            if disp and not any(r[2] >= disp for r in sized):
                biggest = max(sized, key=lambda t: t[2])
                problems.append(("TYPE FLOOR",
                                 "the largest type on this {} is {:.0f}pt, under the {}pt DISPLAY "
                                 "floor — nothing on the board reads from across the hall. "
                                 "deckkit's cover caps titles at 46pt, which is right on a 10in "
                                 "slide and unreadable at 5m on a {:.0f}in board, so a poster title "
                                 "has to be sized explicitly."
                                 .format(fmt.label, biggest[2], disp, W)))
            sec = fl.get("section")
            if sec and not any(sec <= r[2] < (disp or 1e9) for r in sized):
                problems.append(("TYPE FLOOR",
                                 "no run sits between the {}pt SECTION floor and the display size: "
                                 "this board has a headline and body and nothing in between, so a "
                                 "reader standing 2m away has no way to decide whether to come "
                                 "closer. The middle tier is what makes a poster scannable."
                                 .format(sec)))

    # ── FILL — a printed board is composed once, and the whole board is the deliverable ────
    if fmt.fill_range and rects:
        lo, hi = fmt.fill_range
        # Coarse coverage grid rather than a rectangle union: overlapping cards must not be
        # double-counted (that is how an "88% full" board turns out to be one stack of panels).
        NX, NY = 160, 220
        grid = bytearray(NX * NY)
        for (x0, y0, wd, hei) in rects:
            cx0 = max(0, min(NX - 1, int(x0 / W * NX)))
            cx1 = max(0, min(NX, int((x0 + wd) / W * NX) + 1))
            cy0 = max(0, min(NY - 1, int(y0 / H * NY)))
            cy1 = max(0, min(NY, int((y0 + hei) / H * NY) + 1))
            for cy in range(cy0, cy1):
                base = cy * NX
                for cx in range(cx0, cx1):
                    grid[base + cx] = 1
        cover = sum(grid) / float(NX * NY)
        facts["fill"] = "{:.0%} of the board covered (target {:.0%}-{:.0%})".format(cover, lo, hi)
        if cover < lo:
            problems.append(("FILL",
                             "content covers {:.0%} of the board, under the {:.0%} floor for a {}. "
                             "A poster is composed ONCE and read standing: unlike a deck, where "
                             "whitespace across a sequence is rhythm, empty board here is the only "
                             "space this work gets and it is being given away. Either the content "
                             "is thinner than the format, or the layout never grew into it."
                             .format(cover, lo, fmt.label)))
        elif cover > hi:
            problems.append(("FILL",
                             "content covers {:.0%} of the board, over the {:.0%} ceiling for a {}. "
                             "This is the wall people walk past. Cut, do not shrink — the type "
                             "floors above are absolute, so density has to come out of the words."
                             .format(cover, hi, fmt.label)))

    # ── MISSING SECTION — content a surface is not finished without ────────────────────────
    if fmt.required_sections:
        if waive_sections:
            facts["sections_waived"] = waive_sections
        else:
            blob = " ".join(text_all).lower()
            more = {str(k).lower(): list(v) for k, v in (extra_terms or {}).items()}
            for label, keys in fmt.required_sections:
                keys = tuple(keys) + tuple(more.get(label.lower(), ()))
                if not any(str(k).lower() in blob for k in keys):
                    problems.append(("MISSING SECTION",
                                     "nothing on this {} names {} (looked for: {}). A poster is read "
                                     "without its author present: the billboard style that draws "
                                     "people in is also the style that drops the two things a "
                                     "passer-by cannot reconstruct and cannot fairly judge the claim "
                                     "without. If it IS there under a word this list does not "
                                     "know — another language, a field's own term — add it to "
                                     "design_plan.surface_section_terms rather than waiving the "
                                     "check; if this poster genuinely has none, say so with "
                                     "--waive-sections '<why>'."
                                     .format(fmt.label, label, ", ".join(keys[:4]))))
    facts["slides"] = len(prs.slides._sldIdLst)
    return problems, facts


# --------------------------------------------------------------------------- selftest

def _selftest():
    import tempfile
    import formats
    ok, bad = [], []
    tmp = Path(tempfile.mkdtemp(prefix="surface-"))

    def build(fmt_name, runs, extra=None):
        """runs = [(text, pt, x, y, w, h)] on ONE slide of the named format."""
        import deckkit as dk
        f = formats.get(fmt_name)
        prs = dk.blank_deck(f.w_in, f.h_in)
        s = dk.add_slide(prs)
        for txt, pt, x, y, w, h in runs:
            dk.text(s, x, y, w, h, [[(txt, pt, dk.DEEP, False, False)]])
        if extra:
            extra(dk, s, f)
        path = tmp / (fmt_name + "-" + str(len(list(tmp.glob("*.pptx")))) + ".pptx")
        prs.save(str(path))
        return path

    def codes(path, **kw):
        return {c for c, _ in check(path, **kw)[0]}

    # A0 poster that meets every floor, names both required sections, and fills the board.
    good = build("poster_a0", [
        ("A result worth crossing a hall for", 110, 1.6, 1.6, 29.9, 4.4),
        ("Methods", 42, 1.6, 7.5, 14.0, 1.4),
        ("We measured every deck the skill built and compared them.", 26, 1.6, 9.2, 14.0, 12.0),
        ("Results", 42, 17.5, 7.5, 14.0, 1.4),
        ("Two prose rules now fail loudly on real decks.", 26, 17.5, 9.2, 14.0, 12.0),
        ("Limitations", 42, 1.6, 23.0, 14.0, 1.4),
        ("One site, one operator, no held-out cohort.", 26, 1.6, 24.7, 14.0, 19.0),
        ("Next", 42, 17.5, 23.0, 14.0, 1.4),
        ("Record accent hexes so freshness can see more than the ground.", 26, 17.5, 24.7, 14.0, 19.0)])
    probs, facts = check(good)
    got = {c for c, _ in probs}
    (ok if not got else bad).append(
        "a well-set A0 poster passes — {}".format(facts.get("fill")) if not got
        else "good poster flagged: {} ({})".format(got, facts.get("fill")))

    # The defect the eye sees first on a real board and no gate had a name for.
    sparse = build("poster_a0", [
        ("A result worth crossing a hall for", 110, 1.6, 1.6, 29.9, 4.4),
        ("Methods", 42, 1.6, 8.0, 14.0, 1.4),
        ("Short.", 26, 1.6, 9.7, 14.0, 2.0),
        ("Limitations", 42, 17.5, 8.0, 14.0, 1.4),
        ("Also short.", 26, 17.5, 9.7, 14.0, 2.0)])
    probs, facts = check(sparse)
    got = {c for c, _ in probs}
    (ok if "FILL" in got else bad).append(
        "a half-empty board is caught ({}) — measured on a real A0 render whose lower half was "
        "blank and which passed every existing gate".format(facts.get("fill"))
        if "FILL" in got else "sparse poster not caught: {} ({})".format(got, facts.get("fill")))

    # The exact trap the format exists for: a poster typeset like a slide.
    slideish = build("poster_a0", [
        ("A result worth crossing a hall for", 46, 2.0, 2.0, 29.0, 2.0),
        ("Methods", 20, 2.0, 8.0, 14.0, 1.0),
        ("Limitations: one site only.", 14, 2.0, 10.0, 14.0, 1.0)])
    got = codes(slideish)
    (ok if "TYPE FLOOR" in got else bad).append(
        "a poster typeset at SLIDE sizes is caught — deckkit's own 46pt cover cap prints a title "
        "that cannot be read across a hall, and nothing said so before"
        if "TYPE FLOOR" in got else "slide-sized poster not caught: {}".format(got))

    # Billboard poster with no methods and no limitations.
    billboard = build("poster_a0", [
        ("Ours wins", 140, 2.0, 2.0, 29.0, 6.0),
        ("A single number, very large, and a picture of it.", 30, 2.0, 12.0, 29.0, 3.0)])
    got = [m for c, m in check(billboard)[0] if c == "MISSING SECTION"]
    (ok if len(got) == 2 else bad).append(
        "a billboard poster missing BOTH methods and limitations is caught, once per section"
        if len(got) == 2 else "missing sections: {} finding(s)".format(len(got)))
    got = codes(billboard, waive_sections="a purely descriptive display board")
    (ok if "MISSING SECTION" not in got else bad).append(
        "...and a written waiver stands that half down"
        if "MISSING SECTION" not in got else "waiver ignored: {}".format(got))

    # Story: text under the swipe bar, and a side-by-side split.
    story = build("story", [
        ("Swipe up for the rest", 18, 0.5, 9.3, 4.6, 0.5),
        ("left column", 16, 0.45, 4.0, 2.2, 1.2),
        ("right column", 16, 2.9, 4.0, 2.2, 1.2)])
    got = codes(story)
    (ok if "SAFE ZONE" in got and "COLUMNS" in got else bad).append(
        "a 9:16 story with text under the swipe bar AND a side-by-side split is caught — both "
        "rules were prose in canvas-formats.md and nothing downstream read the registry"
        if "SAFE ZONE" in got and "COLUMNS" in got else "story faults missed: {}".format(got))

    # A normal 16:9 deck must not be judged against printed floors or social rules.
    slide = build("wide", [("A normal projected slide", 24, 0.6, 0.6, 8.8, 1.0),
                           ("Body copy at a normal projected size.", 14, 0.6, 2.0, 4.0, 1.0),
                           ("A second column, which 16:9 allows.", 14, 5.2, 2.0, 4.0, 1.0)])
    got = codes(slide)
    (ok if not got else bad).append(
        "a normal 16:9 deck is untouched — a printed floor applied to a projected canvas would be "
        "noise, and the registry says which is which" if not got
        else "16:9 deck wrongly flagged: {}".format(got))

    # An unregistered canvas must report NOT CHECKED rather than clean.
    import deckkit as dk
    prs = dk.blank_deck(12.34, 3.21)
    dk.add_slide(prs)
    odd = tmp / "odd.pptx"
    prs.save(str(odd))
    probs, facts = check(odd)
    (ok if not probs and facts.get("note") else bad).append(
        "an unregistered canvas is reported as unchecked, not passed"
        if not probs and facts.get("note") else "odd canvas: {} {}".format(probs, facts))

    print("\n".join("  ok   " + x for x in ok))
    if bad:
        print("\n".join("  FAIL " + x for x in bad))
    print("\n{} passed, {} failed".format(len(ok), len(bad)))
    return 1 if bad else 0


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0],
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", nargs="?")
    ap.add_argument("--format", help="force a registered format instead of matching the canvas size")
    ap.add_argument("--waive-sections", help="written reason this surface needs no required sections")
    ap.add_argument("--section-terms", metavar="JSON",
                    help='extra words a required section is recognised by, e.g. '
                         '\'{"methods": ["methode", "手法"]}\' — for a poster not written in '
                         'English or Chinese')
    ap.add_argument("--selftest", action="store_true")
    a = ap.parse_args(argv)
    if a.selftest:
        return _selftest()
    if not a.pptx:
        ap.print_help()
        return 2
    if not os.path.exists(a.pptx):
        print("no such file: {}".format(a.pptx), file=sys.stderr)
        return 2
    extra = None
    if a.section_terms:
        import json
        try:
            extra = json.loads(a.section_terms)
        except ValueError as exc:
            print("--section-terms is not valid JSON: {}".format(exc), file=sys.stderr)
            return 2
    probs, facts = check(a.pptx, a.format, a.waive_sections, extra)
    print("canvas {} -> {}".format(facts.get("canvas"), facts.get("format", "unregistered")))
    if facts.get("floors"):
        print("  printed type floors: " + facts["floors"])
    if facts.get("fill"):
        print("  " + facts["fill"])
    if facts.get("sections_waived"):
        print("  required sections waived: " + facts["sections_waived"])
    if facts.get("note"):
        print("  [--] " + facts["note"])
    if not probs:
        print("the surface contract holds.")
        return 0
    print("\n{} finding(s):\n".format(len(probs)))
    for code, msg in probs:
        print("  {}: {}\n".format(code, msg))
    return 1


try:
    from _console import safe_stdio
    safe_stdio()
except Exception:
    pass


if __name__ == "__main__":
    sys.exit(main())

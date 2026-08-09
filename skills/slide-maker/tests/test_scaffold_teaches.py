#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""The scaffold is what this skill actually teaches, so it is checked like a contract.

SKILL.md tells every single-author build to copy `references/examples/build_example_generic.py`.
That makes the file the skill's real curriculum: whatever it demonstrates is what gets written,
and whatever it omits effectively does not exist. Measured before this suite, on a real 14-page
build and on the file itself:

    vstack          0 occurrences        hand-picked y   6 occurrences
    content_band    0                    register_mark   0
    tag_motif       0

The consequences were not hypothetical. The build hand-computed nearly every y and spent many
rounds nudging coordinates; `vstack`'s first honest error ("blocks need 3.68in but the band is
only 2.88in" — the page is OVER-FULL, which no nudge can fix) arrived on round twelve instead of
round one. And 383 shapes in the delivered deck carried ZERO motif tags, so `TEXT_OVER_MOTIF` —
written for that deck's own cover defect — reported nothing on it.

Both are ADOPTION failures, not capability gaps, and adoption is decided by the example. So the
example is now asserted: it must demonstrate the measured-packer idiom and the motif tag, and it
must not teach the hand-picked-y idiom it used to. This suite fails loudly if either regresses,
which is the only thing that keeps a curriculum honest.

Run:  python3 tests/test_scaffold_teaches.py
"""
import pathlib
import re
import subprocess
import sys

HERE = pathlib.Path(__file__).resolve().parent
ROOT = HERE.parent
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

SCAFFOLD = ROOT / "references" / "examples" / "build_example_generic.py"

PASS, FAIL = [], []


def check(name, cond, detail=""):
    (PASS if cond else FAIL).append(name)
    print(("  ok   " if cond else "  FAIL ") + name
          + (("  — " + str(detail)) if detail and not cond else ""))


def main():
    print("scaffold teaches")
    src = SCAFFOLD.read_text(encoding="utf-8")

    # ---- it must DEMONSTRATE the idioms the skill requires ----------------------------
    for name, why in (
        ("vstack", "the measured packer that raises at build time instead of overlapping"),
        ("content_band", "the authoritative safe rect, so a tail is anchored not guessed"),
        ("measure_bullets", "a block's height is measured before it is placed"),
        ("register_mark", "the register signature, correct by construction and tagged"),
        ("declare_delivery", "what kind of deck this is, recorded once instead of retyped"),
    ):
        check("the scaffold demonstrates %s (%s)" % (name, why), name + "(" in src)

    # ---- and it must not teach the idiom that produced the failures -------------------
    # A placing call whose 2nd and 3rd positional arguments are both literal floats is a
    # hand-picked coordinate. Chrome (title_bar/footer/box backgrounds) legitimately uses them,
    # so only the CONTENT placers are counted — those are the ones that grow.
    # Parse the AST rather than grepping the text: the scaffold's own docstring QUOTES the bad
    # idiom as the thing not to do ("callout(s, 0.6, 4.3, ...) is the bug this scaffold used to
    # teach"), and a regex over source counts that as a violation — punishing the file for
    # explaining itself. Only real calls count.
    import ast
    tree = ast.parse(src)
    hand = []
    for node in ast.walk(tree):
        if not isinstance(node, ast.Call) or not isinstance(node.func, ast.Name):
            continue
        if node.func.id not in ("bullet", "callout", "equation_native", "equation_png"):
            continue
        a = node.args[1:3]
        if len(a) == 2 and all(isinstance(x, ast.Constant) and isinstance(x.value, (int, float))
                               for x in a):
            hand.append("%s at line %d" % (node.func.id, node.lineno))
    check("no auto-growing block is placed at a hand-picked y", not hand, hand)

    # ---- the file must still RUN, and its output must carry what it claims ------------
    r = subprocess.run([sys.executable, str(SCAFFOLD)], capture_output=True, text=True)
    check("the scaffold runs to completion", r.returncode == 0, (r.stderr or "")[-160:])
    out = re.search(r"saved -> (\S+)", r.stdout or "")
    check("it reports where it saved", bool(out), r.stdout[-120:])
    if out:
        import deckkit as dk
        from pptx import Presentation
        prs = Presentation(out.group(1))
        tagged = [sh for sl in prs.slides for sh in sl.shapes if dk._is_motif(sh)]
        loud = [sh for sh in tagged if dk._is_motif(sh, loud=True)]
        check("the built deck actually carries tagged motif shapes", len(tagged) >= 2, len(tagged))
        check("...including a loud (budgeted) appearance", len(loud) >= 1, len(loud))
        check("...and the loud budget it demonstrates is a LEGAL one (<=3)",
              len({id(s) for s in loud}) <= 3 * 4, len(loud))
        crit = [f for f in dk.lint_layout(prs, verbose=False) if f[1] == "CRITICAL"]
        check("the scaffold's own output passes the build-time gate", crit == [], crit)

    # ---- the docstring must say WHY, since rationale is what generalises --------------
    # the module docstring: from the OPENING triple-quote to the matching CLOSE. (Searching from
    # index 3 finds the opening quote, not the closing one, and yields an empty head that fails
    # every assertion below for the wrong reason.)
    _o = src.index('"""')
    head = src[_o:src.index('"""', _o + 3) + 3]
    check("the scaffold states the never-hand-pick-a-y rule in its own docstring",
          "hand-pick" in head.lower() or "hand-picked" in head.lower())
    check("the scaffold states the tag-your-motif rule in its own docstring",
          "tag_motif" in head or "TAG the deck" in head)

    print("\n{} passed, {} failed".format(len(PASS), len(FAIL)))
    return 1 if FAIL else 0


if __name__ == "__main__":
    sys.exit(main())

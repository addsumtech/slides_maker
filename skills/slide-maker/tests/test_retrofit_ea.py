#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""`CJK_NO_EA` is a build BLOCKER whose advice used to be unusable at the moment it fired.

`lint_layout(prs, strict=True)` runs at the END of a build script and raises on this CRITICAL. Its
remedy said "set deckkit.EAFONT before building" — which by then is the one thing you cannot do, and
which would not have helped on the two paths where the fault actually lives: a TEMPLATE-branch deck
whose text came in with `open_template()`, and a REDESIGN / surgical fix-pass editing a deck this
skill did not author. Neither routes text through `set_font()`, and `set_font()` is the only reader
of `EAFONT`. So the gate said "you are blocked" and pointed at a lever connected to nothing.

`deckkit.retrofit_ea(prs)` is the lever that is connected. This suite holds it to what the CRITICAL
and the docs now promise, in both directions — a retrofit that silently fixed nothing, or that
quietly rewrote a deliberate per-run face, would each be worse than the fault it clears.

Run:  python3 tests/test_retrofit_ea.py
"""
import copy
import pathlib
import sys
import tempfile

HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

from pptx.oxml.ns import qn                                              # noqa: E402
from pptx.util import Inches                                             # noqa: E402

import deckkit as dk                                                     # noqa: E402

FACE = "Hiragino Sans GB"
PASS, FAIL = [], []


def check(name, cond, detail=""):
    (PASS if cond else FAIL).append(name)
    print(("  ok   " if cond else "  FAIL ") + name
          + (("  — " + str(detail)) if detail and not cond else ""))


def runs(prs):
    """(text, ea-typeface-or-None, child tag order) for every run in the deck."""
    out = []
    for slide in prs.slides:
        for r in slide.shapes._spTree.iter(qn("a:r")):
            t = r.find(qn("a:t"))
            rPr = r.find(qn("a:rPr"))
            ea = rPr.find(qn("a:ea")) if rPr is not None else None
            out.append(((t.text or "") if t is not None else "",
                        ea.get("typeface") if ea is not None else None,
                        [e.tag.split("}")[1] for e in rPr] if rPr is not None else []))
    return out


def ea_of(prs, needle):
    for txt, ea, _ in runs(prs):
        if needle in txt:
            return ea
    return "<no such run>"


def _ea_attr(el):
    """The `typeface` on the first `<a:ea>` anywhere under an element, or None."""
    for ea in el.iter(qn("a:ea")):
        return ea.get("typeface")
    return None


def inherited_face(prs, needle):
    """What the run WILL render with, resolved the way deckkit resolves it."""
    for slide in prs.slides:
        for r in slide.shapes._spTree.iter(qn("a:r")):
            t = r.find(qn("a:t"))
            if t is not None and needle in (t.text or ""):
                return dk._inherited_ea(r)
    return "<no such run>"


def order_of(prs, needle):
    for txt, _, order in runs(prs):
        if needle in txt:
            return order
    return []


def foreign_deck():
    """A deck as the template / fix-pass paths hand it over: nothing built by set_font().

    Every run here is one python-pptx would make, which is the point — these are exactly the runs
    `EAFONT` cannot reach, so they are the ones the retrofit exists for.
    """
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.6))
    tb.text_frame.text = "四维重建 pipeline"                      # plain CJK run

    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4), Inches(0.6))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = "论文链接"                                          # rPr ALREADY has a child
    r2.hyperlink.address = "https://example.org"

    tbl = s.shapes.add_table(2, 2, Inches(0.5), Inches(2.2), Inches(5), Inches(1.2)).table
    tbl.cell(0, 0).text = "指标"                                  # inside a TABLE — gate is blind
    tbl.cell(0, 1).text = "value"                                 # Latin — must stay untouched

    g = s.shapes.add_group_shape()
    inner = g.shapes.add_textbox(Inches(6), Inches(0.5), Inches(3), Inches(0.6))
    inner.text_frame.text = "组内文字"                            # inside a GROUP — gate is blind

    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(4), Inches(0.5))
    tb3.text_frame.text = "pure latin only"
    return prs


def inherited_deck():
    """A supplied CJK template's shape: the EA face lives one level up, not on the run.

    Both routes are real OOXML inheritance and both render correctly, so both must be silent —
    and, just as importantly, retrofit_ea must not overwrite them with the deck's own EAFONT,
    which would retypeset someone else's deck as a side effect of clearing a lint line.
    """
    prs = dk.blank_deck()
    s = dk.add_slide(prs)

    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.6))
    para = tb.text_frame.paragraphs[0]
    para.add_run().text = "段落继承"
    pPr = para._p.get_or_add_pPr()
    d = pPr.makeelement(qn("a:defRPr"), {})
    pPr.append(d)
    d.append(d.makeelement(qn("a:ea"), {"typeface": "思源黑体"}))

    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4), Inches(0.6))
    tf2 = tb2.text_frame
    tf2.paragraphs[0].add_run().text = "形状继承"
    lst = tf2._txBody.makeelement(qn("a:lstStyle"), {})
    lvl = lst.makeelement(qn("a:lvl1pPr"), {})
    lst.append(lvl)
    dd = lvl.makeelement(qn("a:defRPr"), {})
    lvl.append(dd)
    dd.append(dd.makeelement(qn("a:ea"), {"typeface": "方正黑体"}))
    tf2._txBody.insert(1, lst)

    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4), Inches(0.6))
    tb3.text_frame.paragraphs[0].add_run().text = "无任何继承"
    return prs


def cjk_no_ea(prs):
    return [f for f in dk.lint_layout(prs, verbose=False) if f[2] == "CJK_NO_EA"]


def main():
    print("retrofit_ea contract")
    saved_ea, saved_disp = dk.EAFONT, dk.EADISPLAY
    dk.EAFONT = dk.EADISPLAY = None
    try:
        # --- the fault the CRITICAL is about ------------------------------------------------
        prs = foreign_deck()
        check("a template-shaped deck trips CJK_NO_EA before the retrofit", cjk_no_ea(prs) != [])

        # --- with no face anywhere it RAISES; it must never return a quiet 0 -----------------
        try:
            dk.retrofit_ea(prs)
            check("no face set raises instead of no-oping", False, "returned instead of raising")
        except ValueError as exc:
            check("no face set raises instead of no-oping", True)
            check("the raise names EAFONT so the message is actionable", "EAFONT" in str(exc))

        # --- the retrofit clears the gate ---------------------------------------------------
        n = dk.retrofit_ea(prs, FACE, verbose=False)
        check("it clears CJK_NO_EA", cjk_no_ea(prs) == [], cjk_no_ea(prs))
        check("it reports how many runs it fixed", n == 4, n)

        # --- and reaches further than the gate can see ---------------------------------------
        check("a plain CJK run is stamped", ea_of(prs, "四维重建") == FACE)
        check("a run inside a TABLE cell is stamped (the gate cannot see it)",
              ea_of(prs, "指标") == FACE)
        check("a run inside a GROUP is stamped (the gate cannot see it)",
              ea_of(prs, "组内文字") == FACE)

        # --- what it must NOT touch ----------------------------------------------------------
        check("a Latin-only run is left alone", ea_of(prs, "pure latin") is None)
        check("a Latin-only table cell is left alone", ea_of(prs, "value") is None)

        # --- schema order: <a:ea> precedes hlinkClick, or PowerPoint rejects the file ---------
        order = order_of(prs, "论文链接")
        check("<a:ea> is inserted BEFORE hlinkClick, not appended",
              "ea" in order and "hlinkClick" in order
              and order.index("ea") < order.index("hlinkClick"), order)

        # --- idempotent, and never overwrites a deliberate per-run face -----------------------
        check("a second call is a no-op", dk.retrofit_ea(prs, FACE, verbose=False) == 0)
        check("a different face does NOT overwrite an existing slot",
              dk.retrofit_ea(prs, "Noto Sans CJK SC", verbose=False) == 0
              and ea_of(prs, "四维重建") == FACE)

        # --- the file it produces is still openable -------------------------------------------
        with tempfile.TemporaryDirectory() as tmp:
            out = pathlib.Path(tmp) / "retrofit.pptx"
            prs.save(str(out))
            from pptx import Presentation
            check("the saved deck reopens", len(Presentation(str(out)).slides) == 1)

        # --- EAFONT is the default, and EADISPLAY is the fallback ------------------------------
        dk.EAFONT = FACE
        p2 = foreign_deck()
        dk.retrofit_ea(p2, verbose=False)
        check("face defaults to EAFONT", ea_of(p2, "四维重建") == FACE)
        dk.EAFONT, dk.EADISPLAY = None, "Noto Sans CJK SC"
        p3 = foreign_deck()
        dk.retrofit_ea(p3, verbose=False)
        check("face falls back to EADISPLAY", ea_of(p3, "四维重建") == "Noto Sans CJK SC")

        # --- schema position holds for EVERY rPr shape, not just the one we thought of ---------
        # PowerPoint rejects an out-of-order rPr rather than ignoring it, and a foreign deck's runs
        # carry children a deckkit-built run never does. So sweep the whole child vocabulary
        # rather than trusting the hyperlink case above to stand for it.
        SCHEMA = ["ln", "noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill",
                  "effectLst", "effectDag", "highlight", "uLnTx", "uLn", "uFillTx", "uFill",
                  "latin", "ea", "cs", "sym", "hlinkClick", "hlinkMouseOver", "rtl", "extLst"]
        p6 = dk.blank_deck()
        s6 = dk.add_slide(p6)
        seeds = ["solidFill", "highlight", "uLn", "effectLst", "latin", "cs", "sym",
                 "hlinkClick", "hlinkMouseOver", "rtl", "extLst", None]
        for i, seed in enumerate(seeds):
            tb = s6.shapes.add_textbox(Inches(0.3), Inches(0.2 + i * 0.4), Inches(3), Inches(0.35))
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = "中文" + (seed or "bare")
            if seed:
                rPr = run._r.get_or_add_rPr()
                rPr.append(rPr.makeelement(qn("a:" + seed), {}))
        dk.retrofit_ea(p6, FACE, verbose=False)
        bad = []
        for _, _, order in runs(p6):
            idx = [SCHEMA.index(t) for t in order if t in SCHEMA]
            if idx != sorted(idx) or "ea" not in order:
                bad.append(order)
        check("<a:ea> lands in schema position for every rPr child it can meet", not bad, bad)
        with tempfile.TemporaryDirectory() as tmp:
            out = pathlib.Path(tmp) / "order.pptx"
            p6.save(str(out))
            from pptx import Presentation
            check("that deck reopens too", len(Presentation(str(out)).slides) == 1)

        # --- INHERITED faces: not a defect, and not the retrofit's to overwrite ----------------
        dk.EAFONT, dk.EADISPLAY = FACE, None
        p5 = inherited_deck()
        found = cjk_no_ea(p5)
        check("a run inheriting its EA face from the paragraph is not flagged",
              found and "段落继承" not in found[0][3], found)
        check("a run inheriting its EA face from the shape's lstStyle is not flagged",
              found and "形状继承" not in found[0][3], found)
        check("a run inheriting nothing IS still flagged",
              len(found) == 1 and "无任何继承" in found[0][3], found)
        check("the retrofit stamps only the genuinely bare run",
              dk.retrofit_ea(p5, verbose=False) == 1)
        check("the template's paragraph-level face survives the retrofit",
              inherited_face(p5, "段落继承") == "思源黑体")
        check("the template's shape-level face survives the retrofit",
              inherited_face(p5, "形状继承") == "方正黑体")
        check("the bare run took the deck's EAFONT", inherited_face(p5, "无任何继承") == FACE)
        check("the retrofit clears the lint on an inheriting deck too", cjk_no_ea(p5) == [])

        # --- an EMPTY <a:ea> is the fault, not a face --------------------------------------
        # `<a:ea typeface=""/>` is how OOXML spells "no East-Asian font" — the stock Office theme's
        # own fontScheme carries it verbatim, and templates copy it down. A first draft guarded
        # _stamp_ea on the PRESENCE of the element while the gate read the resolved FACE, so the
        # CRITICAL fired, retrofit_ea decided the run needed fixing, stamped nothing, and
        # strict=True raised AFTER the documented remedy had run, with no second lever anywhere.
        for attrs, name in (({"typeface": ""}, 'typeface=""'), ({}, "no typeface attribute")):
            p7 = dk.blank_deck()
            s7 = dk.add_slide(p7)
            tb = s7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.6))
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = "四维重建"
            rPr = run._r.get_or_add_rPr()
            rPr.append(rPr.makeelement(qn("a:ea"), attrs))
            check("<a:ea %s> is reported as missing" % name, cjk_no_ea(p7) != [])
            check("<a:ea %s> is FILLED, not skipped" % name,
                  dk.retrofit_ea(p7, FACE, verbose=False) == 1)
            check("<a:ea %s> — the remedy clears the gate" % name, cjk_no_ea(p7) == [])
            check("<a:ea %s> — exactly one slot afterwards" % name,
                  len(rPr.findall(qn("a:ea"))) == 1)

        # --- coverage the CRITICAL and the docs now promise ---------------------------------
        p8 = dk.blank_deck()
        s8 = dk.add_slide(p8)
        para = s8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.6))
        pel = para.text_frame.paragraphs[0]._p
        fld = pel.makeelement(qn("a:fld"), {"id": "{x}", "type": "datetime1"})
        ft = fld.makeelement(qn("a:t"), {})
        ft.text = "二〇二六年八月"
        fld.append(ft)
        pel.append(fld)
        check("a CJK <a:fld> date/footer field is stamped (a template's own chrome)",
              dk.retrofit_ea(p8, FACE, verbose=False) == 1)
        check("the field's slot carries the face",
              _ea_attr(fld) == FACE, _ea_attr(fld))

        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        def chart_deck(cats, drop_txpr, series="营收"):
            """`drop_txpr` is the FOREIGN shape: a chart with no chartSpace-level <c:txPr>.

            Every chart python-pptx or deckkit builds already carries one, so a suite that only
            builds its own charts never reaches the branch that creates it — which is the branch
            that was writing the element in the wrong place.
            """
            prs = dk.blank_deck()
            s = dk.add_slide(prs)
            cd = CategoryChartData()
            cd.categories = cats
            cd.add_series(series, (3, 5, 4))
            gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(0.5),
                                    Inches(6), Inches(3), cd)
            cs = gf.chart._chartSpace
            if drop_txpr:
                t = cs.find(qn("c:txPr"))
                if t is not None:
                    cs.remove(t)
            return prs, cs

        def tags(cs):
            return [e.tag.split("}")[1] for e in cs]

        p9, cs9 = chart_deck(["一月", "二月", "三月"], drop_txpr=False)
        check("a chart's CJK category text is stamped (another part, no runs of its own)",
              dk.retrofit_ea(p9, FACE, verbose=False) == 1)
        check("the chart's defRPr carries the face",
              any(_ea_attr(d) == FACE for sl in p9.slides for d in dk._chart_ea_parts(sl)))

        # CT_ChartSpace is a SEQUENCE and <c:txPr> is 10th of 14 — BEFORE <c:externalData>.
        # Appending it put the element last and the chart part stopped validating: the same
        # failure _EA_FOLLOWERS prevents for <a:rPr>, on the one code path that skipped the guard.
        p10, cs10 = chart_deck(["一月", "二月", "三月"], drop_txpr=True)
        check("the foreign shape really lacks c:txPr, so this exercises the create branch",
              "txPr" not in tags(cs10), tags(cs10))
        check("a foreign chart is stamped too", dk.retrofit_ea(p10, FACE, verbose=False) == 1)
        _t = tags(cs10)
        check("<c:txPr> is created BEFORE <c:externalData>, not appended after it",
              "txPr" in _t and "externalData" in _t
              and _t.index("txPr") < _t.index("externalData"), _t)

        # ...and a Latin chart is not this function's business at all. Every other branch of the
        # retrofit is guarded by _has_cjk; this one ran on any chart, so an all-Latin deck came
        # back "stamped 1 run(s)" with its chart part rewritten for nothing.
        # The SERIES NAME is a c:v cell too — the first draft of this fixture left it as 营收 and
        # the "Latin" chart was not Latin, which the assertion caught.
        p11, cs11 = chart_deck(["Jan", "Feb", "Mar"], drop_txpr=True, series="revenue")
        check("a LATIN chart is left completely alone",
              dk.retrofit_ea(p11, FACE, verbose=False) == 0 and "txPr" not in tags(cs11),
              tags(cs11))
        p12, cs12 = chart_deck(["Jan", "Feb", "Mar"], drop_txpr=True, series="营收")
        check("a chart whose only CJK is its SERIES NAME is still stamped",
              dk.retrofit_ea(p12, FACE, verbose=False) == 1)

        # --- layouts: the one real gap, and it must be LOUD ----------------------------------
        p10 = dk.blank_deck()
        s10 = dk.add_slide(p10)
        tmp_tb = s10.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4), Inches(0.4))
        tmp_tb.text_frame.text = "版式页脚中文"
        moved = copy.deepcopy(tmp_tb._element)
        s10.shapes._spTree.remove(tmp_tb._element)
        p10.slide_masters[0].slide_layouts[0].shapes._spTree.append(moved)
        check("layout CJK is invisible to the build gate", cjk_no_ea(p10) == [])
        check("the default pass does not touch layouts",
              dk.retrofit_ea(p10, FACE, verbose=False) == 0)
        check("...but it COUNTS them, so the gap is not silent",
              len(dk._layout_cjk_without_ea(p10)) == 1)
        check("layouts=True is the lever that fixes them",
              dk.retrofit_ea(p10, FACE, layouts=True, verbose=False) == 1)
        check("and then there is nothing left to report",
              dk._layout_cjk_without_ea(p10) == [])

        # --- the normal build path is unchanged by the _apply_ea refactor ----------------------
        dk.EAFONT, dk.EADISPLAY = FACE, None
        p4 = dk.blank_deck()
        s4 = dk.add_slide(p4)
        dk.text(s4, 0.5, 0.5, 5, 0.6, [[("中文标题", 20, dk.DEEP, True, False, dk.FONT)]])
        check("set_font still stamps <a:ea> on a deckkit-built run", ea_of(p4, "中文标题") == FACE)
        check("a deckkit-built CJK deck never trips CJK_NO_EA", cjk_no_ea(p4) == [])
        run = list(s4.shapes)[0].text_frame.paragraphs[0].runs[0]
        dk._apply_ea(run, "Noto Sans CJK SC")
        rPr = run._r.find(qn("a:rPr"))
        check("set_font's path DOES overwrite (the author is speaking, not a retrofit)",
              rPr.find(qn("a:ea")).get("typeface") == "Noto Sans CJK SC")
        check("overwriting leaves exactly one <a:ea>", len(rPr.findall(qn("a:ea"))) == 1)
    finally:
        dk.EAFONT, dk.EADISPLAY = saved_ea, saved_disp

    print("\n{} passed, {} failed".format(len(PASS), len(FAIL)))
    return 1 if FAIL else 0


if __name__ == "__main__":
    sys.exit(main())

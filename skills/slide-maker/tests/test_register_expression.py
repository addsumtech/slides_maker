#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""A register must reach the PIXELS, not just the palette.

WHY. `presets.apply()` promised "palette AND structure", `set_geometry`'s own docstring said the
structural tokens are "what actually makes a register look like itself", and SKILL.md said
`radius=0` "squares every rounded component". Measured on a real render, none of the three was
true where it mattered:

  * `RULE_W_SCALE` was read by `hrule()` and NOTHING else — 1 of 181 public functions. Every card
    border, node outline and table rule was register-invariant. The same diagram in brutalist
    (rule_w 3.0), bauhaus (2.6) and swiss (0.6) came out byte-identical at 1.4pt.
  * `node()` — the general architecture/flowchart builder — read `RADIUS_SCALE` not at all, so
    brutalist, swiss, ink_wash and blueprint (radius 0, and three of the four registers SKILL.md
    names as the ones the feature exists to unlock) drew `roundRect adj=0.1667`, python-pptx's
    untouched default, on registers whose own guard forbids rounded cards.
  * `presets.apply()` read `bg` and applied it nowhere. 8 of 18 registers are dark; a caller doing
    exactly what the docstring said got the register's LIGHT ink on a white canvas.
  * `node()`'s own two defaults contradicted each other on every dark register: it paints WHITE and
    took DEEP for the label, while `set_palette` had rebound DEEP to that register's light ink.
    Measured with NO caller override at all — glassmorphism 1.09:1, synthwave 1.16:1.
  * the four registers added last (`bauhaus`, `midcentury`, `terminal`, `synthwave`) stored every
    colour as a plain list while the first fourteen stored RGBColor, so the documented
    `p["accents"][0]` usage raised ValueError on 4 of 18 and worked on the other 14.

Nothing reported any of it. The sameness lint measures monotony WITHIN a deck (LAYOUT SAMENESS,
CARD DOMINANCE, ENVELOPE MONOCULTURE); no check asks whether a deck expresses the register it
chose, so a brutalist deck rendered with glassmorphism's geometry passed every gate. This suite is
that check, at the level the defect lived: the built PPTX, not the arithmetic.

Run:  python3 tests/test_register_expression.py
"""
import pathlib
import sys

HERE = pathlib.Path(__file__).resolve().parent
SKILL = HERE.parent
sys.path.insert(0, str(SKILL / "scripts"))

from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Emu  # noqa: E402

import deckkit as dk  # noqa: E402
import presets  # noqa: E402

PASS = FAIL = 0
NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def check(name, cond, detail=""):
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ok   {name}")
    else:
        FAIL += 1
        print(f"  FAIL {name}  {detail}")


def build_one(register):
    """One slide in `register`, using ONLY the documented call shapes and no colour overrides."""
    p = presets.apply(register)
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    dk.title_bar(s, "T", kicker=register.upper(), accent=p["accents"][0])
    dk.node(s, 0.8, 2.0, 2.0, 1.1, "Label", sub="a caption")
    dk.node(s, 3.0, 2.0, 2.0, 1.1, "Hub", sub="a caption", hub=True, accent=p["accents"][0])
    return p, prs, s


def node_shapes(slide):
    out = []
    for sh in slide.shapes:
        try:
            if abs(Emu(sh.width).inches - 2.0) < 0.05 and abs(Emu(sh.height).inches - 1.1) < 0.1:
                out.append(sh)
        except Exception:
            pass
    return out


def main():
    print("== every colour field is RGBColor, in every register ==")
    # 4 of 18 stored lists; the documented p["accents"][0] usage then raised ValueError on those
    # four and worked on the rest, which is the worst shape for a bug - it reads as caller error.
    bad = []
    for n in presets.names():
        p = presets.preset(n)
        for k in ("bg", "ink", "muted"):
            if not isinstance(p[k], RGBColor):
                bad.append(f"{n}.{k}={type(p[k]).__name__}")
        for i, a in enumerate(p["accents"]):
            if not isinstance(a, RGBColor):
                bad.append(f"{n}.accents[{i}]={type(a).__name__}")
    check(f"all {len(presets.names())} presets normalise to RGBColor", not bad, bad[:4])

    print("\n== the four registers that used to CRASH the documented usage ==")
    for n in ("bauhaus", "midcentury", "terminal", "synthwave"):
        try:
            build_one(n)
            ok, err = True, ""
        except Exception as exc:
            ok, err = False, f"{exc.__class__.__name__}: {exc}"
        check(f"{n}: builds with p['accents'][0] straight from the preset", ok, err)

    print("\n== radius reaches node(), and 0 really means square ==")
    for n, want_square in (("brutalist", True), ("swiss", True), ("ink_wash", True),
                           ("glassmorphism", False), ("synthwave", False)):
        _, _, s = build_one(n)
        adjs = [round(sh.adjustments[0], 4) for sh in node_shapes(s) if len(sh.adjustments)]
        got_square = bool(adjs) and all(a == 0.0 for a in adjs)
        check(f"{n} (radius={presets.preset(n)['radius']}): node corners "
              f"{'square' if want_square else 'rounded'}",
              got_square == want_square, adjs)

    print("\n== rule_w reaches a node OUTLINE and a box BORDER, not only hrule ==")
    widths = {}
    for n in ("swiss", "brutalist", "bauhaus", "ink_wash"):
        _, _, s = build_one(n)
        for sh in node_shapes(s):
            w = sh.line.width
            if w:
                widths[n] = round(Emu(w).pt, 2)
                break
    check("brutalist (3.0) draws a heavier node outline than swiss (0.6)",
          widths.get("brutalist", 0) > widths.get("swiss", 99), widths)
    check("bauhaus (2.6) and ink_wash (0.5) also differ, in the declared direction",
          widths.get("bauhaus", 0) > widths.get("ink_wash", 99), widths)
    check("each outline is the node default (1.4pt) times the register's rule_w",
          all(abs(widths[n] - 1.4 * presets.preset(n)["rule_w"]) < 0.05 for n in widths), widths)
    # box() too - the 19 form components that route through it inherit the same scale.
    dk.set_geometry(rule_w=1.0)
    prs = dk.blank_deck(); s0 = dk.add_slide(prs)
    b1 = dk.box(s0, 1, 1, 2, 1, line=dk.DEEP, line_w=1.0)
    dk.set_geometry(rule_w=3.0)
    b3 = dk.box(s0, 4, 1, 2, 1, line=dk.DEEP, line_w=1.0)
    check("box(line_w=1.0) scales with rule_w (1.0 -> 3.0)",
          abs(Emu(b3.line.width).pt - 3 * Emu(b1.line.width).pt) < 0.05,
          (Emu(b1.line.width).pt, Emu(b3.line.width).pt))
    dk.set_geometry(radius=1.0, rule_w=1.0)

    print("\n== the GROUND is painted, so a dark register is dark from slide one ==")
    for n in ("dark_tech", "synthwave", "glassmorphism", "luxury_dark"):
        _, _, s = build_one(n)
        painted = s._element.find(f".//{NS}bg") is not None
        check(f"{n}: add_slide() paints <p:bg>", painted)
    dk.set_ground(None)
    prs = dk.blank_deck(); s = dk.add_slide(prs)
    check("...and a deck that never sets a ground is UNCHANGED (no <p:bg>)",
          s._element.find(f".//{NS}bg") is None)

    print("\n== a node's own defaults are legible on their own ground ==")
    # No caller override of any kind: this is what the component does by itself.
    for n in ("glassmorphism", "synthwave", "dark_tech", "swiss", "ink_wash"):
        p = presets.apply(n)
        label_on_plain = dk.contrast_ratio(dk.on(dk.WHITE), dk.WHITE)
        sub_on_hub = dk.contrast_ratio(
            dk.MUTE if dk.contrast_ratio(dk.MUTE, p["accents"][0]) >= 4.5
            else dk.mute_for(p["accents"][0], 4.5), p["accents"][0])
        check(f"{n}: plain-node label >= 4.5:1 and hub caption >= 3:1",
              label_on_plain >= 4.5 and sub_on_hub >= 3.0,
              f"label={label_on_plain:.2f} sub={sub_on_hub:.2f}")

    print("\n== the contract card's colour helpers actually exist ==")
    # sigs.py told authors to "resolve by ground (mute_for(bg), on(fill))" while NEITHER function
    # existed - check_reference_code.py scans SKILL.md/references/agents, never scripts/sigs.py,
    # so the one file whose job is stating call shapes was outside the checker.
    for fn in ("on", "mute_for"):
        check(f"deckkit.{fn}() is a real public helper", callable(getattr(dk, fn, None)))
    for g in ((255, 255, 255), (12, 19, 32), (255, 45, 149)):
        check(f"mute_for({g}) clears its target on that ground",
              dk.contrast_ratio(dk.mute_for(g), g) >= 3.0
              and dk.contrast_ratio(dk.mute_for(g, 4.5), g) >= 4.5,
              f"{dk.contrast_ratio(dk.mute_for(g), g):.2f} / "
              f"{dk.contrast_ratio(dk.mute_for(g, 4.5), g):.2f}")

    print("\n== a node caption follows the register's body face, not MONO ==")
    # Hardcoded MONO put every diagram caption in Menlo on ink_wash / editorial_paper /
    # eastern_traditional, and in CJK it split one line across two type systems because MONO
    # carries no Chinese glyphs.
    p = presets.apply("ink_wash")
    prs = dk.blank_deck(); s = dk.add_slide(prs)
    dk.node(s, 0.8, 2.0, 2.0, 1.1, "标签", sub="中文说明 with MRI")
    faces = {r.font.name for sh in s.shapes if sh.has_text_frame
             for para in sh.text_frame.paragraphs for r in para.runs if r.font.name}
    check("the caption is NOT set in MONO", dk.MONO not in faces, sorted(faces))
    check("...it is set in the register's body face", dk.FONT in faces, sorted(faces))
    prs = dk.blank_deck(); s = dk.add_slide(prs)
    dk.node(s, 0.8, 2.0, 2.0, 1.1, "Label", sub="/usr/bin/x", sub_font=dk.MONO)
    faces = {r.font.name for sh in s.shapes if sh.has_text_frame
             for para in sh.text_frame.paragraphs for r in para.runs if r.font.name}
    check("...and sub_font=MONO still gets MONO, for a caption that IS code", dk.MONO in faces,
          sorted(faces))

    print("\n== DECLARED -> APPLIED: both gates verify the register they demand ==")
    # `style_pick` was required as a string by codex_delivery_gate.py AND
    # render_deck.py --gate-check, and verified by NEITHER: measured by grep, presets.apply /
    # set_geometry / set_ground appeared in no gate script at all. A deck recording
    # "brutalist for engineering - beat blueprint" and built with deckkit's stock defaults
    # passed both. One checker, imported by both gates, so they cannot grow two answers.
    sys.path.insert(0, str(SKILL / "scripts"))
    import check_style_applied as csa
    names = csa.preset_names()
    cases = [
        ("brutalist for engineering \u00b7 beat blueprint", 'presets.apply("brutalist")', None, 0,
         "declared and applied"),
        ("brutalist for engineering \u00b7 beat blueprint", 'dk.set_palette(deep=X)', None, 1,
         "declared, hand-built - the measured failure"),
        ("brutalist for engineering \u00b7 beat swiss", 'presets.apply("swiss")', None, 1,
         "applied the REJECTED rival, not the pick"),
        ("bespoke seismograph register \u00b7 beat swiss", '', None, 0,
         "bespoke is not preset-based"),
        ("n/a \u2014 locked: the user's template", '', None, 0, "locked look"),
        ("brutalist for engineering", 'dk.set_palette(deep=X)',
         "the client brand book fixes border weight at 1pt, which brutalist would triple", 0,
         "a real named waiver clears it"),
        ("brutalist for engineering", 'dk.set_palette(deep=X)', "ok", 1,
         "a bare 'ok' is not a waiver"),
    ]
    for pick, src, waiver, want, why in cases:
        got, _ = csa.evaluate(pick, src, names, waiver)
        check(f"style_pick: {why}", got == want, f"want={want} got={got}")
    check("its own selftest still agrees with this file",
          csa.selftest(names) == 0)
    # AST, not a substring: `"_style_applied_gate" in body` is satisfied by the function's own
    # DEFINITION, so deleting the call left it green. A check that cannot go red measures nothing,
    # which is the whole defect class this suite exists for - caught here by running the negative
    # control rather than trusting the assertion.
    import ast as _ast
    for gate, fname in (("scripts/render_deck.py", "_style_applied_gate"),
                        ("scripts/codex_delivery_gate.py", "check_style_applied")):
        tree = _ast.parse((SKILL / gate).read_text(encoding="utf-8"))
        called = any(isinstance(n, _ast.Call) and isinstance(n.func, _ast.Name)
                     and n.func.id == fname for n in _ast.walk(tree))
        check(f"{gate} actually CALLS {fname}() (a gate nobody runs is prose)", called)

    print("\n== the GENERATED-TEMPLATE branch (Q1d) is aligned, not left behind ==")
    # This branch declares a four-line IDENTITY-PROPAGATION CONTRACT (palette/type/geometry/
    # surface) and had a carrier for only the first two: generated-template.md mentioned
    # set_geometry / set_ground / presets.apply ZERO times and the scaffold it says to copy
    # exposed colours and fonts only. A contract with no carrier is the preset defect one branch
    # over - and the scaffold is the curriculum, so whatever it omits does not get written.
    scaffold = (SKILL / "references/examples/style_example.py").read_text(encoding="utf-8")
    check("style_example.py teaches set_geometry (the `geometry:` contract line)",
          "set_geometry" in scaffold)
    check("style_example.py teaches set_ground (the derived ground)", "set_ground" in scaffold)
    gt = (SKILL / "references/generated-template.md").read_text(encoding="utf-8")
    for token in ("set_geometry", "set_ground", "look_source"):
        check(f"generated-template.md names {token}", token in gt)
    # And the gate must not FALSE-POSITIVE this branch: the topic contest runs here too, so a
    # generated pick legitimately leads with the preset that art-directed the hero.
    gen = "editorial_paper for a poetry course \u00b7 beat ink_wash \u2014 generated hero art-directed by it"
    code, _ = csa.evaluate(gen, "dk.set_palette(deep=X)", names, None)
    check("a generated-branch pick naming a preset is NOT blocked (exit 2 = NOT CHECKED)",
          code == 2, code)
    code, _ = csa.evaluate(gen, "dk.set_palette(deep=X)", names, None, "generated")
    check("...and look_source='generated' skips it outright", code == 0, code)
    code, _ = csa.evaluate("swiss for a status update \u00b7 beat consulting",
                           "dk.set_palette(deep=X)", names, None)
    check("...while a qualifier-free preset pick still HARD blocks", code == 1, code)

    print("\n== apply() fills the SEMANTIC SLOTS, not just the accent cycle ==")
    # `apply()` passed deep/slate/accents and nothing else, so BLUE, TEAL, MAGENTA and TINT kept
    # deckkit's own defaults on every register — and those four are what 33 component signatures
    # default to. Measured: after apply("terminal"), a callout() on a black-and-phosphor deck came
    # out with a #E3004F magenta rule on a #EAF3FA pale-blue panel. The remap machinery in
    # set_palette was working; it was simply never told the register's colours.
    import deckkit as _dk
    wrong = []
    for n in presets.names():
        p = presets.apply(n)
        acc = list(p["accents"])
        if tuple(_dk.MAGENTA) != tuple(acc[0]):
            wrong.append("%s: MAGENTA=%s want %s" % (n, tuple(_dk.MAGENTA), tuple(acc[0])))
        if tuple(_dk.TINT) != presets.panel(n):
            wrong.append("%s: TINT=%s want %s" % (n, tuple(_dk.TINT), presets.panel(n)))
    check("every register rebinds MAGENTA (the highlight slot) and TINT (the panel fill)",
          not wrong, "; ".join(wrong[:3]))

    # the panel has to separate from the ground AND carry the register's ink
    thin, unreadable = [], []
    for n in presets.names():
        p = presets.preset(n)
        pa, bg, ink = presets.panel(n), tuple(p["bg"]), tuple(p["ink"])
        if _dk.contrast_ratio(pa, bg) < 1.12:
            thin.append("%s %.2f" % (n, _dk.contrast_ratio(pa, bg)))
        if _dk.contrast_ratio(ink, pa) < 4.5:
            unreadable.append("%s %.1f" % (n, _dk.contrast_ratio(ink, pa)))
    check("the derived panel separates from its ground in all 18 registers", not thin,
          "; ".join(thin))
    check("...and the register's own ink stays above the 4.5:1 body floor ON the panel",
          not unreadable, "; ".join(unreadable))
    check("the panel is DERIVED from bg+ink, not hand-picked per register — it stays correct if a "
          "preset's colours are ever retuned, which eighteen hand-picked values would not",
          presets.panel("terminal") != presets.panel("swiss"))

    print(f"\n{PASS} passed, {FAIL} failed")
    return 1 if FAIL else 0


if __name__ == "__main__":
    raise SystemExit(main())

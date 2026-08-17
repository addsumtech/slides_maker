#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""The counterweight: a deck can now be held for being measurably SAFE, not only for erring.

WHY. Every blocking signal in this skill punished a deck for being too much or too same. Nothing
could hold one for being too timid: `TIMID COVER` and `FLAT TYPE` are deliberately excluded from
SAMENESS_CODES, and the only force that could call a deck forgettable — the critic's
distinctiveness axis — is non-blocking at the default `balanced+` dial and lives inside a review
the user may decline.

Measured on a real 12-page build: ten iterations, every one driven by an advisory, every one making
the deck flatter. The dark pivot page was deleted to satisfy ONE-OFF CANVAS FLIP; content was cut
to satisfy TEXT WALL; the type scale collapsed to satisfy SIZE SPRAWL. Each of those advisories
names the ambitious repair FIRST — "enrich with a second column of substance", "repeat the
treatment as a divider family" — but subtraction is the cheaper way to make a number go away, and
with feedback on one side only, the cheap way always wins. The user's verdict was 设计能力变弱了,
and nothing in the pipeline had said so.

Both directions are asserted, because a miscalibrated gate on THIS side is worse than none — it
would punish exactly the restrained registers the skill protects:

  BLOCKS      a text-only, greyscale, timid-cover deck.
  DOES NOT    fire on a deck whose pages carry real protagonists (charts, figures).
  STANDS DOWN under `boldness: conservative` + a recorded `deliberately restrained` move, and
              under the size/aspect floors, exactly like the sameness gate.
  WAIVABLE    only by NAMING the register, with a codes list bound to this state of the deck.

Run:  python3 tests/test_timidity_gate.py
"""
import json
import pathlib
import subprocess
import sys
import tempfile

HERE = pathlib.Path(__file__).resolve().parent
SKILL = HERE.parent
RENDER = SKILL / "scripts" / "render_deck.py"
sys.path.insert(0, str(SKILL / "scripts"))
sys.path.insert(0, str(HERE))

import lint_deck as LD  # noqa: E402
from test_critic_waiver_gate import fit_content, ARC_OK, DESIGN_OK, GOOD_REASON, PROV_OK, write_proof  # noqa: E402

PASS = FAIL = 0


def check(name, cond, detail=""):
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ok   {name}")
    else:
        FAIL += 1
        print(f"  FAIL {name}\n       {str(detail)[:400]}")


def build(dest, *, visual=False, timid_cover=True, n=12, pics=False):
    """A 12-slide deck. `visual=True` gives half the pages a real non-text protagonist."""
    import deckkit as dk
    prs = dk.blank_deck(13.333, 7.5)
    grey = dk.RGBColor.from_string("444444")
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        head = 26 if (i or not timid_cover) else 26
        if i == 0 and not timid_cover:
            head = 72
        dk.text(s, 0.8, 0.7, 11.5, 1.1, [[(f"Slide {i+1} title here", head, grey, True, False)]])
        # Short on purpose: a long paragraph would trip the DENSITY gate and this suite would be
        # asserting the wrong thing. Timidity is about what a page CONTAINS, not how wordy it is.
        dk.text(s, 0.8, 2.2, 11.5, 3.2,
                [[("One short line of body copy.", 18, grey, False, False)]])
        if pics and 0 < i < n - 1:
            from PIL import Image
            f = dest / f"p{i}.png"
            Image.new("RGB", (400, 300), (120, 140, 160)).save(f)
            dk.picture(s, str(f), 7.6, 2.2, 4.4, 3.0, fit="contain")
        if visual and 0 < i < n - 1 and i % 2:
            dk.native_chart(s, 1.0, 3.2, 6.0, 3.0, ["a", "b", "c"],
                            [("s", [3, 5, 4])], kind="bar")
    out = dest / "t.pptx"
    dk.lint_layout(prs, verbose=False)
    prs.save(str(out))
    return out


def gate(deck, gates, *flags):
    gates = fit_content(gates, deck)
    (deck.parent / ".deck-gates.json").write_text(json.dumps(gates, ensure_ascii=False),
                                                  encoding="utf-8")
    p = subprocess.run([sys.executable, str(RENDER), str(deck), "--gate-check", "--static", *flags],
                       capture_output=True, text=True)
    return p.returncode, p.stdout + p.stderr


def record(**design):
    d = dict(DESIGN_OK)
    d.update(design)
    return {"critic": {"waived": GOOD_REASON, "waived_category": "no-dispatch-on-host",
                       "inline_ran": True},
            "design_plan": d, "content": dict(ARC_OK), "provenance": dict(PROV_OK)}


def main():
    print("== the composite's shape mirrors sameness ==")
    check("4 codes, exactly ONE of them structural",
          len(LD.TIMIDITY_CODES) == 4 and LD.TIMIDITY_STRUCTURAL == ("TEXT-ONLY DECK",))
    check("type drama alone cannot hold a deck (TIMID COVER/FLAT TYPE are not structural)",
          not ({"TIMID COVER", "FLAT TYPE"} & set(LD.TIMIDITY_STRUCTURAL)))
    check("a deliberately monochrome register cannot be held on quietness alone "
          "(MONOTONE INK is not structural — ink-wash / 留白 / mono spec sheets are protected)",
          "MONOTONE INK" not in LD.TIMIDITY_STRUCTURAL)
    check("timidity_codes() reads warn lines like sameness_codes()",
          LD.timidity_codes(["TEXT-ONLY DECK: 9 of 10 …", "CROWDED: slide 2 …"])
          == ("TEXT-ONLY DECK",))

    print("== the measure is DRAWN AREA, not the ink/text difference ==")
    with tempfile.TemporaryDirectory() as td2:
        import deckkit as dk
        d2 = pathlib.Path(td2)
        prs = dk.blank_deck(13.333, 7.5)
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        # A labelled filled form: exactly the shape that broke the first attempt, because the
        # label sits ON the plate and ink coverage is a union, so `ink_cov - text_cov` collapsed
        # a page whose whole geometry was drawn into the same score as bare paragraphs.
        dk.box(sl, 1.0, 2.0, 11.0, 2.4, fill=dk.RGBColor.from_string("0F6E63"))
        dk.text(sl, 1.4, 2.0, 10.2, 2.4,
                [[("a name written on the form it belongs to", 20,
                   dk.RGBColor.from_string("FFFFFF"), True, False)]])
        out2 = d2 / "form.pptx"
        prs.save(str(out2))
        import lint_deck as _L
        from pptx import Presentation
        pr = Presentation(str(out2))
        sw, sh = pr.slide_width / 914400, pr.slide_height / 914400
        bx = _L._boxes(pr.slides[0], sw, sh, record=False)
        drawn = _L._coverage([b for b in bx if not b.get("bg") and not b.get("text")], sw, sh)
        inkgap = (_L._coverage([b for b in bx if not b.get("bg")], sw, sh)
                  - _L._coverage([b for b in bx if b.get("text") and not b.get("bg")], sw, sh))
        check("a labelled filled form reads as DRAWN (>=8% of the page)", drawn >= 0.08,
              f"drawn={drawn:.3f}")
        check("...while the ink/text difference would have called it text-only",
              inkgap < 0.08, f"ink-text={inkgap:.3f}")

    with tempfile.TemporaryDirectory() as td:
        d = pathlib.Path(td)
        deck = build(d, visual=False, timid_cover=True)
        write_proof(d)

        print("== a text-only, timid-cover deck is HELD ==")
        rc, out = gate(deck, record())
        check("blocks", rc != 0, out)
        check("...names it as a finding about ambition, not error",
              "ambition rather than about error" in out, out)
        check("...lists the fired signals", "TEXT-ONLY DECK" in out, out)
        check("...offers the enrich-first repair", "protagonist" in out, out)

        print("== it stands down where restraint is the declared position ==")
        rc, out = gate(deck, record(boldness="conservative",
                                    signature_move="deliberately restrained: 一份法务条款通读稿，"
                                                   "任何视觉噱头都会削弱条文本身的权威"))
        check("conservative + recorded restraint → not applied", rc == 0, out)
        check("...and says why", "restraint IS the position" in out, out)

        print("== the waiver must NAME a register and match this state of the deck ==")
        rec = record()
        rec["timidity"] = {"waived": "短", "waived_category": "register-restrained", "codes": []}
        rc, out = gate(deck, rec)
        check("a two-character reason is not a named register", rc != 0 and "NAME the register" in out, out)

        rec["timidity"] = {"waived": "这是一份逐条通读的法务条款稿，文字本身就是交付物，"
                                     "加视觉只会削弱条文的权威性。",
                           "waived_category": "not-a-kind", "codes": []}
        rc, out = gate(deck, rec)
        check("the category must be one of the named kinds", rc != 0 and "WHICH kind" in out, out)

        rec["timidity"]["waived_category"] = "text-is-the-artifact"
        rec["timidity"]["codes"] = ["TEXT-ONLY DECK"]      # deck also fires TIMID COVER
        rc, out = gate(deck, rec)
        check("a codes list from another state does not certify this one",
              rc != 0 and "does not certify" in out, out)

        fired = sorted(LD.timidity_codes([]) or [])
        rc, out = gate(deck, rec)                          # still wrong; read the real set below
        import re
        m = re.search(r"fires \[(.*?)\]", out)
        real = [x.strip().strip("'\"") for x in (m.group(1).split(",") if m else [])]
        rec["timidity"]["codes"] = real
        rc, out = gate(deck, rec)
        check("a truthful waiver passes and is printed", rc == 0 and "WAIVED" in out, out)

    print("== a deck with real protagonists is NOT held ==")
    with tempfile.TemporaryDirectory() as td:
        d = pathlib.Path(td)
        deck = build(d, visual=True, timid_cover=False)
        write_proof(d)
        rc, out = gate(deck, record())
        check("charts + a display cover → passes", rc == 0, out)
        check("...and TEXT-ONLY DECK does not fire", "TEXT-ONLY DECK" not in out, out)

    print("== a photo/figure deck (pictures, no charts) is NOT held ==")
    with tempfile.TemporaryDirectory() as td:
        d = pathlib.Path(td)
        deck = build(d, visual=False, timid_cover=True, pics=True)
        write_proof(d)
        rc, out = gate(deck, record())
        check("foreground pictures count as a protagonist", rc == 0 and "TEXT-ONLY DECK" not in out, out)

    print("== a signal that could not run says so, rather than passing quietly ==")
    with tempfile.TemporaryDirectory() as td:
        d = pathlib.Path(td)
        deck = build(d, visual=True, timid_cover=False)
        write_proof(d)
        rc, out = gate(deck, record())
        check("render-dependent code reported as NOT CHECKED without renders",
              "NOT CHECKED" in out and "MONOTONE INK" in out, out)

    print("== the size floor mirrors sameness (calibrated for 8+ content slides) ==")
    with tempfile.TemporaryDirectory() as td:
        d = pathlib.Path(td)
        deck = build(d, visual=False, timid_cover=True, n=5)
        write_proof(d)
        rc, out = gate(deck, record())
        check("a 5-slide deck is not judged for ambition", rc == 0 and "not applied" in out, out)

    print(f"\n{PASS} passed, {FAIL} failed")
    return 1 if FAIL else 0


if __name__ == "__main__":
    raise SystemExit(main())

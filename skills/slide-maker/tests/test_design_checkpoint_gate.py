#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Step 2 (design plan + design checkpoint) must run on EVERY Q1 branch — and be ENFORCED before the
first full render, not merely asked for in prose.

The generated-template (Q1=d) branch once routed interview -> build and let its hero checkpoint (a
LOOK gate) stand in for the Step-2 design checkpoint, so a deck shipped with a `design_plan` that was
reconstructed AFTER the build just to pass the hand-off gate. `render_deck.py` now REFUSES a full
render when a content plan (Step 1) is recorded but no design plan + design checkpoint (Step 2) is —
reading BOTH the shared `.deck-gates.json` and the Codex `.codex-deck-evidence.json`.

🔴 The load-bearing half is what it must stay SILENT on: a `--slides` probe (the Step-2 material
probe / hero-checkpoint sample / signature proof all render this way, legitimately BEFORE the plan is
final), a deck with no content plan (a test fixture / ad-hoc render), and a tiny 1–3 slide ask. A
gate that fires on those would break the very flow that produces the plan.

Run:  python3 tests/test_design_checkpoint_gate.py
"""
import json
import os
import pathlib
import subprocess
import sys
import tempfile

HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import render_deck as rd                                               # noqa: E402

PASS, FAIL = [], []
GATE_MSG = "STEP 2 NOT DONE"


def ok(c, m):
    (PASS if c else FAIL).append(m)


def _content(n):
    return {"content": {"slides": [{"slide": i} for i in range(1, n + 1)]}}


PLAN = {"design_plan": {"boldness": "bold", "signature_move": "x",
                        "checkpoint": {"mode": "approved", "record": "approved by the user"}}}

# ── unit: _content_plan_slide_count ──────────────────────────────────────────
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    ok(rd._content_plan_slide_count(str(dd)) == 0, "no evidence file -> content count 0")
    (dd / ".deck-gates.json").write_text(json.dumps(_content(5)), encoding="utf-8")
    ok(rd._content_plan_slide_count(str(dd)) == 5, "reads content.slides from .deck-gates.json")

with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    (dd / ".codex-deck-evidence.json").write_text(json.dumps(_content(7)), encoding="utf-8")
    ok(rd._content_plan_slide_count(str(dd)) == 7, "reads content.slides from Codex evidence")

with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    (dd / ".deck-gates.json").write_text("{ this is not json", encoding="utf-8")
    ok(rd._content_plan_slide_count(str(dd)) == 0, "malformed evidence -> 0, no crash")

# ── unit: _design_plan_and_checkpoint_present ─────────────────────────────────
def _present(gates=None, codex=None):
    with tempfile.TemporaryDirectory() as d:
        dd = pathlib.Path(d)
        if gates is not None:
            (dd / ".deck-gates.json").write_text(json.dumps(gates), encoding="utf-8")
        if codex is not None:
            (dd / ".codex-deck-evidence.json").write_text(json.dumps(codex), encoding="utf-8")
        return rd._design_plan_and_checkpoint_present(str(dd))


ok(not _present(), "nothing recorded -> not present")
ok(_present(gates={"design_plan": {"boldness": "bold", "checkpoint": {"mode": "approved", "record": "x"}}}),
   "design_plan + design_plan.checkpoint(approved) -> present")
ok(_present(gates={"design_plan": {"boldness": "bold"}, "design": {"checkpoint": {"mode": "auto"}}}),
   "design_plan + top-level design.checkpoint(auto) -> present")
ok(not _present(gates={"design_plan": {"boldness": "bold"}}),
   "design_plan with NO checkpoint -> not present")
ok(not _present(gates={"design_plan": {"checkpoint": {"mode": "pending"}}}),
   "checkpoint mode not approved/auto -> not present")
ok(not _present(gates={"design_plan": {}}), "empty design_plan -> not present")
ok(_present(codex={"design": {"concept": "x", "checkpoint": {"mode": "approved"}}}),
   "Codex design + design.checkpoint(approved) -> present")
ok(not _present(codex={"design": {"concept": "x"}}),
   "Codex design with NO checkpoint -> not present")

# ── integration: the full render is refused / allowed through the gate ────────
def _pptx(path, n):
    from pptx import Presentation
    p = Presentation()
    blank = p.slide_layouts[6]
    for _ in range(n):
        p.slides.add_slide(blank)
    p.save(str(path))


def _run(dd, extra=(), fast_fail=False):
    """Run a render. `fast_fail` points SOFFICE at a bogus path so a run that PASSES the design gate
    dies immediately at LibreOffice instead of doing a real (slow) render — we only assert the gate
    message's presence/absence, never that the render itself succeeded."""
    env = dict(os.environ)
    if fast_fail:
        env["SOFFICE"] = str(dd / "no-such-soffice")
    r = subprocess.run([sys.executable, str(SCRIPTS / "render_deck.py"), "deck.pptx", "render", *extra],
                       capture_output=True, text=True, cwd=str(dd), env=env)
    return r.returncode, (r.stdout or "") + (r.stderr or "")


# (1) content plan, NO design plan -> gate fires, nothing rendered
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    _pptx(dd / "deck.pptx", 4)
    (dd / ".deck-gates.json").write_text(json.dumps(_content(4)), encoding="utf-8")
    rc, out = _run(dd)
    ok(GATE_MSG in out and rc != 0, "full render REFUSED when Step 2 is not recorded")
    ok(not (dd / "render" / "slide01.png").exists(), "…and no pixels were produced")

# (2) content plan + design plan + checkpoint -> gate does NOT fire
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    _pptx(dd / "deck.pptx", 4)
    (dd / ".deck-gates.json").write_text(json.dumps({**_content(4), **PLAN}), encoding="utf-8")
    _, out = _run(dd, fast_fail=True)
    ok(GATE_MSG not in out, "full render ALLOWED once the design plan + checkpoint are recorded")

# (3) a --slides PROBE is exempt even with no design plan (material probe / hero sample)
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    _pptx(dd / "deck.pptx", 4)
    (dd / ".deck-gates.json").write_text(json.dumps(_content(4)), encoding="utf-8")
    _, out = _run(dd, extra=["--slides", "1"], fast_fail=True)
    ok(GATE_MSG not in out, "a --slides probe render is EXEMPT (renders before the plan is final)")

# (4) no content plan at all -> exempt (test fixture / ad-hoc render)
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    _pptx(dd / "deck.pptx", 4)
    _, out = _run(dd, fast_fail=True)
    ok(GATE_MSG not in out, "no content plan -> gate does not apply")

# (5) a 1–3 slide tiny ask is exempt (below the _cp >= 4 threshold)
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    _pptx(dd / "deck.pptx", 3)
    (dd / ".deck-gates.json").write_text(json.dumps(_content(3)), encoding="utf-8")
    _, out = _run(dd, fast_fail=True)
    ok(GATE_MSG not in out, "a 1–3 slide tiny ask is exempt")

# (6) Codex path: content plan + no design -> fires; + design.checkpoint -> allowed
with tempfile.TemporaryDirectory() as d:
    dd = pathlib.Path(d)
    _pptx(dd / "deck.pptx", 4)
    (dd / ".codex-deck-evidence.json").write_text(json.dumps(_content(4)), encoding="utf-8")
    rc, out = _run(dd)
    ok(GATE_MSG in out and rc != 0, "Codex evidence: content plan, no design -> REFUSED")
    (dd / ".codex-deck-evidence.json").write_text(json.dumps(
        {**_content(4), "design": {"concept": "x", "checkpoint": {"mode": "approved"}}}), encoding="utf-8")
    _, out = _run(dd, fast_fail=True)
    ok(GATE_MSG not in out, "Codex evidence: content + design.checkpoint -> ALLOWED")

# ── report ────────────────────────────────────────────────────────────────────
print("\n".join("  ok  " + m for m in PASS))
if FAIL:
    print("\n".join("  XX  " + m for m in FAIL))
print("\n{} passed, {} failed".format(len(PASS), len(FAIL)))
sys.exit(1 if FAIL else 0)

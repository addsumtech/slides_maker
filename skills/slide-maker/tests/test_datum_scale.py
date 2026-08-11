#!/usr/bin/env python3
"""DATUM SCALE — a bar whose length no longer matches the number it claims to show.

Every other geometric check in this skill asks whether a page is READABLE. This one asks whether
it is TRUE, which is the only failure where a fully passing deck still misleads the room: a
truncated baseline turns 1.5 vs 2.1 into 1 : 7 while legibility, overflow and contrast all stay
green.

Two things this suite is careful about.

First, the tolerance is not tuned. Measured on a delivered 12-page deck whose two hand-drawn
charts were correct, the inches-per-unit constant held to 1.0003 and 1.0000; the 2% band has two
orders of magnitude of headroom over a correct build's rounding, and both real charts are
reproduced below as live negative controls.

Second, the first version of this check was wrapped in a bare `except: continue` that swallowed a
NameError, so it silently did nothing while looking like it passed — the same "green because it
stopped looking" failure it exists to prevent. The mid-suite assertion that a KNOWN-BAD deck
still fires is there so that can never recur quietly.
"""
import pathlib, sys, tempfile

HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

ok, bad = [], []
try:
    __import__("pptx")
except ImportError:
    print("skip: python-pptx not installed")
    raise SystemExit(0)

import deckkit as dk                                        # noqa: E402
from pptx import Presentation                               # noqa: E402

dk.FONT = "Helvetica Neue"
TMP = pathlib.Path(tempfile.mkdtemp(prefix="datum-"))
W, H = 10.0, 5.625


def _page():
    prs = dk.blank_deck(W, H)
    s = dk.add_slide(prs)
    dk.slide_background(s, "F5F1E6")
    return prs, s


def _hits(prs):
    return [f for f in dk.lint_layout(prs, verbose=False) if f[2] == "DATUM SCALE"]


# ---------------------------------------------------------------- 1. bar_scale is correct
# The delivered deck's slide 3: a GDP split with one negative component.
prs, s = _page()
VALS3 = [67.98, 17.61, 17.08, -2.68]
sc = dk.bar_scale(5.2, VALS3, group="gdp")
for i, v in enumerate(VALS3):
    sc.bar(s, 2.6, 1.0 + i * 0.5, 0.32, v, fill="1F3B2F")
if not _hits(prs):
    ok.append("a signed horizontal group drawn through bar_scale() is clean (the real deck's "
              "67.98 / 17.61 / 17.08 / -2.68 split)")
else:
    bad.append(f"bar_scale's own output was reported as disproportionate: {_hits(prs)}")

ks = [sc.length(v) / abs(v) for v in VALS3]
if max(ks) / min(ks) < 1.0001:
    ok.append(f"bar_scale is proportional to 4 decimal places (spread {max(ks) / min(ks):.6f})")
else:
    bad.append(f"bar_scale is not proportional: spread {max(ks) / min(ks)}")

# the delivered deck's slide 5 columns, reproduced to the same inch values
prs, s = _page()
sc = dk.bar_scale(1.932, [1.5, 2.1], group="prod")
for i, v in enumerate([1.5, 2.1]):
    sc.bar(s, 2.0 + i * 1.6, 3.92, 1.2, v, vertical=True, fill="8A8377")
got = [round(sc.length(v), 3) for v in (1.5, 2.1)]
if not _hits(prs) and got == [1.38, 1.932]:
    ok.append("vertical columns reproduce the delivered deck's measured 1.380 / 1.932 exactly")
else:
    bad.append(f"vertical columns wrong: lengths {got}, findings {_hits(prs)}")

# ---------------------------------------------------------------- 2. the defects it exists for
prs, s = _page()
K = 1.932 / (2.1 - 1.4)                                     # a baseline at 1.4, not zero
for v in (1.5, 2.1):
    shp = dk.box(s, 2.0 + (v > 1.6) * 1.6, 3.92 - (v - 1.4) * K, 1.2, (v - 1.4) * K, fill="8A8377")
    dk.mark_datum(shp, v, group="prod")
h = _hits(prs)
if len(h) == 1 and h[0][1] == "CRITICAL" and "5.00x" in h[0][3]:
    ok.append("a TRUNCATED BASELINE is CRITICAL and the message names the 5.00x distortion")
else:
    bad.append(f"a truncated baseline was not caught as expected: {h}")

# the mid-suite liveness assertion: if this ever returns nothing, the check has stopped looking
if not h:
    bad.append("LIVENESS: the known-bad deck produced no finding at all — the check is dead, "
               "which is exactly how the first version failed (a swallowed NameError)")

prs, s = _page()                                            # a second scale leaking into a group
for i, (v, k) in enumerate(((10.0, 0.30), (20.0, 0.24))):
    dk.mark_datum(dk.box(s, 2.0, 1.0 + i * 0.5, v * k, 0.32, fill="8A8377"), v, group="mix")
if len(_hits(prs)) == 1:
    ok.append("two different scales inside one group are caught")
else:
    bad.append("a second scale leaking into one group passed")

prs, s = _page()                                            # a decline drawn as growth
for i, v in enumerate((17.61, -2.68)):
    dk.mark_datum(dk.box(s, 2.0, 1.0 + i * 0.5, abs(v) * 0.25, 0.32, fill="8A8377"), v,
                  group="sgn")
h = _hits(prs)
if h and "drawn as growth" in h[0][3]:
    ok.append("a NEGATIVE value drawn from the positive edge is caught — length alone cannot see "
              "this, since the magnitude is still perfectly proportional")
else:
    bad.append(f"a decline drawn as growth passed the length check unchallenged: {h}")

prs, s = _page()                                            # …the same, as columns
for i, v in enumerate((2.1, -0.6)):
    dk.mark_datum(dk.box(s, 2.0 + i * 1.5, 3.9 - abs(v) * 0.9, 1.0, abs(v) * 0.9, fill="8A8377"),
                  v, group="col")
if _hits(prs):
    ok.append("the same sign error is caught on vertical columns")
else:
    bad.append("a negative column drawn upward from the baseline passed")

prs, s = _page()                                            # bar_scale gets both right
sc = dk.bar_scale(4.0, [17.61, -2.68], group="sgn")
for i, v in enumerate([17.61, -2.68]):
    sc.bar(s, 2.0, 1.0 + i * 0.5, 0.32, v, fill="8A8377")
sc2 = dk.bar_scale(2.7, [2.1, -0.6], group="col")
for i, v in enumerate([2.1, -0.6]):
    sc2.bar(s, 6.0 + i * 1.2, 3.9, 0.6, v, vertical=True, fill="8A8377")
if not _hits(prs):
    ok.append("bar_scale places the zero line so signed groups are correct in both orientations")
else:
    bad.append(f"bar_scale produced a sign fault on its own output: {_hits(prs)}")

# ---------------------------------------------------------------- 3. it refuses to guess
prs, s = _page()
for i, v in enumerate((1.0, 2.0)):                          # both dimensions vary
    dk.mark_datum(dk.box(s, 2.0, 1.0 + i * 0.6, v * 0.8, 0.2 + i * 0.3, fill="8A8377"), v,
                  group="amb")
h = _hits(prs)
if h and "no single dimension varies" in h[0][3]:
    ok.append("when NEITHER dimension is constant the encoding axis is unknowable, and it says "
              "so instead of picking one")
else:
    bad.append("an ambiguous group was silently assigned an encoding axis")

prs, s = _page()                                            # one bar is not a proportion claim
dk.mark_datum(dk.box(s, 2.0, 1.0, 3.0, 0.32, fill="8A8377"), 42.0, group="solo")
if not _hits(prs):
    ok.append("a single tagged bar is silent — a ratio needs two bars")
else:
    bad.append("a lone bar was reported")

prs, s = _page()                                            # zero values carry no ratio
sc = dk.bar_scale(4.0, [0.0, 10.0, 20.0], group="z")
for i, v in enumerate([0.0, 10.0, 20.0]):
    sc.bar(s, 2.0, 1.0 + i * 0.5, 0.32, v, fill="8A8377")
if not _hits(prs):
    ok.append("a zero value is excluded from the ratio rather than dividing by it")
else:
    bad.append("a zero-valued bar broke the proportionality test")

prs, s = _page()                                            # untagged bars are UNCHECKED
K = 1.932 / (2.1 - 1.4)
for v in (1.5, 2.1):
    dk.box(s, 2.0 + (v > 1.6) * 1.6, 3.92 - (v - 1.4) * K, 1.2, (v - 1.4) * K, fill="8A8377")
if not _hits(prs):
    ok.append("untagged bars are UNCHECKED, not assumed correct — the check cannot recover a "
              "number the author never wrote down")
else:
    bad.append("an untagged shape was judged against a value it never declared")

prs, s = _page()                                            # a foreign shape borrowing the prefix
shp = dk.box(s, 2.0, 1.0, 3.0, 0.32, fill="8A8377")
shp.name = "deckkit-datum-something-else"
if not _hits(prs):
    ok.append("a shape whose name merely starts with the prefix does not crash the parser")
else:
    bad.append("a malformed datum name produced a finding")

# ---------------------------------------------------------------- 3b. groups and rotation
# `slide.shapes` stops at the group, and composing related marks into one is ordinary practice.
# Measured: a mis-scaled bar pair went unreported the moment it was grouped.
prs, s = _page()
K = 1.932 / (2.1 - 1.4)
for v in (1.5, 2.1):
    dk.mark_datum(dk.box(s, 2.0 + (v > 1.6) * 1.6, 3.92 - (v - 1.4) * K, 1.2, (v - 1.4) * K,
                         fill="8A8377"), v, group="prod")
s.shapes.add_group_shape([x for x in list(s.shapes)][-2:])
if len(_hits(prs)) == 1:
    ok.append("a truncated baseline INSIDE a group is still caught (the walk descends)")
else:
    bad.append("grouping a mis-scaled bar pair hid it from the check")

prs, s = _page()
sc = dk.bar_scale(1.932, [1.5, 2.1], group="ok")
for i, v in enumerate([1.5, 2.1]):
    sc.bar(s, 2.0 + i * 1.6, 3.92, 1.2, v, vertical=True, fill="8A8377")
s.shapes.add_group_shape([x for x in list(s.shapes)][-2:])
if not _hits(prs):
    ok.append("a CORRECT grouped pair stays silent — descending did not break the ratios")
else:
    bad.append("a correct pair inside a group was reported")

prs, s = _page()                                            # a group's children live in the
sc = dk.bar_scale(3.0, [1.0, 2.0], group="x")               # GROUP's coordinate space
sc.bar(s, 2.0, 1.0, 0.3, 1.0, fill="1F3B2F")
b2 = sc.bar(s, 2.0, 1.6, 0.3, 2.0, fill="1F3B2F")
s.shapes.add_group_shape([b2])
h = _hits(prs)
if h and "spans a group boundary" in h[0][3]:
    ok.append("a datum group split across a group boundary REFUSES to compare rather than "
              "comparing lengths measured in different coordinate spaces")
else:
    bad.append(f"bars in different coordinate spaces were compared anyway: {h}")

prs, s = _page()
sc = dk.bar_scale(3.0, [1.0, 2.0], group="r")
sc.bar(s, 2.0, 1.0, 0.3, 1.0, fill="1F3B2F")
b = sc.bar(s, 2.0, 1.6, 0.3, 2.0, fill="1F3B2F")
b.rotation = 90
if not _hits(prs):
    ok.append("a ROTATED bar is skipped — neither dimension is the length the reader sees, so "
              "measuring it would compare the wrong one")
else:
    bad.append("a rotated bar was measured on its unrotated box")

# ---------------------------------------------------------------- 4. survives the round trip
prs, s = _page()
K = 1.932 / (2.1 - 1.4)
for v in (1.5, 2.1):
    dk.mark_datum(dk.box(s, 2.0 + (v > 1.6) * 1.6, 3.92 - (v - 1.4) * K, 1.2, (v - 1.4) * K,
                         fill="8A8377"), v, group="prod")
p = TMP / "rt.pptx"
prs.save(str(p))
if len(_hits(Presentation(str(p)))) == 1:
    ok.append("the tag survives save/reopen, so the check works on a file it did not build")
else:
    bad.append("the datum tag did not survive the round trip")

print("\n".join("  ok   " + x for x in ok))
if bad:
    print("\n".join("  FAIL " + x for x in bad))
print(f"\n{len(ok)} passed, {len(bad)} failed")
raise SystemExit(1 if bad else 0)

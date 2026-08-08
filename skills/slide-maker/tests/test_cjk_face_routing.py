#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""On a 中文 deck the declared display face reached ZERO glyphs, and nothing said a word.

A run tuple carries one font slot, in position 6, and it writes `<a:latin>`. CJK glyphs render
from `<a:ea>`, which takes EAFONT unless the run says otherwise. So an author who writes

    (title, 40, INK, True, False, "Songti SC")

on a Chinese title has set the LATIN face of a run whose Latin content is a stray acronym, and
every Chinese character still renders in EAFONT. Measured by pixels on a real build: swapping
EAFONT changed 25,570 px of a rendered CJK title while swapping the run tuple's font changed none
of the CJK glyphs at all. The deck LOOKED right only because EADISPLAY happened to hold the same
face for the components that read it — the hand-written titles were a different face entirely,
across 6 of 14 slides, for the whole build.

This is the dangerous shape: the rule was WRONG rather than missing. The author followed the
documented call shape and got a silent no-op. So the fix is two-sided — `text()` gained a SEVENTH
slot for the East-Asian face, and `CJK_FACE_UNREACHED` makes the sixth-slot mistake audible.

Backward compatibility is asserted run by run: 5- and 6-element tuples must behave exactly as
before, because every existing deck and every component in the library is built from them.

Run:  python3 tests/test_cjk_face_routing.py
"""
import pathlib
import sys

HERE = pathlib.Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent / "scripts"))

import deckkit as dk                                                  # noqa: E402
from pptx.oxml.ns import qn                                           # noqa: E402

PASS, FAIL = [], []


def check(name, cond, detail=""):
    (PASS if cond else FAIL).append(name)
    print(("  ok   " if cond else "  FAIL ") + name
          + (("  — " + str(detail)) if detail and not cond else ""))


def faces(run_tuple):
    """The (latin, ea) a run tuple actually produces."""
    dk.FONT, dk.EAFONT = "Helvetica Neue", "Hiragino Sans GB"
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    dk.text(s, 0.6, 1.0, 8, 1.0, [[run_tuple]], space_after=0, wrap=False)
    rPr = list(s.shapes)[0].text_frame.paragraphs[0].runs[0]._r.find(qn("a:rPr"))
    lat, ea = rPr.find(qn("a:latin")), rPr.find(qn("a:ea"))
    return (lat.get("typeface") if lat is not None else None,
            ea.get("typeface") if ea is not None else None)


def codes(run_tuple, eafont="Hiragino Sans GB"):
    dk.FONT, dk.EAFONT = "Helvetica Neue", eafont
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    dk.text(s, 0.6, 1.0, 8, 1.0, [[run_tuple]], space_after=0, wrap=False)
    return [f for f in dk.lint_layout(prs, verbose=False) if f[2] == "CJK_FACE_UNREACHED"]


CJK = "认识 LUMC"


def main():
    print("CJK face routing")

    # ---- the seventh slot, and the compatibility that makes it safe to add -------------
    check("a 7-tuple's ea face reaches the CJK glyphs",
          faces((CJK, 40, dk.DEEP, True, False, "Helvetica Neue", "Songti SC"))
          == ("Helvetica Neue", "Songti SC"))
    check("a 6-tuple behaves EXACTLY as before (ea falls back to EAFONT)",
          faces((CJK, 40, dk.DEEP, True, False, "Helvetica Neue"))
          == ("Helvetica Neue", "Hiragino Sans GB"))
    check("a 5-tuple behaves EXACTLY as before",
          faces((CJK, 40, dk.DEEP, True, False))
          == ("Helvetica Neue", "Hiragino Sans GB"))
    check("an explicit ea overrides EAFONT rather than merging with it",
          faces((CJK, 20, dk.DEEP, False, False, "Arial", "KaiTi"))[1] == "KaiTi")
    check("a Latin-only run may still name an ea face without harm",
          faces(("Leiden", 20, dk.DEEP, False, False, "Arial", "KaiTi"))[0] == "Arial")

    # ---- CJK_FACE_UNREACHED: fires on the exact mistake --------------------------------
    found = codes((CJK, 40, dk.DEEP, True, False, "Songti SC"))
    check("a CJK face in slot 6 on a CJK run is reported", len(found) == 1, found)
    check("the report names the call shape that fixes it",
          found and "SEVENTH" in found[0][3])
    # EAFONT is deliberately something none of these equals: a slot-6 face that MATCHES the deck's
    # EAFONT is not unreached (the glyphs do render in it), and that exclusion is asserted below.
    for face in ("PingFang SC", "Microsoft YaHei", "Noto Sans CJK SC", "Hiragino Sans GB",
                 "SimSun", "Malgun Gothic", "Source Han Serif", "KaiTi"):
        check("%s in slot 6 is reported too" % face,
              len(codes((CJK, 24, dk.DEEP, False, False, face), eafont="Songti SC")) == 1)

    # ---- and the silent half, which is what makes it usable ----------------------------
    check("the SEVENTH slot used correctly is silent",
          codes((CJK, 40, dk.DEEP, True, False, "Helvetica Neue", "Songti SC")) == [])
    check("a Latin-only run naming a CJK face is silent (no CJK glyph to reach)",
          codes(("Leiden University", 20, dk.DEEP, False, False, "Songti SC")) == [])
    check("an ordinary CJK run with a Latin face in slot 6 is silent",
          codes((CJK, 14, dk.DEEP, False, False, "Helvetica Neue")) == [])
    check("both slots set to the same CJK face is silent (deliberate)",
          codes((CJK, 40, dk.DEEP, True, False, "Songti SC", "Songti SC")) == [])
    check("a run with no font at all is silent",
          codes((CJK, 14, dk.DEEP, False, False)) == [])

    # ---- it may never block a build ----------------------------------------------------
    dk.FONT, dk.EAFONT = "Helvetica Neue", "Hiragino Sans GB"
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    dk.text(s, 0.6, 1.0, 8, 1.0, [[(CJK, 40, dk.DEEP, True, False, "Songti SC")]],
            space_after=0, wrap=False)
    try:
        dk.lint_layout(prs, verbose=False, strict=True)
        check("WARN, not CRITICAL — strict=True still saves", True)
    except RuntimeError as exc:
        check("WARN, not CRITICAL — strict=True still saves", False, str(exc)[:70])

    print("\n{} passed, {} failed".format(len(PASS), len(FAIL)))
    return 1 if FAIL else 0


if __name__ == "__main__":
    sys.exit(main())

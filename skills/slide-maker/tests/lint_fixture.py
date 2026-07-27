#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Two-sided regression fixture for the lint changes.

PASS deck  — ordinary, correctly-built slides. Every one of these passes the gates TODAY and
             must still pass afterwards. If a change breaks one of these, the change is wrong:
             it is catching craft rather than defects.
FAIL deck  — one slide per defect the changes are supposed to start catching. Each is clean
             today (that is the bug) and must be caught afterwards.

Run before and after; diff the two reports.
"""
import sys, pathlib
SK = pathlib.Path(__file__).resolve().parent.parent / "scripts"
sys.path.insert(0, str(SK))
import deckkit as dk
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C = lambda h: RGBColor.from_string(h)
BG, INK, GREY, ACC = C("FFFFFF"), C("14181C"), C("5A6470"), C("1F6FB2")
FAINT = C("9AA6B4")          # 2.9:1 on white — under the 4.5 body floor
LINE = C("D6DCE2")
OUT = pathlib.Path.cwd()          # write beside the caller, not beside this file
# A fixture that hard-codes a macOS-only face is not a portable test: on Linux the name
# silently falls back to whatever matplotlib finds (DejaVu Sans is ~13.6% wider than
# Helvetica Neue at the same size), so the PASS deck's margins shift under the assertions.
# Pick the first face that genuinely resolves here, and say which one, so a platform-specific
# failure is readable instead of mysterious.
dk.FONT = next((f for f in ("Helvetica Neue", "Helvetica", "Arial", "Liberation Sans",
                            "DejaVu Sans") if dk._font_file(f, False)), "DejaVu Sans")
print("fixture font: %s" % dk.FONT)


def head(s, t):
    dk.text(s, 0.6, 0.5, 8.8, 0.5, [[(t, 22, INK, True, False, dk.FONT)]])


def _panel_png():
    """A four-panel composite whose panels are UNEQUAL in width and separated by page-coloured
    gutters — the shape a real multi-panel figure takes when its panels keep their own aspect
    ratios. Gutters match the page so the render segments them as ground, not as ink."""
    path = OUT / "_fixture_panels.png"
    if not path.exists():
        from PIL import Image, ImageDraw
        im = Image.new("RGB", (1200, 260), (0xFF, 0xFF, 0xFF))
        d = ImageDraw.Draw(im)
        for x0, x1 in ((10, 260), (300, 520), (560, 780), (820, 1040)):
            d.rectangle([x0, 8, x1, 252], fill=(0x22, 0x28, 0x30))
        im.save(path)
    return path


def body(s, y, txt, size=13, color=None):
    dk.text(s, 0.6, y, 6.4, 0.9, [[(txt, size, color or GREY, False, False, dk.FONT)]])


# ───────────────────────────── PASS deck ─────────────────────────────
def build_pass():
    prs = dk.blank_deck()
    L = prs.slide_layouts[6]

    # 1 ordinary slide: rule sits BELOW the text block with a real gap
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "An ordinary content slide")
    body(s, 1.3, "A normal paragraph of body copy at a readable size and a contrast ratio "
                 "well above the floor, sitting where it belongs.")
    dk.box(s, 0.6, 2.6, 8.8, 0.02, fill=LINE, line=None)          # rule BELOW, clear gap
    body(s, 2.85, "A second block under the divider.")
    dk.footer(s, tag="fixture", page=1)

    # 2 an UNDERLINE — the case RULE_THROUGH_TEXT's pad legitimately spares
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "An underline must stay legal")
    dk.text(s, 0.6, 1.4, 4.0, 0.4, [[("Underlined heading", 15, INK, True, False, dk.FONT)]])
    dk.box(s, 0.6, 1.72, 2.1, 0.015, fill=ACC, line=None)          # just below the baseline
    body(s, 2.2, "The rule above is an underline, not a strike-through. It must not be flagged.")
    dk.footer(s, tag="fixture", page=2)

    # 3 text on a solid card, good contrast
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "Text on a card, comfortably above the floor")
    dk.box(s, 0.6, 1.3, 5.0, 1.4, fill=C("EEF3F8"), line=None)
    dk.text(s, 0.8, 1.5, 4.6, 1.0, [[("Dark text on a pale card — roughly 12:1.",
                                      13, INK, False, False, dk.FONT)]])
    dk.footer(s, tag="fixture", page=3)

    # 4-6 a DECLARED rhymed triptych: identical skeleton on purpose
    for i, label in enumerate(["Speed", "Cost", "Risk"]):
        s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
        head(s, f"{label}: the frame repeats on purpose")
        dk.box(s, 0.6, 1.4, 8.8, 1.6, fill=C("F4F7FA"), line=None)
        dk.text(s, 0.9, 1.7, 8.2, 1.0, [[(f"{label} — one variable changes, the frame does not.",
                                          16, INK, True, False, dk.FONT)]])
        body(s, 3.3, "Small multiples at deck scale: identical structure is what makes the "
                     "difference the only thing the eye can see.")
        dk.design_intent(s, rhyme=1, reason="deliberate Speed/Cost/Risk triptych")
        dk.footer(s, tag="fixture", page=4 + i)

    # 7 a DECLARED quiet page — one sentence, lots of air
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    dk.text(s, 0.6, 2.2, 8.0, 1.2,
            [[("This is the pause, and nothing else belongs anywhere on it — the air is the "
               "entire point of the page.", 20, INK, True, False, dk.FONT)]])
    dk.design_intent(s, envelope="quiet", reason="deliberate pause slide before the turn")
    dk.footer(s, tag="fixture", page=7)

    # 8 the SAME composite figure as FAIL slide 11, captioned CORRECTLY — each caption centred
    #   on the panel it names, derived from where the panel actually is. Without this the
    #   alignment check could drift into firing on every captioned figure and only the FAIL
    #   deck would notice, which is the failure mode a one-sided fixture always has.
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "Captions centred on their panels")
    pic = dk.picture(s, str(_panel_png()), 0.6, 2.00, 8.8, 8.8 * 260 / 1200.0, fit="contain",
                     alt="four panels of unequal width")
    px, pw = pic.left / 914400.0, pic.width / 914400.0
    for (a, b), lab in zip(((10, 260), (300, 520), (560, 780), (820, 1040)),
                           ("first", "second", "third", "fourth")):
        c = px + (a + b) / 2.0 / 1200.0 * pw
        dk.text(s, c - 0.75, 3.98, 1.5, 0.26, [[(lab, 11, GREY, True, False, dk.FONT)]],
                align=PP_ALIGN.CENTER, space_after=0)
    dk.footer(s, tag="fixture", page=8)

    dk.lint_layout(prs, strict=True)
    prs.save(str(OUT / "fx_pass.pptx"))
    print("built fx_pass.pptx (8 slides — must stay clean)")


# ───────────────────────────── FAIL deck ─────────────────────────────
def build_fail():
    prs = dk.blank_deck()
    L = prs.slide_layouts[6]

    # 1 RULE THROUGH TEXT — hairline drawn AFTER the text, across its x-height (Run 2b)
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "A hairline painted over the type")
    dk.text(s, 0.6, 1.40, 8.0, 0.34,
            [[("This sentence has a divider drawn straight through it after the fact.",
               11, GREY, False, False, dk.FONT)]])
    dk.box(s, 0.6, 1.485, 8.8, 0.012, fill=LINE, line=None)        # painted AFTER, mid-glyph
    dk.footer(s, tag="fixture", page=1)

    # 2 OCCLUSION — a solid panel drawn after the text, covering it entirely
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "A panel that swallows a sentence")
    dk.text(s, 0.8, 1.6, 5.0, 0.5,
            [[("You cannot read this — a panel is drawn on top of it.",
               14, INK, False, False, dk.FONT)]])
    dk.box(s, 0.7, 1.5, 5.4, 0.8, fill=C("E3E9EF"), line=None)     # painted AFTER, opaque
    dk.footer(s, tag="fixture", page=2)

    # 3 BODY CONTRAST — body copy at ~2.9:1 on a resolvable solid fill
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "Body copy under the contrast floor")
    dk.text(s, 0.6, 1.4, 8.0, 0.9,
            [[("This paragraph is set in a pale grey that does not clear 4.5:1 against the "
               "white it sits on, at body size.", 13, FAINT, False, False, dk.FONT)]])
    dk.footer(s, tag="fixture", page=3)

    # 4-6 an UNDECLARED rhymed run — must still trip LAYOUT SAMENESS
    for i, label in enumerate(["One", "Two", "Three"]):
        s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
        head(s, f"{label}: same skeleton, nothing declared")
        dk.box(s, 0.6, 1.4, 8.8, 1.6, fill=C("F4F7FA"), line=None)
        dk.text(s, 0.9, 1.7, 8.2, 1.0, [[(f"{label} — and no rhyme was declared.",
                                          16, INK, True, False, dk.FONT)]])
        body(s, 3.3, "Three slides sharing one skeleton with no declared intent is sameness, "
                     "not rhythm, and should still be flagged.")
        dk.footer(s, tag="fixture", page=4 + i)

    # 7 an UNDECLARED thin page — must still trip UNDERFILLED
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "A thin page nobody declared")
    body(s, 1.4, "Sixteen or so words of body copy and nothing else at all on the whole canvas.")
    dk.footer(s, tag="fixture", page=7)

    # 8-10 the PAINT-ORDER faults that a per-shape threshold cannot see. All three shipped on
    # real decks with the gate reporting clean, and each was found by a human looking at a PNG.
    # The shared cause: the old test asked whether ONE shape covered >=60% of a text block.

    # 8 a field of 150 tiles over a caption — no single tile covers 1% of the text
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "Composite coverage")
    dk.text(s, 0.6, 2.60, 8.8, 0.30,
            [[("COMPOSITE TILES MUST NOT HIDE THIS CAPTION", 10, GREY, True, False, dk.FONT)]],
            space_after=0)
    cell, gp = 0.352, 0.05
    for r in range(6):
        for c in range(25):
            dk.box(s, c * (cell + gp), 1.60 + r * (cell + gp), cell, cell, fill=LINE, line=None)
    dk.footer(s, tag="fixture", page=8)

    # 9 a DASHED rule — 40 separate boxes — struck through a footnote
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "A dashed rule is still a rule")
    dk.text(s, 0.6, 2.40, 8.8, 0.44,
            [[("a dashed rule of many boxes must be caught exactly like a solid one, because "
               "the reader cannot tell them apart", 10, GREY, False, False, dk.FONT)]],
            space_after=0)
    seg, g2 = 0.10, 0.07
    for k in range(int(8.8 // (seg + g2))):
        dk.box(s, 0.6 + k * (seg + g2), 2.47, seg, 0.012, fill=LINE, line=None)
    dk.footer(s, tag="fixture", page=9)

    # 10 an opaque PICTURE over the first line only — unknowable from the XML, since nothing
    #    in the file records a picture's alpha. Only the pixels can answer this one.
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "A picture over one line")
    dk.text(s, 0.6, 2.30, 5.0, 0.34,
            [[("an opaque picture over line one of this sentence hides it completely",
               13, GREY, False, False, dk.FONT)]], space_after=0)
    plate = OUT / "_fixture_plate.png"
    if not plate.exists():
        from PIL import Image
        Image.new("RGB", (600, 90), (0x9A, 0xA6, 0xB4)).save(plate)
    dk.picture(s, str(plate), 0.6, 2.26, 5.0, 0.30, fit="cover", alt="opaque plate")
    dk.footer(s, tag="fixture", page=10)

    # 11 captions laid out on the TEXT grid under a composite figure whose panels are not on
    #    that grid. Nothing overlaps and nothing overflows, so every geometric rule is silent;
    #    only the pixels know where the panels actually landed. Panel widths differ here on
    #    purpose — that is what makes "column width over four" wrong in the first place.
    s = prs.slides.add_slide(L); dk.box(s, 0, 0, 10, 5.625, fill=BG, line=None)
    head(s, "Captions off the panels they name")
    fig = _panel_png()
    dk.picture(s, str(fig), 0.6, 2.00, 8.8, 8.8 * 260 / 1200.0, fit="contain",
               alt="four panels of unequal width")
    for i, lab in enumerate(("first", "second", "third", "fourth")):
        dk.text(s, 0.6 + i * 2.2, 3.98, 2.2, 0.26,
                [[(lab, 11, GREY, True, False, dk.FONT)]], align=PP_ALIGN.CENTER, space_after=0)
    dk.footer(s, tag="fixture", page=11)

    dk.lint_layout(prs, strict=False)     # this deck is SUPPOSED to be defective
    prs.save(str(OUT / "fx_fail.pptx"))
    print("built fx_fail.pptx (11 slides — each carries one defect)")


if __name__ == "__main__":
    build_pass()
    build_fail()

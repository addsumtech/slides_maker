#!/usr/bin/env python3
"""slide_background() — a REAL <p:bg> backdrop, and the linter's ability to still see it.

Moving the backdrop out of the shape tree is a one-line change with a silent failure mode:
`_backing_fill` documented "Slide bg unknown -> None", and its callers treat None as "skip this
check". So a deck painted with <p:bg> would have passed contrast checking by NOT BEING CHECKED —
green gates, unreadable slides, no signal anywhere. The load-bearing test in this file is
therefore not "does <p:bg> exist" but the EQUIVALENCE property: a <p:bg> deck and the
full-canvas-rect deck it replaces must produce the same findings, including the failures.
"""
import pathlib, sys, tempfile

HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

ok, bad = [], []


def _need():
    for m in ("pptx",):
        try:
            __import__(m)
        except ImportError:
            print("skip: python-pptx not installed")
            raise SystemExit(0)


_need()

import deckkit as dk                                        # noqa: E402
import lint_deck as L                                       # noqa: E402
from pptx import Presentation                               # noqa: E402
from pptx.oxml.ns import qn                                 # noqa: E402
from pptx.dml.color import RGBColor                          # noqa: E402


def _c(h):
    """run tuples take RGBColor, never a hex string (one of the three call-shape contracts)."""
    return RGBColor.from_string(h)

dk.FONT = "Helvetica Neue"
TMP = pathlib.Path(tempfile.mkdtemp(prefix="bgtest-"))
W, H = 10.0, 5.625
PAPER, INK = "F5F1E6", "1A1A1A"


def _save(prs, name):
    p = TMP / name
    prs.save(str(p))
    return p


def _content(s, tcolor):
    """The same content on either backdrop: a title on the bare canvas, a filled card, a run of
    body text. Every element whose contrast is judged against the BACKGROUND, not a card."""
    dk.text(s, 0.6, 0.5, 8.8, 0.8, [[("A title on the bare canvas", 30, _c(tcolor), True, False)]])
    dk.box(s, 0.6, 1.8, 4.0, 2.0, fill="2B4C7E")
    dk.text(s, 0.9, 2.2, 3.4, 0.6, [[("on a card", 16, _c("FFFFFF"), False, False)]])
    dk.text(s, 5.2, 1.8, 4.2, 1.2, [[("body text, also on the bare canvas", 14, _c(tcolor), False, False)]])


def _deck(bg_mode, tcolor, name):
    prs = dk.blank_deck(W, H)
    s = dk.add_slide(prs)
    if bg_mode == "rect":
        dk.box(s, 0, 0, W, H, fill=PAPER)
    else:
        dk.slide_background(s, PAPER)
    _content(s, tcolor)
    return _save(prs, name)


# ---------------------------------------------------------------- 1. the XML is what it claims
p = _deck("bg", INK, "basic.pptx")
sl = Presentation(str(p)).slides[0]
cSld = sl._element.find(qn("p:cSld"))
kids = [e.tag.split("}")[-1] for e in cSld]

if kids and kids[0] == "bg":
    ok.append("<p:bg> is the FIRST child of <p:cSld> (CT_CommonSlideData orders bg before spTree)")
else:
    bad.append(f"<p:bg> is not first in <p:cSld>: {kids} — PowerPoint will offer to repair the file")

if kids.count("bg") == 1:
    ok.append("exactly one <p:bg> element")
else:
    bad.append(f"{kids.count('bg')} <p:bg> elements — the schema allows one")

vals = [e.get("val") for e in cSld.find(qn("p:bg")).findall(".//" + qn("a:srgbClr"))]
if vals == [PAPER]:
    ok.append("the background colour round-trips through save/reopen")
else:
    bad.append(f"background colour did not round-trip: {vals!r} != ['{PAPER}']")

# ---------------------------------------------------------------- 2. it is NOT a selectable shape
n_bg = len(list(sl.shapes))
r = Presentation(str(_deck("rect", INK, "rect.pptx"))).slides[0]
n_rect = len(list(r.shapes))
if n_bg == n_rect - 1:
    ok.append(f"the backdrop left the shape tree ({n_rect} shapes with a rect -> {n_bg} with p:bg): "
              f"it can no longer be selected, dragged or deleted by accident")
else:
    bad.append(f"shape counts say the backdrop is still selectable: p:bg={n_bg}, rect={n_rect}")

# ---------------------------------------------------------------- 3. the linter still sees it
bx = L._boxes(sl, W, H)
bgrec = [b for b in bx if b["bg"]]
if len(bgrec) == 1 and bgrec[0]["fill"] == PAPER:
    ok.append("_boxes synthesises ONE bg record carrying the real colour")
else:
    bad.append(f"_boxes did not synthesise a usable bg record: {[(b['bg'], b['fill']) for b in bx]}")

if bx and bx[0]["bg"]:
    ok.append("the bg record sits at index 0 — below every real shape, which is what the "
              "index-walking _backing_fill assumes")
else:
    bad.append("the bg record is not first in bx; _backing_fill walks range(ti) and would miss it")

if bgrec and bgrec[0]["w"] * bgrec[0]["h"] >= 0.95 * W * H:
    ok.append("the bg record covers the canvas, so every `not s['bg']` density filter drops it")
else:
    bad.append("the bg record is not canvas-sized; ink-coverage denominators will shift")

# _backing_fill on the title (index 1 = first real shape) must resolve to the background
ti = next(i for i, b in enumerate(bx) if b["text"] and "bare canvas" in b["full"])
got = L._backing_fill(bx, ti)
if got == PAPER:
    ok.append("_backing_fill resolves text on the bare canvas to the p:bg colour (not None) — "
              "the 62 contrast call sites stay live")
else:
    bad.append(f"_backing_fill returned {got!r} for text on the bare canvas; None means every "
               f"contrast check SKIPS and the deck passes by not being looked at")

# ---------------------------------------------------------------- 4. THE equivalence property
def _findings(path):
    prs = Presentation(str(path))
    return sorted((f[1], f[2]) for f in dk.lint_layout(prs, verbose=False))


for label, tcol in (("readable ink", INK), ("unreadable pale grey", "EDE8DC")):
    a = _findings(_deck("rect", tcol, f"eq_rect_{tcol}.pptx"))
    b = _findings(_deck("bg", tcol, f"eq_bg_{tcol}.pptx"))
    if a == b:
        ok.append(f"identical build-time findings on rect vs p:bg ({label}): {len(a)} finding(s)")
    else:
        bad.append(f"findings DIVERGE between rect and p:bg ({label}): rect={a} p:bg={b}")

# the pale-grey deck must actually FAIL somewhere — an equivalence of two silences proves nothing
pale_bx = L._boxes(Presentation(str(TMP / "eq_bg_EDE8DC.pptx")).slides[0], W, H)
pti = next(i for i, b in enumerate(pale_bx) if b["text"] and "bare canvas" in b["full"])
pback = L._backing_fill(pale_bx, pti)
if pback == PAPER:
    ok.append("pale-grey-on-paper still resolves a backing colour, so the contrast maths runs "
              "on it (equivalence above is not two silences agreeing)")
else:
    bad.append(f"pale-grey text resolved backing={pback!r}; the negative control is not live")

# ---------------------------------------------------------------- 5. idempotent + input forms
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
dk.slide_background(s, "112233")
dk.slide_background(s, "AABBCC")
p2 = _save(prs, "twice.pptx")
c2 = Presentation(str(p2)).slides[0]._element.find(qn("p:cSld"))
k2 = [e.tag.split("}")[-1] for e in c2]
v2 = [e.get("val") for e in c2.find(qn("p:bg")).findall(".//" + qn("a:srgbClr"))]
if k2.count("bg") == 1 and v2 == ["AABBCC"]:
    ok.append("calling it twice REPAINTS rather than stacking a second <p:bg>")
else:
    bad.append(f"a second call stacked or failed to repaint: kids={k2} vals={v2}")

prs = dk.blank_deck(W, H)
forms = {}
for i, val in enumerate(("#f5f1e6", "f5f1e6", RGBColor.from_string("F5F1E6"))):
    sx = dk.add_slide(prs)
    dk.slide_background(sx, val)
p3 = _save(prs, "forms.pptx")
for i, sx in enumerate(Presentation(str(p3)).slides):
    cs = sx._element.find(qn("p:cSld")).find(qn("p:bg"))
    forms[i] = cs.findall(".//" + qn("a:srgbClr"))[0].get("val")
if set(forms.values()) == {PAPER}:
    ok.append("'#hex', 'hex' and RGBColor all normalise to the same value")
else:
    bad.append(f"colour input forms disagree: {forms}")

# ---------------------------------------------------------------- 6. schema gate covers p:cSld
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
dk.slide_background(s, PAPER)
dk.text(s, 1, 1, 4, 0.5, [[("x", 14, _c(INK), False, False)]])
crit = [f for f in dk.lint_layout(prs, verbose=False) if f[1] == "CRITICAL"]
if not crit:
    ok.append("a p:bg deck is CRITICAL-clean (the new cSld order check does not false-fire)")
else:
    bad.append(f"p:bg deck reports CRITICAL findings: {crit}")

# deliberately break the order: append <p:bg> AFTER spTree, the mistake the check exists for
from pptx.oxml import parse_xml                              # noqa: E402
from pptx.oxml.ns import nsdecls                             # noqa: E402
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
cs = s._element.find(qn("p:cSld"))
cs.append(parse_xml('<p:bg %s><p:bgPr><a:solidFill><a:srgbClr val="112233"/></a:solidFill>'
                    '<a:effectLst/></p:bgPr></p:bg>' % nsdecls("p", "a")))
f = [x for x in dk.lint_layout(prs, verbose=False) if x[2] == "OOXML_SHAPE"]
if f and f[0][1] == "CRITICAL":
    ok.append("OOXML_SHAPE fires CRITICAL when <p:bg> is appended after <p:spTree>")
else:
    bad.append("appending <p:bg> after <p:spTree> was NOT caught — the file PowerPoint offers "
               "to repair passes the gate")

# two <p:bg> elements
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
cs = s._element.find(qn("p:cSld"))
for _ in range(2):
    cs.insert(0, parse_xml('<p:bg %s><p:bgPr><a:solidFill><a:srgbClr val="112233"/></a:solidFill>'
                           '<a:effectLst/></p:bgPr></p:bg>' % nsdecls("p", "a")))
f = [x for x in dk.lint_layout(prs, verbose=False) if x[2] == "OOXML_SHAPE"]
if f and f[0][1] == "CRITICAL":
    ok.append("OOXML_SHAPE fires CRITICAL on a duplicated <p:bg>")
else:
    bad.append("two <p:bg> elements passed the schema gate")

# ---------------------------------------------------------------- 7. a dark p:bg reads as dark
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
dk.slide_background(s, "0E1B2A")
dk.text(s, 1, 1, 6, 0.6, [[("light on dark", 20, _c("FFFFFF"), True, False)]])
p4 = _save(prs, "dark.pptx")
dbx = L._boxes(Presentation(str(p4)).slides[0], W, H)
dark_plate = any(b["bg"] and b["fill"] and L._lum(b["fill"]) < 0.45 for b in dbx)
if dark_plate:
    ok.append("a dark p:bg is detected as a dark plate (canvas-flip and one-off checks keep working)")
else:
    bad.append("a dark p:bg is not seen as a dark plate; ONE-OFF CANVAS FLIP goes blind")

# ---------------------------------------------------------------- 8. no bg at all -> no record
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
dk.text(s, 1, 1, 4, 0.5, [[("bare", 14, _c(INK), False, False)]])
p5 = _save(prs, "none.pptx")
nbx = L._boxes(Presentation(str(p5)).slides[0], W, H)
if not any(b["bg"] for b in nbx):
    ok.append("a slide with no <p:bg> gets no synthetic record (unknown backdrop stays unknown)")
else:
    bad.append("a slide without <p:bg> was given a background record out of thin air")

# a theme/gradient background must be UNKNOWN, never a guess
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
s._element.find(qn("p:cSld")).insert(0, parse_xml(
    '<p:bg %s><p:bgPr><a:gradFill><a:gsLst>'
    '<a:gs pos="0"><a:srgbClr val="112233"/></a:gs>'
    '<a:gs pos="100000"><a:srgbClr val="AABBCC"/></a:gs>'
    '</a:gsLst><a:lin ang="0"/></a:gradFill><a:effectLst/></p:bgPr></p:bg>' % nsdecls("p", "a")))
p6 = _save(prs, "grad.pptx")
gbx = L._boxes(Presentation(str(p6)).slides[0], W, H)
g = [b for b in gbx if b["bg"]]
if g and g[0]["fill"] is None and g[0]["unk"]:
    ok.append("a gradient <p:bg> is recorded UNKNOWN, not guessed at one of its stops")
else:
    bad.append(f"a gradient background was resolved to a definite colour: {g}")

# ---------------------------------------------------------------- 9. inheritance + robustness
# Backgrounds resolve slide -> layout -> master, the way OOXML does. Reading only the slide meant
# a deck built on the USER'S TEMPLATE — which states its canvas once on the master, not per page —
# resolved every run's backing to None, and None tells all 62 contrast callers to skip. Contrast
# checking was silently off for the whole template branch.
SOLID = ('<p:bg %s><p:bgPr><a:solidFill><a:srgbClr val="0E1B2A"/></a:solidFill>'
         '<a:effectLst/></p:bgPr></p:bg>' % nsdecls("p", "a"))


def _at(level, xml):
    prs = dk.blank_deck(W, H)
    s = dk.add_slide(prs)
    part = {"slide": s, "layout": s.slide_layout, "master": prs.slide_masters[0]}[level]
    part._element.find(qn("p:cSld")).insert(0, parse_xml(xml))
    dk.text(s, 1, 2, 6, 0.6, [[("a line of text", 18, _c("14243A"), False, False)]])
    out = _save(prs, "inh_%s.pptx" % level)
    bxx = L._boxes(Presentation(str(out)).slides[0], W, H)
    ti2 = next(i for i, b in enumerate(bxx) if b["text"])
    return L._backing_fill(bxx, ti2), bxx


for lvl in ("slide", "layout", "master"):
    back, bxx = _at(lvl, SOLID)
    dark = any(b["bg"] and b["fill"] and L._lum(b["fill"]) < 0.45 for b in bxx)
    if back == "0E1B2A" and dark:
        ok.append(f"a solid background declared on the {lvl} resolves, and reads as a dark plate")
    else:
        bad.append(f"{lvl}-level background did not resolve: backing={back!r} dark={dark}")

# …but an INHERITED background we cannot resolve must NOT become a record. python-pptx's own
# default master carries <p:bgRef idx="1001"><a:schemeClr/>; synthesising an `unk` record for it
# would hand every deck an unknowable plate and flip `unk_plate` deck-wide.
prs = dk.blank_deck(W, H)
s = dk.add_slide(prs)
dk.text(s, 1, 2, 6, 0.6, [[("x", 14, _c(INK), False, False)]])
pth = _save(prs, "themeref.pptx")
bx2 = L._boxes(Presentation(str(pth)).slides[0], W, H)
if not any(b["bg"] for b in bx2):
    ok.append("an inherited THEME-reference background stays unknown — no record invented, so "
              "`unk_plate` does not become true on every deck ever built")
else:
    bad.append("a theme-reference master background was turned into a background record")

# a gradient the SLIDE declares is still a record (a fact about that page), just unresolvable
GRAD = ('<p:bg %s><p:bgPr><a:gradFill><a:gsLst><a:gs pos="0"><a:srgbClr val="112233"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="AABBCC"/></a:gs></a:gsLst><a:lin ang="0"/>'
        '</a:gradFill><a:effectLst/></p:bgPr></p:bg>') % nsdecls("p", "a")
_b, bx3 = _at("slide", GRAD)
if _b == "UNKNOWN" and any(x["bg"] for x in bx3):
    ok.append("a gradient the SLIDE declares keeps its record and reports UNKNOWN")
else:
    bad.append(f"a slide-declared gradient lost its record: backing={_b!r}")

_b, _ = _at("master", GRAD)
if _b is None:
    ok.append("an INHERITED gradient stays None — inheritance only ever adds resolved colour")
else:
    bad.append(f"an inherited gradient produced {_b!r} instead of staying unknown")

# the colour must come from inside <a:solidFill>, not the first srgbClr anywhere under <p:bg>:
# a legal <a:effectLst> declared first would otherwise be read as the page colour
EFFECT_FIRST = ('<p:bg %s><p:bgPr>'
                '<a:effectLst><a:glow><a:srgbClr val="FF0000"/></a:glow></a:effectLst>'
                '<a:solidFill><a:srgbClr val="F5F1E6"/></a:solidFill></p:bgPr></p:bg>'
                % nsdecls("p", "a"))
_b, _ = _at("slide", EFFECT_FIRST)
if _b == PAPER:
    ok.append("a glow colour declared before the fill is not mistaken for the page colour")
else:
    bad.append(f"an effect colour was read as the background: {_b!r}")

# every canvas format, not just 16:9
for label, (cw, ch) in (("4:3", (10.0, 7.5)), ("story 9:16", (5.625, 10.0)),
                        ("A4 portrait", (8.27, 11.69))):
    prs = dk.blank_deck(cw, ch)
    s = dk.add_slide(prs)
    dk.slide_background(s, PAPER)
    dk.text(s, 0.6, 0.5, cw - 1.2, 0.6, [[("a title", 22, _c(INK), True, False)]])
    pth = _save(prs, "fmt_%s.pptx" % label.replace(":", "x").replace(" ", "_"))
    bxf = L._boxes(Presentation(str(pth)).slides[0], cw, ch)
    rec = [b for b in bxf if b["bg"]]
    tif = next(i for i, b in enumerate(bxf) if b["text"])
    if len(rec) == 1 and rec[0]["w"] * rec[0]["h"] >= 0.95 * cw * ch \
            and L._backing_fill(bxf, tif) == PAPER:
        ok.append(f"{label}: the background covers the canvas and resolves")
    else:
        bad.append(f"{label}: background record wrong on a non-16:9 canvas")

# ---------------------------------------------------------------- report
print("\n".join("  ok   " + x for x in ok))
if bad:
    print("\n".join("  FAIL " + x for x in bad))
print(f"\n{len(ok)} passed, {len(bad)} failed")
raise SystemExit(1 if bad else 0)

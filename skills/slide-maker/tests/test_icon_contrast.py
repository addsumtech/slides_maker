#!/usr/bin/env python3
"""ICON CONTRAST — the WCAG 1.4.11 floor applied to recolored monochrome icons.

The non-text contrast check skips `pic`, and this skill's icons ARE pictures, so icon contrast
was unchecked on every deck ever built here while icons are default-on for category-rich content.
Measured on a real 12-page deck: a gold icon at 2.69:1 passed every DETERMINISTIC gate and was
caught only by the design critic — a model that costs a dispatch and might not look next time.
Arithmetic with one right answer belongs in a check that runs every time, for free.

A contrast check is easy; a contrast check that does not cry wolf is the work. So most of this
file is NEGATIVE: a photograph, a two-colour illustration, an opaque PNG and an unresolvable
backdrop must all produce SILENCE, because a confident wrong finding about a colour the image
does not actually have is worse than the miss it replaces.
"""
import io, pathlib, sys, tempfile

HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

ok, bad = [], []

for _m in ("pptx", "PIL"):
    try:
        __import__(_m)
    except ImportError:
        print(f"skip: {_m} not installed")
        raise SystemExit(0)

import deckkit as dk                                        # noqa: E402
import lint_deck as L                                       # noqa: E402
from pptx import Presentation                               # noqa: E402
from pptx.dml.color import RGBColor                         # noqa: E402
from PIL import Image                                       # noqa: E402

dk.FONT = "Helvetica Neue"
TMP = pathlib.Path(tempfile.mkdtemp(prefix="iconc-"))
W, H = 10.0, 5.625
PAPER = "F5F1E6"                                            # the real deck's paper
GOLD = "C08A2E"                                             # the real deck's gold: 2.69:1 on paper
INKG = "1F3B2F"                                             # its ink green: 10.79:1


def _blob(im):
    b = io.BytesIO()
    im.save(b, "PNG")
    return b.getvalue()


def _glyph(hexv, size=192, coverage=0.30):
    """A monochrome glyph on transparency — what icons.py's recolor+rasterise produces."""
    r, g, b = int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16)
    im = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    px = im.load()
    n = int(size * size * coverage)
    for i in range(n):                                      # a contiguous block: shape is irrelevant
        px[i % size, i // size] = (r, g, b, 255)
    return im


def _photo(size=192):
    """An opaque photograph: no alpha, and no colour owning its pixels."""
    im = Image.new("RGB", (size, size))
    px = im.load()
    for y in range(size):
        for x in range(size):
            px[x, y] = ((x * 7) % 256, (y * 11) % 256, ((x + y) * 13) % 256)
    return im


def _two_colour(size=192):
    """A two-colour illustration — neither colour is 'the' ink."""
    im = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    px = im.load()
    for y in range(size):
        for x in range(size // 2):
            px[x, y] = (0xC0, 0x8A, 0x2E, 255)
    for y in range(size):
        for x in range(size // 2, size):
            px[x, y] = (0x1F, 0x3B, 0x2F, 255)
    return im


# ---------------------------------------------------------------- 1. _icon_ink classification
cases = [
    ("monochrome gold icon", _blob(_glyph(GOLD)), GOLD, True),
    ("monochrome ink icon", _blob(_glyph(INKG)), INKG, True),
    ("sparse icon (5% ink)", _blob(_glyph(GOLD, coverage=0.05)), GOLD, True),
]
for name, blob, want, should in cases:
    got = L._icon_ink(blob)
    if should and got and got[0] == want:
        ok.append(f"{name}: ink read as #{got[0]} (purity {got[1]:.0%}, opaque {got[2]:.0%})")
    else:
        bad.append(f"{name}: expected ink #{want}, got {got!r}")

neg = [
    ("an opaque photograph", _blob(_photo())),
    ("a two-colour illustration", _blob(_two_colour())),
    ("a fully opaque RGBA image", _blob(Image.new("RGBA", (64, 64), (200, 30, 30, 255)))),
    ("a fully transparent image", _blob(Image.new("RGBA", (64, 64), (0, 0, 0, 0)))),
    ("not an image at all", b"not-an-image"),
]
for name, blob in neg:
    got = L._icon_ink(blob)
    if got is None:
        ok.append(f"{name}: correctly NOT given an ink colour")
    else:
        bad.append(f"{name}: was assigned ink {got!r} — a confident wrong finding")

# a JPEG photo must be rejected from the HEADER, without decoding pixels
jb = io.BytesIO()
_photo().save(jb, "JPEG", quality=85)
if L._icon_ink(jb.getvalue()) is None:
    ok.append("a JPEG is rejected on mode alone (no alpha channel => cannot be an icon)")
else:
    bad.append("a JPEG was treated as an icon")

# purity threshold: 80% one colour is an icon, 60% is not
mixed = _glyph(GOLD, coverage=0.30)
mp = mixed.load()
n_flip = int(192 * 192 * 0.30 * 0.35)                       # flip 35% of the ink to another hue
for i in range(n_flip):
    mp[i % 192, i // 192] = (0x1F, 0x3B, 0x2F, 255)
if L._icon_ink(_blob(mixed)) is None:
    ok.append("a 65/35 two-hue mark stays below the purity floor and is skipped")
else:
    bad.append("a 65/35 two-hue mark was given a single ink colour")


# ---------------------------------------------------------------- 2. the check on a real deck
def _deck(icon_hex, back=PAPER, name="d.pptx", backing="bg"):
    prs = dk.blank_deck(W, H)
    s = dk.add_slide(prs)
    if backing == "bg":
        dk.slide_background(s, back)
    elif backing == "picture":                              # unresolvable backdrop
        ph = TMP / "photo.png"
        ph.write_bytes(_blob(_photo(400)))
        dk.picture(s, str(ph), 0, 0, W, H)
    p = TMP / f"icon_{icon_hex}.png"
    p.write_bytes(_blob(_glyph(icon_hex)))
    dk.icon(s, str(p), 1.0, 2.0, 0.8, alt="a category icon")
    dk.text(s, 2.4, 2.0, 6.0, 0.6,
            [[("a label", 14, RGBColor.from_string("1A1A1A"), False, False)]])
    out = TMP / name
    prs.save(str(out))
    return out


def _warns(path, renders=None):
    import contextlib
    buf = io.StringIO()
    with contextlib.redirect_stdout(buf), contextlib.redirect_stderr(buf):
        try:
            L.lint(str(path), renders_dir=renders, static_ok=True)
        except SystemExit:
            pass
    return buf.getvalue()


out = _warns(_deck(GOLD, name="gold.pptx"))
hits = [l for l in out.splitlines() if "ICON CONTRAST" in l]
if len(hits) == 1 and "2.69:1" in hits[0] and GOLD in hits[0].upper():
    ok.append("gold-on-paper reproduces the real deck's finding EXACTLY: 2.69:1")
else:
    bad.append(f"gold-on-paper did not reproduce 2.69:1 — got {hits!r}")

out = _warns(_deck(INKG, name="ink.pptx"))
if not [l for l in out.splitlines() if "ICON CONTRAST" in l]:
    ok.append("ink-green-on-paper (10.79:1) is silent — the check does not fire on good icons")
else:
    bad.append("a 10.79:1 icon was reported as low contrast")

out = _warns(_deck(GOLD, name="onpic.pptx", backing="picture"))
if not [l for l in out.splitlines() if "ICON CONTRAST" in l]:
    ok.append("an icon over a PHOTO stays silent — an unresolvable backdrop is not guessed at")
else:
    bad.append("an icon over a photo was judged against a backdrop the linter cannot know")


# one warn per colour pair per slide, however many icons repeat it
def _multi(n, name):
    prs = dk.blank_deck(W, H)
    s = dk.add_slide(prs)
    dk.slide_background(s, PAPER)
    p = TMP / f"g_{GOLD}.png"
    p.write_bytes(_blob(_glyph(GOLD)))
    for i in range(n):
        dk.icon(s, str(p), 0.8 + i * 1.1, 2.0, 0.8, alt=f"icon {i}")
    out = TMP / name
    prs.save(str(out))
    return out


out = _warns(_multi(5, "five.pptx"))
n = len([l for l in out.splitlines() if "ICON CONTRAST" in l])
if n == 1:
    ok.append("five icons sharing one bad colour pair produce ONE warning, not five")
else:
    bad.append(f"colour-pair dedup failed: {n} warnings for one repeated pair")

# ---------------------------------------------------------------- 2b. page devices are not icons
# A watermark or backdrop motif is DELIBERATELY faint — that is the whole point — and by every
# pixel test it IS a monochrome icon. Measured: a pale full-page motif PNG was reported for
# failing a floor it was never meant to meet. Both helpers tag their shape, so the tag is read
# rather than guessed; the area cap catches an untagged page device.
def _faint_png(name, w_in, h_in, tag=None):
    im = Image.new("RGBA", (400, 400), (0, 0, 0, 0))
    px = im.load()
    for y in range(0, 400, 3):
        for x in range(400):
            px[x, y] = (0xEA, 0xE4, 0xD6, 255)              # barely there on paper
    p = TMP / f"{name}.png"
    p.write_bytes(_blob(im))
    prs = dk.blank_deck(W, H)
    s = dk.add_slide(prs)
    dk.slide_background(s, PAPER)
    shp = dk.picture(s, str(p), 0.2, 0.2, w_in, h_in)
    if tag:
        try:
            shp.name = tag
        except Exception:
            pass
    dk.text(s, 1, 2, 6, 0.6, [[("body", 15, RGBColor.from_string("1F3B2F"), False, False)]])
    out = TMP / f"{name}.pptx"
    prs.save(str(out))
    return out


for nm, wi, hi, tag, want, why in (
        ("wm_full", 9.6, 5.2, None, 0, "a page-scale faint image is a backdrop, not an icon"),
        ("wm_tagged", 1.0, 1.0, "deckkit-watermark", 0, "a tagged watermark is exempt at any size"),
        ("wm_motif", 1.0, 1.0, "deckkit-motif-quiet", 0, "a tagged motif is exempt at any size")):
    hits = [l for l in _warns(_faint_png(nm, wi, hi, tag)).splitlines() if "ICON CONTRAST" in l]
    if len(hits) == want:
        ok.append(why)
    else:
        bad.append(f"{nm}: expected {want} finding(s), got {len(hits)}")

# ---------------------------------------------------------------- 3. it is a WARN, not a blocker
out = _warns(_deck(GOLD, name="sev.pptx"))
if "[warn]" in out and "ICON CONTRAST" in out:
    line = [l for l in out.splitlines() if "ICON CONTRAST" in l][0]
    if "[warn]" in line:
        ok.append("ICON CONTRAST is a WARN — decorative-vs-essential is a judgment lint can't make")
    else:
        bad.append("ICON CONTRAST is not emitted as a warn")
else:
    bad.append("ICON CONTRAST did not appear in the warn stream")

# the message must say what to DO, not just what is wrong
line = [l for l in out.splitlines() if "ICON CONTRAST" in l][0]
if "recolor" in line and "plate" in line:
    ok.append("the message names both remedies (recolor to a darker tone / lift it on a plate)")
else:
    bad.append("the message reports the defect without naming a remedy")

# ---------------------------------------------------------------- 4. cost stays negligible
import time                                                 # noqa: E402
prs = Presentation(str(TMP / "five.pptx"))
sl = prs.slides[0]
t = time.perf_counter()
for _ in range(3):
    L._boxes(sl, W, H)
on = (time.perf_counter() - t) / 3
_orig = L._icon_ink
L._icon_ink = lambda b: None
t = time.perf_counter()
for _ in range(3):
    L._boxes(sl, W, H)
off = (time.perf_counter() - t) / 3
L._icon_ink = _orig
per_icon = (on - off) / 5 * 1000
if per_icon < 15:
    ok.append(f"icon inspection costs {per_icon:.1f} ms per icon (sampled at 96x96, NEAREST)")
else:
    bad.append(f"icon inspection costs {per_icon:.1f} ms per icon — too slow for every lint run")

print("\n".join("  ok   " + x for x in ok))
if bad:
    print("\n".join("  FAIL " + x for x in bad))
print(f"\n{len(ok)} passed, {len(bad)} failed")
raise SystemExit(1 if bad else 0)

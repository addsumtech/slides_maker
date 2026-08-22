#!/usr/bin/env python3
"""Regression fixtures for the strict Codex-only presentation delivery gate.

Run with: python3 tests/test_codex_delivery_gate.py
"""

from __future__ import annotations

import pathlib

import hashlib
import importlib.util
import json
import struct
import subprocess
import sys
import tempfile
import zlib
from pathlib import Path


HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def write_png(path: Path, width: int = 640, height: int = 360, *, alpha: bool = False) -> None:
    def chunk(kind: bytes, payload: bytes) -> bytes:
        return (
            struct.pack(">I", len(payload))
            + kind
            + payload
            + struct.pack(">I", zlib.crc32(kind + payload) & 0xFFFFFFFF)
        )

    pixel = b"\xff\xff\xff\xff" if alpha else b"\xff\xff\xff"
    pixels = b"\x00" + pixel * width
    raw = pixels * height
    color_type = 6 if alpha else 2
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, color_type, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw, 9))
        + chunk(b"IEND", b"")
    )


def load_gate():
    spec = importlib.util.spec_from_file_location(
        "codex_delivery_gate_test", SCRIPTS / "codex_delivery_gate.py"
    )
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def fresh_component_audit(build: Path, deck: Path) -> dict:
    spec = importlib.util.spec_from_file_location("component_audit_test", SCRIPTS / "component_audit.py")
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module.audit(str(build), str(deck))


def critic(lens: str) -> dict:
    probes = {"memory_sentence": "The remembered takeaway remains precise and defensible."}
    if lens == "design":
        probes["per_slide"] = [
            {
                "slide": 1,
                "first_read": "A clear title and visual focal point.",
                "takeaway_guess": "The workflow makes evidence review mandatory.",
            }
        ]
    return {
        "purpose": "Validate a one-slide presentation fixture for the final delivery gate.",
        # Which of review-rubrics.md's nine per-purpose overlays this review applied. Required
        # since the scoping rule got a gate — see validate_review._rubric_overlay.
        "rubric_overlay": "work status update",
        "reviewer": {
            "origin": "isolated",
            "identity": f"fixture-{lens}-reviewer",
            "fresh_context": True,
        },
        "coverage": {
            "slides_opened": [1],
            "passes": [f"{lens} lens (full deck)"],
            "stats_block_seen": True,
            "contract_card_seen": True,
        },
        # A real contract-card audit, not a placeholder. This used to read
        # {"lens_a": {"status": "pass"}, "lens_b": {"status": "pass"}} — `status` is not a field
        # in agents/critic.md's output block at all, and it validated clean because nothing
        # checked the subfields. validate_review.py now requires the audit each declared lens
        # owes, so the fixture has to be the shape a critic actually returns.
        "plan_audit": {
            "lens_a": {
                "memory_sentence": "One fixture slide, one claim, bound to its evidence.",
                "matches_deck_message": True,
                "curve_visible": "single slide — no curve to pace",
                "takeaway_titles": "title matches the takeaway table",
                "motion_manifest": "kept",
            },
            "lens_b": {
                "concept_landed": "carried: the evidence-binding idea is the slide's own geometry",
                "skeleton_rhythm": "kept — single slide",
                "signature_move": {"verdict": "landed", "why": "the binding is the visual",
                                   "carried": [], "proof": "matches"},
                "memorable_one_thing": "evidence travels with the deck",
                "composition": {"cover_archetype": "kept", "home_skeleton_plurality": "kept"},
                "register_interiors": "kept",
                "money_slide": "landed on slide 1",
                "semantic_colour": "ledger kept",
                "type_tokens": "sizes drawn from the declared tokens",
            },
        },
        "probes": probes,
        "verdict": "consent",
        "summary": "The final fixture meets the named review lens.",
        "strengths": ["Evidence and delivery artifacts are bound to the final deck."],
        "findings": [],
    }


def write_json(path: Path, value: dict) -> None:
    path.write_text(json.dumps(value, indent=2) + "\n", encoding="utf-8")


def fixture(root: Path) -> tuple[dict, dict, dict, Path]:
    try:
        from pptx import Presentation
    except ImportError:
        print("SKIPPED: this suite needs python-pptx")
        sys.exit(0)

    deck = root / "deck.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(deck)
    source = root / "README.md"
    source.write_text("# Source\n\nThe gate verifies workflow evidence.\n", encoding="utf-8")
    direction = root / "directions.html"
    direction.write_text(
        "<html><body>" + "direction preview " * 40 + "A B C D</body></html>",
        encoding="utf-8",
    )
    proof = root / "slide-1.png"
    icon = root / "feature.png"
    write_png(proof)
    write_png(icon, 512, 512, alpha=True)
    build = root / "build_deck.py"
    # APPLIES the register the evidence declares (`design.style_pick` = editorial_report). The
    # fixture used to declare a register and hand-build without it — which is exactly the defect
    # check_style_applied.py was written for, and the scaffolding teaching the defect is how it
    # spreads: whatever the example demonstrates is what gets written.
    build.write_text(
        "import deckkit as dk\n"
        "import presets\n\n"
        "p = presets.apply('editorial_report')\n\n"
        "def slide_01(slide):\n"
        "    dk.icon_card(slide, 0, 0, 1, 1, 'Feature')\n\n"
        "dk.lint_layout(prs, strict=True)\n",
        encoding="utf-8",
    )
    visual_manifest = root / "visual-contract.json"
    write_json(
        visual_manifest,
        {
            "schema": "slide-maker-codex-visual-contract/v1",
            "zones": [],
            "icons": [],
        },
    )
    visual_result = root / "visual-contract-final.json"
    write_json(
        visual_result,
        {
            "schema": "slide-maker-codex-visual-contract-result/v1",
            "pptx_sha256": sha256(deck),
            "manifest_sha256": sha256(visual_manifest),
            "passed": True,
            "zones": [],
            "icons": [],
            "errors": [],
        },
    )
    content_review = root / "critic-content.json"
    design_review = root / "critic-design.json"
    write_json(content_review, critic("content"))
    write_json(design_review, critic("design"))
    evidence = {
        "schema": "slide-maker-codex-evidence/v2",
        "runtime": "codex",
        "delivery": "presented",
        "review_effort": "standard",
        "deck": {"pptx": deck.name, "sha256": sha256(deck), "slide_count": 1},
        "interview": {
            "mode": "answered",
            "record": "The user requested a concise evidence-focused deck.",
            # The length axis, recorded. It is the question that disappears on a runtime with no
            # choice UI, and a deck silently built at one page is its symptom.
            "length": "user asked for a single-slide record fixture",
        },
        "content": {
            "source_mode": "provided",
            "sources": [{"kind": "provided", "path": source.name, "sha256": sha256(source)}],
            "slides": [
                {
                    "slide": 1,
                    "role": "cover",
                    "takeaway": "The gate makes workflow evidence verifiable.",
                    "evidence": ["README.md:1-3"],
                }
            ],
            "claim_ledger": [
                {
                    "claim": "The gate validates final workflow evidence.",
                    "source": "README.md:3",
                    "verified": True,
                }
            ],
            # The CANDIDATES, not a verdict about them: this gate now runs
            # arc_divergence.check() over them, because `"divergence": "ok"` is a string the run
            # writes about itself and a delivered deck passed that way with the script never run.
            "arc": {
                "chosen": "evidence-first",
                "shape": "evidence-build",
                "candidates": [
                    {"name": "evidence-first",
                     "shape": "evidence-build",
                     "roles": ["problem", "evidence", "conclusion"],
                     "audience_question": "does the measurement actually hold up",
                     "objection": "one record is not a result",
                     "closing_ask": "accept the single-record fixture as evidence",
                     "evidence": ["c1", "c2"]},
                    {"name": "recommendation-first",
                     "shape": "recommendation-first",
                     "roles": ["conclusion", "evidence", "roadmap"],
                     "audience_question": "should the fixture ship as the CI default",
                     "objection": "nobody has run it on the other host",
                     "closing_ask": "make it the default gate fixture",
                     "evidence": ["c1", "c3"]},
                ],
                "rejected": [
                    {"name": "recommendation-first",
                     "why_lost": "there is no decision to lead with on a one-slide record"},
                ],
            },
            "checkpoint": {
                "mode": "approved",
                "record": "The one-slide narrative and claim were approved.",
            },
        },
        "design": {
            "direction": {
                "branch": "clean",
                "artifact": direction.name,
                "sha256": sha256(direction),
                "directions": [
                    {
                        # A — the selected BESPOKE register (carries its own motif → real DNA)
                        "id": "A",
                        "name": "A",
                        "bg": "#071820",
                        "accent": "#54D9D0",
                        "font_display": "Inter",
                        "font_body": "Inter",
                        "density": "minimal",
                        "cover": "low-left",
                        "skeleton": "rail",
                        "cover_motif": "<div class='sonar'/>",
                        "ambient_motif": "<i class='sonar-echo'/>",
                    },
                    {
                        # B — a best-fit DNA preset
                        "id": "B",
                        "name": "B",
                        "bg": "#FFFFFF",
                        "accent": "#E4572E",
                        "font_display": "Georgia",
                        "font_body": "Arial",
                        "density": "normal",
                        "cover": "centred",
                        "skeleton": "split",
                        "dna": "editorial_paper",
                    },
                    {
                        # C — a best-fit DNA preset
                        "id": "C",
                        "name": "C",
                        "bg": "#182033",
                        "accent": "#F2C14E",
                        "font_display": "Menlo",
                        "font_body": "Inter",
                        "density": "dense",
                        "cover": "full-bleed-type",
                        "skeleton": "dashboard",
                        "dna": "terminal",
                    },
                    {
                        # D — the ONE allowed motif-less colour-scheme option
                        "id": "D",
                        "name": "D",
                        "bg": "#F4E9D8",
                        "accent": "#355C7D",
                        "font_display": "Rockwell",
                        "font_body": "Arial",
                        "density": "spacious",
                        "cover": "split-vertical",
                        "skeleton": "gallery",
                    },
                ],
                "decision": "user-approved",
                "record": "Direction A was selected after reviewing four distinct previews.",
            },
            "type_scale": {"display": 34, "title": 24, "body": 14},
            "concept": {
                "chosen": "a chain of custody — every claim carries its receipt",
                "via": "provenance, verification -> a rail that threads every page",
                "rejected": [
                    {"concept": "a laboratory bench",
                     "why_lost": "it pictures the work, not the guarantee the deck is selling"},
                    {"concept": "a set of scales",
                     "why_lost": "it implies a tradeoff, and nothing here is being traded off"},
                ],
            },
            "boldness": "balanced+",
            "style_pick": "editorial_report for AI/ML research · beat blueprint because the deck is a "
                          "data readout, not a schematic · anti-pick avoided: dark_tech neon cliché",
            "palette": "ink #101820 on paper #F6F4EF (14.9:1); accent #1F5FA9 FILL-only, "
                       "text-safe variant #17457C (5.4:1) — per palette_audit.py",
            "motif_generates": {
                "background": "a faint rule field the evidence rail sits on",
                "markers": "numbered evidence stamps, one per claim",
                "page": "slide 1 — the rail IS the slide's geometry",
            },
            "signature_move": "A visible evidence rail connects source, build, and review.",
            "carried_by": [1],
            "signature_proof": {
                "slide": 1,
                "path": proof.name,
                "sha256": sha256(proof),
                "pptx_sha256": sha256(deck),
            },
            "slides": [
                {
                    "slide": 1,
                    "function": "slide_01",
                    "form": "cover",
                    "runner_up": "editorial opener",
                    "reason": "The cover provides a single memorable entry point for the workflow.",
                    "categorical": True,
                    "components": [],
                }
            ],
            "checkpoint": {
                "mode": "approved",
                "record": "The design direction and form ledger were approved.",
            },
        },
        "build": {"script": build.name, "sha256": sha256(build), "strict_layout": True},
        "icons": [
            {
                "slide": 1,
                "family": "lucide",
                "asset": icon.name,
                "sha256": sha256(icon),
                "rasterizer": "scripts/icons.py",
            }
        ],
        "visual_contract": {
            "manifest": visual_manifest.name,
            "sha256": sha256(visual_manifest),
            "result": visual_result.name,
            "result_sha256": sha256(visual_result),
            "pptx_sha256": sha256(deck),
        },
        "critics": [
            {
                "lens": "content",
                "review": content_review.name,
                "sha256": sha256(content_review),
                "pptx_sha256": sha256(deck),
            },
            {
                "lens": "design",
                "review": design_review.name,
                "sha256": sha256(design_review),
                "pptx_sha256": sha256(deck),
            },
        ],
        # The actor's render look — one verdict per slide (this fixture deck is a single page).
        "render_selfcheck": {"slides": [{"n": 1, "verdict": "ok — cover reads clean"}]},
        "waivers": [],
    }
    lint = {
        "findings": [],
        "pixel_checks": [{"pass": True}],
        "text_runs": [{"role": "body", "size_pt": 14}],
        "stats": {"warnings": []},
    }
    components = fresh_component_audit(build, deck)
    return evidence, lint, components, build


def main() -> int:
    gate = load_gate()
    with tempfile.TemporaryDirectory(prefix="codex-gate-") as name:
        root = Path(name)
        evidence, lint, components, build = fixture(root)
        icon = root / evidence["icons"][0]["asset"]
        errors = gate.evaluate(lint, components, build, evidence, root)
        failures = []
        if errors:
            failures.append("valid evidence unexpectedly blocked:\n" + "\n".join(errors))

        evidence_path = root / ".codex-deck-evidence.json"
        lint_path = root / "lint-final.json"
        components_path = root / "components-final.json"
        receipt_path = root / ".codex-delivery-receipt.json"
        write_json(evidence_path, evidence)
        write_json(lint_path, lint)
        write_json(components_path, components)
        run = subprocess.run(
            [
                sys.executable, str(SCRIPTS / "codex_delivery_gate.py"),
                "--lint", str(lint_path),
                "--components", str(components_path),
                "--build-script", str(build),
                "--evidence", str(evidence_path),
                "--receipt", str(receipt_path),
            ],
            text=True, capture_output=True,
        )
        if run.returncode or not receipt_path.exists():
            failures.append("a valid gate run did not emit a PASS receipt:\n" + run.stdout + run.stderr)
        deck_path = root / evidence["deck"]["pptx"]
        guard = subprocess.run(
            [sys.executable, str(SCRIPTS / "codex_handoff_guard.py"), "--receipt", str(receipt_path), "--deck", str(deck_path)],
            text=True, capture_output=True,
        )
        if guard.returncode:
            failures.append("a matching PASS receipt was rejected by the hand-off guard:\n" + guard.stdout + guard.stderr)
        tampered = root / "tampered-deck.pptx"
        tampered.write_bytes(deck_path.read_bytes() + b"tampered")
        guard = subprocess.run(
            [sys.executable, str(SCRIPTS / "codex_handoff_guard.py"), "--receipt", str(receipt_path), "--deck", str(tampered)],
            text=True, capture_output=True,
        )
        if guard.returncode == 0:
            failures.append("a receipt for a different PPTX passed the hand-off guard")

        evidence["runtime"] = "openai-gpt-bridged"
        errors = gate.evaluate(lint, components, build, evidence, root)
        if errors:
            failures.append("a bridged GPT Store runtime was incorrectly blocked:\n" + "\n".join(errors))
        evidence["runtime"] = "codex"

        lint["stats"]["warnings"] = ["card_dominance"]
        evidence["waivers"] = [
            {
                "kind": "stats",
                "warning": "card_dominance",
                "reason": "The deliberate scorecard composition is the approved signature move.",
            }
        ]
        errors = gate.evaluate(lint, components, build, evidence, root)
        if errors:
            failures.append("a documented, deliberate statistics waiver was blocked:\n" + "\n".join(errors))
        lint["stats"]["warnings"] = []
        evidence["waivers"] = []

        # THE LENGTH AXIS. It is the interview question that disappears on a runtime with no
        # choice UI — nothing downstream demands it, so nothing notices — and the observed symptom
        # is a deck silently built at ONE page. `interview.record` cannot stand in for it: a
        # 12-char free-text floor passes a stub.
        _len = evidence["interview"].pop("length")
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("interview.length" in error for error in errors):
            failures.append("an interview with no recorded deck length passed the gate")
        evidence["interview"]["length"] = _len
        errors = gate.evaluate(lint, components, build, evidence, root)
        if errors:
            failures.append("restoring interview.length did not clear the gate:\n" + "\n".join(errors))

        evidence["critics"][0]["pptx_sha256"] = "0" * 64
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("not bound to the final PPTX" in error for error in errors):
            failures.append("a critic review for a stale deck passed the gate")
        evidence["critics"][0]["pptx_sha256"] = evidence["deck"]["sha256"]

        bad_review = root / "critic-content.json"
        write_json(
            bad_review,
            {"purpose": "This artifact is intentionally incomplete but not empty.", "verdict": "consent"},
        )
        evidence["critics"][0]["sha256"] = sha256(bad_review)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("invalid critic schema" in error for error in errors):
            failures.append("an incomplete critic review passed the gate")

        write_json(bad_review, critic("content"))
        evidence["critics"][0]["sha256"] = sha256(bad_review)
        components = {
            "clusters": [{"slide": 1, "pattern": "equal tiles", "suggest": ["scorecard"]}],
            "actionable": [],
            "suppressed_by": ["scorecard"],
        }
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("does not match a fresh audit" in error for error in errors):
            failures.append("a stale or fabricated component audit passed the gate")

        component_errors = []
        gate.check_components(
            evidence,
            {
                "clusters": [{"slide": 1, "pattern": "tile row", "suggest": ["scorecard"]}],
                "suppressed_by": ["org_tree"],
            },
            {1: {"function": "slide_01", "components": ["org_tree"]}},
            {"slide_01": {"org_tree"}},
            component_errors,
        )
        if component_errors:
            failures.append("a same-slide, registered component emitter was incorrectly blocked")

        components = fresh_component_audit(build, root / evidence["deck"]["pptx"])
        proof = root / evidence["design"]["signature_proof"]["path"]
        proof.write_bytes(b"")
        evidence["design"]["signature_proof"]["sha256"] = sha256(proof)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("signature_proof" in error for error in errors):
            failures.append("an empty signature proof passed the gate")

        write_png(icon, 240, 240, alpha=True)
        evidence["icons"][0]["sha256"] = sha256(icon)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("thumbnail blur" in error for error in errors):
            failures.append("a low-resolution icon thumbnail passed the gate")

        write_png(icon, 512, 512, alpha=False)
        evidence["icons"][0]["sha256"] = sha256(icon)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("transparent alpha" in error for error in errors):
            failures.append("a matted icon asset without alpha passed the gate")

        write_png(icon, 512, 512, alpha=True)
        evidence["icons"][0]["sha256"] = sha256(icon)
        build.write_text(
            build.read_text(encoding="utf-8") + "\nimport subprocess\nsubprocess.run(['qlmanage', '-t'], check=False)\n",
            encoding="utf-8",
        )
        evidence["build"]["sha256"] = sha256(build)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("Quick Look thumbnail generation" in error for error in errors):
            failures.append("a Quick Look icon rasterization workaround passed the gate")

        # Every STRICT_STATS name must be a code lint_deck.py can actually EMIT. The test above
        # feeds the gate `["card_dominance"]` -- the gate's own expected shape -- so it proves the
        # logic and nothing about the vocabulary. Measured before this check existed: `cjk_risk`
        # and `color_envelope` matched no linter output at all, so two of the seven entries were
        # a gate believing it enforced something it could never see. The real codes are
        # `cjk_tight_leading` and `envelope_monoculture`. This is a cross-file agreement, which is
        # decidable by a program and therefore should never again be left to someone remembering.
        import re as _re

        _src = (pathlib.Path(__file__).resolve().parent.parent
                / "scripts" / "lint_deck.py").read_text(encoding="utf-8")
        emitted = {"_".join(m.group(1).strip().lower().split())
                   for m in _re.finditer(r"[\"\']([A-Z][A-Z0-9 \-/&]{2,40}):", _src)}
        # the accessibility tier: an objective floor on the per-slide `warnings` stream, which
        # previously had NO strict path at all (errors blocked, stats could be waived, per-slide
        # warnings were unreachable from here). Both directions, plus the rubber-stamp guard.
        _a11y = {"stats_warnings": [], "warnings": [
            {"slide": 8, "text": "ICON CONTRAST: icon ink #C08A2E on #F5F1E6 — 2.69:1 "
                                 "(<3:1, WCAG 1.4.11)."}]}

        def _a11y_errs(ev):
            errs: list[str] = []
            gate.check_lint(_a11y, "presented", ev, errs)
            return [e for e in errs if "1.4.11" in e]

        if not _a11y_errs({}):
            failures.append("an icon below the WCAG 1.4.11 3:1 floor passed the codex gate")
        if _a11y_errs({"waivers": [{"kind": "a11y", "warning": "ICON CONTRAST",
                                    "reason": "decorative flourish; meaning is carried by the "
                                              "label beside it"}]}):
            failures.append("a properly reasoned a11y waiver was still blocked")
        if not _a11y_errs({"waivers": [{"kind": "a11y", "warning": "ICON CONTRAST",
                                        "reason": "ok"}]}):
            failures.append("a rubber-stamp waiver (<12 chars of reason) was accepted")
        _unrelated = {"stats_warnings": [], "warnings": [
            {"slide": 3, "text": "TEXT WALL: slide 3 carries a reading load of ~95 words"}]}
        _e: list[str] = []
        gate.check_lint(_unrelated, "presented", _e, [])
        if [x for x in _e if "1.4.11" in x]:
            failures.append("an unrelated per-slide warning was treated as an a11y floor")

        unreachable = sorted(name for name in gate.STRICT_STATS if name not in emitted)
        if unreachable:
            failures.append(
                "STRICT_STATS names that lint_deck.py never emits, so the gate can never fire "
                "on them: " + ", ".join(unreachable))

        # ADVERSARY-FOUND, two ways the arc/slides/checkpoint changes broke on their own skeleton:
        #
        # (1) `--init` wrote `"candidates": ["<one descriptive string>"]` and check_content then did
        #     `str(c.get("name") …)` over it — an uncaught AttributeError on the exact evidence file
        #     the gate itself emits. A gate that crashes on `--init` output cannot be run at all.
        # (2) The unedited skeleton, with its `<placeholder>` fields, cleared every width floor and
        #     PASSED — the same blindness the shared gate had.
        init_path = root / "init-evidence.json"
        init_run = subprocess.run(
            [sys.executable, str(SCRIPTS / "codex_delivery_gate.py"), "--init", str(init_path)],
            text=True, capture_output=True)
        if init_run.returncode or not init_path.exists():
            failures.append("codex_delivery_gate.py --init did not write a skeleton:\n"
                            + init_run.stdout + init_run.stderr)
        else:
            skeleton = json.loads(init_path.read_text(encoding="utf-8"))
            crash = subprocess.run(
                [sys.executable, str(SCRIPTS / "codex_delivery_gate.py"),
                 "--lint", str(lint_path), "--components", str(components_path),
                 "--build-script", str(build), "--evidence", str(init_path)],
                text=True, capture_output=True)
            if "Traceback" in crash.stderr or "AttributeError" in (crash.stdout + crash.stderr):
                failures.append("the gate CRASHES on its own --init skeleton:\n" + crash.stderr)
            # The skeleton is placeholder-filled, so it must be BLOCKED (cleanly), never accepted.
            init_errors = gate.evaluate(lint, components, build, skeleton, root)
            if not init_errors:
                failures.append("the unedited --init skeleton PASSED the gate — placeholders are "
                                "clearing the width floors")

        # And on the real fixture: swap a filled candidate field for a raw `--template` placeholder
        # and confirm it is refused, so the guard is doing work rather than being always-on.
        placeholder_ev = json.loads(json.dumps(evidence))
        placeholder_ev["content"]["arc"]["candidates"][0]["audience_question"] = \
            "<the question this room is actually asking>"
        errs = gate.evaluate(lint, components, build, placeholder_ev, root)
        if not any("placeholder" in e for e in errs):
            failures.append("a raw --template placeholder in a candidate field passed the codex gate")

        if failures:
            print("\n".join("FAIL: " + failure for failure in failures))
            return 1
    print("ok - Codex delivery gate rejects icon-rasterization and historical bypasses, and "
          "every STRICT_STATS name is a code the linter really emits")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

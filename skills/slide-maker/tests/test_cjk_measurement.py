#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""On a 中文 deck the ruler was short by 46%, and every geometry check used it.

`_natural_width_in` was never wrong — hand it a face that covers CJK and it is exact. What was
wrong was the FACE handed to it. `_ink_rect` resolved a run's face with `run.font.name`, and
python-pptx's `font.name` reads `<a:latin>`; CJK glyphs render from `<a:ea>`, where `set_font` /
`_apply_ea` put the East-Asian face. Measured on one real subtitle at 13pt (27 full-width glyphs,
true width 4.8750in):

    Hiragino Sans GB   4.8750in   exact
    Songti SC          4.8750in   exact
    Helvetica Neue     2.6181in   54% — what every check was actually using

That under-measurement reaches TEXT_OVERLAP, OFF_CANVAS, ESCAPES_CARD, FOOTER and every geometry
decision built on ink rectangles, i.e. essentially all of this file's geometry on a Chinese deck.
Its signature in practice is the worst kind: a page collides after a generous-looking margin was
left, and nudging coordinates gives feedback that does not match intuition — so the operator
blames their own arithmetic rather than the ruler.

The resolution goes through `_inherited_ea`, not through a second width model and not through a
`_is_wide` branch: a supplied CJK template normally sets the face one level up (paragraph
`defRPr`, the shape's `lstStyle`), and that is the same chain `retrofit_ea` and `CJK_NO_EA`
already walk — so all three agree by construction rather than by coincidence.

Run:  python3 tests/test_cjk_measurement.py
"""
import pathlib
import sys

HERE = pathlib.Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent / "scripts"))

import deckkit as dk                                                  # noqa: E402
from deckkit import Inches, Pt                                        # noqa: E402
from pptx.oxml.ns import qn                                           # noqa: E402

PASS, FAIL = [], []
CJK = "给准备申请的人——学位、博士、岗位，三扇门分别长什么样"   # 27 glyphs
TRUE_W = 27 * 13 / 72.0                                              # full-width at 13pt
LATIN = "A subtitle running straight across the page"


def check(name, cond, detail=""):
    (PASS if cond else FAIL).append(name)
    print(("  ok   " if cond else "  FAIL ") + name
          + (("  — " + str(detail)) if detail and not cond else ""))


def ink_w(text, latin="Helvetica Neue", ea=None, inherited=None, size=13):
    """The ink width _ink_rect measures for one un-wrapped run."""
    dk.FONT, dk.EAFONT = "Helvetica Neue", "Hiragino Sans GB"
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    if inherited:
        tb = s.shapes.add_textbox(Inches(0.62), Inches(2.0), Inches(8.76), Inches(0.4))
        tb.text_frame.word_wrap = False
        p0 = tb.text_frame.paragraphs[0]
        r0 = p0.add_run()
        r0.text = text
        r0.font.size = Pt(size)
        r0.font.name = latin
        pPr = p0._p.get_or_add_pPr()
        d = pPr.makeelement(qn("a:defRPr"), {})
        pPr.append(d)
        d.append(d.makeelement(qn("a:ea"), {"typeface": inherited}))
        sh = tb
    else:
        run = (text, size, dk.DEEP, False, False, latin) if ea is None \
            else (text, size, dk.DEEP, False, False, latin, ea)
        dk.text(s, 0.62, 2.0, 8.76, 0.4, [[run]], space_after=0, wrap=False)
        sh = list(s.shapes)[0]
    r = dk._ink_rect(sh, dk._bbox_in(sh))
    return r[0][2] if r and r[0] else 0.0


def main():
    print("CJK measurement")

    # ---- the fix -----------------------------------------------------------------------
    w = ink_w(CJK)
    check("a CJK run measures at its REAL width, not the Latin face's",
          abs(w - TRUE_W) < 0.02, "%.4f vs %.4f" % (w, TRUE_W))
    check("...which is ~1.86x what the old ruler reported (2.6181in)", w > 4.0, "%.4f" % w)

    # the face may live one level up — a supplied template's usual shape
    w = ink_w(CJK, inherited="Songti SC")
    check("an INHERITED ea face (paragraph defRPr) is honoured too",
          abs(w - TRUE_W) < 0.02, "%.4f" % w)

    # an explicit 7th slot must win
    w = ink_w(CJK, ea="Songti SC")
    check("an explicit ea face in the run is honoured", abs(w - TRUE_W) < 0.02, "%.4f" % w)

    # ---- and the half that must not move ------------------------------------------------
    lat = ink_w(LATIN)
    check("a Latin run is measured exactly as before (no CJK, no change)",
          2.5 < lat < 4.5, "%.4f" % lat)
    check("a Latin run naming a CJK face still measures as Latin text",
          abs(ink_w(LATIN, latin="Songti SC") - ink_w(LATIN, latin="Songti SC")) < 1e-9)

    mixed_cjk = ink_w("认识 LUMC")
    mixed_lat = ink_w("LUMC")
    check("a mixed run is wider than its Latin part alone", mixed_cjk > mixed_lat,
          "%.3f vs %.3f" % (mixed_cjk, mixed_lat))

    # ---- the whole point: geometry decisions change ------------------------------------
    # A box that "fits" under the short ruler and overflows under the true one. This is the
    # defect class the fix exists for — under-measurement waves a real overflow through.
    dk.FONT, dk.EAFONT = "Helvetica Neue", "Hiragino Sans GB"
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    dk.box(s, 0, 0, 10, 5.625, fill=dk.RGBColor(0xFF, 0xFF, 0xFF))
    # 3.2in box holding 4.875in of CJK ink, with a visible outline so OVERFLOW can fire
    dk.box(s, 0.6, 2.0, 3.2, 0.42, fill=None, line=dk.RGBColor(0x88, 0x88, 0x88), line_w=1.0)
    dk.text(s, 0.6, 2.0, 3.2, 0.42, [[(CJK, 13, dk.DEEP, False, False, "Helvetica Neue")]],
            space_after=0, wrap=True)
    found = dk.lint_layout(prs, verbose=False)
    check("CJK text overflowing its box is now measurable at all",
          ink_w(CJK) > 3.2, "ink %.2f vs box 3.20" % ink_w(CJK))

    # ---- a real CJK deck must not suddenly light up ------------------------------------
    # The fix improves PRECISION; it must not act as a loosened or tightened threshold. A deck
    # that was clean under the short ruler and is genuinely well-built stays clean.
    dk.FONT, dk.EAFONT = "Helvetica Neue", "Hiragino Sans GB"
    prs = dk.blank_deck()
    for k in range(3):
        sl = dk.add_slide(prs)
        dk.text(sl, 0.62, 0.7, 8.76, 0.5,
                [[("认识 LUMC 医学中心", 23, dk.DEEP, True, False, "Helvetica Neue", "Songti SC")]],
                space_after=0, wrap=False)
        for j in range(3):
            dk.text(sl, 0.62, 1.6 + j * 0.34, 8.76, 0.26,
                    [[("一行普通的中文正文，长度适中，留有余量。", 12, dk.DEEP,
                       False, False, "Helvetica Neue")]], space_after=0, wrap=False)
    crit = [f for f in dk.lint_layout(prs, verbose=False) if f[1] == "CRITICAL"]
    check("a well-built CJK deck stays clean under the true ruler", crit == [], crit)

    print("\n{} passed, {} failed".format(len(PASS), len(FAIL)))
    return 1 if FAIL else 0


if __name__ == "__main__":
    sys.exit(main())

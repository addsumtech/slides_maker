#!/usr/bin/env python3
"""The critic waiver on the shared path must be CLASSIFIED, not just written.

Script-style (`main()` + explicit exit) to match test_lint_regressions.py and the codex
suites; pytest collects nothing from this file by design, and ci.yml invokes it directly.

History this guards: the Codex delivery gate required a distinct schema-valid review artifact
per lens, while the shared path accepted `{"critic": {"waived": "<any string>"}}` and printed
one line. A hand-typed waiver carried a real 10-slide deck through "all hand-off gates pass"
with no independent critic ever involved.
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SKILL = HERE.parent
RENDER = SKILL / "scripts" / "render_deck.py"

GOOD_REASON = ("No subagent dispatch on this host, so the content and design lenses were "
               "run inline in the author's own context.")

DESIGN_OK = {
    "concept": {"chosen": "the deck is a picture of a signal being recovered from noise",
                "rejected": [{"concept": "a clock running out", "why_lost": "makes time the subject, not fidelity"},
                             {"concept": "a pair of hands", "why_lost": "no room for the data-consistency argument"}]},
    "boldness": "balanced+",
    "signature_move": "s" * 30,
    "carried_by": ["slide 3", "slide 4"],
    "form_ledger": "f" * 30,
    # TRUTHFUL for this fixture: build_deck() makes three slides of one text box each and places
    # no icons at all. This said "tabler" and passed, because the declared-but-absent direction
    # only printed. It dies now, and the fixture is the first thing it caught.
    "icon_family": "none — a three-slide text fixture, no categorical content to mark",
    "palette": "FILL E2543A / TEXT BD4630 on cream, A3341F on tint",
    # Required from ~6 content slides up; always allowed. A truthful reason for every fixture here.
    "build_shape": "solo — CI fixture, no subagent dispatch",
    "type_scale": {"display": 34, "title": 24, "body": 14},
    # The 🔴 design checkpoint's mode + record. Same two values codex_delivery_gate accepts
    # (`approved` | `auto`); on the shared path there is no top-level `design` section, so it
    # rides beside the plan it belongs to.
    "checkpoint": {"mode": "approved", "record": "CI fixture — design plan presented and approved"},
    # A motif that only RECURS is an ornament with a schedule; these are the three things it makes
    # besides itself. Carved out under a conservative dial with a `deliberately restrained:` move,
    # exactly as signature_proof is.
    "motif_generates": {"background": "a scanline field the recovered signal rides",
                        "markers": "sample-index numerals in the corner",
                        "page": "slide 3 — the recovery plot IS the motif"},
    # THE ANCHOR PROOF — three rendered pages, three different failures (scripts/anchor_proof.py).
    # One page only ever proved the aesthetic risk; the design that looks right on it can still fail
    # to hold the deck's densest page, and the charts can still obey none of the palette/type
    # decisions that were made against text.
    # Step 2 opens by MAKING one real slide and looking at it; this is the one field on the plan
    # that cannot be written without an artifact.
    "material_probe": {"png": "probe.png",
                       "safe_version": "安全版本会是同一句话配一张通用图表——那不是同一类东西"},
    "signature_proof": [
        {"role": "signature", "slide": 3, "png": "proof.png"},
        {"role": "complex", "slide": 2, "png": "proof_complex.png"},
        {"role": "data", "slide": 1, "png": "proof_data.png"},
    ],
}
PROV_OK = {"claims": [{"claim": "c", "verdict": "CONFIRMED", "url": "https://example.org"}]}
# The arc competition (Step 1). Gated on the shared path since the same competition was already
# bound on the Codex one — a wrong form costs one slide, a wrong arc costs the build under it.
# 🔴 This fixture USED to be the exact shape the gate now refuses:
#     {"chosen": ..., "rejected": [...], "divergence": "ok"}
# — every field authored by the run, and `"ok"` standing in for a competition that never happened.
# A delivered deck passed that way while arc_divergence.py had never been invoked for it. The gate
# now runs `arc_divergence.check()` over `candidates` itself, so the fixture has to carry two arcs
# that actually diverge: different shape, different opening role, different ask, different stance,
# and comparable ledger evidence (a candidate under half the winner's evidence is a strawman).
ARC_CANDIDATES = [
    {"name": "problem-turn-evidence",
     "shape": "problem-turn-evidence",
     "roles": ["problem", "diagnosis", "evidence", "conclusion"],
     "audience_question": "why did the recovery fall over on this cohort",
     "objection": "the model is overfit to the training scanner",
     "closing_ask": "fund a second acquisition round on the other scanner",
     "evidence": ["c1", "c2"]},
    {"name": "recommendation-first",
     "shape": "recommendation-first",
     "roles": ["conclusion", "evidence", "roadmap"],
     "audience_question": "should the pipeline ship this quarter or next",
     "objection": "nobody has measured the latency budget end to end",
     "closing_ask": "approve the December ship date",
     "evidence": ["c1", "c3"]},
]


def content_ok(n_slides: int = 3) -> dict:
    """A complete `content` record for an n-slide fixture deck.

    Three artifacts, not one: the arc COMPETITION (scored here by the gate), the per-slide plan
    (`content.slides` — the content checkpoint's own table, the field codex_delivery_gate has
    always required and the shared path did not), and the checkpoint MODE. Sized to the deck
    because `content.slides` must cover every slide exactly once.
    """
    rows = []
    for i in range(1, n_slides + 1):
        if i == 1:
            role, take = "cover", "this deck is about recovering a signal from noise"
        elif i == n_slides:
            role, take = "conclusion", "the recovery holds on the held-out scanner too"
        else:
            role, take = "evidence", f"measurement {i} is what rules out the overfit explanation"
        rows.append({"slide": i, "role": role, "takeaway": take,
                     "evidence": [f"fixture ledger c{i}"], "units": 1})
    return {"arc": {"chosen": "problem-turn-evidence",
                    "candidates": [dict(c) for c in ARC_CANDIDATES],
                    "rejected": [{"name": "recommendation-first",
                                  "why_lost": "there is no decision to lead with here"}]},
            "slides": rows,
            "checkpoint": {"mode": "approved",
                           "record": "CI fixture — table presented, approved as-is"}}


ARC_OK = content_ok(3)


def fit_content(gates: dict, pptx) -> dict:
    """Resize a fixture's `content.slides` to the deck actually under test.

    The per-slide plan must cover every slide exactly once — that IS the gate — so a fixture with
    three hardcoded rows cannot be shared by tests that build 3-, 8- and 16-slide decks. This
    regenerates only the rows, only when they do not already match, and touches no other field, so
    a test that deliberately corrupts `content` still sees its own corruption rather than this.
    """
    c = gates.get("content")
    if not isinstance(c, dict) or "slides" not in c:
        return gates
    try:
        from pptx import Presentation
        n = len(Presentation(str(pptx)).slides)
    except Exception:                                              # noqa: BLE001
        return gates
    if not n or len(c.get("slides") or []) == n:
        return gates
    out = dict(gates)
    out["content"] = dict(c, slides=content_ok(n)["slides"])
    return out


def build_deck(dest: Path, n_slides: int = 3) -> Path:
    sys.path.insert(0, str(SKILL / "scripts"))
    import deckkit as dk
    prs = dk.blank_deck(10, 5.625)
    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        dk.text(s, 1, 1, 8, 1, [[(f"Slide {i+1}", 28, dk.DEEP, True, False)]])
    out = dest / "t.pptx"
    prs.save(str(out))
    return out


def write_proof(dest: Path) -> None:
    """Real PNGs next to the deck — the anchors point at rendered evidence, not a promise."""
    sys.path.insert(0, str(SKILL / "scripts"))
    from PIL import Image
    # NOT a flat rectangle: the gate now refuses a single-colour PNG as an anchor, because a
    # 960x540 field of one grey used to satisfy the ANCHOR PROOF — evidence that proves nothing.
    for k, name in enumerate(("proof.png", "proof_complex.png", "proof_data.png",
                              "probe.png")):
        im = Image.new("RGB", (960, 540), (240, 240, 245))
        for x in range(80, 880):
            for y in range(120 + k * 20, 300 + k * 20):
                im.putpixel((x, y), (16 + k * 20, 110, 99))
        im.save(dest / name)


def run_gate(deck: Path, gates: dict, *flags: str) -> tuple[int, str]:
    gates = fit_content(gates, deck)
    (deck.parent / ".deck-gates.json").write_text(json.dumps(gates))
    p = subprocess.run(
        [sys.executable, str(RENDER), str(deck), "--gate-check", "--static", *flags],
        capture_output=True, text=True)
    return p.returncode, p.stdout + p.stderr


# The delivery a gate enforces comes from two places — a recorded `delivery` key and the CLI mode
# flags — and they used to be read INDEPENDENTLY: the type-scale floor read the key, the density
# gate read the flag. One run could therefore enforce two different deliveries, and --selfread was
# INERT for the floor (a self-read deck with 12pt body died citing the *presented* floor).
# The rule now: a recorded key wins, the flag is the fallback, an unrecognised recorded value DIES
# rather than falling back to a floor it was never meant to be held to, and both gates read the one
# resolved value. Body is 12.0 throughout, which is legal for selfread/surface and illegal for
# presented/textheavy — so each cell's outcome is decided purely by the resolved delivery.
DELIVERY_CASES = [
    # (name,                          recorded,      flags,           expect_floor_death)
    ("no key + no flag = presented",  None,          (),              True),
    ("no key + --selfread applies",   None,          ("--selfread",), False),
    ("no key + --surface applies",    None,          ("--surface",),  False),
    ("recorded presented wins",       "presented",   (),              True),
    ("recorded presented beats flag", "presented",   ("--selfread",), True),
    ("recorded selfread applies",     "selfread",    (),              False),
    ("recorded selfread beats flag",  "selfread",    ("--textheavy",), False),
    ("unknown recorded value dies",   "briefing",    (),              None),
]


def check_delivery(deck: Path) -> tuple[int, int]:
    ok = bad = 0
    for name, recorded, flags, want_death in DELIVERY_CASES:
        g = {"critic": {"verdict": "consent", "rounds": 2},
             "design_plan": dict(DESIGN_OK, type_scale={"display": 34, "title": 24, "body": 12}),
             "provenance": PROV_OK, "content": ARC_OK}
        if recorded:
            g["delivery"] = recorded
        _, out = run_gate(deck, g, *flags)
        if want_death is None:                       # unknown value must be REFUSED, not defaulted
            good = "not a delivery mode" in out
            why = "died on the unknown delivery" if good else "silently accepted an unknown delivery"
        else:
            died = "legibility floor, not a style choice" in out
            good = died == want_death
            why = ("enforced the floor" if died else "let 12pt through")
        if good:
            ok += 1
            print("  ok   delivery: {} -> {}".format(name, why))
        else:
            bad += 1
            print("  FAIL delivery: {} -> {}\n       {}".format(
                name, why, out.strip().splitlines()[-1][:150] if out.strip() else "(no output)"))
    return ok, bad


CASES = [
    ("a placeholder waiver is refused",
     {"critic": {"waived": "auto mode"}},
     False, "written reason"),

    ("a substantive but UNCLASSIFIED waiver is refused",
     {"critic": {"waived": GOOD_REASON}},
     False, "waived_category"),

    ("an unknown category is refused",
     {"critic": {"waived": GOOD_REASON, "waived_category": "because-i-said-so"}},
     False, "waived_category"),

    ("no-dispatch-on-host without inline_ran is refused",
     {"critic": {"waived": GOOD_REASON, "waived_category": "no-dispatch-on-host"}},
     False, "inline_ran"),

    ("a fully classified waiver passes, labelled NOT INDEPENDENT",
     {"critic": {"waived": GOOD_REASON, "waived_category": "no-dispatch-on-host",
                 "inline_ran": True}},
     True, "NOT INDEPENDENTLY REVIEWED"),

    ("a legitimate minor-edit waiver passes",
     {"critic": {"waived": "One-slide typo fix to a deck that already passed its full loop.",
                 "waived_category": "already-reviewed-minor-edit"}},
     True, "critic WAIVED"),

    ("the consent path is unchanged",
     {"critic": {"verdict": "consent", "rounds": 2}},
     True, "critic consented"),

    # `cap-reached-majors-open`. Every other kind describes a loop that was SKIPPED, so a deck that
    # RAN to its round cap with majors still open had no honest route: `user-waived` asserts "the
    # user was asked and chose to ship over it", a claim about a conversation that did not happen.
    # Measured on a real build, the builder's words were "the four options force either a lie or a
    # red gate" — and it correctly left the gate red rather than invent one. This kind claims the
    # work HAPPENED, so it owes the reader what the work ran into; both refusals below are the
    # difference between a record and an excuse.
    ("cap-reached-majors-open without the open list is refused",
     {"critic": {"waived": "Two full rounds ran; three design majors remain open at the cap.",
                 "waived_category": "cap-reached-majors-open"}},
     False, '"open"'),

    ("cap-reached-majors-open with an EMPTY open list is refused",
     {"critic": {"waived": "Two full rounds ran; three design majors remain open at the cap.",
                 "waived_category": "cap-reached-majors-open", "open": []}},
     False, '"open"'),

    ("cap-reached-majors-open without surfaced_to_user is refused",
     {"critic": {"waived": "Two full rounds ran; three design majors remain open at the cap.",
                 "waived_category": "cap-reached-majors-open",
                 "open": ["envelope monoculture", "signature move reaches too few pages"]}},
     False, "surfaced_to_user"),

    ("a complete cap-reached record passes, still labelled NOT INDEPENDENT",
     {"critic": {"waived": "Two full rounds ran; three design majors remain open at the cap.",
                 "waived_category": "cap-reached-majors-open",
                 "open": ["envelope monoculture", "signature move reaches too few pages"],
                 "surfaced_to_user": True}},
     True, "NOT INDEPENDENTLY REVIEWED"),

    # The two-token contrast rule was declared in a design plan and then broken four times on
    # the same deck, each in a pair nobody was computing contrast for. `palette` is a required
    # field so the split has to be resolved (palette_audit.py) rather than remembered.
    ("a design plan with no resolved palette is refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "palette"}},
     False, "palette"),

    ("a design plan carrying the palette split passes",
     {"critic": {"verdict": "consent", "rounds": 2}, "design_plan": DESIGN_OK},
     True, "design plan: boldness"),

    # type_scale and signature_proof were gated on the CODEX path only. Typography was then the one
    # pillar of the visual language the shared path never made anyone resolve, and the signature move
    # was accepted as a sentence with nothing showing it survived the build. Same asymmetry as the
    # critic-waiver bug above, which is why both directions are pinned here.
    ("a design plan with no type_scale is refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "type_scale"}},
     False, "type_scale"),

    ("a type_scale whose tiers do not rank is refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, type_scale={"display": 18, "title": 24, "body": 14})},
     False, "not a scale"),

    ("a body size under the legibility floor is refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, type_scale={"display": 34, "title": 24, "body": 9})},
     False, "legibility floor"),

    ("a design plan with no signature_proof is refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "signature_proof"}},
     False, "signature_proof"),

    # An OpenAI/Codex-bridged run keeps both .codex-deck-evidence.json and .deck-gates.json, and its
    # own gate spells this key "path". Rejecting that spelling here would fail the same evidence for
    # its key name alone.
    ("the Codex spelling signature_proof[].path is accepted",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof=[
          {"role": "signature", "slide": 3, "path": "proof.png"},
          {"role": "complex", "slide": 2, "path": "proof_complex.png"},
          {"role": "data", "slide": 1, "path": "proof_data.png"}])},
     True, "design plan: boldness"),

    # ---- the ANCHOR PROOF rules. Every one of these passed as recently as the single-anchor
    # shape, which is the point: a deck could prove its bravest page and ship a design that had
    # never been seen holding real content or a real chart.
    ("the LEGACY single-anchor dict is refused, and told what to add",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof={"slide": 3, "png": "proof.png"})},
     False, "only 1 anchor"),

    ("...and the refusal names the two anchors that are missing, not just a count",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof={"slide": 3, "png": "proof.png"})},
     False, '"complex"'),

    ("two anchors claiming the same ROLE are refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof=[
          {"role": "signature", "slide": 3, "png": "proof.png"},
          {"role": "complex", "slide": 2, "png": "proof_complex.png"},
          {"role": "complex", "slide": 1, "png": "proof_data.png"}])},
     False, "same role"),

    ("two anchors pointing at the same SLIDE are refused — proving one page three times "
     "proves one page",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof=[
          {"role": "signature", "slide": 3, "png": "proof.png"},
          {"role": "complex", "slide": 3, "png": "proof_complex.png"},
          {"role": "data", "slide": 3, "png": "proof_data.png"}])},
     False, "same slide"),

    ("an anchor set with no `signature` role is refused — the risk is the anchor most likely to "
     "have been sanded away",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof=[
          {"role": "data", "slide": 3, "png": "proof.png"},
          {"role": "complex", "slide": 2, "png": "proof_complex.png"},
          {"role": "data", "slide": 1, "png": "proof_data.png"}])},
     False, "no `signature` anchor"),

    ("an anchor pointing at a slide the deck does not have is refused",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof=[
          {"role": "signature", "slide": 3, "png": "proof.png"},
          {"role": "complex", "slide": 2, "png": "proof_complex.png"},
          {"role": "data", "slide": 99, "png": "proof_data.png"}])},
     False, "not a final slide"),

    # THE RESTRAINT CARVE, on the escape agents/slide-design.md already documents: under a
    # *conservative* dial the risk is OPTIONAL — take a modest move, or write the one-clause
    # "deliberately restrained: <why>" so the field is never blank. That existed only in prose, so an
    # honest 5-minute lab-meeting plan was rejected for lacking a rendered proof of a risk it never
    # took, and the only escape ({"waived": …}) also switches off palette/type_scale/icon_family.
    # The carve must stay narrow, so every direction is pinned.
    ("conservative + a 'deliberately restrained:' move drops only signature_proof",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "signature_proof"}
      | {"boldness": "conservative",
         "signature_move": "deliberately restrained: 5-minute working update; one accent is "
                           "reserved for the new result and nothing competes with it"}},
     True, "signature_proof not required"),

    # at balanced+ and above a real signature move is required, not optional — the phrase alone
    # must not buy the exemption
    ("the 'deliberately restrained:' phrase does NOT work above the conservative dial",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "signature_proof"}
      | {"signature_move": "deliberately restrained: trying to dodge the proof"}},
     False, "signature_proof"),

    ("a conservative deck that took a REAL move still owes the proof",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "signature_proof"}
      | {"boldness": "conservative"}},
     False, "signature_proof"),

    ("the carved plan still needs a NON-BLANK signature_move",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items()
                      if k not in ("signature_proof", "signature_move")}
      | {"boldness": "conservative"}},
     False, "signature_move"),

    ("the carve is not a blanket exemption — type_scale is still required",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": {k: v for k, v in DESIGN_OK.items()
                      if k not in ("signature_proof", "type_scale")}
      | {"boldness": "conservative",
         "signature_move": "deliberately restrained: sober status readout"}},
     False, "type_scale"),

    ("an anchor pointing at a MISSING png is refused — including when it is not the first one, "
     "so the loop cannot stop at the signature anchor and call the set proven",
     {"critic": {"verdict": "consent", "rounds": 2},
      "design_plan": dict(DESIGN_OK, signature_proof=[
          {"role": "signature", "slide": 3, "png": "proof.png"},
          {"role": "complex", "slide": 2, "png": "proof_complex.png"},
          {"role": "data", "slide": 1, "png": "nope.png"}])},
     False, "does not exist"),
]


def build_icon_deck(dest: Path, *, logo_every=False, icon_slides=(), label_row=False, n=8) -> None:
    """A deck shaped to probe one branch of the icon waiver at a time."""
    sys.path.insert(0, str(SKILL / "scripts"))
    import deckkit as dk
    from PIL import Image
    lg = dest.parent / "iconfx_logo.png"
    Image.new("RGB", (64, 64), (30, 60, 120)).save(lg)
    prs = dk.blank_deck(10, 5.625)
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        dk.text(s, 0.6, 0.4, 8.8, 0.6, [[(f"Title of slide {i+1}", 28, dk.DEEP, True, False)]])
        for k in range(6):                                   # ordinary body copy: NOT a category set
            dk.text(s, 0.6, 1.3 + k * 0.34, 8.8, 0.32,
                    [[("body copy line carrying several real words", 14, dk.DEEP, False, False)]],
                    space_after=0)
        if logo_every:
            dk.logo(s, str(lg), corner="tr", h=0.6)
        if label_row and i in (1, 3, 5):                     # 4 short labels ACROSS the page
            for j, lab in enumerate(("Portrait", "Genre", "Still life", "Landscape")):
                dk.text(s, 0.6 + j * 2.2, 4.6, 2.0, 0.34, [[(lab, 15, dk.DEEP, True, False)]],
                        space_after=0)
        if i + 1 in icon_slides:
            for j in range(3):
                p = dest.parent / f"iconfx_{i}{j}.png"
                Image.new("RGB", (64, 64), (200, 40, 40)).save(p)
                dk.picture(s, str(p), 0.7 + j * 2.6, 3.9, 0.5, 0.5)
    prs.save(str(dest))


ICON_CASES = [
    # (name, deck shape, icon_family, must the waiver fire?)
    ("icon waiver: a LOGO repeated on every slide is not an icon family",
     {"tag": "ic_logo", "logo_every": True}, "none - brand allows only the logo", False),
    ("icon waiver: plain body copy is not a category set",
     {"tag": "ic_plain"}, "none - narrative deck, no entities", False),
    ("icon waiver: a real icon set contradicts a `none` record",
     {"tag": "ic_real", "icon_slides": (2, 4, 6)}, "none - conceptual content", True),
    ("icon waiver: a logo does not mask a real icon set",
     {"tag": "ic_both", "logo_every": True, "icon_slides": (2, 4, 6)}, "none - x", True),
    ("icon waiver: label ROWS across the page do contradict `none`",
     {"tag": "ic_rows", "label_row": True}, "none - concepts, icons would decorate", True),
    ("icon waiver: a declared family is never second-guessed",
     {"tag": "ic_decl", "icon_slides": (2, 4, 6)}, "tabler outline 1.75px", False),
]


# ── The consent record must be evidence about the DECK, not only about the FILE ──────────────
# Until the coverage gate existed, "verified against <review>.json" meant the artifact was present
# and still hashed to what was recorded. Measured on a 15-slide deck: a SCHEMA-VALID review
# declaring `slides_opened: [1]` was accepted by validate_review.py, recorded with a sha256, and
# printed as verified, and `--gate-check` exited 0 under "all hand-off gates pass". `slides_opened`
# is the anti-skim field and nothing compared it to the deck it claims to have read.
# Both directions matter. A per-section critic legitimately opens a RANGE, so the escape has to
# work — it just has to be DECLARED (`coverage.scope`) rather than inferred from a short list,
# because "short list" and "skimmed" are the same bytes.
def _record_review(deck: Path, review: dict, name: str = "review.json") -> dict:
    """A critic block as `validate_review.py --record` writes it: path + sha256 + verdict.

    `reviews_seen`, not `rounds`: the recorder counts review FILES, and leaves `rounds` to the
    coordinator who is the only one who knows it (it preserves a `rounds` already present and
    invents none). This fixture used to write `rounds: 1`, i.e. it claimed to be a tool-written
    record while carrying a field the tool does not write — so the shape these tests exercise was
    not the shape the gate meets in the field.
    """
    import hashlib
    rp = deck.parent / name
    rp.write_text(json.dumps(review))
    h = hashlib.sha256(rp.read_bytes()).hexdigest()
    return {"verdict": "consent", "reviews_seen": 1, "blockers": 0, "majors": 0,
            "source": str(rp), "sha256": h, "reviews": [str(rp)],
            "recorded_by": "validate_review.py"}


def _review(opened, scope=None) -> dict:
    cov = {"slides_opened": list(opened), "passes": ["content lens", "design lens"],
           "stats_block_seen": True, "contract_card_seen": True}
    if scope is not None:
        cov["scope"] = list(scope)
    return {"purpose": "lab meeting", "coverage": cov, "plan_audit": None,
            "verdict": "consent", "summary": "ok", "strengths": ["clear"], "findings": []}


# (name, slides_opened, scope, should_pass, needle) — deck built by build_deck() has 3 slides
COVERAGE_CASES = [
    ("whole-deck review covering every slide passes", [1, 2, 3], None, True,
     "opened 3/3 slides"),
    ("a review that opened 1 of 3 is refused", [1], None, False,
     "never lists slide(s) 2, 3"),
    ("a review missing one middle slide is refused", [1, 3], None, False,
     "never lists slide(s) 2"),
    ("a DECLARED per-section scope passes", [2, 3], (2, 3), True,
     "opened 2/2 slides"),
    ("a declared scope with a hole inside it is still refused", [2], (2, 3), False,
     "never lists slide(s) 3"),
]


def check_coverage(deck: Path) -> tuple[int, int]:
    ok = bad = 0
    for name, opened, scope, should_pass, needle in COVERAGE_CASES:
        g = {"critic": _record_review(deck, _review(opened, scope)),
             "design_plan": DESIGN_OK, "provenance": PROV_OK, "content": ARC_OK}
        code, out = run_gate(deck, g)
        good = (code == 0) == should_pass and needle in out
        if good:
            ok += 1
            print(f"  ok   {name}")
        else:
            bad += 1
            print(f"  FAIL {name}: exit={code} (wanted pass={should_pass}), missing {needle!r}")
            print("       " + out.strip().replace("\n", "\n       ")[:400])
    return ok, bad


# ── An arbiter pass is corroboration only when it CORROBORATES ────────────────────────────────
# `dulled_reopened` was written by one line of validate_review.py and read by nothing, so a Job-2
# payload reporting resolved=False + dulled=True + a regressed neighbour printed as "consent
# corroborated by 1 arbiter pass(es)" and exited 0. A failed verification round became a hand-off
# credential. The clean pass must keep working, or the fix would just make arbitration unusable.
ARBITER_CASES = [
    ("an arbiter pass with nothing open corroborates", [], True, "no open items"),
    ("an unresolved finding is not corroboration",
     [{"finding_ref": "F1", "resolved": False, "still_wrong": "callout still on the footer",
       "regressions": [], "from": "a.json"}], False, "NOT resolved"),
    ("a dulled strength is not corroboration",
     [{"finding_ref": "F2", "resolved": True, "dulled": True, "regressions": [],
       "from": "a.json"}], False, "dulled a named strength"),
    ("a regressed neighbour is not corroboration",
     [{"finding_ref": "F3", "resolved": True, "dulled": False,
       "regressions": ["slide 4 lost its hero"], "from": "a.json"}], False, "regressed"),
]


def check_arbiter(deck: Path) -> tuple[int, int]:
    ok = bad = 0
    for name, open_items, should_pass, needle in ARBITER_CASES:
        critic = _record_review(deck, _review([1, 2, 3]))
        critic["corroborated_by"] = [str(deck.parent / "a.json")]
        critic["dulled_reopened"] = len(open_items)
        critic["arbiter_open"] = open_items
        g = {"critic": critic, "design_plan": DESIGN_OK, "provenance": PROV_OK, "content": ARC_OK}
        code, out = run_gate(deck, g)
        good = (code == 0) == should_pass and needle in out
        if good:
            ok += 1
            print(f"  ok   {name}")
        else:
            bad += 1
            print(f"  FAIL {name}: exit={code} (wanted pass={should_pass}), missing {needle!r}")
            print("       " + out.strip().replace("\n", "\n       ")[:400])
    return ok, bad


# ── The three answers the skill itself disowns are not signature moves ────────────────────────
# SKILL.md Step 2, slide-design self-verify (h) and review-rubrics all name "a big number / a nice
# gradient / a full-bleed photo" as the SAFE CATALOGUE. Nothing checked, so the literal example
# passed the gate. Paraphrase evades this trivially and that is fine — the case it closes is the
# example copied verbatim because it was the nearest words to hand.
SIGNATURE_CASES = [
    ("the literal safe-catalogue answer is refused", "a big number", False),
    ("...with a trailing period too", "A Big Number.", False),
    ("a nice gradient is refused", "a nice gradient", False),
    ("a full-bleed photo is refused", "a full-bleed photo", False),
    ("a real move passes", "an off-grid k-space band that runs past its own axis", True),
    ("a move that merely CONTAINS the phrase passes",
     "a big number set into the negative space the sampling mask leaves open", True),
]


def check_signature(deck: Path) -> tuple[int, int]:
    ok = bad = 0
    for name, move, should_pass in SIGNATURE_CASES:
        g = {"critic": _record_review(deck, _review([1, 2, 3])),
             "design_plan": dict(DESIGN_OK, signature_move=move), "provenance": PROV_OK, "content": ARC_OK}
        code, out = run_gate(deck, g)
        good = (code == 0) == should_pass and (should_pass or "SAFE CATALOGUE" in out)
        if good:
            ok += 1
            print(f"  ok   {name}")
        else:
            bad += 1
            print(f"  FAIL {name}: exit={code} (wanted pass={should_pass})")
            print("       " + out.strip().replace("\n", "\n       ")[:300])
    return ok, bad


# ── A bypass may skip the gates; it may not assert they passed ────────────────────────────────
def check_skip_env(deck: Path) -> tuple[int, int]:
    import os
    env = dict(os.environ, SLIDE_MAKER_SKIP_GATES="1")
    (deck.parent / ".deck-gates.json").unlink(missing_ok=True)   # no record AT ALL
    p = subprocess.run([sys.executable, str(RENDER), str(deck), "--gate-check", "--static"],
                       capture_output=True, text=True, env=env)
    out = p.stdout + p.stderr
    ok = bad = 0
    for name, cond in (
            ("SKIP_GATES says it skipped", "SKIPPED" in out),
            ("SKIP_GATES never claims a pass", "all hand-off gates pass" not in out),
            ("SKIP_GATES fails --gate-check rather than blessing it", p.returncode != 0)):
        if cond:
            ok += 1
            print(f"  ok   {name}")
        else:
            bad += 1
            print(f"  FAIL {name}: exit={p.returncode}\n       "
                  + out.strip().replace("\n", "\n       ")[:300])
    return ok, bad


# ── SAMENESS: the deck-level monotony the [stats] block measured and nobody read ──────────────
# Calibrated against 11 decks BUILT AND LINTED in the registers this skill prescribes. The raw
# signal count is a bad gate: a 6-slide template-locked status update with ZERO hard findings
# already reaches 3 families, a 9-card 小红书 carousel built exactly to canvas-formats.md's DNA
# reaches 7, an appendix-heavy defense deck reaches 5. Every one is a deck the skill tells you to
# build, so the scope is three deterministic properties (body_n>=8, landscape, not surface) rather
# than a taxonomy of registers — and the threshold is >=4 DISTINCT codes with >=1 structural.
# The false-positive direction matters more than the true-positive one here: a blocking gate that
# fires on a good deck is worse than no gate ("a report that is always wrong is a report everyone
# learns to ignore" — ci.yml).
def _R(t, size=20, c=None, b=True):
    import deckkit as _dk
    return (t, size, c or _dk.DEEP, b, False, _dk.FONT)


def build_samey(dest: Path, n=12, portrait=False, appendix_at=None) -> Path:
    """cover + identical body pages + closer: one skeleton, one card grid, one strip, one rule."""
    import deckkit as dk
    prs = dk.blank_deck(7.5, 10.0) if portrait else dk.blank_deck()
    W = 7.5 if portrait else 10.0
    s = dk.add_slide(prs)
    dk.text(s, 0.5, 2.0, W - 1, 1.0, [[_R("Cover", 44)]])
    for i in range(n - 2):
        s = dk.add_slide(prs)
        if appendix_at is not None and i + 2 == appendix_at:
            dk.design_intent(s, role="appendix", reason="backup slides")
        dk.text(s, 0.5, 0.35, W - 1, 0.5, [[_R("Section %d" % (i + 1), 20)]])
        dk.hrule(s, 0.5, 0.95, W - 1)
        for (cx, cy, cw, _h) in dk.columns(3, slide=s, top=1.2, bottom=1.1):
            dk.box(s, cx, cy, cw, 2.2, fill="F2F4F7", round=True)
            dk.text(s, cx + 0.15, cy + 0.2, cw - 0.3, 1.6, [[_R("card body text", 12, dk.SLATE, False)]])
        dk.bottom_callout(s, 0.5, W - 1, "TAKEAWAY", "the same strip on every page")
    s = dk.add_slide(prs)
    dk.text(s, 0.5, 2.0, W - 1, 1.0, [[_R("Closing", 40)]])
    out = dest / "samey.pptx"
    prs.save(str(out))
    return out


def build_varied(dest: Path) -> Path:
    """cover + 8 distinct skeletons + closer — what the skill actually asks for."""
    import deckkit as dk
    prs = dk.blank_deck()
    s = dk.add_slide(prs); dk.box(s, 0, 0, 10, 5.625, fill=dk.DEEP)
    dk.text(s, 0.7, 2.0, 8.6, 1.2, [[_R("Cover", 44, dk.WHITE)]])
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Hero", 20)]]); dk.big_numeral(s, 3.5, 1.8, "34")
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Diagram", 20)]])
    dk.node(s, 1, 2, 2, 0.9, "In"); dk.node(s, 4, 2, 2, 0.9, "Model", hub=True); dk.node(s, 7, 2, 2, 0.9, "Out")
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Chart", 20)]])
    dk.native_chart(s, 0.8, 1.2, 8.4, 3.4, ["A", "B", "C"], [("v", [3, 5, 9])], kind="column")
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Table", 20)]])
    dk.table(s, 0.8, 1.4, 8.0, [["a", "b"], ["1", "2"], ["3", "4"]])
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Timeline", 20)]])
    dk.timeline(s, 0.8, 2.4, 8.4, [("Q1", "a", "x"), ("Q2", "b", "y"), ("Q3", "c", "z")])
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Two up", 20)]])
    L, Rr = dk.columns(2, slide=s, bottom=0.8)
    dk.bullet(s, L[0], L[1], L[2], [("one ", "x"), ("two ", "y")], size=16)
    dk.box(s, Rr[0], Rr[1], Rr[2], 2.4, fill="EEF2F7", round=True)
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Quote", 20)]])
    dk.pull_quote(s, 1.0, 1.8, 8.0, "A line that matters.")
    s = dk.add_slide(prs); dk.text(s, 0.5, 0.35, 9, 0.5, [[_R("Meter", 20)]])
    dk.meter_bar(s, 1.0, 2.2, 7.5, 0.62, label="share", value="62%")
    s = dk.add_slide(prs); dk.box(s, 0, 0, 10, 5.625, fill=dk.DEEP)
    dk.text(s, 0.7, 2.3, 8.6, 0.9, [[_R("Closing", 38, dk.WHITE)]])
    out = dest / "varied.pptx"
    prs.save(str(out))
    return out


def _fired(deck: Path):
    """The codes this deck actually fires — the same call the gate makes."""
    import io, contextlib
    import lint_deck as ld
    stats, buf = {}, io.StringIO()
    with contextlib.redirect_stdout(buf):
        ld.lint(str(deck), mode="presented", static_ok=True, stats_out=stats)
    return tuple(stats.get("sameness_codes") or ()), int(stats.get("body_n") or 0)


# ── CONCEPT: what the deck's idea is a PICTURE of, and the two pictures it beat ───────────────
# The pipeline diverged on STYLE (the direction gate renders "the same four slide types … only the
# style differs") and on LAYOUT (form-selection's per-slide runner-up) and never on the IDEA. Three
# governing images for one argument is the missing divergence; one picture with no alternatives is
# the first thing that came to mind, which is the default the field exists to interrupt.
CONCEPT_CASES = [
    ("a complete concept passes", DESIGN_OK["concept"], True, "concept:"),
    ("a bare string is not a choice", "an intelligence network", False, "must name the governing image"),
    ("one rejected alternative is not enough", {
        "chosen": "a network", "rejected": [{"concept": "an organism", "why_lost": "too soft"}]},
     False, "must name the governing image"),
    ("the same picture relabelled is refused", {
        "chosen": "an intelligence network",
        "rejected": [{"concept": "An Intelligence Network", "why_lost": "x"},
                     {"concept": "a network of intelligence", "why_lost": "y"}]},
     False, "same picture more than once"),
    ("a rejected concept with no reason is refused", {
        "chosen": "a network", "rejected": [{"concept": "an organism", "why_lost": ""},
                                            {"concept": "two hands", "why_lost": "z"}]},
     False, "why_lost"),
]


def check_concept(deck_path: Path) -> tuple[int, int]:
    ok = bad = 0
    deck = deck_path
    for name, con, should_pass, needle in CONCEPT_CASES:
        g = {"critic": _record_review(deck, _review(list(range(1, 4)))),
             "design_plan": dict(DESIGN_OK, concept=con), "provenance": PROV_OK, "content": ARC_OK}
        code, out = run_gate(deck, g)
        good = (code == 0) == should_pass and needle in out
        if good:
            ok += 1
            print("  ok   concept: " + name)
        else:
            bad += 1
            print("  FAIL concept: %s: exit=%d (wanted pass=%s), missing %r" %
                  (name, code, should_pass, needle))
            print("       " + out.strip().replace("\n", "\n       ")[:300])
    return ok, bad


def check_sameness(deck_path: Path) -> tuple[int, int]:
    td = deck_path.parent
    import lint_deck as ld
    ok_n = bad_n = 0

    def one(name, cond, detail=""):
        nonlocal ok_n, bad_n
        if cond:
            ok_n += 1
            print("  ok   " + name)
        else:
            bad_n += 1
            print("  FAIL " + name + (("\n       " + str(detail)[:400]) if detail else ""))

    # These sameness fixtures ARE repeated label rows with no pictures, so the icon gate fires on
    # them correctly — they simply predate it. Record the re-decision the gate asks for rather than
    # weakening it; this suite is about the SAMENESS composite and must not be the place an icon
    # rule gets quietly relaxed.
    _dp = dict(DESIGN_OK)
    _dp["icon_none_checked"] = ["slide 2", "slide 3", "slide 4"]
    # The icon `none` waiver now needs a CLASSIFIED high-bar reason too (not just named slides) —
    # naming the class is what a casual "not category-rich" cannot do. This synthetic sameness deck
    # is a repeated-label fixture, so any valid category clears it for the purpose of THIS suite
    # (which tests the sameness composite, not the icon rule).
    _dp["icon_none_category"] = "editorial-register"
    base = {"critic": {"waived": "No subagent dispatch on this host; both lenses ran inline.",
                       "waived_category": "no-dispatch-on-host", "inline_ran": True},
            "design_plan": _dp, "provenance": PROV_OK, "content": ARC_OK}

    d = td / "sm"; d.mkdir()
    samey = build_samey(d, n=12)
    write_proof(d)
    fired, body_n = _fired(samey)
    one("samey deck fires >=4 codes with a structural member (%d: %s)" % (len(fired), ",".join(fired)),
        len(fired) >= 4 and any(c in ld.SAMENESS_STRUCTURAL for c in fired), fired)
    code, out = run_gate(samey, dict(base))
    one("samey deck BLOCKS the hand-off", code != 0 and "SAMENESS:" in out, out[-400:])
    one("...and the message names the fired codes", all(c in out for c in fired), out[-400:])

    dv = td / "vr"; dv.mkdir()
    varied = build_varied(dv)
    write_proof(dv)
    vf, _ = _fired(varied)
    code, out = run_gate(varied, dict(base))
    one("a varied deck PASSES (fired %d)" % len(vf), code == 0 and "hand-off gates pass" in out, out[-400:])

    # THE load-bearing false-positive test: repetition ALONE must not block.
    one("a deck firing <4 codes passes even though it repeats", len(vf) < 4, vf)

    ds = td / "sml"; ds.mkdir()
    small = build_samey(ds, n=8)
    write_proof(ds)
    code, out = run_gate(small, dict(base))
    one("too small => not applied (body_n<8)", code == 0 and "not applied" in out, out[-300:])

    dp = td / "prt"; dp.mkdir()
    port = build_samey(dp, n=12, portrait=True)
    write_proof(dp)
    code, out = run_gate(port, dict(base))
    one("portrait/carousel => not applied", code == 0 and "not applied" in out, out[-300:])

    code, out = run_gate(samey, dict(base, delivery="surface"), "--surface")
    one("surface => not applied", code == 0 and "not applied" in out, out[-300:])

    da = td / "apx"; da.mkdir()
    apx = build_samey(da, n=12, appendix_at=6)
    write_proof(da)
    _, apx_body = _fired(apx)
    code, out = run_gate(apx, dict(base))
    one("declared appendix shrinks the body run (body_n=%d)" % apx_body,
        apx_body < 8 and code == 0 and "not applied" in out, out[-300:])

    # waivers
    good = "This is a registered corporate template whose grid the deck may not break."
    code, out = run_gate(samey, dict(base, sameness={
        "waived": good, "waived_category": "template-locked", "codes": list(fired)}))
    one("a complete waiver passes and is announced", code == 0 and "sameness: WAIVED" in out, out[-400:])
    code, out = run_gate(samey, dict(base, sameness={
        "waived": good, "waived_category": "template-locked", "codes": list(fired)[:-1]}))
    one("a waiver written for a DIFFERENT code set dies", code != 0 and "does not certify" in out, out[-300:])
    code, out = run_gate(samey, dict(base, sameness={
        "waived": "too short", "waived_category": "template-locked", "codes": list(fired)}))
    one("a thin waiver reason dies", code != 0, out[-200:])
    code, out = run_gate(samey, dict(base, sameness={"waived": good, "codes": list(fired)}))
    one("an uncategorised waiver dies", code != 0 and "waived_category" in out, out[-200:])
    code, out = run_gate(varied, dict(base, sameness={
        "waived": good, "waived_category": "template-locked", "codes": list(vf)}))
    one("an unneeded waiver is NOISY, not fatal", code == 0 and "NOT needed" in out, out[-300:])

    # the stringly-typed contract this whole gate rests on
    src = (SKILL / "scripts" / "lint_deck.py").read_text(encoding="utf-8")
    missing = [c for c in ld.SAMENESS_CODES if ('"%s: ' % c) not in src and ("'%s: " % c) not in src]
    one("every counted code is still emitted by lint_deck's own f-strings", not missing, missing)
    one("structural codes are a subset of the counted set",
        set(ld.SAMENESS_STRUCTURAL) <= set(ld.SAMENESS_CODES))
    one("every render-dependent code is declared in _PIXEL_CHECKS",
        set(ld.SAMENESS_RENDER_DEPENDENT) <= set(ld._PIXEL_CHECKS))

    # The deck this repo asserts is good must score ZERO. lint_fixture's own docstring: these
    # slides "pass the gates TODAY and must still pass afterwards. If a change breaks one of
    # these, the change is wrong." It is also the empirical reason TIMID COVER and FLAT TYPE are
    # NOT counted — this fixture emits both, so a composite that counted type drama would start
    # the asserted-good deck at 2 of 7 before looking at a single repeated page.
    import contextlib, io
    sys.path.insert(0, str(HERE))
    import lint_fixture
    # lint_fixture binds OUT = Path.cwd() at IMPORT time, so chdir-ing afterwards does not move
    # where it writes. Ask it where it put the file instead of assuming.
    with contextlib.redirect_stdout(io.StringIO()):
        lint_fixture.build_pass()
    fx_pass = lint_fixture.OUT / "fx_pass.pptx"
    st, buf = {}, io.StringIO()
    try:
        with contextlib.redirect_stdout(buf):
            ld.lint(str(fx_pass), mode="presented", static_ok=True, stats_out=st)
    finally:
        try:                      # do not leave a fixture deck in whatever dir CI ran from
            fx_pass.unlink()
        except OSError:
            pass
    one("the repo's must-stay-clean PASS fixture scores ZERO sameness codes",
        not st.get("sameness_codes"), st.get("sameness_codes"))
    one("...and it DOES emit TIMID COVER / FLAT TYPE, which is why they are excluded",
        any(w.startswith("TIMID COVER") for w in st.get("warns", []))
        and any(w.startswith("FLAT TYPE") for w in st.get("warns", [])),
        [w.split(":")[0] for w in st.get("warns", [])])
    return ok_n, bad_n



def check_form_reach(deck: Path) -> tuple[int, int]:
    """A report that never asks for an answer is a line people learn to scroll past.

    `form reach` printed `1 of 23 named components; the rest is raw box/text` on a DELIVERED deck
    and let it through — and three of that deck's review findings were defects the unused
    components prevent by construction (a label grazing its bar, a value floating off a track's
    centreline, a reference line drawn three different ways).

    🔴 It must never block on the NUMBER. Bespoke composition is legitimate and is often the
    signature move itself — a Mondrian page cannot come from a catalogue. It blocks on the absence
    of a DECISION, which is the failure that actually happened: nobody looked. Both directions are
    asserted, and the silent cases are the load-bearing half.
    """
    ok = bad = 0

    def one(name, cond, detail=""):
        nonlocal ok, bad
        if cond:
            ok += 1
            print("  ok   form_reach: " + name)
        else:
            bad += 1
            print("  FAIL form_reach: " + name + ("  — " + str(detail) if detail else ""))

    base = {"critic": {"verdict": "consent", "rounds": 2}, "provenance": PROV_OK, "content": ARC_OK,
            "density": {"waived": "a fixture deck is not a real density decision"}}

    # the fixture deck has no build*.py beside it, so the reporter cannot even measure reach —
    # it must stay silent rather than inventing a blocker for a deck it knows nothing about.
    code, out = run_gate(deck, dict(base, design_plan=DESIGN_OK))
    one("a deck with no build script is unaffected (nothing to measure)",
        "form reach is" not in out, out[-140:])

    # now give it a build script that is all raw box/text — the measured failure
    script = deck.parent / "build_fixture.py"
    script.write_text("import deckkit as dk\n"
                      "def main(s):\n"
                      "    dk.box(s, 0, 0, 1, 1)\n"
                      "    dk.text(s, 0, 0, 1, 1, [[('x', 12, dk.DEEP, False, False)]])\n",
                      encoding="utf-8")
    try:
        code, out = run_gate(deck, dict(base, design_plan=DESIGN_OK))
        one("a raw box/text build with no recorded reason BLOCKS", code != 0, code)
        one("...and the message names the one command that answers it",
            "sigs.py --list" in out, out[-160:])
        one("...and says it blocks on the decision, not the number",
            "blocks on there being no decision" in out, out[-160:])

        d = dict(DESIGN_OK)
        d["form_reach"] = {"waived": "the deck's whole argument is one bespoke datum line"}
        code, out = run_gate(deck, dict(base, design_plan=d))
        one("a recorded reason passes", code == 0, out[-160:])
        one("...and the waiver is printed, not swallowed", "form reach WAIVED" in out, out[-160:])

        d["form_reach"] = {"waived": "bespoke"}
        code, out = run_gate(deck, dict(base, design_plan=d))
        one("a one-word reason is not a reason", code != 0, code)

        d["form_reach"] = "bespoke composition throughout this deck"
        code, out = run_gate(deck, dict(base, design_plan=d))
        one("a bare string is not the recorded shape either", code != 0, code)
    finally:
        script.unlink(missing_ok=True)

    # ---- plan files in the deliverable folder ------------------------------------------
    # checkpoint-convention.md says the checkpoint artifact is pasted into the CONVERSATION and
    # that content-plan.md / design-plan.md must not be written into the deck folder. That lived
    # as one sentence inside a reference file — the position most likely to be skipped — and
    # nothing anywhere reported one, so the user found the clutter instead of a gate.
    # Advisory, not blocking: the same sentence carves out "unless the user asked for plan files".
    code, out = run_gate(deck, dict(base, design_plan=DESIGN_OK))
    one("a clean folder says nothing about plan files", "plan files in the deliverable" not in out)
    plan = deck.parent / "design-plan.md"
    plan.write_text("# design plan\n", encoding="utf-8")
    try:
        code, out = run_gate(deck, dict(base, design_plan=DESIGN_OK))
        one("a stray design-plan.md is reported", "plan files in the deliverable" in out, out[-140:])
        one("...and it is ADVISORY — the hand-off is not blocked over clutter", code == 0, code)
        one("...and it names the file", "design-plan.md" in out)
    finally:
        plan.unlink(missing_ok=True)
    return ok, bad


def main() -> int:
    passed = failed = 0
    with tempfile.TemporaryDirectory() as td:
        deck = build_deck(Path(td))
        write_proof(Path(td))
        for name, critic_block, should_pass, needle in CASES:
            gates = dict(critic_block)
            gates.setdefault("design_plan", DESIGN_OK)
            gates.setdefault("provenance", PROV_OK)
            gates.setdefault("content", ARC_OK)
            code, out = run_gate(deck, gates)
            ok = (code == 0) == should_pass and needle in out
            if ok:
                passed += 1
                print(f"  ok   {name}")
            else:
                failed += 1
                print(f"  FAIL {name}: exit={code} (wanted pass={should_pass}), "
                      f"missing {needle!r}")
                print("       " + out.strip().replace("\n", "\n       ")[:400])
        # ── the `build_shape` gate. The build step is 40-71% of a session's model-active minutes
        # and the fan-out rule that addresses it lived only in prose: a 13-slide deck was built
        # solo at 241 round-trips against a ~125 budget, and nothing anywhere asked why. From ~6
        # content slides the DECISION is recorded — fanout, or solo with a reason. Never blocks
        # the choice (solo is mandatory on hosts without dispatch), only the absence of one.
        # The 3-slide deck above is under the threshold, which is itself asserted: a tiny deck
        # must never be asked this question.
        big_dir = Path(td) / "big"
        big_dir.mkdir()
        big = build_deck(big_dir, 8)
        write_proof(big_dir)
        base_big = {"critic": {"waived": GOOD_REASON, "waived_category": "no-dispatch-on-host",
                               "inline_ran": True},
                    "design_plan": {k: v for k, v in DESIGN_OK.items() if k != "build_shape"},
                    "provenance": PROV_OK, "content": content_ok(8)}
        code, out = run_gate(big, base_big)
        if code != 0 and "build_shape" in out:
            passed += 1; print("  ok   an 8-slide deck with no build_shape is refused")
        else:
            failed += 1; print(f"  FAIL an 8-slide deck with no build_shape passed (exit={code})")
        base_big["design_plan"] = dict(DESIGN_OK, build_shape="solo — host has no subagent dispatch")
        code, out = run_gate(big, base_big)
        if code == 0 and "build shape: solo" in out:
            passed += 1; print("  ok   solo with a reason passes, and is echoed")
        else:
            failed += 1; print(f"  FAIL solo-with-reason was refused (exit={code})\n       " + out.strip()[-300:])
        code, out = run_gate(deck, {"critic": {"waived": GOOD_REASON,
                                               "waived_category": "no-dispatch-on-host",
                                               "inline_ran": True},
                                    "design_plan": dict(DESIGN_OK), "provenance": PROV_OK,
                                    "content": ARC_OK})
        if code == 0:
            passed += 1; print("  ok   a 3-slide deck is never asked for build_shape")
        else:
            failed += 1; print(f"  FAIL the under-threshold deck was asked for build_shape (exit={code})")

        # ── the `icon waiver` gate. `icon_family: "none - <reason>"` is free text written at PLAN
        # time, before any slide exists, and nothing revisited it: one real build shipped ZERO icons
        # past every gate on a deck of category slides. It must stay satisfiable for a genuinely
        # icon-free deck, so both FALSE-POSITIVE cases below matter at least as much as the true
        # ones — every one of them was a live bug in the first cut of this check.
        for name, kw, fam, want in ICON_CASES:
            d2 = Path(td) / (kw["tag"] + ".pptx")
            build_icon_deck(d2, **{k: v for k, v in kw.items() if k != "tag"})
            g = {"critic": {"verdict": "consent", "rounds": 2},
                 "design_plan": dict(DESIGN_OK, icon_family=fam, carried_by=[2, 3]),
                 "provenance": PROV_OK, "content": ARC_OK}
            code, out = run_gate(d2, g)
            if "hand-off gates pass" not in out and "icon waiver" not in out:
                failed += 1
                print(f"  FAIL {name}: the gate aborted before the icon check ran, so this "
                      f"assertion means nothing\n       " + out.strip().splitlines()[-1][:160])
                continue
            fired = "icon waiver" in out
            if fired == want:
                passed += 1
                print(f"  ok   {name}")
            else:
                failed += 1
                print(f"  FAIL {name}: icon waiver {'fired' if fired else 'stayed silent'}, "
                      f"wanted the opposite")
        o, b = check_delivery(deck)
        passed += o
        failed += b
        for fn in (check_coverage, check_arbiter, check_signature, check_skip_env,
                   check_concept, check_sameness, check_form_reach):
            o, b = fn(deck)
            passed += o
            failed += b
    print(f"\n{passed} passed, {failed} failed")
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(main())

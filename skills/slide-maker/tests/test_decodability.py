#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""A viewer must be able to DECODE what is on the page — and that must be measured, not claimed.

WHY. Three pieces of feedback from the first human reader of a deck this repo built, all of which
every automated gate had passed:

  「icon 应该是一个必须包含的东西但是这个却没有」   the icon gate was cleared by a prose waiver
  「1/5/10/12 的设计我没有看懂」                    the signature move's stranger test was asserted
  「第一页的那些横线是什么意思」                     nine repeated marks nobody had named

One root cause under all three: **the skill could SAY a visual element was readable, and nothing
checked whether the sentence was true.** Each fix below turns a rule the skill had already written
in prose into a measurement over the built file.

  MOTIF_UNEXPLAINED_AT_FIRST_USE  SKILL.md says, in as many words, that "a reading that defers to
                                  a later slide is a FAILED test written as a passing sentence" —
                                  and the check tested the opposite, clearing on a legend ANYWHERE.
                                  Measured: a loud motif debuted on the cover with nothing to read
                                  it by, the legend arrived later, the deck passed, and the first
                                  reader asked what it meant.
  UNNAMED_REPEATED_MARK           a set of identical marks with no text near it was invisible to
                                  every check: not text, not tagged as a motif, and trivially
                                  clearing contrast and overlap. Nine unlabelled rules on a cover
                                  shipped, and the reader's question was "什么意思".
  icon waiver categories          `motif-dominant` / `tiny-deck` / `template-locked` were compared
                                  against a list of strings. Any of the four words cleared the
                                  gate, so the WORD did the work and not the fact.

Run: python3 tests/test_decodability.py
"""
from __future__ import annotations

import contextlib
import io
import pathlib
import sys
import tempfile
import warnings

warnings.filterwarnings("ignore")
HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import deckkit as dk            # noqa: E402
import render_deck as rd        # noqa: E402

OKS: list[str] = []
FAILS: list[str] = []
TMP = pathlib.Path(tempfile.mkdtemp(prefix="decodability-"))


def check(cond, msg, detail=""):
    (OKS if cond else FAILS).append(msg if cond else "{} — {}".format(msg, detail))


def codes(prs):
    buf = io.StringIO()
    with contextlib.redirect_stdout(buf):
        dk.lint_layout(prs)
    out = buf.getvalue()
    return {w for w in ("MOTIF_UNEXPLAINED_AT_FIRST_USE", "MOTIF_UNEXPLAINED",
                        "UNNAMED_REPEATED_MARK", "MOTIF_BUDGET") if w in out}


# ── 1. the stranger test is about FIRST appearance ───────────────────────────────────────────
prs = dk.blank_deck()
s1 = dk.add_slide(prs)
dk.tag_motif(dk.box(s1, 0.7, 2.0, 8.6, 0.02, fill=dk.MAGENTA), loud=True)
dk.text(s1, 0.7, 2.6, 8.6, 0.6, [[("A cover", 28, dk.DEEP, True, False)]])
s2 = dk.add_slide(prs)
dk.tag_motif(dk.box(s2, 0.7, 2.0, 8.6, 0.02, fill=dk.MAGENTA), loud=True)
dk.motif_legend(s2, "the rule marks a boundary", x=0.7, y=2.3)
check("MOTIF_UNEXPLAINED_AT_FIRST_USE" in codes(prs),
      "a loud motif that debuts unexplained and is keyed LATER is caught — deferring the reading "
      "to a later slide is the exact failure SKILL.md names, and the old check cleared it",
      str(codes(prs)))

prs = dk.blank_deck()
s1 = dk.add_slide(prs)
dk.tag_motif(dk.box(s1, 0.7, 2.0, 8.6, 0.02, fill=dk.MAGENTA), loud=True)
dk.motif_legend(s1, "the rule marks a boundary", x=0.7, y=2.3)
s2 = dk.add_slide(prs)
dk.tag_motif(dk.box(s2, 0.7, 2.0, 8.6, 0.02, fill=dk.MAGENTA), loud=True)
check("MOTIF_UNEXPLAINED_AT_FIRST_USE" not in codes(prs),
      "...and keying it AT its debut passes — the rule is about where the reader meets the device",
      str(codes(prs)))

prs = dk.blank_deck()
for _ in range(2):
    s = dk.add_slide(prs)
    dk.tag_motif(dk.box(s, 0.7, 2.0, 8.6, 0.02, fill=dk.MAGENTA), loud=True)
check("MOTIF_UNEXPLAINED" in codes(prs),
      "a loud motif with NO legend anywhere is still caught (the pre-existing check is intact)")

# ── 2. repetition nobody named ───────────────────────────────────────────────────────────────
def _marks(named):
    prs = dk.blank_deck()
    s = dk.add_slide(prs)
    for i in range(9):
        dk.box(s, 3.0, 1.5 + i * 0.13, 6.0, 0.01, fill=dk.DEEP)
    if named:
        dk.text(s, 1.2, 2.0, 1.5, 0.4, [[("TWELVE FLOORS", 11, dk.DEEP, True, False)]])
    dk.text(s, 0.7, 4.6, 8.6, 0.6, [[("A cover", 28, dk.DEEP, True, False)]])
    return prs


check("UNNAMED_REPEATED_MARK" in codes(_marks(False)),
      "nine identical rules with no text near them are caught — not text, not tagged, and they "
      "clear contrast and overlap trivially, so every other check was blind to them",
      str(codes(_marks(False))))
check("UNNAMED_REPEATED_MARK" not in codes(_marks(True)),
      "...and the SAME nine rules pass once a word sits beside them. The fix a reader needs is a "
      "label, and that is exactly what clears it",
      str(codes(_marks(True))))

# A LABEL, not merely text that happens to be nearby. The first version asked "is there text
# within reach", and on the real cover this check was written for, a 28pt HEADLINE sat 0.22in below
# the nine rules and cleared them — a page title does not name your diagram. A label is small and
# close; a headline is large and is titling the page.
def _cover(named):
    prs = dk.blank_deck()
    sl = dk.add_slide(prs)
    for i in range(9):
        dk.box(sl, 2.55 if named else 0.72, 3.30 + i * 0.132,
               6.73 if named else 8.56, 0.010, fill=dk.DEEP)
    if named:
        dk.text(sl, 0.72, 3.30 + 4 * 0.132 - 0.30, 1.40, 0.30,
                [[("TWELVE", 11, dk.DEEP, True, False)]])
        dk.text(sl, 0.72, 3.30 + 4 * 0.132 - 0.04, 1.40, 0.30,
                [[("FLOORS", 11, dk.DEEP, True, False)]])
    dk.text(sl, 0.72, 4.58, 8.0, 0.46, [[("Every check was green.", 28, dk.DEEP, True, False)]])
    return prs


check("UNNAMED_REPEATED_MARK" in codes(_cover(False)),
      "a 28pt HEADLINE 0.22in below the marks does not count as naming them — this is the real "
      "cover the check was written for, and 'text is nearby' cleared it",
      str(codes(_cover(False))))
check("UNNAMED_REPEATED_MARK" not in codes(_cover(True)),
      "...and an 11pt label beside the group does count. Small and close is a label; large is a "
      "page title", str(codes(_cover(True))))

prs = dk.blank_deck()
s = dk.add_slide(prs)
for i in range(3):
    dk.box(s, 1.0 + i * 3.0, 2.0, 2.4, 1.2, fill=dk.TINT)
check("UNNAMED_REPEATED_MARK" not in codes(prs),
      "three cards are not reported — the floor is 4, so an ordinary card row is untouched")

# It must not fire on LAWFUL composition. Probed against the shapes this toolkit itself builds —
# the check is on every deck, so a false fire here is the reflex-to-waive problem this batch is
# otherwise about. TEXTURE is separated by CONTRAST, not by a count: a backdrop is faint by
# definition ("keep it near #EEE so it never fights body content" is its own docstring), and
# anything faint enough to be texture sits under the 3:1 floor where NON-TEXT CONTRAST already
# owns it. The two checks divide the space instead of both firing on a ground.
def _page(build):
    prs = dk.blank_deck()
    sl = dk.add_slide(prs)
    dk.box(sl, 0, 0, 10, 5.625, fill=dk.WHITE)
    build(sl)
    return prs


for _label, _build in (
        ("unit_grid (isotype)", lambda sl: dk.unit_grid(sl, 0.7, 1.5, 8.6, 2.5, 100, "people",
                                                        filled=62)),
        ("dot_strip", lambda sl: dk.dot_strip(sl, 0.6, 2.0, 8.0,
                                              [("A", 70), ("B", 100), ("C", 180)], 0, 200)),
        ("small_multiples", lambda sl: dk.small_multiples(
            sl, 0.6, 1.4, 8.8, 3.2, [("a", [1, 2, 3]), ("b", [3, 2, 1]), ("c", [2, 2, 2])])),
        ("timeline", lambda sl: dk.timeline(sl, 0.7, 2.0, 8.6,
                                            [("Q1", "a"), ("Q2", "b"), ("Q3", "c")])),
        ("backdrop_motif texture", lambda sl: (dk.backdrop_motif(sl, kind="grid"),
                                               dk.text(sl, 0.7, 2.0, 8.6, 0.8,
                                                       [[("Cover", 40, dk.DEEP, True, False)]]))),
):
    check("UNNAMED_REPEATED_MARK" not in codes(_page(_build)),
          "{} is not reported — a check that fires on lawful composition teaches the waive reflex"
          .format(_label), str(codes(_page(_build))))


def _bento(filled, fill=None):
    prs = dk.blank_deck()
    sl = dk.add_slide(prs)
    dk.box(sl, 0, 0, 10, 5.625, fill=dk.WHITE)
    for i, (x, y, w, h) in enumerate(dk.bento(sl, 0.6, 1.2, 8.8, 3.9, [(1, 1)] * 8, cols=4)):
        dk.box(sl, x, y, w, h, fill=fill or dk.TINT)
        if filled:
            dk.text(sl, x + 0.12, y + 0.10, w - 0.24, 0.4,
                    [[("tile %d" % i, 12, dk.DEEP, False, False)]])
    return prs


check("UNNAMED_REPEATED_MARK" not in codes(_bento(True)),
      "a bento grid with content in its tiles passes")
check("UNNAMED_REPEATED_MARK" not in codes(_bento(False)),
      "...and PALE empty tiles are not reported either — TINT reads 1.12:1 on white, which is "
      "texture by the same contrast test, and NON-TEXT CONTRAST owns anything that faint if it "
      "was meant to be read")
check("UNNAMED_REPEATED_MARK" in codes(_bento(False, fill=dk.DEEP)),
      "...but eight SOLID empty tiles are reported — a skeleton with nothing in it is not a page, "
      "and a reader cannot decode it either")

# ── 3. the icon waiver's category must be TRUE of the built file ─────────────────────────────
def _deck(n_slides, loud=False):
    prs = dk.blank_deck()
    for _ in range(n_slides):
        s = dk.add_slide(prs)
        if loud:
            dk.tag_motif(dk.box(s, 0.7, 2.0, 8.6, 0.02, fill=dk.MAGENTA), loud=True)
    p = TMP / ("d%d%s.pptx" % (n_slides, "m" if loud else ""))
    prs.save(str(p))
    return p


ok, why = rd._icon_none_category_holds("motif-dominant", _deck(12), [2, 3])
check(not ok and "loud motif" in why,
      "`motif-dominant` on a deck with NO loud motif is rejected — the strongest word on the list "
      "was clearing the gate on decks that had no motif to dilute", why[:90])
ok, _ = rd._icon_none_category_holds("motif-dominant", _deck(4, loud=True), [2])
check(ok, "...and it holds on a deck that really carries one")
ok, why = rd._icon_none_category_holds("tiny-deck", _deck(12), [2])
check(not ok and "tiny-deck" in why, "`tiny-deck` on twelve slides is rejected", why[:80])
ok, _ = rd._icon_none_category_holds("tiny-deck", _deck(2), [2])
check(ok, "...and holds on a 2-slide ask")
ok, why = rd._icon_none_category_holds("template-locked", _deck(6), [2])
check(not ok and "stock layouts" in why,
      "`template-locked` on a blank deck is rejected — python-pptx ships eleven named layouts, so "
      "'has layout names' proved nothing and this cleared on the first try", why[:90])
ok, _ = rd._icon_none_category_holds("editorial-register", _deck(12), [2])
check(ok, "`editorial-register` is left as declared — it is a taste claim about a look, and "
          "inventing a measurement for it would be worse than admitting there is none")

# ── 4. the waiver must name EVERY flagged slide ──────────────────────────────────────────────
src = (SCRIPTS / "render_deck.py").read_text(encoding="utf-8")
check("_icon_none_category_holds" in src and "icon waiver REJECTED" in src,
      "the gate calls the category check and says REJECTED out loud when it fails")
check("not re-decided" in src,
      "...and REPORTS the flagged slides a waiver did not name. Deliberately a report and not a "
      "hold: as a hold it rejected a lawful `editorial-register` waiver that had named two of a "
      "dozen pages, and the defect this batch actually had — a category FALSE of the built file — "
      "is caught precisely by the check above. A blunt second rule that fires on lawful use is how "
      "a gate earns the reflex to waive it.")

# ── 5. the two runtimes must hold the SAME icon bar ──────────────────────────────────────────
# Found by auditing this batch: the shared path learned to CHECK each category against the built
# file while the codex path kept comparing a string to a HAND-COPIED list — under a comment saying
# to keep them identical, which is the shape that drifts. A deck was motif-dominant-and-fine on one
# runtime and rejected on the other.
import importlib.util                                                        # noqa: E402

_spec = importlib.util.spec_from_file_location("cdg_icons", SCRIPTS / "codex_delivery_gate.py")
_cdg = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(_cdg)
check(tuple(_cdg.ICON_NONE_CATEGORIES) == tuple(rd._ICON_NONE_CATEGORIES),
      "both runtimes take the icon categories from ONE source, not a hand-copied mirror",
      "{} vs {}".format(_cdg.ICON_NONE_CATEGORIES, rd._ICON_NONE_CATEGORIES))
_ev = {"waivers": [{"kind": "icon", "scope": "undeclared-categorical",
                    "reason": "a reason long enough to clear the width floor",
                    "category": "motif-dominant"}]}
check(not _cdg._icon_waiver_ok(_ev, _deck(12), [2, 3]),
      "the codex path also REJECTS `motif-dominant` on a deck with no loud motif — a floor that "
      "one runtime enforces and the other does not is not a floor")
check(_cdg._icon_waiver_ok(_ev, _deck(4, loud=True), [2]),
      "...and accepts it where the motif is really there")

# ── 6. a declared register must be OBEYED, not merely paletted ───────────────────────────────
# Measured on this toolchain: the same content through all 18 `presets.apply()` calls produced 18
# pages differing only in ground colour, corner radius and rule weight — memphis without its
# coloured header bands, bauhaus without a primitive. `check_style_applied` verifies the CALL and
# `check_register_pixels` says in its own docstring that it judges COLOUR only, so "declare
# brutalist, ship its palette on rounded cards" cleared every gate.
import check_register_guard as crg                                           # noqa: E402
import presets as _presets                                                   # noqa: E402
from pptx.util import Inches as _In                                          # noqa: E402
from pptx.enum.shapes import MSO_SHAPE as _MS                                # noqa: E402

_gn = [0]


def _guard_deck(fn):
    prs = dk.blank_deck()
    sl = dk.add_slide(prs)
    fn(sl)
    _gn[0] += 1
    p = TMP / ("guard%d.pptx" % _gn[0])
    prs.save(str(p))
    return p


def _gcodes(path, reg):
    return {c for c, _m in crg.check(path, register=reg)[0]}


check("rounded" in _gcodes(_guard_deck(
          lambda sl: dk.box(sl, 1, 1, 3, 2, fill=dk.DEEP, round=True)), "swiss"),
      "a rounded card under `swiss` is caught — its own guard forbids them, which is why the "
      "preset sets radius=0")
check("rounded" not in _gcodes(_guard_deck(
          lambda sl: dk.box(sl, 1, 1, 3, 2, fill=dk.DEEP)), "swiss"),
      "...and a square one passes")


def _raw(sl):
    c = sl.shapes.add_shape(_MS.OVAL, _In(5), _In(0), _In(4.6), _In(4.6))
    c.fill.solid()
    c.fill.fore_color.rgb = dk.DEEP
    c.line.fill.background()          # shadow.inherit left True on purpose


check("soft-shadow" in _gcodes(_guard_deck(_raw), "bauhaus"),
      "a raw add_shape() that never switched shadow.inherit off is caught — the author of this "
      "check made exactly that mistake on the page that motivated it")


def _hero_and_title(sl):
    c = sl.shapes.add_shape(_MS.OVAL, _In(5), _In(0), _In(4.6), _In(4.6))
    c.fill.solid()
    c.fill.fore_color.rgb = dk.DEEP
    c.line.fill.background()
    c.shadow.inherit = False
    dk.text(sl, 0.6, 1.0, 5.1, 1.5, [[("three things", 40, dk.DEEP, True, False)]])


check("confetti" not in _gcodes(_guard_deck(_hero_and_title), "bauhaus"),
      "ONE hero primitive beside a large TITLE BLOCK is not confetti — a text box's prstGeom is "
      "`rect` like any square, and counting geometry alone reported a page carrying exactly one "
      "circle. Found on a real page: every fixture here was text-free, which is why they missed it")

_v, _f = crg.check(_guard_deck(lambda sl: dk.box(sl, 1, 1, 3, 2, fill=dk.DEEP, round=True)),
                   register="consulting")
check(not _v and _f.get("note"),
      "a register whose guard needs judgement of meaning is REPORTED as unchecked, never passed — "
      "only 7 of 18 declare machine-settleable prohibitions and the checker says so")
check(len(_presets.FORBIDS) < len(_presets.PRESETS),
      "FORBIDS is deliberately a SUBSET — a check that fired on lawful composition would teach the "
      "waive reflex, which is what this whole area is about")

# The register is read by `check_style_applied.declared_preset`, never by a substring search.
# Found on the real delivered deck: a BESPOKE look whose pick read "…beat blueprint-the-preset
# because…" was checked as `blueprint` — a substring cannot tell the register a deck DECLARES from
# the rival it says it BEAT, and checking a deck against a register it never chose is worse than
# not checking it.
_bespoke = {"design_plan": {"style_pick": "bespoke Section Drawing for a toolchain domain - beat "
                                          "blueprint-the-preset because a true blueprint is a plan "
                                          "view - anti-pick avoided: the terminal reflex"}}
_v, _f = crg.check(_guard_deck(lambda sl: dk.box(sl, 1, 1, 3, 2, fill=dk.DEEP, round=True)),
                   None, _bespoke)
check(_f.get("register") is None and not _v,
      "a bespoke pick that NAMES a preset as the rival it beat is not checked against that preset "
      "— one parser owns this, and it is the one that already handles the generated-template branch",
      "register={!r}".format(_f.get("register")))

_rd_src = (SCRIPTS / "render_deck.py").read_text(encoding="utf-8")
_cx_src = (SCRIPTS / "codex_delivery_gate.py").read_text(encoding="utf-8")
check("_gate_section('register_guard')" in _rd_src,
      "render_deck registers a `register_guard` section")
check("check_register_guard" in _cx_src and "check_register_guard.py" in _cx_src,
      "the codex path runs the SAME module — a register enforced on one runtime only is the drift "
      "this file has already been through twice on icons")

print("\n".join("  ok   " + m for m in OKS))
if FAILS:
    print("\n".join("  FAIL " + m for m in FAILS))
print("\n{} passed, {} failed".format(len(OKS), len(FAILS)))
raise SystemExit(1 if FAILS else 0)
